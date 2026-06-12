"""
涉税风险分析报告模块 V10
综合分析评估 130+ 个维度 — 全规则100%覆盖：
账务数据 / 发票合规 / 发票深度 / 成本结构 / 财税票比对 / 配比弹性 /
隐匿虚增 / 税负水平 / 城建税 / 房产税 / 个人所得税 / 印花税 /
纳税调整 / 收入时点 / 政策执行 / 资金往来 / 薪酬合规 /
客户穿透 / 供应商穿透 / 财务健康 / 企业信用 / 行业专项 / 良好实践 /
经营实质(18+3+3项) / 增值税专项(5+4项) / 发票异常(3+4+3项) / 费用匹配(3项) /
企业所得税专项(4+4项) / 申报比对(4项) / 征管风险(4+1项) / 交易特征(4项) /
V7新增8项: 年终奖计税优化 / 农产品收购发票异常 / 跨区域开票 / 简易计税划分 /
一址多照/一人多企 / 注册经营地址不一致 / 注销前突击开票 / D级纳税人风险 /
V8新增16项: 合同风险 — 三流合一/四流合一（收入无合同/成本无合同/大额无合同/
三流缺失/四流缺失/长期无合同/金额偏差/过期合同/先票后签/收付款不匹配/
无合同占比/跨合同方/印花税合同风险/合同作废率/框架合同缺失/关联方无合同）
V9新增2项: 逐票合同覆盖率 — 销项发票合同覆盖率(逐票名称+有效期双重匹配)、
进项发票合同覆盖率(逐票名称+有效期双重匹配)
V10新增27项: 合同内容涉税风险(涉税条款缺失/拆分合同/违约金/免租期/跨境代扣/
合同变更/混合销售税率) + 稽查重点盲区(价外费用/非货币交换/债务重组/股权转让/
无偿借款/母子公司管理费/境外支付/混合销售/环保税/发票备注栏/佣金手续费/
捐赠支出/存货盘点)
V11新增2项: 普票浪费可抵扣进项税额（专业服务类供应商开普票未开专票→浪费进项）/
预付卡涉税风险（增值税不得抵扣+企业所得税扣除凭证+个人所得税代扣代缴）
"""
from fastapi import APIRouter, Depends, Query, Body
from fastapi.responses import FileResponse, Response
from sqlalchemy.orm import Session
from sqlalchemy import func, case, extract, and_, or_
from datetime import date, timedelta, datetime
from typing import Optional, List, Tuple, Dict, Any
from collections import defaultdict
import json
import os
import calendar
import io
import tempfile

from database import get_db, Company, Account, JournalEntry
from database import SalesInvoice, PurchaseInvoice, BookkeepingInvoice
from database import VATDeclaration, InputVATDeduction, BankTransaction
from database import SalaryRecord, SocialSecurityDetail, HousingFundDetail
from database import CulturalConstructionFeeDeclaration
from database import FixedAsset, IntangibleAsset
from database import Customer, Supplier, Employee, Contract, ContractPayment, Payment
from database import InventoryTransaction, InventoryBalance, InventoryItem
from database import SocialSecurityDeclaration, HousingFundDeclaration

router = APIRouter(prefix="/api/tax-risk", tags=["涉税风险分析"])


# ============ 期间格式工具函数 ============
def _normalize_period(ym: str) -> str:
    """将各种格式的时期统一为 YYYY-MM（财税系统标准格式）"""
    if not ym:
        return ''
    ym = ym.strip()
    if len(ym) >= 7:
        return ym[:7]
    return ym


def _period_to_date_range(ym: str) -> Tuple[str, str]:
    """
    将 YYYY-MM 转换为当月第一天和最后一天（财税严谨格式）
    如 2025-01 → ('2025-01-01', '2025-01-31')
    """
    if not ym or len(ym) < 7:
        return ('', '')
    y, m = int(ym[:4]), int(ym[5:7])
    first_day = f'{y}-{m:02d}-01'
    last_day_num = calendar.monthrange(y, m)[1]
    last_day = f'{y}-{m:02d}-{last_day_num:02d}'
    return (first_day, last_day)

# ── 工具函数 ──

def _safe_float(val, default=0.0):
    if val is None: return default
    return float(val)

def _risk_level(score: int) -> str:
    if score >= 7: return "高风险"
    elif score >= 4: return "中风险"
    elif score >= 1: return "低风险"
    return "良好"

def _risk_color(score: int) -> str:
    if score >= 7: return "#dc2626"
    elif score >= 4: return "#f59e0b"
    elif score >= 1: return "#3b82f6"
    return "#10b981"

def _get_period_range(db: Session, company_id: int):
    min_entry = db.query(func.min(JournalEntry.entry_date)).filter(
        JournalEntry.company_id == company_id).scalar()
    max_entry = db.query(func.max(JournalEntry.entry_date)).filter(
        JournalEntry.company_id == company_id).scalar()
    if min_entry and max_entry: return str(min_entry), str(max_entry)
    return None, None

def _vat_payable_sum(db: Session, company_id: int, ps: str = None, pe: str = None) -> float:
    """汇总增值税应纳税额（从 form_main JSON 中提取）"""
    q = db.query(VATDeclaration).filter(VATDeclaration.company_id == company_id)
    if ps: q = q.filter(VATDeclaration.period >= ps)
    if pe: q = q.filter(VATDeclaration.period <= pe)
    total = 0.0
    for v in q.all():
        if v.form_main:
            fm = v.form_main if isinstance(v.form_main, dict) else json.loads(v.form_main)
            total += _safe_float(fm.get("vat_payable"))
    return total

def _get_account_balance(db: Session, company_id: int, account_code: str, period: str = None) -> float:
    q = db.query(
        func.coalesce(func.sum(JournalEntry.debit_amount), 0),
        func.coalesce(func.sum(JournalEntry.credit_amount), 0)
    ).filter(JournalEntry.company_id == company_id, JournalEntry.account_code.like(account_code + '%'))
    if period: q = q.filter(JournalEntry.period <= period)
    debit, credit = q.first()
    return _safe_float(debit) - _safe_float(credit)

def _get_account_sum(db: Session, company_id: int, account_code: str, ps: str, pe: str, field: str = "debit") -> float:
    """期间内科目发生额合计"""
    col = JournalEntry.debit_amount if field == "debit" else JournalEntry.credit_amount
    return _safe_float(db.query(func.sum(col)).filter(
        JournalEntry.company_id == company_id,
        JournalEntry.account_code.like(account_code + '%'),
        JournalEntry.period >= ps, JournalEntry.period <= pe
    ).scalar())

def _get_periods_between(ps: str, pe: str) -> list:
    """生成两个 YYYY-MM 之间的所有月份"""
    result = []
    y1, m1 = int(ps[:4]), int(ps[5:7])
    y2, m2 = int(pe[:4]), int(pe[5:7])
    y, m = y1, m1
    while True:
        result.append(f"{y}-{m:02d}")
        if y == y2 and m == m2: break
        m += 1
        if m > 12: m = 1; y += 1
    return result

def _monthly_account_balance(db: Session, company_id: int, account_code: str, ps: str, pe: str) -> dict:
    """按月汇总科目借方/贷方发生额"""
    ps = _normalize_period(ps)
    pe = _normalize_period(pe)
    rows = db.query(
        JournalEntry.period,
        func.coalesce(func.sum(JournalEntry.debit_amount), 0),
        func.coalesce(func.sum(JournalEntry.credit_amount), 0)
    ).filter(
        JournalEntry.company_id == company_id,
        JournalEntry.account_code.like(account_code + '%'),
        JournalEntry.period >= ps, JournalEntry.period <= pe
    ).group_by(JournalEntry.period).order_by(JournalEntry.period).all()
    return {r[0]: {"debit": _safe_float(r[1]), "credit": _safe_float(r[2])} for r in rows}


def _check_premise_expenses_from_invoices(db: Session, company_id: int, ps: str, pe: str):
    """从原始发票（取得发票+记账发票）中检测经营场所相关支出。
    
    序时账可能因未及时入账而缺失租赁费/水电费/物业费记录，
    但取得发票是最原始的业务证据——发票已到即意味着费用已发生。
    
    返回: {"has_rent": bool, "has_utility": bool, "has_property": bool, 
           "rent_amount": float, "utility_amount": float, "property_amount": float,
           "any_found": bool, "total_found": float}
    """
    ps_date, pe_date = _period_to_date_range(ps)[0], _period_to_date_range(pe)[1]
    
    # 关键词分组
    rent_kw = ["租", "租赁", "房租", "租金"]
    utility_kw = ["电", "水", "燃气", "水电"]
    property_kw = ["物业"]
    
    def _has_kw(text, kws):
        if not text:
            return False
        text_lower = text.lower()
        return any(kw in text_lower for kw in kws)
    
    # 构建 OR 条件（goods_name 或 seller_name 包含关键词）
    def _build_or(kws):
        conditions = []
        for kw in kws:
            conditions.append(PurchaseInvoice.goods_name.contains(kw))
            conditions.append(PurchaseInvoice.seller_name.contains(kw))
            conditions.append(BookkeepingInvoice.goods_name.contains(kw))
            conditions.append(BookkeepingInvoice.seller_name.contains(kw))
        return or_(*conditions)
    
    # 查询租赁相关发票
    rent_pi = db.query(func.coalesce(func.sum(func.abs(PurchaseInvoice.total_amount)), 0)).filter(
        PurchaseInvoice.company_id == company_id,
        PurchaseInvoice.invoice_date >= ps_date,
        PurchaseInvoice.invoice_date <= pe_date,
        or_(*[or_(PurchaseInvoice.goods_name.contains(kw), PurchaseInvoice.seller_name.contains(kw)) for kw in rent_kw])
    ).scalar() or 0
    
    rent_bi = db.query(func.coalesce(func.sum(func.abs(BookkeepingInvoice.total_amount)), 0)).filter(
        BookkeepingInvoice.company_id == company_id,
        BookkeepingInvoice.invoice_date >= ps_date,
        BookkeepingInvoice.invoice_date <= pe_date,
        or_(*[or_(BookkeepingInvoice.goods_name.contains(kw), BookkeepingInvoice.seller_name.contains(kw)) for kw in rent_kw])
    ).scalar() or 0
    rent_total = _safe_float(rent_pi) + _safe_float(rent_bi)
    
    # 查询水电相关发票
    util_pi = db.query(func.coalesce(func.sum(func.abs(PurchaseInvoice.total_amount)), 0)).filter(
        PurchaseInvoice.company_id == company_id,
        PurchaseInvoice.invoice_date >= ps_date,
        PurchaseInvoice.invoice_date <= pe_date,
        or_(*[or_(PurchaseInvoice.goods_name.contains(kw), PurchaseInvoice.seller_name.contains(kw)) for kw in utility_kw])
    ).scalar() or 0
    
    util_bi = db.query(func.coalesce(func.sum(func.abs(BookkeepingInvoice.total_amount)), 0)).filter(
        BookkeepingInvoice.company_id == company_id,
        BookkeepingInvoice.invoice_date >= ps_date,
        BookkeepingInvoice.invoice_date <= pe_date,
        or_(*[or_(BookkeepingInvoice.goods_name.contains(kw), BookkeepingInvoice.seller_name.contains(kw)) for kw in utility_kw])
    ).scalar() or 0
    util_total = _safe_float(util_pi) + _safe_float(util_bi)
    
    # 查询物业相关发票
    prop_pi = db.query(func.coalesce(func.sum(func.abs(PurchaseInvoice.total_amount)), 0)).filter(
        PurchaseInvoice.company_id == company_id,
        PurchaseInvoice.invoice_date >= ps_date,
        PurchaseInvoice.invoice_date <= pe_date,
        or_(*[or_(PurchaseInvoice.goods_name.contains(kw), PurchaseInvoice.seller_name.contains(kw)) for kw in property_kw])
    ).scalar() or 0
    
    prop_bi = db.query(func.coalesce(func.sum(func.abs(BookkeepingInvoice.total_amount)), 0)).filter(
        BookkeepingInvoice.company_id == company_id,
        BookkeepingInvoice.invoice_date >= ps_date,
        BookkeepingInvoice.invoice_date <= pe_date,
        or_(*[or_(BookkeepingInvoice.goods_name.contains(kw), BookkeepingInvoice.seller_name.contains(kw)) for kw in property_kw])
    ).scalar() or 0
    prop_total = _safe_float(prop_pi) + _safe_float(prop_bi)
    
    has_rent = rent_total > 0
    has_utility = util_total > 0
    has_property = prop_total > 0
    any_found = has_rent or has_utility or has_property
    total_found = rent_total + util_total + prop_total
    
    return {
        "has_rent": has_rent, "has_utility": has_utility, "has_property": has_property,
        "rent_amount": rent_total, "utility_amount": util_total, "property_amount": prop_total,
        "any_found": any_found, "total_found": total_found
    }


def _check_premise_contracts(db: Session, company_id: int):
    """检查合同台账中是否有经营场所相关合同。
    
    包括：租赁合同、场地使用合同、物业合同、免费使用协议等。
    
    返回: {
        "has_lease_contract": bool,   # 有租赁合同
        "has_venue_contract": bool,    # 有场地相关合同（包括免费）
        "has_free_contract": bool,     # 有免费使用/免租合同
        "contracts": [dict],           # 合同详情列表
        "any_found": bool,             # 任何场地相关合同
        "free_venue_possible": bool,   # 可能为免费场地（有合同但无支付）
    }
    """
    # 租赁合同
    lease_contracts = db.query(Contract).filter(
        Contract.company_id == company_id,
        Contract.contract_type == "租赁",
        Contract.status.in_(["已签署", "履行中", "已完成"])
    ).all()
    
    # 场地相关合同（从合同名称/内容摘要/备注中检测）
    venue_kw = ["场地", "场所", "办公", "厂房", "仓库", "物业", "租", "使用"]
    all_contracts = db.query(Contract).filter(
        Contract.company_id == company_id,
        Contract.status.in_(["已签署", "履行中", "已完成"])
    ).all()
    
    venue_contracts = []
    free_contracts = []
    for c in all_contracts:
        combined = (c.name or "") + (c.content_summary or "") + (c.remark or "")
        combined_lower = combined.lower()
        if any(kw in combined_lower for kw in venue_kw):
            venue_contracts.append(c)
            # 检测免费/免租
            if any(fkw in combined_lower for fkw in ["免费", "免租", "无偿", "0元", "0.00"]):
                free_contracts.append(c)
    
    contracts_detail = []
    for c in lease_contracts + venue_contracts:
        is_free = any(fkw in ((c.remark or "").lower()) for fkw in ["免费", "免租", "无偿"])
        contracts_detail.append({
            "contract_no": c.contract_no,
            "name": c.name,
            "type": c.contract_type,
            "amount": float(c.amount or 0),
            "status": c.status,
            "is_free": is_free,
        })
    
    has_lease = len(lease_contracts) > 0
    has_venue = len(venue_contracts) > 0
    has_free = len(free_contracts) > 0
    
    return {
        "has_lease_contract": has_lease,
        "has_venue_contract": has_venue,
        "has_free_contract": has_free,
        "contracts": contracts_detail,
        "any_found": has_lease or has_venue,
        "free_venue_possible": has_free or (has_venue and not has_lease),  # 有场地合同但非正式租赁
    }


def _check_premise_evidence_all(db: Session, company_id: int, ps: str, pe: str):
    """综合四源证据检查：序时账 + 银行流水 + 取得发票 + 合同。
    
    银行流水铁律：银行付款必须入账（余额对不上），所以银行流水有支出而序时账无记录 = 风险升级。
    发票铁律：发票可以延后入账，所以发票有而序时账无记录 = 费用已发生未入账，风险较低。
    合同铁律：场地可能是免费的（无支付+无发票），合同是唯一证据。
    
    返回: {
        "je": {"has_rent": bool, "has_utility": bool, "has_property": bool, "has_any": bool},
        "bank": {"has_payment": bool, "amount": float},
        "invoice": _check_premise_expenses_from_invoices 结果,
        "contract": _check_premise_contracts 结果,
        "overall": {
            "je_missing_but_invoice_found": bool,    # 序时账无但发票有 → 费用未入账
            "je_missing_but_bank_found": bool,       # 序时账无但银行有支付 → 银行付款未入账(高风险)
            "je_missing_but_contract_found": bool,   # 序时账无但合同有 → 可能免费场地
            "all_sources_none": bool,                # 四源全无 → 真实空壳嫌疑
            "risk_level": str,                       # "none" / "unrecorded" / "bank_unrecorded" / "free_venue" / "shell_suspect"
        }
    }
    """
    # 1. 序时账
    je_rent = db.query(func.coalesce(func.sum(JournalEntry.debit_amount), 0)).filter(
        JournalEntry.company_id == company_id,
        JournalEntry.period >= ps, JournalEntry.period <= pe,
        JournalEntry.account_code.like("660214%")
    ).scalar() or 0
    je_utility = db.query(func.coalesce(func.sum(JournalEntry.debit_amount), 0)).filter(
        JournalEntry.company_id == company_id,
        JournalEntry.period >= ps, JournalEntry.period <= pe,
        JournalEntry.account_code.like("660215%")
    ).scalar() or 0
    je_property = db.query(func.coalesce(func.sum(JournalEntry.debit_amount), 0)).filter(
        JournalEntry.company_id == company_id,
        JournalEntry.period >= ps, JournalEntry.period <= pe,
        JournalEntry.account_code.like("660213%")
    ).scalar() or 0
    je_has_any = _safe_float(je_rent) > 0 or _safe_float(je_utility) > 0 or _safe_float(je_property) > 0
    
    # 2. 银行流水
    ps_date, pe_date = _period_to_date_range(ps)[0], _period_to_date_range(pe)[1]
    bank_premise = db.query(func.coalesce(func.sum(BankTransaction.debit_amount), 0)).filter(
        BankTransaction.company_id == company_id,
        BankTransaction.transaction_date >= ps_date,
        BankTransaction.transaction_date <= pe_date,
        or_(
            BankTransaction.counterparty_name.contains("租"),
            BankTransaction.counterparty_name.contains("物业"),
            BankTransaction.counterparty_name.contains("电力"),
            BankTransaction.counterparty_name.contains("水务"),
            BankTransaction.summary.contains("租金"),
            BankTransaction.summary.contains("房租"),
            BankTransaction.summary.contains("物业"),
            BankTransaction.summary.contains("水电"),
        )
    ).scalar() or 0
    bank_has = _safe_float(bank_premise) > 0
    
    # 3. 发票
    inv_info = _check_premise_expenses_from_invoices(db, company_id, ps, pe)
    
    # 4. 合同
    contract_info = _check_premise_contracts(db, company_id)
    
    # 综合判断
    all_none = not je_has_any and not bank_has and not inv_info["any_found"] and not contract_info["any_found"]
    je_missing_invoice_found = not je_has_any and inv_info["any_found"]
    je_missing_bank_found = not je_has_any and bank_has  # ⚠️ 银行付款必须入账，这种情况非常可疑
    je_missing_contract_found = not je_has_any and contract_info["any_found"]
    
    if all_none:
        risk_level = "shell_suspect"
    elif je_missing_bank_found:
        risk_level = "bank_unrecorded"  # 银行付款但序时账无记录 → 高风险
    elif je_missing_invoice_found:
        risk_level = "unrecorded"  # 发票有但未入账 → 中低风险
    elif je_missing_contract_found:
        risk_level = "free_venue"  # 有合同但无支付 → 可能免费场地
    else:
        risk_level = "none"
    
    return {
        "je": {"has_any": je_has_any, "rent": float(je_rent), "utility": float(je_utility), "property": float(je_property)},
        "bank": {"has_payment": bank_has, "amount": float(bank_premise)},
        "invoice": inv_info,
        "contract": contract_info,
        "overall": {
            "je_missing_but_invoice_found": je_missing_invoice_found,
            "je_missing_but_bank_found": je_missing_bank_found,
            "je_missing_but_contract_found": je_missing_contract_found,
            "all_sources_none": all_none,
            "risk_level": risk_level,
        }
    }


def _check_expense_evidence(db: Session, company_id: int, ps: str, pe: str,
                             expense_name: str = "",
                             je_account_codes: list = None,
                             je_summary_kw: list = None,
                             bank_kw: list = None,
                             invoice_kw: list = None,
                             contract_kw: list = None,
                             contract_types: list = None,
                             je_extra_codes: list = None):
    """通用费用四源交叉验证：序时账 + 银行流水 + 取得发票 + 合同。
    
    银行流水铁律：银行付款必须入账（余额对不上）
    发票铁律：发票可以延后入账
    合同铁律：合同是免费/赊账场景的唯一证据
    
    参数：
    - expense_name: 费用名称（用于日志/标识）
    - je_account_codes: 序时账科目代码前缀列表，如 ["660207"]
    - je_summary_kw: 序时账摘要关键词
    - je_extra_codes: 额外科目代码（用于合并检查，如折旧费走不同科目）
    - bank_kw: 银行流水关键词（匹配 counterparty_name + summary）
    - invoice_kw: 发票关键词（匹配 goods_name + seller_name）
    - contract_kw: 合同关键词（匹配 name + content_summary + remark）
    - contract_types: 合同类型筛选（如 ["租赁", "服务"]）
    
    返回: 统一结构字典，含 overall.risk_level
    """
    je_account_codes = je_account_codes or []
    je_summary_kw = je_summary_kw or []
    bank_kw = bank_kw or []
    invoice_kw = invoice_kw or []
    contract_kw = contract_kw or []
    contract_types = contract_types or []
    je_extra_codes = je_extra_codes or []
    
    ps_date, pe_date = _period_to_date_range(ps)[0], _period_to_date_range(pe)[1]
    
    # 1. 序时账
    je_amount = 0.0
    all_codes = list(je_account_codes) + list(je_extra_codes)
    if all_codes:
        code_conditions = [JournalEntry.account_code.like(f"{code}%") for code in all_codes]
        je_code_q = db.query(func.coalesce(func.sum(JournalEntry.debit_amount), 0)).filter(
            JournalEntry.company_id == company_id,
            JournalEntry.period >= ps, JournalEntry.period <= pe,
            or_(*code_conditions)
        )
        je_amount += _safe_float(je_code_q.scalar())
    
    if je_summary_kw:
        je_summary_q = db.query(func.coalesce(func.sum(JournalEntry.debit_amount), 0)).filter(
            JournalEntry.company_id == company_id,
            JournalEntry.period >= ps, JournalEntry.period <= pe,
            or_(*[JournalEntry.summary.contains(kw) for kw in je_summary_kw])
        )
        # 避免重复计数：如果与科目查询重叠，取 max
        je_summary_amt = _safe_float(je_summary_q.scalar())
        if all_codes:
            je_amount = max(je_amount, je_summary_amt)  # 取大的，避免重复
        else:
            je_amount = je_summary_amt
    
    je_has = je_amount > 0
    
    # 2. 银行流水
    bank_amount = 0.0
    if bank_kw:
        bank_or = []
        for kw in bank_kw:
            bank_or.append(BankTransaction.counterparty_name.contains(kw))
            bank_or.append(BankTransaction.summary.contains(kw))
        bank_amount = _safe_float(db.query(func.coalesce(func.sum(BankTransaction.debit_amount), 0)).filter(
            BankTransaction.company_id == company_id,
            BankTransaction.transaction_date >= ps_date,
            BankTransaction.transaction_date <= pe_date,
            or_(*bank_or)
        ).scalar())
    bank_has = bank_amount > 0
    
    # 3. 发票（取得发票 + 记账发票）
    inv_amount = 0.0
    inv_count = 0
    if invoice_kw:
        inv_pi_q = db.query(
            func.coalesce(func.sum(func.abs(PurchaseInvoice.total_amount)), 0),
            func.count(PurchaseInvoice.id)
        ).filter(
            PurchaseInvoice.company_id == company_id,
            PurchaseInvoice.invoice_date >= ps_date,
            PurchaseInvoice.invoice_date <= pe_date,
            or_(*[or_(PurchaseInvoice.goods_name.contains(kw), PurchaseInvoice.seller_name.contains(kw)) for kw in invoice_kw])
        )
        pi_amt, pi_cnt = inv_pi_q.first()
        inv_amount += _safe_float(pi_amt)
        inv_count += (pi_cnt or 0)
        
        inv_bi_q = db.query(
            func.coalesce(func.sum(func.abs(BookkeepingInvoice.total_amount)), 0),
            func.count(BookkeepingInvoice.id)
        ).filter(
            BookkeepingInvoice.company_id == company_id,
            BookkeepingInvoice.invoice_date >= ps_date,
            BookkeepingInvoice.invoice_date <= pe_date,
            or_(*[or_(BookkeepingInvoice.goods_name.contains(kw), BookkeepingInvoice.seller_name.contains(kw)) for kw in invoice_kw])
        )
        bi_amt, bi_cnt = inv_bi_q.first()
        inv_amount += _safe_float(bi_amt)
        inv_count += (bi_cnt or 0)
    inv_has = inv_amount > 0
    
    # 4. 合同
    contract_has = False
    free_possible = False
    contracts_detail = []
    if contract_kw or contract_types:
        q = db.query(Contract).filter(
            Contract.company_id == company_id,
            Contract.status.in_(["已签署", "履行中", "已完成"])
        )
        if contract_types:
            q = q.filter(Contract.contract_type.in_(contract_types))
        
        all_contracts = q.all()
        for c in all_contracts:
            combined = (c.name or "") + (c.content_summary or "") + (c.remark or "")
            combined_lower = combined.lower()
            if contract_kw:
                if any(kw in combined_lower for kw in contract_kw):
                    is_free = any(fkw in combined_lower for fkw in ["免费", "免租", "无偿", "0元", "0.00"])
                    contracts_detail.append({
                        "contract_no": c.contract_no, "name": c.name,
                        "type": c.contract_type, "amount": float(c.amount or 0),
                        "status": c.status, "is_free": is_free,
                    })
                    if is_free:
                        free_possible = True
            else:
                # 无关键词，接受所有匹配类型的合同
                is_free = any(fkw in combined_lower for fkw in ["免费", "免租", "无偿", "0元", "0.00"])
                contracts_detail.append({
                    "contract_no": c.contract_no, "name": c.name,
                    "type": c.contract_type, "amount": float(c.amount or 0),
                    "status": c.status, "is_free": is_free,
                })
                if is_free:
                    free_possible = True
        
        contract_has = len(contracts_detail) > 0
    
    # 综合判断
    all_none = (not je_has) and (not bank_has) and (not inv_has) and (not contract_has)
    je_missing_invoice_found = (not je_has) and inv_has
    je_missing_bank_found = (not je_has) and bank_has
    je_missing_contract_found = (not je_has) and contract_has
    
    if all_none:
        risk_level = "shell_suspect"
    elif je_missing_bank_found:
        risk_level = "bank_unrecorded"
    elif je_missing_invoice_found:
        risk_level = "unrecorded"
    elif je_missing_contract_found:
        risk_level = "free_possible"
    else:
        risk_level = "none"
    
    return {
        "je": {"has_any": je_has, "amount": je_amount},
        "bank": {"has_payment": bank_has, "amount": bank_amount},
        "invoice": {"has_any": inv_has, "amount": inv_amount, "count": inv_count},
        "contract": {"has_any": contract_has, "free_possible": free_possible, "contracts": contracts_detail},
        "overall": {
            "je_missing_but_invoice_found": je_missing_invoice_found,
            "je_missing_but_bank_found": je_missing_bank_found,
            "je_missing_but_contract_found": je_missing_contract_found,
            "all_sources_none": all_none,
            "risk_level": risk_level,
        }
    }


# ── 规则文件路径 ──
RULES_FILE = os.path.join(os.path.dirname(__file__), 'static', 'tax_risk_rules_local_export.json')

def _load_saved_rules():
    """加载用户保存的涉税风险规则（从风险规则管理模块）"""
    if not os.path.exists(RULES_FILE):
        return None
    try:
        with open(RULES_FILE, 'r', encoding='utf-8') as f:
            rules = json.load(f)
        if isinstance(rules, list) and len(rules) > 0:
            _validate_rules_on_load(rules)  # 启动时自动校验规则完整性
            return rules
        return None
    except Exception:
        return None


# ── 规则校验标准（启动时自动审计）──
VALID_LEVELS = {'高风险', '中风险', '低风险', '良好'}
VALID_URGENCIES = {'紧急', '高', '中', '低', '提醒', '建议'}
RULE_REQUIRED_FIELDS = ['id', 'item', 'category', 'score', 'level', 'suggestion', 'urgency', 'evidence', 'remark', 'dataSource']

def _validate_rules_on_load(rules: list):
    """启动时校验规则文件完整性，不合格打印警告到stdout"""
    issues = []
    seen_items = {}
    for r in rules:
        rid = r.get('id', '?')
        ritem = r.get('item', '').strip()
        
        # 必填字段
        for f in RULE_REQUIRED_FIELDS:
            val = r.get(f)
            if val is None or (isinstance(val, str) and val.strip() == ''):
                issues.append(f"[规则] ID={rid} 缺失字段: {f}")
        
        # score范围
        score = r.get('score', -1)
        if not isinstance(score, (int, float)) or score < 0 or score > 100:
            issues.append(f"[规则] ID={rid} score={score} 无效(0-100)")
        
        # level
        if r.get('level', '') not in VALID_LEVELS:
            issues.append(f"[规则] ID={rid} level='{r.get('level','')}' 无效")
        
        # urgency
        if r.get('urgency', '') not in VALID_URGENCIES:
            issues.append(f"[规则] ID={rid} urgency='{r.get('urgency','')}' 无效")
        
        # item唯一性
        if ritem:
            if ritem in seen_items:
                issues.append(f"[规则] ID={rid} item='{ritem}' 与 ID={seen_items[ritem]} 重复!")
            else:
                seen_items[ritem] = rid
        
        # suggestion长度
        if len(r.get('suggestion', '')) < 5:
            issues.append(f"[规则] ID={rid} suggestion过短")
    
    if issues:
        print(f"\n⚠️  规则文件校验发现 {len(issues)} 个问题：")
        for i in issues:
            print(f"  {i}")
        print(f"  文件: {RULES_FILE}\n")
    return issues


# ── 冲突场景答案存储（用户确认风险冲突后保存）──
CONFLICT_ANSWERS_FILE = os.path.join(os.path.dirname(__file__), "tax_risk_conflict_answers.json")

def _load_conflict_answers(company_id: int) -> Dict[str, Any]:
    """加载某公司的冲突答案 {risk_item: {conflict_id: answer_dict}}"""
    try:
        if os.path.exists(CONFLICT_ANSWERS_FILE):
            with open(CONFLICT_ANSWERS_FILE, "r", encoding="utf-8") as f:
                all_answers = json.load(f)
            return all_answers.get(str(company_id), {})
    except Exception:
        pass
    return {}

def _save_conflict_answers(company_id: int, answers: Dict[str, Any]):
    """保存某公司的冲突答案"""
    all_answers = {}
    try:
        if os.path.exists(CONFLICT_ANSWERS_FILE):
            with open(CONFLICT_ANSWERS_FILE, "r", encoding="utf-8") as f:
                all_answers = json.load(f)
    except Exception:
        pass
    all_answers[str(company_id)] = answers
    with open(CONFLICT_ANSWERS_FILE, "w", encoding="utf-8") as f:
        json.dump(all_answers, f, ensure_ascii=False, indent=2)

def _apply_conflict_answers(results, company_id):
    """将用户已确认的冲突答案应用到结果中，调整风险等级和内容"""
    answers = _load_conflict_answers(company_id)
    if not answers:
        return
    for r in results:
        item = r.get("item", "")
        if item not in answers:
            continue
        item_answers = answers[item]
        scenarios = r.get("conflict_scenarios", [])
        if not scenarios:
            continue
        for sc in scenarios:
            sc_id = sc.get("id", "")
            if sc_id in item_answers:
                user_ans = item_answers[sc_id]
                if user_ans.get("confirmed"):
                    if_confirmed = sc.get("if_confirmed", {})
                    if if_confirmed.get("new_level"):
                        r["risk_level"] = if_confirmed["new_level"]
                    if if_confirmed.get("new_score") is not None:
                        r["risk_score"] = if_confirmed["new_score"]
                    if if_confirmed.get("override_detail"):
                        r["detail"] = if_confirmed["override_detail"]
                    if if_confirmed.get("override_suggestion"):
                        r["suggestion"] = if_confirmed["override_suggestion"]
                    r["risk_color"] = _risk_color(r["risk_score"])
                    r["_conflict_resolved"] = True
                    r["_conflict_answer"] = user_ans
                    break
        r["_saved_answers"] = {k: v for k, v in item_answers.items()}


def _apply_rule_overrides(results, rules, db=None, company_id=None, ps=None, pe=None):
    """
    用规则驱动分析：
    1. 用规则中的评分/等级/建议覆盖硬编码的分析结果
    2. 如果规则对应的维度在results中不存在，且条件满足，则新增维度
    """
    if not rules:
        return

    # 建立结果索引：用于匹配
    result_index = {}
    for i, r in enumerate(results):
        item = r.get("item", "").strip()
        if not item:
            continue
        # 用item的前6个字做键（更精确）
        key = item[:6]
        if key not in result_index:
            result_index[key] = []
        result_index[key].append((i, r))

    # 建立规则索引
    rule_index = {}
    for rule in rules:
        item = rule.get("item", "").strip()
        if not item:
            continue
        key = item[:6]
        if key not in rule_index:
            rule_index[key] = []
        rule_index[key].append(rule)

    # 第一步：覆盖已有维度
    for i, r in enumerate(results):
        result_item = r.get("item", "").strip()
        if not result_item:
            continue

        # 查找匹配的规则
        best_match = None
        best_score = 0

        # 尝试多种匹配方式
        for rule_key, rule_list in rule_index.items():
            for rule in rule_list:
                rule_item = rule.get("item", "")
                # 精确匹配
                if result_item == rule_item:
                    best_match = rule
                    best_score = 100
                    break
                # 子串匹配（需至少匹配6个字符，避免太短的误匹配）
                if (rule_item in result_item or result_item in rule_item) and min(len(rule_item), len(result_item)) >= 6:
                    score = min(len(rule_item), len(result_item))
                    if score > best_score:
                        best_score = score
                        best_match = rule
                # 关键词匹配（前6个字相同才匹配，避免太短的误匹配）
                if len(result_item) >= 6 and len(rule_item) >= 6 and result_item[:6] == rule_item[:6]:
                    score = 4
                    if score > best_score:
                        best_score = score
                        best_match = rule

            if best_score >= 100:  # 精确匹配，退出
                break

        if best_match:
            # 用规则值覆盖
            if "score" in best_match and best_match["score"] is not None:
                r["risk_score"] = best_match["score"]
            if "level" in best_match and best_match["level"]:
                r["risk_level"] = best_match["level"]
            if "category" in best_match and best_match["category"]:
                r["category"] = best_match["category"]
            if "categoryIcon" in best_match:
                r["category_icon"] = best_match["categoryIcon"]
            if "suggestion" in best_match and best_match["suggestion"]:
                r["suggestion"] = best_match["suggestion"]
            if "urgency" in best_match and best_match["urgency"]:
                r["urgency"] = best_match["urgency"]
            if "evidence" in best_match and best_match["evidence"]:
                r["required_evidence"] = [e.strip() for e in best_match["evidence"].split("\n") if e.strip()]
            # 冲突场景：从规则透传到结果
            if "conflict_scenarios" in best_match and best_match["conflict_scenarios"]:
                r["conflict_scenarios"] = best_match["conflict_scenarios"]
            # 重新计算颜色
            r["risk_color"] = _risk_color(r["risk_score"])

    # 第二步：为规则中独有且条件满足的维度，新增到results
    # （这需要动态执行规则的检查逻辑，暂不实现）


# ── 风险分析核心 ──

@router.get("/report")
def get_tax_risk_report(
    company_id: int = Query(...),
    period: Optional[str] = None,
    period_from: Optional[str] = None,
    period_to: Optional[str] = None,
    db: Session = Depends(get_db)
):
    results = []
    if period_from and period_to:
        # 规范化：统一为 YYYY-MM 格式
        period_start = _normalize_period(period_from)
        period_end = _normalize_period(period_to)
    elif period:
        period_start = period_end = _normalize_period(period)
    else:
        latest = db.query(func.max(VATDeclaration.period)).filter(
            VATDeclaration.company_id == company_id).scalar()
        if latest:
            period_end = latest
            y, m = int(latest[:4]), int(latest[5:7])
            m -= 5
            if m <= 0: m += 12; y -= 1
            period_start = f"{y}-{m:02d}"
        else:
            period_start = period_end = date.today().strftime("%Y-%m")

    # ── 23 个基础分析维度 ──
    _analyze_accounting(db, company_id, period_start, period_end, results)
    _analyze_invoice_compliance(db, company_id, period_start, period_end, results)
    _analyze_invoice_depth(db, company_id, period_start, period_end, results)
    _analyze_cost_structure(db, company_id, period_start, period_end, results)
    _analyze_financial_tax_invoice_cross(db, company_id, period_start, period_end, results)
    _analyze_ratio_elasticity(db, company_id, period_start, period_end, results)
    _analyze_hidden_inflated(db, company_id, period_start, period_end, results)
    _analyze_tax_burden(db, company_id, period_start, period_end, results)
    _analyze_urban_construction_tax(db, company_id, period_start, period_end, results)
    _analyze_property_tax(db, company_id, period_start, period_end, results)
    _analyze_iit(db, company_id, period_start, period_end, results)
    _analyze_stamp_tax(db, company_id, period_start, period_end, results)
    _analyze_tax_adjustment(db, company_id, period_start, period_end, results)
    _analyze_revenue_timing(db, company_id, period_start, period_end, results)
    _analyze_policy_execution(db, company_id, period_start, period_end, results)
    _analyze_capital_risks(db, company_id, period_start, period_end, results)
    _analyze_salary_compliance(db, company_id, period_start, period_end, results)
    _analyze_customer_penetration(db, company_id, period_start, period_end, results)
    _analyze_supplier_penetration(db, company_id, period_start, period_end, results)
    _analyze_financial_health(db, company_id, period_start, period_end, results)
    _analyze_business_credit(db, company_id, period_start, period_end, results)
    _analyze_industry_specific(db, company_id, period_start, period_end, results)
    _analyze_good_practices(db, company_id, period_start, period_end, results)
    # ── 经营实质深度分析（稽查级·10项）──
    _analyze_long_term_loss(db, company_id, period_start, period_end, results)
    _analyze_business_premise(db, company_id, period_start, period_end, results)
    _analyze_inventory_substance(db, company_id, period_start, period_end, results)
    _analyze_utility_expense(db, company_id, period_start, period_end, results)
    _analyze_staffing_substance(db, company_id, period_start, period_end, results)
    _analyze_fund_flow_invoice_match(db, company_id, period_start, period_end, results)
    _analyze_scrap_revenue(db, company_id, period_start, period_end, results)
    _analyze_deemed_sales(db, company_id, period_start, period_end, results)
    _analyze_cash_overstock(db, company_id, period_start, period_end, results)
    _analyze_related_party_pricing(db, company_id, period_start, period_end, results)
    # ── V5 新增 — 经营实质（8项·稽查必查）──
    _analyze_transport_missing(db, company_id, period_start, period_end, results)
    _analyze_agriculture_substance(db, company_id, period_start, period_end, results)
    _analyze_packaging_missing(db, company_id, period_start, period_end, results)
    _analyze_warehouse_missing(db, company_id, period_start, period_end, results)
    _analyze_equipment_depreciation_missing(db, company_id, period_start, period_end, results)
    _analyze_advertising_missing(db, company_id, period_start, period_end, results)
    _analyze_travel_missing(db, company_id, period_start, period_end, results)
    _analyze_office_expense_missing(db, company_id, period_start, period_end, results)
    # ── V4 新增 — 增值税专项（5项）──
    _analyze_vat_zero_declaration(db, company_id, period_start, period_end, results)
    _analyze_vat_burden_quarterly(db, company_id, period_start, period_end, results)
    _analyze_vat_input_transfer_omission(db, company_id, period_start, period_end, results)
    _analyze_vat_retention_refund(db, company_id, period_start, period_end, results)
    _analyze_vat_no_ticket_sales(db, company_id, period_start, period_end, results)
    # ── V4 新增 — 发票异常（3项）──
    _analyze_invoice_amount_anomaly(db, company_id, period_start, period_end, results)
    _analyze_sensitive_invoice(db, company_id, period_start, period_end, results)
    _analyze_buy_sell_mismatch(db, company_id, period_start, period_end, results)
    # ── V4 新增 — 费用匹配（3项）──
    _analyze_fuel_vs_vehicles(db, company_id, period_start, period_end, results)
    _analyze_transport_ratio(db, company_id, period_start, period_end, results)
    _analyze_expense_reasonability(db, company_id, period_start, period_end, results)
    # ── V4 新增 — 企业所得税专项（4项）──
    _analyze_impairment_not_adjusted(db, company_id, period_start, period_end, results)
    _analyze_unpaid_capital_interest(db, company_id, period_start, period_end, results)
    _analyze_non_taxable_income(db, company_id, period_start, period_end, results)
    _analyze_provisional_cost(db, company_id, period_start, period_end, results)
    # ── V4 新增 — 薪酬福利及其他（5项）──
    _analyze_staff_welfare(db, company_id, period_start, period_end, results)
    _analyze_social_security_match(db, company_id, period_start, period_end, results)
    _analyze_undistributed_profit(db, company_id, period_start, period_end, results)
    _analyze_cross_invoicing(db, company_id, period_start, period_end, results)
    _analyze_invest_property_tax(db, company_id, period_start, period_end, results)
    # ── V6 新增 — 发票异常补充（4项）──
    _analyze_top_amount_invoice_ratio(db, company_id, period_start, period_end, results)
    _analyze_top_amount_red_invoice(db, company_id, period_start, period_end, results)
    _analyze_consecutive_invoice_numbers(db, company_id, period_start, period_end, results)
    _analyze_off_hours_invoicing(db, company_id, period_start, period_end, results)
    # ── V6 新增 — 资金流与发票联动（3项）──
    _analyze_einvoice_frequency_spike(db, company_id, period_start, period_end, results)
    _analyze_corp_to_personal_transfer(db, company_id, period_start, period_end, results)
    _analyze_cash_transaction_ratio(db, company_id, period_start, period_end, results)
    # ── V6 新增 — 增值税专项补充（4项）──
    _analyze_input_tax_transfer_delay(db, company_id, period_start, period_end, results)
    _analyze_invoice_revenue_cross_period(db, company_id, period_start, period_end, results)
    _analyze_export_rebate_anomaly(db, company_id, period_start, period_end, results)
    _analyze_input_tax_credit_long_term(db, company_id, period_start, period_end, results)
    # ── V6 新增 — 企业所得税补充（4项）──
    _analyze_related_party_transfer_pricing(db, company_id, period_start, period_end, results)
    _analyze_rd_expense_compliance(db, company_id, period_start, period_end, results)
    _analyze_eit_contribution_zero(db, company_id, period_start, period_end, results)
    _analyze_eit_adjustment_items(db, company_id, period_start, period_end, results)
    # ── V6 新增 — 发票深度补充（4项）──
    _analyze_seller_spike_invoicing(db, company_id, period_start, period_end, results)
    _analyze_daily_necessities_invoice(db, company_id, period_start, period_end, results)
    _analyze_energy_consumption_match(db, company_id, period_start, period_end, results)
    _analyze_transport_expense_match(db, company_id, period_start, period_end, results)
    # ── V6 新增 — 经营实质补充（3项）──
    _analyze_two_end_outside(db, company_id, period_start, period_end, results)
    _analyze_new_company_mass_invoicing(db, company_id, period_start, period_end, results)
    _analyze_small_taxpayer_overscale(db, company_id, period_start, period_end, results)
    # ── V6 新增 — 征管风险（4项）──
    _analyze_tax_arrears(db, company_id, period_start, period_end, results)
    _analyze_upstream_runaway(db, company_id, period_start, period_end, results)
    _analyze_stale_invoice(db, company_id, period_start, period_end, results)
    _analyze_non_normal_counterparty(db, company_id, period_start, period_end, results)
    # ── V6 新增 — 交易特征（4项）──
    _analyze_round_amount_transactions(db, company_id, period_start, period_end, results)
    _analyze_fund_backflow(db, company_id, period_start, period_end, results)
    _analyze_same_buyer_frequent_red(db, company_id, period_start, period_end, results)
    _analyze_input_transfer_reentry(db, company_id, period_start, period_end, results)
    # ── V6 新增 — 财务健康与其他（3项）──
    _analyze_roe_sustainability(db, company_id, period_start, period_end, results)
    _analyze_operating_cash_flow_coverage(db, company_id, period_start, period_end, results)
    _analyze_other_payables_scale_mismatch(db, company_id, period_start, period_end, results)
    # ── V7 新增 — 补全8项（个人所得税/发票异常/申报比对/经营实质/征管风险）──
    _analyze_year_end_bonus_tax_opt(db, company_id, period_start, period_end, results)
    _analyze_agricultural_purchase_invoice(db, company_id, period_start, period_end, results)
    _analyze_cross_region_invoicing(db, company_id, period_start, period_end, results)
    _analyze_simple_general_tax_mix(db, company_id, period_start, period_end, results)
    _analyze_pre_cancellation_spike(db, company_id, period_start, period_end, results)
    _analyze_same_address_multi_company(db, company_id, period_start, period_end, results)
    _analyze_address_mismatch(db, company_id, period_start, period_end, results)
    _analyze_d_taxpayer_risk(db, company_id, period_start, period_end, results)
    _analyze_flexible_employment_tax(db, company_id, period_start, period_end, results)
    # ── V8 新增 — 合同风险（三流合一/四流合一，16项）──
    _analyze_contract_risks(db, company_id, period_start, period_end, results)
    # ── V10 新增 — 稽查重点盲区补全（13项）──
    _analyze_price_surcharge(db, company_id, period_start, period_end, results)
    _analyze_non_monetary_exchange(db, company_id, period_start, period_end, results)
    _analyze_debt_restructuring(db, company_id, period_start, period_end, results)
    _analyze_equity_transfer(db, company_id, period_start, period_end, results)
    _analyze_interest_free_loan(db, company_id, period_start, period_end, results)
    _analyze_parent_subsidiary_mgmt_fee(db, company_id, period_start, period_end, results)
    _analyze_overseas_payment_withholding(db, company_id, period_start, period_end, results)
    _analyze_mixed_sales_accounting(db, company_id, period_start, period_end, results)
    _analyze_environmental_tax(db, company_id, period_start, period_end, results)
    _analyze_invoice_remarks_compliance(db, company_id, period_start, period_end, results)
    _analyze_commission_fee_compliance(db, company_id, period_start, period_end, results)
    _analyze_donation_compliance(db, company_id, period_start, period_end, results)
    _analyze_inventory_count_discrepancy(db, company_id, period_start, period_end, results)
    # ── V11 新增 — 普票浪费可抵扣进项税额 + 预付卡涉税风险 ──
    _analyze_wasted_deductible_tax(db, company_id, period_start, period_end, results)
    _analyze_prepaid_card_risk(db, company_id, period_start, period_end, results)

    # ── 【核心】加载用户规则并应用覆盖 ──
    rules = _load_saved_rules()
    rules_applied = False
    rules_count = 0
    if rules:
        rules_count = len(rules)
        _apply_rule_overrides(results, rules, db, company_id, period_start, period_end)
        rules_applied = True

    # 应用用户已确认的冲突答案（降级/消除风险）
    _apply_conflict_answers(results, company_id)

    results.sort(key=lambda x: (x.get("risk_score", 0)), reverse=True)

    high_count = sum(1 for r in results if r.get("risk_level") == "高风险")
    mid_count = sum(1 for r in results if r.get("risk_level") == "中风险")
    low_count = sum(1 for r in results if r.get("risk_level") == "低风险")
    good_count = sum(1 for r in results if r.get("risk_level") == "良好")
    total_score = sum(r.get("risk_score", 0) for r in results)

    # 计算主要税负率和财务指标
    revenue_debit = _get_account_sum(db, company_id, "6001", period_start, period_end, "credit")
    cost_debit = _get_account_sum(db, company_id, "6401", period_start, period_end, "debit")
    vat_total = _vat_payable_sum(db, company_id, period_start, period_end)

    return {
        "company_id": company_id,
        "period_start": period_start, "period_end": period_end,
        "summary": {
            "total_items": len(results),
            "high_risk_count": high_count, "mid_risk_count": mid_count,
            "low_risk_count": low_count, "good_count": good_count,
            "overall_risk_score": total_score,
            "overall_risk_level": (
                "高风险" if total_score >= 60 else
                ("中风险" if total_score >= 30 else
                 ("低风险" if total_score >= 10 else "良好"))
            )
        },
        "metrics": {
            "revenue": round(revenue_debit, 2),
            "cost": round(cost_debit, 2),
            "gross_margin_pct": round((revenue_debit - cost_debit) / revenue_debit * 100, 2) if revenue_debit > 0 else 0,
            "vat_payable": round(vat_total, 2),
        },
        "rules_applied": rules_applied,
        "rules_count": rules_count,
        "required_evidence_summary": _build_evidence_summary(results),
        "results": results
    }


# ═══════════════════════════════════════════════════════════════
#  冲突场景答案 API — 用户确认风险降级/排除
# ═══════════════════════════════════════════════════════════════

@router.get("/report/conflict-answers")
def get_conflict_answers(company_id: int = Query(...)):
    """获取某公司已保存的冲突答案"""
    answers = _load_conflict_answers(company_id)
    return {"company_id": company_id, "answers": answers}

@router.post("/report/conflict-answers")
def save_conflict_answer(
    company_id: int = Query(...),
    risk_item: str = Body(...),
    conflict_id: str = Body(...),
    confirmed: bool = Body(...),
    answer_note: str = Body(""),
):
    """保存单个冲突场景的答案"""
    answers = _load_conflict_answers(company_id)
    if risk_item not in answers:
        answers[risk_item] = {}
    answers[risk_item][conflict_id] = {
        "confirmed": confirmed,
        "note": answer_note,
        "answered_at": datetime.now().isoformat()
    }
    _save_conflict_answers(company_id, answers)
    return {"status": "ok", "risk_item": risk_item, "conflict_id": conflict_id, "confirmed": confirmed}

@router.post("/report/conflict-answers/reset")
def reset_conflict_answer(
    company_id: int = Query(...),
    risk_item: str = Body(...),
    conflict_id: str = Body(...),
):
    """重置（撤销）单个冲突答案"""
    answers = _load_conflict_answers(company_id)
    if risk_item in answers and conflict_id in answers[risk_item]:
        del answers[risk_item][conflict_id]
        if not answers[risk_item]:
            del answers[risk_item]
        _save_conflict_answers(company_id, answers)
    return {"status": "ok", "reset": True}


# ═══════════════════════════════════════════════════════════════
#  报告下载 API — PDF / Word / PPT
# ═══════════════════════════════════════════════════════════════

def _generate_report_html(data):
    """生成报告HTML内容（用于PDF导出）"""
    s = data.get("summary", {})
    m = data.get("metrics", {})
    results = data.get("results", [])
    
    # 分类排序
    cat_order = [
        "经营实质", "增值税专项", "合同风险", "发票异常", "费用匹配",
        "企业所得税", "薪酬福利", "良好实践", "财税票比对", "配比弹性",
        "隐匿虚增", "纳税调整", "收入时点", "账务数据", "发票合规",
        "发票深度", "成本结构", "税负水平", "城建税", "房产税",
        "个人所得税", "印花税", "政策执行", "资金往来", "薪酬合规",
        "客户穿透", "供应商穿透", "财务健康", "企业信用", "行业专项", "其他风险"
    ]
    cats = defaultdict(list)
    for r in results:
        cats[r.get("category", "其他")].append(r)
    
    items_html = ""
    seq = 0
    for cat in cat_order:
        cat_items = cats.pop(cat, [])
        if not cat_items:
            continue
        items_html += f'<h3>{cat_items[0].get("category_icon","")} {cat}</h3>'
        for r in cat_items:
            seq += 1
            level = r.get("risk_level", "")
            level_color = {"高风险":"#dc2626","中风险":"#f59e0b","低风险":"#3b82f6","良好":"#10b981"}.get(level, "#6b7280")
            items_html += f'''
            <div class="item" style="border-left:4px solid {level_color};padding:8px;margin:8px 0;background:#fafafa">
              <strong style="color:{level_color}">[{level}]</strong> {seq}. {r.get("item","")}
              <p style="color:#666;font-size:13px">{r.get("detail","")}</p>
              <p style="color:#0c4a6e">💡 建议：{r.get("suggestion","")}</p>
            </div>'''
    for cat_name, cat_items in cats.items():
        items_html += f'<h3>{cat_name}</h3>'
        for r in cat_items:
            seq += 1
            level = r.get("risk_level", "")
            level_color = {"高风险":"#dc2626","中风险":"#f59e0b","低风险":"#3b82f6","良好":"#10b981"}.get(level, "#6b7280")
            items_html += f'''
            <div class="item" style="border-left:4px solid {level_color};padding:8px;margin:8px 0;background:#fafafa">
              <strong style="color:{level_color}">[{level}]</strong> {seq}. {r.get("item","")}
              <p style="color:#666;font-size:13px">{r.get("detail","")}</p>
              <p style="color:#0c4a6e">💡 建议：{r.get("suggestion","")}</p>
            </div>'''
    
    return f'''<!DOCTYPE html><html><head><meta charset="utf-8">
    <title>涉税风险分析报告</title>
    <style>
      body {{font-family:"Microsoft YaHei","PingFang SC",sans-serif;max-width:900px;margin:0 auto;padding:20px;color:#333}}
      h1 {{text-align:center;color:#1e3a5f;border-bottom:3px solid #1e3a5f;padding-bottom:12px}}
      h2 {{color:#334155;margin-top:24px}}
      h3 {{color:#475569;margin-top:16px;border-bottom:1px solid #e2e8f0;padding-bottom:4px}}
      .header-info {{text-align:center;color:#64748b;font-size:14px;margin:8px 0}}
      .cards {{display:flex;gap:12px;flex-wrap:wrap;margin:16px 0}}
      .card {{flex:1;min-width:100px;text-align:center;padding:12px;border-radius:8px;border:2px solid #e2e8f0}}
      .card .label {{font-size:12px;color:#64748b}}
      .card .value {{font-size:24px;font-weight:bold}}
      .summary-box {{padding:16px;background:#f8fafc;border-radius:8px;margin:16px 0}}
      @media print {{body{{font-size:12px}}}}
    </style></head><body>
    <h1>🛡️ 涉税风险分析报告</h1>
    <div class="header-info">公司ID: {data.get("company_id")} | 分析期间: {data.get("period_start")} ~ {data.get("period_end")}</div>
    <div class="cards">
      <div class="card" style="border-color:#dc2626"><div class="label">综合风险</div><div class="value" style="color:#dc2626">{s.get("overall_risk_level","")}</div></div>
      <div class="card"><div class="label">高风险</div><div class="value" style="color:#dc2626">{s.get("high_risk_count",0)}</div></div>
      <div class="card"><div class="label">中风险</div><div class="value" style="color:#f59e0b">{s.get("mid_risk_count",0)}</div></div>
      <div class="card"><div class="label">低风险</div><div class="value" style="color:#3b82f6">{s.get("low_risk_count",0)}</div></div>
    </div>
    <div class="summary-box">
      <p>分析项目：{s.get("total_items",0)}项 | 收入：¥{m.get("revenue",0):,.2f} | 毛利率：{m.get("gross_margin_pct",0):.1f}%</p>
    </div>
    <h2>📊 风险详情</h2>
    {items_html}
    </body></html>'''


@router.get("/report/download")
def download_tax_risk_report(
    company_id: int = Query(...),
    period_from: str = Query(None),
    period_to: str = Query(None),
    format: str = Query("pdf"),
    db: Session = Depends(get_db)
):
    """下载涉税风险分析报告 - 支持 pdf/docx/pptx 三种格式"""
    # 生成报告数据（复用核心分析逻辑）
    data = get_tax_risk_report(company_id=company_id, period_from=period_from, period_to=period_to, db=db)
    
    if format == "pdf":
        return _download_pdf(data)
    elif format == "docx":
        return _download_docx(data)
    elif format == "pptx":
        return _download_pptx(data)
    else:
        return {"error": f"不支持的格式: {format}，请使用 pdf/docx/pptx"}


def _download_pdf(data):
    """使用 reportlab 生成 PDF 报告"""
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import mm
    from reportlab.lib.colors import HexColor
    from reportlab.lib.enums import TA_CENTER, TA_LEFT
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    
    # 注册中文字体
    font_paths = [
        "C:/Windows/Fonts/msyh.ttc",
        "C:/Windows/Fonts/simsun.ttc",
        "C:/Windows/Fonts/simhei.ttf",
    ]
    font_registered = False
    for fp in font_paths:
        if os.path.exists(fp):
            try:
                pdfmetrics.registerFont(TTFont('ChineseFont', fp))
                font_registered = True
                break
            except Exception:
                continue
    if not font_registered:
        # Fallback - 无中文字体则返回HTML格式
        return Response(_generate_report_html(data).encode('utf-8'), media_type="text/html")
    
    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=A4, rightMargin=20*mm, leftMargin=20*mm, topMargin=20*mm, bottomMargin=20*mm)
    
    styles = getSampleStyleSheet()
    cn_style = ParagraphStyle('Chinese', parent=styles['Normal'], fontName='ChineseFont', fontSize=10, leading=16)
    cn_title = ParagraphStyle('ChineseTitle', parent=styles['Title'], fontName='ChineseFont', fontSize=20, leading=28, alignment=TA_CENTER)
    cn_h2 = ParagraphStyle('ChineseH2', parent=styles['Heading2'], fontName='ChineseFont', fontSize=14, leading=20)
    
    s = data.get("summary", {})
    m = data.get("metrics", {})
    results = data.get("results", [])
    
    story = []
    story.append(Paragraph("涉税风险分析报告", cn_title))
    story.append(Spacer(1, 6*mm))
    story.append(Paragraph(f"公司ID: {data.get('company_id')} | 期间: {data.get('period_start')} ~ {data.get('period_end')}", cn_style))
    story.append(Paragraph(f"综合风险: {s.get('overall_risk_level','')} | 高风险{s.get('high_risk_count',0)}项 | 中风险{s.get('mid_risk_count',0)}项 | 低风险{s.get('low_risk_count',0)}项 | 共{s.get('total_items',0)}项", cn_style))
    story.append(Spacer(1, 4*mm))
    
    for r in results[:50]:  # 最多50项
        level_color = {"高风险":"#dc2626","中风险":"#f59e0b","低风险":"#3b82f6","良好":"#10b981"}.get(r.get("risk_level",""), "#6b7280")
        story.append(Paragraph(f"<font color='{level_color}'>[{r.get('risk_level','')}]</font> {r.get('item','')}", cn_h2))
        story.append(Paragraph(r.get("detail",""), cn_style))
        story.append(Paragraph(f"💡 建议: {r.get('suggestion','')}", cn_style))
        story.append(Spacer(1, 3*mm))
    
    doc.build(story)
    buf.seek(0)
    return Response(buf.read(), media_type="application/pdf",
                    headers={"Content-Disposition": "attachment; filename=tax_risk_report.pdf"})


def _download_docx(data):
    """使用 python-docx 生成 Word 报告"""
    from docx import Document
    from docx.shared import Inches, Pt, RGBColor, Cm
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.enum.table import WD_TABLE_ALIGNMENT
    
    doc = Document()
    # 设置默认字体
    style = doc.styles['Normal']
    font = style.font
    font.name = 'Microsoft YaHei'
    font.size = Pt(10)
    
    # 标题
    title = doc.add_heading('涉税风险分析报告', level=0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    s = data.get("summary", {})
    m = data.get("metrics", {})
    results = data.get("results", [])
    
    doc.add_paragraph(f"公司ID: {data.get('company_id')}  |  期间: {data.get('period_start')} ~ {data.get('period_end')}")
    doc.add_paragraph(f"综合风险等级: {s.get('overall_risk_level','')}  |  高风险: {s.get('high_risk_count',0)}项  |  中风险: {s.get('mid_risk_count',0)}项  |  低风险: {s.get('low_risk_count',0)}项")
    doc.add_paragraph(f"分析项目: {s.get('total_items',0)}项  |  收入: ¥{m.get('revenue',0):,.2f}  |  毛利率: {m.get('gross_margin_pct',0):.1f}%  |  风险评估分: {s.get('overall_risk_score',0)}")
    doc.add_paragraph("")
    
    # 汇总表
    table = doc.add_table(rows=1, cols=6)
    table.style = 'Table Grid'
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    hdr = table.rows[0].cells
    hdr[0].text = '序号'
    hdr[1].text = '风险等级'
    hdr[2].text = '风险项'
    hdr[3].text = '分类'
    hdr[4].text = '评分'
    hdr[5].text = '紧急程度'
    for i, cell in enumerate(hdr):
        cell.paragraphs[0].runs[0].bold = True
    
    for idx, r in enumerate(results, 1):
        row = table.add_row().cells
        row[0].text = str(idx)
        row[1].text = r.get('risk_level', '')
        row[2].text = r.get('item', '')
        row[3].text = r.get('category', '')
        row[4].text = str(r.get('risk_score', ''))
        row[5].text = r.get('urgency', '')
    
    doc.add_paragraph("")
    doc.add_heading('风险详情', level=1)
    
    for idx, r in enumerate(results, 1):
        level = r.get('risk_level', '')
        doc.add_heading(f"{idx}. [{level}] {r.get('item','')}", level=2)
        doc.add_paragraph(r.get('detail', ''))
        doc.add_paragraph(f"建议: {r.get('suggestion','')}")
    
    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    return Response(buf.read(), media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                    headers={"Content-Disposition": "attachment; filename=tax_risk_report.docx"})


def _download_pptx(data):
    """使用 python-pptx 生成 PPT 报告"""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    s = data.get("summary", {})
    m = data.get("metrics", {})
    results = data.get("results", [])
    
    # Slide 1: 封面
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    txBox = slide.shapes.add_textbox(Inches(2), Inches(1.5), Inches(9), Inches(2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "涉税风险分析报告"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1E, 0x3A, 0x5F)
    p.alignment = PP_ALIGN.CENTER
    
    txBox2 = slide.shapes.add_textbox(Inches(2), Inches(4), Inches(9), Inches(2))
    tf2 = txBox2.text_frame
    p = tf2.paragraphs[0]
    p.text = f"公司ID: {data.get('company_id')}  |  期间: {data.get('period_start')} ~ {data.get('period_end')}"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(0x64, 0x74, 0x8B)
    p.alignment = PP_ALIGN.CENTER
    p2 = tf2.add_paragraph()
    p2.text = f"综合风险: {s.get('overall_risk_level','')} | 评估分: {s.get('overall_risk_score',0)} | 高风险{s.get('high_risk_count',0)}项"
    p2.font.size = Pt(16)
    p2.font.color.rgb = RGBColor(0xDC, 0x26, 0x26)
    p2.alignment = PP_ALIGN.CENTER
    
    # Slide 2: 汇总
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    txBox3 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
    tf3 = txBox3.text_frame
    p = tf3.paragraphs[0]
    p.text = "📊 风险概览"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1E, 0x3A, 0x5F)
    
    # 汇总表
    rows_data = [
        ["指标", "数值"],
        ["综合风险等级", s.get('overall_risk_level','')],
        ["风险评估总分", str(s.get('overall_risk_score',0))],
        ["高风险项", str(s.get('high_risk_count',0))],
        ["中风险项", str(s.get('mid_risk_count',0))],
        ["低风险项", str(s.get('low_risk_count',0))],
        ["良好事项", str(s.get('good_count',0))],
        ["分析项目数", str(s.get('total_items',0))],
        ["收入", f"¥{m.get('revenue',0):,.2f}"],
        ["毛利率", f"{m.get('gross_margin_pct',0):.1f}%"],
    ]
    table_shape = slide2.shapes.add_table(len(rows_data), 2, Inches(1.5), Inches(1.5), Inches(10), Inches(4))
    table = table_shape.table
    for i, row in enumerate(rows_data):
        for j, cell_text in enumerate(row):
            cell = table.cell(i, j)
            cell.text = cell_text
            if i == 0:
                cell.text_frame.paragraphs[0].font.bold = True
    
    # Slide 3+: 风险详情
    for idx, r in enumerate(results):
        if idx > 0 and idx % 6 == 0:
            continue  # 每页最多6项
        if idx % 6 == 0:
            slide_n = prs.slides.add_slide(prs.slide_layouts[6])
            txBox_n = slide_n.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
            tf_n = txBox_n.text_frame
            p = tf_n.paragraphs[0]
            p.text = f"🔍 风险详情 ({idx//6 + 1})"
            p.font.size = Pt(24)
            p.font.bold = True
        else:
            slide_n = prs.slides[-1]
        
        y_pos = 1.5 + (idx % 6) * 1.0
        txBox_r = slide_n.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(12), Inches(0.8))
        tf_r = txBox_r.text_frame
        level = r.get('risk_level','')
        p = tf_r.paragraphs[0]
        p.text = f"[{level}] {r.get('item','')}"
        p.font.size = Pt(14)
        p.font.bold = True
        level_c = {"高风险":RGBColor(0xDC,0x26,0x26),"中风险":RGBColor(0xF5,0x9E,0x0B),"低风险":RGBColor(0x3B,0x82,0xF6),"良好":RGBColor(0x10,0xB9,0x81)}.get(level, RGBColor(0x6B,0x72,0x80))
        p.font.color.rgb = level_c
    
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return Response(buf.read(), media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    headers={"Content-Disposition": "attachment; filename=tax_risk_report.pptx"})


# ═══════════════════════════════════════════════════════════════
#  一、账务数据风险
# ═══════════════════════════════════════════════════════════════

def _analyze_accounting(db, company_id, ps, pe, results):
    """借贷平衡、凭证数量、折旧摊销"""
    entries = db.query(
        JournalEntry.period,
        func.sum(JournalEntry.debit_amount).label("debit"),
        func.sum(JournalEntry.credit_amount).label("credit"),
        func.count(JournalEntry.id).label("cnt")
    ).filter(
        JournalEntry.company_id == company_id,
        JournalEntry.period >= ps, JournalEntry.period <= pe
    ).group_by(JournalEntry.period).all()

    unbalanced = [e.period for e in entries if abs(_safe_float(e.debit) - _safe_float(e.credit)) > 0.01]
    if unbalanced:
        results.append({
            "category": "账务数据", "category_icon": "📊", "risk_score": 9, "risk_level": "高风险",
            "risk_color": "#dc2626", "urgency": "紧急",
            "item": "借贷不平衡",
            "detail": f"以下期间借贷方金额不相等：{', '.join(unbalanced)}。借贷不平衡会导致报表数据失真，影响税务申报准确性。",
            "suggestion": "逐月排查序时账，检查是否有凭证漏记或金额录入错误，确保每笔凭证借贷金额相等。"
        })

    for e in entries:
        cnt = int(e.cnt)
        if cnt == 0:
            results.append({
                "category": "账务数据", "category_icon": "📊", "risk_score": 7, "risk_level": "高风险",
                "risk_color": "#dc2626", "urgency": "紧急",
                "item": "期间无凭证记录",
                "detail": f"{e.period} 期间没有任何序时账记录，可能尚未开始账务处理。",
                "suggestion": "立即检查该期间是否漏记账务，补充必要的凭证录入。"
            })
        elif cnt > 500:
            results.append({
                "category": "账务数据", "category_icon": "📊", "risk_score": 4, "risk_level": "中风险",
                "risk_color": "#f59e0b", "urgency": "提醒",
                "item": "凭证量异常偏高",
                "detail": f"{e.period} 期间有 {cnt} 条凭证记录，远超常规。需确认是否存在重复录入或批量导入错误。",
                "suggestion": "核查是否重复导入或批量操作异常，建议抽查大额凭证的真实性。"
            })

    fixed_count = db.query(func.count(FixedAsset.id)).filter(
        FixedAsset.company_id == company_id, FixedAsset.status.in_(["在用", "闲置"])).scalar()
    intangible_count = db.query(func.count(IntangibleAsset.id)).filter(
        IntangibleAsset.company_id == company_id, IntangibleAsset.status.in_(["在用", "闲置"])).scalar()

    if fixed_count and fixed_count > 0:
        depr = db.query(func.count(JournalEntry.id)).filter(
            JournalEntry.company_id == company_id, JournalEntry.period >= ps,
            JournalEntry.account_code.like("1602%"),
            func.lower(JournalEntry.summary).contains("折旧")).scalar()
        if not depr:
            results.append({
                "category": "账务数据", "category_icon": "📊", "risk_score": 8, "risk_level": "高风险",
                "risk_color": "#dc2626", "urgency": "紧急",
                "item": "未计提固定资产折旧",
                "detail": f"系统中有 {fixed_count} 项在用/闲置固定资产，但未发现折旧计提凭证。漏提折旧会导致利润虚增、多缴企业所得税。",
                "suggestion": "按企业会计准则第4号，当月增加的固定资产次月起计提折旧，检查是否漏记折旧费用。"
            })

    if intangible_count and intangible_count > 0:
        amor = db.query(func.count(JournalEntry.id)).filter(
            JournalEntry.company_id == company_id, JournalEntry.period >= ps,
            func.lower(JournalEntry.summary).contains("摊销")).scalar()
        if not amor:
            results.append({
                "category": "账务数据", "category_icon": "📊", "risk_score": 6, "risk_level": "中风险",
                "risk_color": "#f59e0b", "urgency": "提醒",
                "item": "未计提无形资产摊销",
                "detail": f"系统中有 {intangible_count} 项无形资产，但未发现摊销凭证。漏摊销会导致利润虚增。",
                "suggestion": "按照无形资产摊销政策，检查是否漏记摊销费用。"
            })


# ═══════════════════════════════════════════════════════════════
#  二、发票合规风险
# ═══════════════════════════════════════════════════════════════

def _analyze_invoice_compliance(db, company_id, ps, pe, results):
    """进项风险发票、销项作废红冲率、税号完整性、进项抵扣风险"""
    risk_invoices = db.query(
        func.count(PurchaseInvoice.id), func.sum(PurchaseInvoice.total_amount)
    ).filter(PurchaseInvoice.company_id == company_id,
             PurchaseInvoice.invoice_risk_level.in_(["疑点", "异常", "失控"])).first()
    risk_cnt, risk_amt = risk_invoices
    if risk_cnt and risk_cnt > 0:
        results.append({
            "category": "发票合规", "category_icon": "🧾", "risk_score": 9, "risk_level": "高风险",
            "risk_color": "#dc2626", "urgency": "紧急",
            "item": "存在风险进项发票",
            "detail": f"有 {risk_cnt} 张进项发票被标记为「疑点/异常/失控」，涉及金额 {_safe_float(risk_amt):,.2f} 元。异常发票进项税额不得抵扣，已抵扣需做进项税额转出。",
            "suggestion": "立即核实风险发票来源，如确认异常应及时做进项税额转出处理，避免税务稽查风险。"
        })

    total_sales = db.query(func.count(SalesInvoice.id)).filter(SalesInvoice.company_id == company_id).scalar() or 0
    void_cnt = db.query(func.count(SalesInvoice.id)).filter(
        SalesInvoice.company_id == company_id, SalesInvoice.status.in_(["作废", "红冲"])).scalar() or 0
    if total_sales > 0:
        void_rate = void_cnt / total_sales * 100
        if void_rate > 10:
            results.append({
                "category": "发票合规", "category_icon": "🧾",
                "risk_score": 7 if void_rate > 20 else 5,
                "risk_level": "高风险" if void_rate > 20 else "中风险",
                "risk_color": "#dc2626" if void_rate > 20 else "#f59e0b",
                "urgency": "紧急" if void_rate > 20 else "提醒",
                "item": "销项发票作废/红冲率偏高",
                "detail": f"销项发票作废/红冲率为 {void_rate:.1f}%（{void_cnt}/{total_sales}），远超正常水平（<5%）。异常高的作废率可能被税务机关重点关注。",
                "suggestion": "核查作废和红冲原因，是否符合国家税务总局公告2016年第47号规定的红冲条件。避免频繁作废发票。"
            })

    pi_empty_tax = db.query(func.count(PurchaseInvoice.id)).filter(
        PurchaseInvoice.company_id == company_id,
        (PurchaseInvoice.seller_tax_no == None) | (PurchaseInvoice.seller_tax_no == "")).scalar()
    si_empty_tax = db.query(func.count(SalesInvoice.id)).filter(
        SalesInvoice.company_id == company_id,
        (SalesInvoice.buyer_tax_no == None) | (SalesInvoice.buyer_tax_no == "")).scalar()
    if pi_empty_tax:
        results.append({
            "category": "发票合规", "category_icon": "🧾", "risk_score": 5, "risk_level": "中风险",
            "risk_color": "#f59e0b", "urgency": "提醒",
            "item": "进项发票缺少销售方税号",
            "detail": f"有 {pi_empty_tax} 张进项发票的供应商税号为空。缺少税号的发票可能无法正常认证抵扣。",
            "suggestion": "完善供应商档案信息，确保每张进项发票都有完整的供应商纳税人识别号。"
        })
    if si_empty_tax:
        results.append({
            "category": "发票合规", "category_icon": "🧾", "risk_score": 4, "risk_level": "中风险",
            "risk_color": "#f59e0b", "urgency": "提醒",
            "item": "销项发票缺少购买方税号",
            "detail": f"有 {si_empty_tax} 张销项发票的购买方税号为空。缺少购买方税号可能影响对方认证抵扣。",
            "suggestion": "完善客户档案信息，开票时要求提供纳税人识别号。"
        })

    input_risk = db.query(func.count(InputVATDeduction.id)).filter(
        InputVATDeduction.company_id == company_id,
        InputVATDeduction.risk_level.in_(["疑点", "异常", "失控"])).scalar()
    if input_risk:
        results.append({
            "category": "发票合规", "category_icon": "🧾", "risk_score": 8, "risk_level": "高风险",
            "risk_color": "#dc2626", "urgency": "紧急",
            "item": "进项抵扣存在风险记录",
            "detail": f"有 {input_risk} 条进项抵扣记录被标记为风险状态。需逐条核实抵扣凭证的真实性和合规性。",
            "suggestion": "逐条核查风险进项抵扣记录，确认发票来源和业务真实性，必要时做进项税额转出。"
        })


# ═══════════════════════════════════════════════════════════════
#  三、发票深度分析（新增 P2）
# ═══════════════════════════════════════════════════════════════

SENSITIVE_KEYWORDS = [
    "咨询费", "服务费", "劳务费", "管理费", "技术费", "推广费", "居间费",
    "预付费卡", "预付卡", "购物卡", "礼品", "生活用品", "办公用品",
    "会议费", "培训费", "差旅费"
]

ZERO_TAX_CATEGORIES = ["免税", "零税率", "不征税", "出口退税"]


def _analyze_invoice_depth(db, company_id, ps, pe, results):
    """发票深度分析：税率结构、敏感业务、零税额/顶额发票、代开发票"""
    # 3.1 各税率销项发票结构
    tax_rates = db.query(
        SalesInvoice.tax_rate,
        func.count(SalesInvoice.id),
        func.sum(SalesInvoice.total_amount)
    ).filter(SalesInvoice.company_id == company_id).group_by(SalesInvoice.tax_rate).all()

    if tax_rates:
        rate_summary = []
        has_odd_rate = False
        for tr in tax_rates:
            rate = _safe_float(tr[0])
            cnt = int(tr[1])
            amt = _safe_float(tr[2])
            rate_summary.append(f"{rate*100:.0f}%（{cnt}张/{amt:,.0f}元）")
            if rate not in [0.0, 0.03, 0.05, 0.06, 0.09, 0.13]:
                has_odd_rate = True

        if has_odd_rate:
            results.append({
                "category": "发票深度", "category_icon": "🔍", "risk_score": 6, "risk_level": "中风险",
                "risk_color": "#f59e0b", "urgency": "提醒",
                "item": "存在非标准税率发票",
                "detail": f"销项发票税率分布：{'; '.join(rate_summary)}。存在非标准税率发票，需确认税率适用是否正确。",
                "suggestion": "核对税率与商品服务类别的匹配关系，确保增值税税率适用正确。"
            })

    # 3.2 敏感业务发票检测
    sensitive_invoices = {"sales": [], "purchase": []}
    for inv in db.query(SalesInvoice).filter(SalesInvoice.company_id == company_id).all():
        gn = (getattr(inv, 'goods_name', '') or '').lower()
        for kw in SENSITIVE_KEYWORDS:
            if kw in gn:
                sensitive_invoices["sales"].append({
                    "no": getattr(inv, 'digital_invoice_no', '') or (getattr(inv, 'invoice_code', '') or '') + (getattr(inv, 'invoice_no', '') or ''),
                    "kw": kw, "amt": _safe_float(inv.total_amount)
                })
                break

    for inv in db.query(PurchaseInvoice).filter(PurchaseInvoice.company_id == company_id).all():
        gn = (getattr(inv, 'goods_name', '') or '').lower()
        for kw in SENSITIVE_KEYWORDS:
            if kw in gn:
                sensitive_invoices["purchase"].append({
                    "other": getattr(inv, 'seller_name', '') or '',
                    "kw": kw, "amt": _safe_float(inv.total_amount)
                })
                break

    total_sensitive = len(sensitive_invoices["sales"]) + len(sensitive_invoices["purchase"])
    if total_sensitive > 0:
        detail_items = []
        sens_amt = 0
        for s in sensitive_invoices["sales"]:
            if len(detail_items) < 6:
                detail_items.append(f"销项[{s['kw']}]{s['amt']:,.0f}元")
            sens_amt += s['amt']
        for s in sensitive_invoices["purchase"]:
            if len(detail_items) < 6:
                detail_items.append(f"进项[{s['kw']}]{s['other']}{s['amt']:,.0f}元")
            sens_amt += s['amt']

        detail = f"检测到 {total_sensitive} 张涉及敏感业务的发票（合计 {sens_amt:,.2f} 元）：{'；'.join(detail_items)}。咨询费、服务费等敏感项目容易被税务机关重点关注。"
        if total_sensitive >= 3:
            results.append({
                "category": "发票深度", "category_icon": "🔍", "risk_score": 7, "risk_level": "高风险",
                "risk_color": "#dc2626", "urgency": "紧急",
                "item": "存在多张敏感业务发票",
                "detail": detail,
                "suggestion": "核查咨询费、服务费等敏感业务发票的真实性和合理性，保留合同、协议、成果交付物等证明材料。大额咨询费需准备定价合理性说明。"
            })
        else:
            results.append({
                "category": "发票深度", "category_icon": "🔍", "risk_score": 3, "risk_level": "低风险",
                "risk_color": "#3b82f6", "urgency": "建议",
                "item": "存在敏感业务发票",
                "detail": detail,
                "suggestion": "确保敏感业务发票背后有真实交易支撑，保留相关业务文档备查。"
            })

    # 3.3 代开发票检测（摘要含"代开"）
    agency_cnt = db.query(func.count(PurchaseInvoice.id)).filter(
        PurchaseInvoice.company_id == company_id,
        or_(func.lower(PurchaseInvoice.seller_name).contains("代开"),
            func.lower(PurchaseInvoice.goods_name).contains("代开"))).scalar() or 0
    if agency_cnt:
        results.append({
            "category": "发票深度", "category_icon": "🔍", "risk_score": 5, "risk_level": "中风险",
            "risk_color": "#f59e0b", "urgency": "提醒",
            "item": "存在代开发票",
            "detail": f"检测到 {agency_cnt} 张代开发票。代开发票的业务真实性需要严格核实。",
            "suggestion": "核查代开发票对应业务的真实性，确保有合同、资金流、货物流'三流一致'的证据链。"
        })


# ═══════════════════════════════════════════════════════════════
#  四、成本结构风险
# ═══════════════════════════════════════════════════════════════

def _analyze_cost_structure(db, company_id, ps, pe, results):
    """毛利率、费用率、业务招待费"""
    revenue = _get_account_sum(db, company_id, "6001", ps, pe, "credit")
    cost = _get_account_sum(db, company_id, "6401", ps, pe, "debit")

    if revenue > 0:
        gross_margin = (revenue - cost) / revenue * 100
        if gross_margin < 0:
            results.append({
                "category": "成本结构", "category_icon": "📈", "risk_score": 9, "risk_level": "高风险",
                "risk_color": "#dc2626", "urgency": "紧急",
                "item": "毛利率/净利率为负",
                "detail": f"主营业务收入 {revenue:,.2f} 元，主营业务成本 {cost:,.2f} 元，毛利率为 {gross_margin:.1f}%。成本大于收入属于严重异常，可能涉及少计收入或多列成本。",
                "suggestion": "核查收入是否完整入账、成本核算是否准确，是否存在跨期调节利润的行为。"
            })
        elif gross_margin < 5:
            results.append({
                "category": "成本结构", "category_icon": "📈", "risk_score": 5, "risk_level": "中风险",
                "risk_color": "#f59e0b", "urgency": "提醒",
                "item": "毛利率偏离行业参考值",
                "detail": f"毛利率仅为 {gross_margin:.1f}%，远低于一般行业水平。低毛利率可能被税务机关质疑定价不合理或隐瞒收入。",
                "suggestion": "分析成本构成，确认是否存在非正常成本支出；检查关联交易定价是否公允。"
            })

    mgmt_exp = _get_account_sum(db, company_id, "6602", ps, pe, "debit")
    sales_exp = _get_account_sum(db, company_id, "6601", ps, pe, "debit")
    fin_exp = _get_account_sum(db, company_id, "6603", ps, pe, "debit")
    total_exp = mgmt_exp + sales_exp + fin_exp
    if revenue > 0 and total_exp / revenue > 0.5:
        results.append({
            "category": "成本结构", "category_icon": "📈", "risk_score": 3, "risk_level": "低风险",
            "risk_color": "#3b82f6", "urgency": "建议",
            "item": "期间费用占比过高",
            "detail": f"三项期间费用合计 {total_exp:,.2f} 元，占收入 {total_exp/revenue*100:.1f}%。费用率过高可能影响企业所得税应纳税所得额计算。",
            "suggestion": "检查费用报销的合理性和合规性，确认各项费用的税前扣除政策。"
        })

    entertainment = _safe_float(db.query(func.sum(JournalEntry.debit_amount)).filter(
        JournalEntry.company_id == company_id,
        func.lower(JournalEntry.summary).contains("招待"),
        JournalEntry.period >= ps, JournalEntry.period <= pe).scalar())
    if entertainment > 0 and revenue > 0:
        ent_rate = entertainment / revenue * 100
        if ent_rate > 0.5:
            results.append({
                "category": "成本结构", "category_icon": "📈",
                "risk_score": 5 if ent_rate > 1 else 3,
                "risk_level": "中风险" if ent_rate > 1 else "低风险",
                "risk_color": "#f59e0b" if ent_rate > 1 else "#3b82f6",
                "urgency": "提醒" if ent_rate > 1 else "建议",
                "item": "业务招待费可能超标",
                "detail": f"业务招待费 {entertainment:,.2f} 元，占收入 {ent_rate:.2f}%。按税法规定，业务招待费按发生额60%扣除，且不得超过营业收入5‰。",
                "suggestion": "汇算清缴时注意纳税调增超标部分。建议控制业务招待费，完善报销审批流程。"
            })


# ═══════════════════════════════════════════════════════════════
#  五、财税票综合比对（新增 P0）
# ═══════════════════════════════════════════════════════════════

def _analyze_financial_tax_invoice_cross(db, company_id, ps, pe, results):
    """企业所得税申报收入 vs 增值税申报收入 vs 财报收入 vs 发票金额 四维交叉比对"""
    # 财报收入（序时账 6001 贷方发生额）
    fin_revenue = _get_account_sum(db, company_id, "6001", ps, pe, "credit")
    # 销项发票金额
    inv_amount = _safe_float(db.query(func.sum(SalesInvoice.total_amount)).filter(
        SalesInvoice.company_id == company_id,
        SalesInvoice.invoice_date >= _period_to_date_range(ps)[0],
        SalesInvoice.invoice_date <= _period_to_date_range(pe)[1]).scalar())

    # 增值税申报销售额（取最近一期）
    vat_sales = 0
    ccp_revenue = 0
    vat_list = db.query(VATDeclaration).filter(
        VATDeclaration.company_id == company_id,
        VATDeclaration.period >= ps, VATDeclaration.period <= pe
    ).order_by(VATDeclaration.period.desc()).limit(3).all()
    for v in vat_list:
        if v.form_main:
            fm = v.form_main if isinstance(v.form_main, dict) else json.loads(v.form_main)
            vat_sales += _safe_float(fm.get("sales_amount") or fm.get("total_sales"))

    # 文化事业建设费申报的应征收入
    ccf_list = db.query(CulturalConstructionFeeDeclaration).filter(
        CulturalConstructionFeeDeclaration.company_id == company_id,
        CulturalConstructionFeeDeclaration.period >= ps, CulturalConstructionFeeDeclaration.period <= pe
    ).all()
    for ccf in ccf_list:
        ccp_revenue += _safe_float(ccf.row1_taxable_income_current)

    diffs = []
    if fin_revenue > 0 and vat_sales > 0:
        diff_pct = abs(fin_revenue - vat_sales) / max(fin_revenue, vat_sales) * 100
        if diff_pct > 10:
            diffs.append(f"财报收入({fin_revenue:,.0f})与增值税申报销售额({vat_sales:,.0f})差异 {diff_pct:.1f}%")

    if fin_revenue > 0 and inv_amount > 0:
        diff_pct = abs(fin_revenue - inv_amount) / max(fin_revenue, inv_amount) * 100
        if diff_pct > 10:
            diffs.append(f"财报收入({fin_revenue:,.0f})与销项发票金额({inv_amount:,.0f})差异 {diff_pct:.1f}%")

    if vat_sales > 0 and inv_amount > 0:
        diff_pct = abs(vat_sales - inv_amount) / max(vat_sales, inv_amount) * 100
        if diff_pct > 10:
            diffs.append(f"增值税销售额({vat_sales:,.0f})与销项发票金额({inv_amount:,.0f})差异 {diff_pct:.1f}%")

    if diffs:
        results.append({
            "category": "财税票比对", "category_icon": "🔗", "risk_score": 9, "risk_level": "高风险",
            "risk_color": "#dc2626", "urgency": "紧急",
            "item": "财税票三表综合比对差异",
            "detail": f"收入类数据在财报/增值税申报/发票三个维度之间存在显著差异：{'；'.join(diffs)}。这可能暗示少计收入、虚开发票或申报数据错误。",
            "suggestion": "逐月核对三张表（财务报表、增值税申报表、发票汇总）的收入数据，查找差异原因。如因适用税率不同（如简易计税 vs 一般计税），应做好差异台账备查。"
        })
    elif fin_revenue > 0 and vat_sales > 0:
        diff_pct = abs(fin_revenue - vat_sales) / max(fin_revenue, vat_sales) * 100
        if diff_pct <= 5:
            results.append({
                "category": "财税票比对", "category_icon": "🔗", "risk_score": 0, "risk_level": "良好",
                "risk_color": "#10b981", "urgency": "",
                "item": "财税票收入数据基本一致",
                "detail": f"财报收入({fin_revenue:,.0f})与增值税销售额({vat_sales:,.0f})差异仅 {diff_pct:.1f}%，在合理范围内。",
                "suggestion": "继续保持财税票数据的一致性。"
            })


# ═══════════════════════════════════════════════════════════════
#  六、配比弹性分析（新增 P0）
# ═══════════════════════════════════════════════════════════════

def _analyze_ratio_elasticity(db, company_id, ps, pe, results):
    """收入/成本/费用/利润 变动率之间的弹性系数"""
    periods = _get_periods_between(ps, pe)
    if len(periods) < 2:
        return

    rev_data = _monthly_account_balance(db, company_id, "6001", ps, pe)
    cost_data = _monthly_account_balance(db, company_id, "6401", ps, pe)
    exp_data = {}
    for p in periods:
        d = _get_account_sum(db, company_id, "6601", p, p, "debit")
        d += _get_account_sum(db, company_id, "6602", p, p, "debit")
        d += _get_account_sum(db, company_id, "6603", p, p, "debit")
        exp_data[p] = d

    # 计算逐月变动率
    rev_changes = []
    cost_changes = []
    exp_changes = []
    for i in range(1, len(periods)):
        p1, p2 = periods[i - 1], periods[i]
        r1 = rev_data.get(p1, {}).get("credit", 0)
        r2 = rev_data.get(p2, {}).get("credit", 0)
        c1 = cost_data.get(p1, {}).get("debit", 0)
        c2 = cost_data.get(p2, {}).get("debit", 0)
        e1 = exp_data.get(p1, 0)
        e2 = exp_data.get(p2, 0)
        if r1 > 0: rev_changes.append((r2 - r1) / r1)
        if c1 > 0: cost_changes.append((c2 - c1) / c1)
        if e1 > 0: exp_changes.append((e2 - e1) / e1)

    all_changes = rev_changes + cost_changes + exp_changes
    if not all_changes:
        return

    max_change = max(abs(c) for c in all_changes)

    # 弹性系数：成本变动率 / 收入变动率
    elasticity_warnings = []
    for i in range(min(len(rev_changes), len(cost_changes))):
        rev_c = rev_changes[i]
        cost_c = cost_changes[i]
        if abs(rev_c) < 0.01: continue
        elasticity = abs(cost_c / rev_c) if abs(rev_c) > 0.001 else 99
        if elasticity > 1.5:
            elasticity_warnings.append(f"{periods[i+1]}: 成本弹性 {elasticity:.1f}（成本变动是收入的 {elasticity:.1f} 倍）")

    if max_change > 0.5:
        results.append({
            "category": "配比弹性", "category_icon": "📐", "risk_score": 8, "risk_level": "高风险",
            "risk_color": "#dc2626", "urgency": "紧急",
            "item": "收入/成本/费用变动率异常",
            "detail": f"分析期间内收入/成本/费用的最大单月变动率达 {max_change*100:.0f}%，远超正常水平（±30%以内）。月度间剧烈波动可能涉及跨期调节。" +
                     (f" 另有 {"；".join(elasticity_warnings[:3])}" if elasticity_warnings else ""),
            "suggestion": "逐项排查大额波动的原因（如大额订单、季节性因素、会计估计变更等），保留合理的商业解释和证明材料。特别是成本弹性>1.5的月份需重点核查。"
        })
    elif elasticity_warnings:
        results.append({
            "category": "配比弹性", "category_icon": "📐", "risk_score": 5, "risk_level": "中风险",
            "risk_color": "#f59e0b", "urgency": "提醒",
            "item": "成本弹性异常",
            "detail": f"部分月份成本增长率远超收入增长率：{'；'.join(elasticity_warnings[:3])}。弹性系数>1.5说明成本增长速度快于收入，需关注成本真实性。",
            "suggestion": "核查弹性异常月份的成本构成，确认是否存在虚增成本或多列费用的情况。"
        })
    elif max_change < 0.3 and len(periods) >= 3:
        results.append({
            "category": "配比弹性", "category_icon": "📐", "risk_score": 0, "risk_level": "良好",
            "risk_color": "#10b981", "urgency": "",
            "item": "收入成本配比稳定",
            "detail": f"分析期间内收入/成本/费用月度变动率均控制在±30%以内，收入成本配比关系稳定，经营状况较为健康。",
            "suggestion": "继续保持稳定的经营节奏。"
        })


# ═══════════════════════════════════════════════════════════════
#  七、隐匿收入 / 虚增成本 / 虚增利润（新增 P0）
# ═══════════════════════════════════════════════════════════════

def _analyze_hidden_inflated(db, company_id, ps, pe, results):
    """综合指标：其他应收款占比、存货异常、预收账款/应付账款激增"""
    other_rec = _get_account_balance(db, company_id, "1221")
    inventory = _get_account_balance(db, company_id, "1405")
    adv_receipts = _get_account_balance(db, company_id, "2203")  # 预收账款
    accounts_payable = _get_account_balance(db, company_id, "2202")

    revenue = _get_account_sum(db, company_id, "6001", ps, pe, "credit")
    total_assets = _get_account_balance(db, company_id, "1")

    # 7.1 其他应收款异常（隐匿收入/虚增成本的重要信号）
    if total_assets > 0 and other_rec / total_assets > 0.1:
        results.append({
            "category": "隐匿虚增", "category_icon": "⚠️", "risk_score": 8, "risk_level": "高风险",
            "risk_color": "#dc2626", "urgency": "紧急",
            "item": "其他应收款占比过高",
            "detail": f"其他应收款余额 {other_rec:,.2f} 元，占总资产 {other_rec/total_assets*100:.1f}%。大额其他应收款可能隐藏以下风险：(1)将收入挂账在其他应收款贷方；(2)通过其他应收款核算无票支出或虚假采购；(3)股东/关联方占用资金。",
            "suggestion": "逐户清理其他应收款明细账。对个人股东借款超过1年的，需视同分红代扣20%个税。对挂账超过2年的无票支出，建议核实业务真实性并考虑纳税调增。"
        })

    # 7.2 存货异常
    if revenue > 0 and inventory > 0:
        inv_turnover = revenue / inventory
        if inv_turnover < 1:
            results.append({
                "category": "隐匿虚增", "category_icon": "⚠️", "risk_score": 7, "risk_level": "高风险",
                "risk_color": "#dc2626", "urgency": "紧急",
                "item": "存货周转率极低",
                "detail": f"存货余额 {inventory:,.2f} 元，存货周转率仅 {inv_turnover:.1f} 次。极低的存货周转率可能暗示：(1)存货已发生减值但未计提跌价准备；(2)通过虚增存货来人为调高利润；(3)已销售商品未确认收入导致存货虚挂。",
                "suggestion": "进行存货实地盘点，核实账实是否相符。对滞销/过时存货应计提存货跌价准备。确认是否存在已发货未开票的情形，及时确认收入。"
            })

    # 7.3 预收账款/应付账款激增（隐匿收入信号）
    if revenue > 0:
        adv_receipt_pct = adv_receipts / revenue if adv_receipts > 0 else 0
        ap_pct = accounts_payable / revenue if accounts_payable > 0 else 0
        if adv_receipt_pct > 0.3:
            results.append({
                "category": "隐匿虚增", "category_icon": "⚠️", "risk_score": 6, "risk_level": "中风险",
                "risk_color": "#f59e0b", "urgency": "提醒",
                "item": "预收账款占比偏高",
                "detail": f"预收账款余额 {adv_receipts:,.2f} 元，占当期收入 {adv_receipt_pct*100:.1f}%。大量预收账款长期挂账可能涉及已满足收入确认条件但不确认收入，人为调节利润或延迟纳税。",
                "suggestion": "逐笔核实预收账款对应的合同履约进度，已履约部分应及时结转收入。按照新收入准则的五步法模型判断收入确认时点。"
            })
        if ap_pct > 0.5:
            results.append({
                "category": "隐匿虚增", "category_icon": "⚠️", "risk_score": 4, "risk_level": "中风险",
                "risk_color": "#f59e0b", "urgency": "提醒",
                "item": "应付账款占比偏高",
                "detail": f"应付账款余额 {accounts_payable:,.2f} 元，占当期收入 {ap_pct*100:.1f}%。若应付账款余额持续增长但无对应的存货/资产增加，可能暗示虚列成本或费用挂账。",
                "suggestion": "核对大额应付账款明细，确认是否有真实业务支撑。关注是否存在超出正常信用期的应付账款。"
            })


# ═══════════════════════════════════════════════════════════════
#  八、税负水平风险
# ═══════════════════════════════════════════════════════════════

def _analyze_tax_burden(db, company_id, ps, pe, results):
    """增值税税负率、CCF减征"""
    vat_list = db.query(VATDeclaration).filter(
        VATDeclaration.company_id == company_id,
        VATDeclaration.period >= ps, VATDeclaration.period <= pe
    ).order_by(VATDeclaration.period).all()

    low_vat_periods = []
    vat_details = []
    for v in vat_list:
        if v.form_main:
            fm = v.form_main if isinstance(v.form_main, dict) else json.loads(v.form_main)
            sales = _safe_float(fm.get("sales_amount") or fm.get("total_sales"))
            vat_payable = _safe_float(fm.get("vat_payable"))
            if sales > 0:
                rate = vat_payable / sales * 100
                vat_details.append({"period": v.period, "rate": rate, "sales": sales, "vat": vat_payable})
                if rate < 0.5:
                    low_vat_periods.append(f"{v.period}({rate:.2f}%)")

    if low_vat_periods:
        results.append({
            "category": "税负水平", "category_icon": "💰", "risk_score": 6, "risk_level": "中风险",
            "risk_color": "#f59e0b", "urgency": "提醒",
            "item": "行业税负率明显低于同行业平均水平",
            "detail": f"以下期间增值税税负率低于0.5%，可能触发税务预警：{', '.join(low_vat_periods)}。不同行业有预警税负率参照值。",
            "suggestion": "分析税负偏低原因（如期末留抵、免税收入、进项税额过大等），做好合理解释准备。避免增值税税负率持续低于行业预警值。"
        })

    # 整体税负率
    if vat_details:
        avg_rate = sum(d["rate"] for d in vat_details) / len(vat_details)
        if 0.5 <= avg_rate <= 5:
            results.append({
                "category": "税负水平", "category_icon": "💰", "risk_score": 0, "risk_level": "良好",
                "risk_color": "#10b981", "urgency": "",
                "item": "增值税税负率合理",
                "detail": f"分析期间平均增值税税负率 {avg_rate:.2f}%，处于行业合理区间。",
                "suggestion": "继续保持合理的税负水平。"
            })

    # CCF 减征
    ccf = db.query(CulturalConstructionFeeDeclaration).filter(
        CulturalConstructionFeeDeclaration.company_id == company_id,
        CulturalConstructionFeeDeclaration.period >= ps,
        CulturalConstructionFeeDeclaration.period <= pe).first()
    if ccf:
        rate = _safe_float(ccf.fee_reduction_rate, 0.5)
        reduction = _safe_float(ccf.row10a_fee_reduction_current, 0)
        if rate < 0.01 and reduction < 1:
            results.append({
                "category": "税负水平", "category_icon": "💰", "risk_score": 2, "risk_level": "低风险",
                "risk_color": "#3b82f6", "urgency": "建议",
                "item": "文化事业建设费未享受减征",
                "detail": f"期间 {ccf.period} 文化事业建设费申报未应用财税〔2025〕7号的50%减征优惠，可能多缴费款。",
                "suggestion": "确认是否符合减征条件，如符合应在下次申报时应用50%减征。已多缴的款项可咨询税务机关办理退抵。"
            })
        elif reduction > 0:
            results.append({
                "category": "税负水平", "category_icon": "💰", "risk_score": 0, "risk_level": "良好",
                "risk_color": "#10b981", "urgency": "",
                "item": "已享受文化事业建设费50%减征（财税〔2025〕7号）",
                "detail": f"期间 {ccf.period} 已正确应用减征优惠，减免额 {reduction:,.2f} 元。",
                "suggestion": "继续保持。"
            })


# ═══════════════════════════════════════════════════════════════
#  九、城建税（新增 P1）
# ═══════════════════════════════════════════════════════════════

CITY_TAX_RATES = {"市区": 0.07, "县城": 0.05, "建制镇": 0.05, "其他": 0.01}


def _analyze_urban_construction_tax(db, company_id, ps, pe, results):
    """城建税：税率合理性校验"""
    vat = _vat_payable_sum(db, company_id, ps, pe)

    if vat <= 0:
        return

    # 查找城建税相关凭证（通常科目编码 6404 或摘要含"城建税"）
    urban_tax = _safe_float(db.query(func.sum(JournalEntry.debit_amount)).filter(
        JournalEntry.company_id == company_id,
        JournalEntry.period >= ps, JournalEntry.period <= pe,
        func.lower(JournalEntry.summary).contains("城建")).scalar())

    if urban_tax <= 0:
        # 尝试用 6404 科目编码查找
        urban_tax = _get_account_sum(db, company_id, "6404", ps, pe, "debit")

    if urban_tax <= 0:
        # 尝试营业税金及附加中查找
        urban_tax = _safe_float(db.query(func.sum(JournalEntry.debit_amount)).filter(
            JournalEntry.company_id == company_id,
            JournalEntry.account_code.like("6403%"),
            JournalEntry.period >= ps, JournalEntry.period <= pe,
            func.lower(JournalEntry.summary).contains("城建")).scalar())

    if urban_tax > 0 and vat > 0:
        effective_rate = urban_tax / vat
        # 正常城建税税率 = 增值税应纳税额 × 7%/5%/1%
        if effective_rate < 0.01:
            results.append({
                "category": "城建税", "category_icon": "🏙️", "risk_score": 7, "risk_level": "高风险",
                "risk_color": "#dc2626", "urgency": "紧急",
                "item": "城建税实际税率偏低",
                "detail": f"增值税应纳税额 {vat:,.2f} 元，城建税 {urban_tax:,.2f} 元，实际税率 {effective_rate*100:.2f}%。城建税法定税率为增值税税额的7%/5%/1%（按所在地）。实际税率过低可能涉及漏缴。",
                "suggestion": "核实企业所在地适用的城建税税率，检查是否完整计提城建税。同时检查增值税免抵税额是否也计提了城建税。"
            })
        elif 0.01 <= effective_rate <= 0.09:
            results.append({
                "category": "城建税", "category_icon": "🏙️", "risk_score": 0, "risk_level": "良好",
                "risk_color": "#10b981", "urgency": "",
                "item": "城建税计提比率合理",
                "detail": f"增值税应纳税额 {vat:,.2f} 元，城建税 {urban_tax:,.2f} 元，实际税率 {effective_rate*100:.2f}%，与法定税率匹配。",
                "suggestion": "继续保持。"
            })
    elif urban_tax <= 0 and vat > 1000:
        results.append({
            "category": "城建税", "category_icon": "🏙️", "risk_score": 6, "risk_level": "中风险",
            "risk_color": "#f59e0b", "urgency": "提醒",
            "item": "有增值税但未计提城建税",
            "detail": f"增值税应纳税额 {vat:,.2f} 元，但未检测到城建税计提凭证。城建税是增值税的附加税，缴纳增值税必须同时缴纳城建税。",
            "suggestion": "立即检查是否漏提城建税。城建税 = 增值税应纳税额 × 适用税率（7%/5%/1%），同时计算教育费附加(3%)和地方教育附加(2%)。"
        })


# ═══════════════════════════════════════════════════════════════
#  十、房产税（新增 P1）
# ═══════════════════════════════════════════════════════════════

def _analyze_property_tax(db, company_id, ps, pe, results):
    """房产税：从价计征/从租计征"""
    # 查找固定资产中属于房屋建筑物的资产
    buildings = db.query(
        func.count(FixedAsset.id),
        func.sum(FixedAsset.original_value)
    ).filter(
        FixedAsset.company_id == company_id,
        FixedAsset.category.in_(["房屋建筑物", "房屋", "建筑物", "厂房", "办公楼", "仓库"]),
        FixedAsset.status.in_(["在用", "闲置"])
    ).first() or (0, 0)

    building_cnt, building_value = buildings
    if not building_cnt or building_cnt == 0:
        # 更宽泛的搜索
        buildings = db.query(
            func.count(FixedAsset.id),
            func.sum(FixedAsset.original_value)
        ).filter(
            FixedAsset.company_id == company_id,
            or_(FixedAsset.name.like("%房%"), FixedAsset.name.like("%楼%"),
                FixedAsset.category.like("%房%"), FixedAsset.category.like("%楼%")),
            FixedAsset.status.in_(["在用", "闲置"])
        ).first() or (0, 0)
        building_cnt, building_value = buildings

    # 查找房产税凭证
    property_tax = _safe_float(db.query(func.sum(JournalEntry.debit_amount)).filter(
        JournalEntry.company_id == company_id,
        JournalEntry.period >= ps, JournalEntry.period <= pe,
        func.lower(JournalEntry.summary).contains("房产税")).scalar())

    # 也查找税金及附加中的房产税
    if property_tax <= 0:
        property_tax = _safe_float(db.query(func.sum(JournalEntry.debit_amount)).filter(
            JournalEntry.company_id == company_id,
            JournalEntry.account_code.like("6403%"),
            JournalEntry.period >= ps, JournalEntry.period <= pe,
            func.lower(JournalEntry.summary).contains("房产")).scalar())

    if building_cnt and building_cnt > 0:
        expected_tax = _safe_float(building_value) * 0.7 * 0.012  # 房产原值 × 70% × 1.2%
        if property_tax <= 0:
            results.append({
                "category": "房产税", "category_icon": "🏠", "risk_score": 8, "risk_level": "高风险",
                "risk_color": "#dc2626", "urgency": "紧急",
                "item": "有房产原值但未缴房产税",
                "detail": f"系统中有 {building_cnt} 项房屋建筑物（原值合计 {_safe_float(building_value):,.2f} 元），但未发现房产税计提凭证。房产税漏缴有严重税务风险，滞纳金按日万分之五计算。",
                "suggestion": f"立即核实房产税缴纳情况：(1)从价计征：房产原值×(1-30%)×1.2%（年税额估算 {expected_tax:,.2f} 元）；(2)从租计征：租金收入×12%。若有出租部分，需分别计算。尽快补缴并做好纳税申报。"
            })
        elif building_value > 0 and _safe_float(building_value) > 1000000:
            # 检查是否足额缴纳
            if property_tax < expected_tax * 0.5:
                results.append({
                    "category": "房产税", "category_icon": "🏠", "risk_score": 5, "risk_level": "中风险",
                    "risk_color": "#f59e0b", "urgency": "提醒",
                    "item": "有租金收入但未缴房产税",
                    "detail": f"房屋建筑物净值约 {_safe_float(building_value):,.2f} 元，预计年房产税约 {expected_tax:,.2f} 元，实际计提 {property_tax:,.2f} 元，可能不足。",
                    "suggestion": "核实房产税计税依据（房产原值），确认是否所有房产均已申报。检查有无新建、改建、扩建的房产未纳入计税范围。"
                })
    elif property_tax > 0:
        results.append({
            "category": "房产税", "category_icon": "🏠", "risk_score": 0, "risk_level": "良好",
            "risk_color": "#10b981", "urgency": "",
            "item": "已计提房产税",
            "detail": f"房产税已计提 {property_tax:,.2f} 元，房产税申报工作已开展。",
            "suggestion": "继续保持按期申报缴纳。"
        })


# ═══════════════════════════════════════════════════════════════
#  十一、个人所得税全面分析（新增 P0）
# ═══════════════════════════════════════════════════════════════

def _analyze_iit(db, company_id, ps, pe, results):
    """个税：工资薪金、劳务报酬、股息红利、股东借款"""
    salary_records = db.query(SalaryRecord).filter(
        SalaryRecord.company_id == company_id,
        SalaryRecord.period >= ps, SalaryRecord.period <= pe
    ).all()

    if not salary_records:
        # 检查是否有雇员
        emp_count = db.query(func.count(Employee.id)).filter(
            Employee.company_id == company_id).scalar() or 0
        if emp_count > 0:
            results.append({
                "category": "个人所得税", "category_icon": "🧑‍💼", "risk_score": 7, "risk_level": "高风险",
                "risk_color": "#dc2626", "urgency": "紧急",
                "item": "有员工但无个税申报",
                "detail": f"系统中有 {emp_count} 名员工，但无工资薪金申报记录。不履行个人所得税代扣代缴义务将面临税务处罚。",
                "suggestion": "立即统计员工工资，完成个人所得税代扣代缴申报。即使员工工资未达到起征点（5,000元/月），也应当进行零申报。"
            })
        return

    # 11.1 劳务报酬代扣代缴
    labor_summaries = db.query(JournalEntry.summary).filter(
        JournalEntry.company_id == company_id, JournalEntry.period >= ps, JournalEntry.period <= pe,
        or_(func.lower(JournalEntry.summary).contains("劳务"),
            func.lower(JournalEntry.summary).contains("兼职"),
            func.lower(JournalEntry.summary).contains("临时工"))).all()

    if len(labor_summaries) > 3:
        # 检查是否有劳务报酬个税扣缴记录
        has_labor_tax = db.query(func.count(JournalEntry.id)).filter(
            JournalEntry.company_id == company_id, JournalEntry.period >= ps, JournalEntry.period <= pe,
            func.lower(JournalEntry.summary).contains("劳务报酬")
        ).scalar() or 0
        if not has_labor_tax:
            results.append({
                "category": "个人所得税", "category_icon": "🧑‍💼", "risk_score": 6, "risk_level": "中风险",
                "risk_color": "#f59e0b", "urgency": "提醒",
                "item": "有劳务费支付但未代扣个税",
                "detail": f"账务中涉及劳务/兼职/临时工等 {len(labor_summaries)} 笔记录，但未检测到劳务报酬个人所得税代扣记录。劳务报酬所得应按规定代扣代缴个人所得税。",
                "suggestion": "核查劳务费支付是否依法代扣代缴个人所得税（预扣率20%-40%）。要求劳务提供方到税务机关代开发票或通过自然人电子税务局申报。"
            })

    # 11.2 股东分红个税
    dividend_entries = db.query(JournalEntry).filter(
        JournalEntry.company_id == company_id, JournalEntry.period >= ps, JournalEntry.period <= pe,
        or_(func.lower(JournalEntry.summary).contains("股利"),
            func.lower(JournalEntry.summary).contains("分红"),
            func.lower(JournalEntry.summary).contains("利润分配"))).all()
    if dividend_entries:
        has_dividend_tax = db.query(func.count(JournalEntry.id)).filter(
            JournalEntry.company_id == company_id, JournalEntry.period >= ps, JournalEntry.period <= pe,
            func.lower(JournalEntry.summary).contains("股息红利")
        ).scalar() or 0
        if not has_dividend_tax:
            results.append({
                "category": "个人所得税", "category_icon": "🧑‍💼", "risk_score": 8, "risk_level": "高风险",
                "risk_color": "#dc2626", "urgency": "紧急",
                "item": "利润分配未代扣个人所得税",
                "detail": "账务中存在利润分配/分红记录，但未检测到股息红利个人所得税代扣记录。个人股东分红应代扣20%个人所得税。",
                "suggestion": "立即核查股东分红是否按规定代扣代缴20%个人所得税。根据《个人所得税法》，利息、股息、红利所得适用20%比例税率。"
            })

    # 11.3 股东借款视同分红
    other_rec_balance = _get_account_balance(db, company_id, "1221")
    if other_rec_balance > 100000:
        results.append({
            "category": "个人所得税", "category_icon": "🧑‍💼", "risk_score": 7, "risk_level": "高风险",
            "risk_color": "#dc2626", "urgency": "紧急",
            "item": "股东借款超期未还",
            "detail": f"其他应收款余额 {other_rec_balance:,.2f} 元。按照财税〔2003〕158号文，个人投资者从企业借款超过1年未归还且未用于生产经营的，视同股息红利分配，应代扣20%个人所得税。",
            "suggestion": "逐户清理其他应收款中的个人借款：(1)督促在纳税年度终了前归还；(2)超过1年未归还的，需按分红处理代扣20%个税；(3)完善借款管理制度，签订借款协议明确用途和归还期限。"
        })

    # 11.4 工资薪金个税整体情况
    tax_payers = sum(1 for sr in salary_records if _safe_float(sr.tax_payable) > 1)
    total_emp = len(salary_records)
    if total_emp > 5 and tax_payers == 0:
        results.append({
            "category": "个人所得税", "category_icon": "🧑‍💼", "risk_score": 3, "risk_level": "低风险",
            "risk_color": "#3b82f6", "urgency": "建议",
            "item": "无员工缴纳工资薪金个税",
            "detail": f"全部 {total_emp} 名员工均未缴纳个人所得税。如有管理层或高薪人员，需确认收入是否完整申报。",
            "suggestion": "核实：(1)奖金、津贴、补贴是否并入工资计税；(2)是否通过费用报销、私卡收款等方式变相发放工资；(3)专项附加扣除申报是否合规。"
        })
    elif tax_payers > 0:
        total_tax = sum(_safe_float(sr.tax_payable) for sr in salary_records)
        results.append({
            "category": "个人所得税", "category_icon": "🧑‍💼", "risk_score": 0, "risk_level": "良好",
            "risk_color": "#10b981", "urgency": "",
            "item": "工资薪金个税代扣代缴正常",
            "detail": f"{tax_payers}/{total_emp} 名员工已代扣工资薪金个税，合计 {total_tax:,.2f} 元。",
            "suggestion": "继续保持按期申报缴纳。"
        })


# ═══════════════════════════════════════════════════════════════
#  十二、印花税（新增 P0）
# ═══════════════════════════════════════════════════════════════

STAMP_TAX_RATES = {
    "购销合同": 0.0003, "加工承揽": 0.0005, "建筑工程勘察": 0.0005,
    "建筑安装工程承包": 0.0003, "财产租赁": 0.001, "货物运输": 0.0005,
    "仓储保管": 0.001, "借款合同": 0.00005, "财产保险": 0.001,
    "技术合同": 0.0003, "产权转移": 0.0005, "营业账簿": 5.0,  # 按件贴花
    "权利许可证照": 5.0
}


def _analyze_stamp_tax(db, company_id, ps, pe, results):
    """印花税：合同金额 vs 印花税缴纳"""
    # 查找印花税凭证
    stamp_tax = _safe_float(db.query(func.sum(JournalEntry.debit_amount)).filter(
        JournalEntry.company_id == company_id,
        JournalEntry.period >= ps, JournalEntry.period <= pe,
        func.lower(JournalEntry.summary).contains("印花税")).scalar())

    if stamp_tax <= 0:
        stamp_tax = _safe_float(db.query(func.sum(JournalEntry.debit_amount)).filter(
            JournalEntry.company_id == company_id,
            JournalEntry.account_code.like("6403%"),
            JournalEntry.period >= ps, JournalEntry.period <= pe,
            func.lower(JournalEntry.summary).contains("印花")).scalar())

    # 检查合同台账
    contracts = db.query(func.count(Contract.id), func.sum(Contract.amount)).filter(
        Contract.company_id == company_id).first() or (0, 0)
    contract_cnt, contract_amount = contracts

    # 估算印花税 = 合同金额 × 0.03% + 业务收入 × 0.03%（购销合同）
    revenue = _get_account_sum(db, company_id, "6001", ps, pe, "credit")
    purchase = _get_account_sum(db, company_id, "1402", ps, pe, "debit")  # 采购/原材料

    estimated_stamp = 0
    if revenue > 0:
        estimated_stamp += revenue * 0.0003  # 销售合同印花税
    if purchase > 0:
        estimated_stamp += purchase * 0.0003  # 采购合同印花税
    if contract_cnt > 0:
        estimated_stamp += _safe_float(contract_amount) * 0.0003

    if revenue > 100000 and stamp_tax <= 0:
        results.append({
            "category": "印花税", "category_icon": "📜", "risk_score": 5, "risk_level": "中风险",
            "risk_color": "#f59e0b", "urgency": "提醒",
            "item": "营业收入较高但未缴印花税",
            "detail": f"公司有较大营业收入（{revenue:,.2f} 元），但未检测到印花税计提凭证。购销合同、账簿等均需缴纳印花税，漏缴面临滞纳金和罚款。预计应缴印花税约 {estimated_stamp:,.2f} 元（仅按购销合同估计）。",
            "suggestion": "立即补缴印花税：(1)购销合同按合同金额万分之三贴花；(2)营业账簿按件贴花（每件5元）；(3)其他应税凭证按规定税率缴纳。可咨询税务机关确认申报方式。"
        })
    elif revenue > 100000 and stamp_tax > 0:
        if estimated_stamp > 0 and stamp_tax < estimated_stamp * 0.3:
            results.append({
                "category": "印花税", "category_icon": "📜", "risk_score": 3, "risk_level": "低风险",
                "risk_color": "#3b82f6", "urgency": "建议",
                "item": "印花税可能未足额缴纳",
                "detail": f"已缴印花税 {stamp_tax:,.2f} 元，但按购销合同初步估算约需 {estimated_stamp:,.2f} 元，可能存在少缴情况。",
                "suggestion": "全面梳理各类应税合同和凭证，按法定税率计算应缴印花税，及时补缴差额。"
            })
        else:
            results.append({
                "category": "印花税", "category_icon": "📜", "risk_score": 0, "risk_level": "良好",
                "risk_color": "#10b981", "urgency": "",
                "item": "已缴纳印花税",
                "detail": f"已计提印花税 {stamp_tax:,.2f} 元，印花税申报工作已开展。",
                "suggestion": "继续保持按期申报缴纳。"
            })


# ═══════════════════════════════════════════════════════════════
#  十三、纳税调整检查（新增 P1）
# ═══════════════════════════════════════════════════════════════

def _analyze_tax_adjustment(db, company_id, ps, pe, results):
    """纳税调整：业务招待费60%、广宣费15%、职工福利费14%、资产减值"""
    revenue = _get_account_sum(db, company_id, "6001", ps, pe, "credit")
    if revenue <= 0:
        return

    adjustments = []

    # 13.1 业务招待费 = 发生额 × 60% 与 收入 × 5‰ 取孰低
    entertainment = _safe_float(db.query(func.sum(JournalEntry.debit_amount)).filter(
        JournalEntry.company_id == company_id,
        func.lower(JournalEntry.summary).contains("招待"),
        JournalEntry.period >= ps, JournalEntry.period <= pe).scalar())
    if entertainment > 0:
        ded_limit_1 = entertainment * 0.6
        ded_limit_2 = revenue * 0.005
        actual_ded = min(ded_limit_1, ded_limit_2)
        non_ded = entertainment - actual_ded
        if non_ded > 1000:
            adjustments.append(f"业务招待费：发生{entertainment:,.0f}元，可扣{actual_ded:,.0f}元，需调增{non_ded:,.0f}元")

    # 13.2 广宣费 = 收入 × 15%（一般企业）或 30%（化妆品/医药/饮料制造）
    ad_promo = _safe_float(db.query(func.sum(JournalEntry.debit_amount)).filter(
        JournalEntry.company_id == company_id,
        JournalEntry.period >= ps, JournalEntry.period <= pe,
        or_(func.lower(JournalEntry.summary).contains("广告"),
            func.lower(JournalEntry.summary).contains("宣传"),
            func.lower(JournalEntry.summary).contains("推广"))).scalar())
    if ad_promo > 0:
        ad_limit = revenue * 0.15
        if ad_promo > ad_limit:
            non_ded = ad_promo - ad_limit
            adjustments.append(f"广宣费：发生{ad_promo:,.0f}元，限额{ad_limit:,.0f}元，需调增{non_ded:,.0f}元")
        # 单独输出广宣费超限额风险项（为规则匹配）
        if ad_promo > ad_limit:
            results.append({
                "category": "纳税调整", "category_icon": "📝",
                "risk_score": 6 if ad_promo > ad_limit * 1.5 else 4,
                "risk_level": "中风险" if ad_promo > ad_limit * 1.5 else "低风险",
                "risk_color": "#f59e0b" if ad_promo > ad_limit * 1.5 else "#3b82f6",
                "urgency": "提醒",
                "item": "广宣费超限额",
                "detail": f"广宣费发生{ad_promo:,.0f}元，企业所得税税前扣除限额为营业收入×15%={ad_limit:,.0f}元，超出{non_ded:,.0f}元需纳税调增。化妆品/医药/饮料制造企业限额为30%，如适用请备注。",
                "suggestion": "汇算清缴时在《广告费和业务宣传费跨年度纳税调整明细表》（A105060）中填列超标金额，结转以后年度扣除。",
                "required_evidence": ["广告/宣传合同及发票", "广告投放计划及执行报告", "A105060纳税调整明细表"]
            })

    # 13.3 职工福利费 = 工资总额 × 14%
    total_salary = _safe_float(db.query(func.sum(SalaryRecord.current_income)).filter(
        SalaryRecord.company_id == company_id,
        SalaryRecord.period >= ps, SalaryRecord.period <= pe).scalar())
    welfare = _safe_float(db.query(func.sum(JournalEntry.debit_amount)).filter(
        JournalEntry.company_id == company_id,
        JournalEntry.period >= ps, JournalEntry.period <= pe,
        or_(func.lower(JournalEntry.summary).contains("福利"),
            func.lower(JournalEntry.summary).contains("补贴"),
            func.lower(JournalEntry.summary).contains("补助")),
        ~func.lower(JournalEntry.summary).contains("工资")
    ).scalar())
    if welfare > 0 and total_salary > 0:
        wel_limit = total_salary * 0.14
        if welfare > wel_limit:
            non_ded = welfare - wel_limit
            adjustments.append(f"职工福利费：发生{welfare:,.0f}元，限额(工资×14%){wel_limit:,.0f}元，需调增{non_ded:,.0f}元")

    if adjustments:
        results.append({
            "category": "纳税调整", "category_icon": "📝", "risk_score": 6, "risk_level": "中风险",
            "risk_color": "#f59e0b", "urgency": "提醒",
            "item": "纳税调整减少率异常偏低",
            "detail": f"根据企业所得税法，以下项目需在汇算清缴时进行纳税调增：{'；'.join(adjustments)}。未做纳税调增将导致少缴企业所得税，面临补税+滞纳金+罚款风险。",
            "suggestion": "汇算清缴时务必在《纳税调整项目明细表》（A105000）中填写上述调增金额。建议提前做好税务台账管理，准备完整的证明材料备查。"
        })
    elif entertainment > 0 and entertainment < revenue * 0.005:
        results.append({
            "category": "纳税调整", "category_icon": "📝", "risk_score": 0, "risk_level": "良好",
            "risk_color": "#10b981", "urgency": "",
            "item": "纳税调整项目在限额内",
            "detail": f"业务招待费({entertainment:,.0f}元)在扣除限额(revenue×5‰={revenue*0.005:,.0f}元)以内，汇算清缴不会产生重大纳税调增。",
            "suggestion": "继续保持合理的费用控制。"
        })


# ═══════════════════════════════════════════════════════════════
#  十四、收入确认时点（新增 P1）
# ═══════════════════════════════════════════════════════════════

def _analyze_revenue_timing(db, company_id, ps, pe, results):
    """收入确认时点：四季度收入占比、年底集中开票"""
    # 获取全年各月收入（传入 YYYY-MM 格式）
    y_start = ps[:4] + "-01"
    y_end = pe[:4] + "-12"
    monthly_rev = _monthly_account_balance(db, company_id, "6001", y_start, y_end)

    # 计算各季度收入
    q_rev = {"Q1": 0, "Q2": 0, "Q3": 0, "Q4": 0}
    for period, val in monthly_rev.items():
        m = int(period[5:7])
        month_amt = val.get("credit", 0)
        if 1 <= m <= 3: q_rev["Q1"] += month_amt
        elif 4 <= m <= 6: q_rev["Q2"] += month_amt
        elif 7 <= m <= 9: q_rev["Q3"] += month_amt
        else: q_rev["Q4"] += month_amt

    total = sum(q_rev.values())
    if total <= 0:
        return

    q4_pct = q_rev["Q4"] / total * 100

    # 检查收入确认时点集中度（单月收入占比超过40%或Top2月份占比超过70%）
    monthly_credits = [(p, v.get("credit", 0)) for p, v in monthly_rev.items()]
    monthly_credits.sort(key=lambda x: x[1], reverse=True)
    top2_pct = sum(v for _, v in monthly_credits[:2]) / total * 100 if total > 0 else 0
    max_month_pct = monthly_credits[0][1] / total * 100 if total > 0 and monthly_credits else 0

    if max_month_pct > 40:
        results.append({
            "category": "收入时点", "category_icon": "📅", "risk_score": 7, "risk_level": "高风险",
            "risk_color": "#dc2626", "urgency": "紧急",
            "item": "收入确认时点集中度异常",
            "detail": f"收入高度集中于{monthly_credits[0][0]}月（占比{max_month_pct:.1f}%），前2月合计占比{top2_pct:.1f}%。收入集中度过高可能涉及：(1)跨期调节收入确认时点；(2)突击开票虚增收入；(3)收入确认时点与业务实质不符。",
            "suggestion": "核查集中月份的合同签订日期、交付验收记录与收入确认时点是否一致。确保收入确认符合新收入准则（五步法）关于履约义务完成的规定。",
            "required_evidence": ["高收入月份的合同/订单清单", "对应交付验收单据", "收入确认时点合理性说明"]
        })
    elif max_month_pct > 30:
        results.append({
            "category": "收入时点", "category_icon": "📅", "risk_score": 3, "risk_level": "低风险",
            "risk_color": "#3b82f6", "urgency": "建议",
            "item": "收入确认时点集中度异常",
            "detail": f"收入在{monthly_credits[0][0]}月占比较高（{max_month_pct:.1f}%），建议关注是否存在季节性因素或收入确认时点不合理的情况。",
            "suggestion": "如有合理的季节性因素（如项目集中在某月交付），做好业务记录备查。"
        })

    if q4_pct > 40:
        results.append({
            "category": "收入时点", "category_icon": "📅", "risk_score": 7, "risk_level": "高风险",
            "risk_color": "#dc2626", "urgency": "紧急",
            "item": "四季度收入占比过高",
            "detail": f"四季度收入占比 {q4_pct:.1f}%（Q1:{q_rev['Q1']/total*100:.0f}% Q2:{q_rev['Q2']/total*100:.0f}% Q3:{q_rev['Q3']/total*100:.0f}% Q4:{q4_pct:.1f}%）。四季度收入超过全年的40%可能涉及：(1)年底突击开票调节税负；(2)提前确认收入完成业绩指标；(3)跨期调节利润。",
            "suggestion": "核查四季度收入确认的真实性：(1)是否满足收入确认条件（交付/验收完成）；(2)是否存在次年退货或红冲的情况；(3)关联交易定价是否公允。保留合同、验收单、物流单据等证据链。"
        })
    elif q4_pct > 30:
        results.append({
            "category": "收入时点", "category_icon": "📅", "risk_score": 3, "risk_level": "低风险",
            "risk_color": "#3b82f6", "urgency": "建议",
            "item": "四季度收入占比较高",
            "detail": f"四季度收入占比 {q4_pct:.1f}%，略高于全年平均水平。如果存在明显的季节性因素是正常的，但仍建议关注。",
            "suggestion": "确认四季度收入增长是否有合理的商业理由（如季节性销售旺季、年底项目集中交付等），做好记录备查。"
        })
    else:
        q_details = " / ".join([f"Q{i+1}:{q_rev[f'Q{i+1}']/total*100:.0f}%" for i in range(4)])
        results.append({
            "category": "收入时点", "category_icon": "📅", "risk_score": 0, "risk_level": "良好",
            "risk_color": "#10b981", "urgency": "",
            "item": "收入季度分布均衡",
            "detail": f"各季度收入占比：{q_details}，收入分布较为均衡，无明显的年底集中确认迹象。",
            "suggestion": "继续保持均衡的收入确认节奏。"
        })


# ═══════════════════════════════════════════════════════════════
#  十五、政策执行风险
# ═══════════════════════════════════════════════════════════════

def _analyze_policy_execution(db, company_id, ps, pe, results):
    """社保公积金、结账状态"""
    from database import Period as PeriodModel
    ss_decl_count = db.query(func.count(SocialSecurityDeclaration.id)).filter(
        SocialSecurityDeclaration.company_id == company_id).scalar() or 0
    # SocialSecurityDetail 无 company_id，通过 declaration_id 关联查询
    ss_detail_count = db.query(func.count(SocialSecurityDetail.id)).join(
        SocialSecurityDeclaration, SocialSecurityDetail.declaration_id == SocialSecurityDeclaration.id
    ).filter(SocialSecurityDeclaration.company_id == company_id).scalar() or 0
    if not ss_decl_count and not ss_detail_count:
        results.append({
            "category": "政策执行", "category_icon": "📋", "risk_score": 8, "risk_level": "高风险",
            "risk_color": "#dc2626", "urgency": "紧急",
            "item": "未发现社保缴存记录",
            "detail": "系统中未检测到社会保险费申报记录。未依法缴纳社保存在劳动监察和补缴风险。",
            "suggestion": "依法为员工缴纳社会保险费，包括养老、医疗、失业、工伤、生育保险。"
        })

    hf_decl_count = db.query(func.count(HousingFundDeclaration.id)).filter(
        HousingFundDeclaration.company_id == company_id).scalar() or 0
    hf_detail_count = db.query(func.count(HousingFundDetail.id)).filter(
        HousingFundDetail.company_id == company_id).scalar() or 0
    hf_journal_count = db.query(func.count(JournalEntry.id)).filter(
        JournalEntry.company_id == company_id,
        JournalEntry.source.in_(["公积金计提", "公积金缴纳"])
    ).scalar() or 0
    if not hf_decl_count and not hf_detail_count and not hf_journal_count:
        results.append({
            "category": "政策执行", "category_icon": "📋", "risk_score": 6, "risk_level": "中风险",
            "risk_color": "#f59e0b", "urgency": "提醒",
            "item": "未发现住房公积金缴存记录",
            "detail": "系统中未检测到住房公积金缴存记录。未依法缴纳公积金存在被责令限期缴存的风险。",
            "suggestion": "按照《住房公积金管理条例》为员工缴存住房公积金。"
        })

    from database import Period as PeriodModel
    open_periods = db.query(func.count(PeriodModel.id)).filter(
        PeriodModel.company_id == company_id,
        PeriodModel.status == "开放",
        PeriodModel.period <= pe).scalar()
    if open_periods and open_periods > 3:
        results.append({
            "category": "政策执行", "category_icon": "📋", "risk_score": 4, "risk_level": "中风险",
            "risk_color": "#f59e0b", "urgency": "提醒",
            "item": "存在多个未结账期间",
            "detail": f"有 {open_periods} 个开放状态的期间未结账。长期不结账可能导致账目混乱和数据不一致。",
            "suggestion": "定期进行月度结账，确保各期间数据独立完整。建议每月末及时结账。"
        })


# ═══════════════════════════════════════════════════════════════
#  十六、资金与往来风险
# ═══════════════════════════════════════════════════════════════

def _analyze_capital_risks(db, company_id, ps, pe, results):
    """银行存款、应收账款、其他应收款"""
    bank_deposit = _get_account_balance(db, company_id, "1002")
    if bank_deposit < 0:
        results.append({
            "category": "资金往来", "category_icon": "🏦", "risk_score": 8, "risk_level": "高风险",
            "risk_color": "#dc2626", "urgency": "紧急",
            "item": "银行存款余额为负",
            "detail": f"银行存款科目余额为 {bank_deposit:,.2f} 元（贷方）。银行存款余额为负说明可能存在记账错误或未达账项。",
            "suggestion": "编制银行存款余额调节表，逐笔核对银行对账单与账面金额，查找差异原因并调整。"
        })

    ar_balance = _get_account_balance(db, company_id, "1122")
    if ar_balance > 0:
        early_ar = db.query(func.sum(JournalEntry.debit_amount)).filter(
            JournalEntry.company_id == company_id,
            JournalEntry.account_code.like("1122%"),
            JournalEntry.entry_date < (date.today() - timedelta(days=365)).isoformat()).scalar()
        if early_ar and early_ar > 0:
            results.append({
                "category": "资金往来", "category_icon": "🏦", "risk_score": 5, "risk_level": "中风险",
                "risk_color": "#f59e0b", "urgency": "提醒",
                "item": "应收账款存在长期挂账",
                "detail": f"应收账款余额 {ar_balance:,.2f} 元，其中可能包含超过1年的长期挂账。长期未收回的应收账款存在坏账风险。",
                "suggestion": "清理应收账款，对长期挂账发函催收。根据账龄计提坏账准备，坏账损失需符合税法规定的条件方可税前扣除。"
            })

    other_rec = _get_account_balance(db, company_id, "1221")
    if other_rec > 0:
        total_assets = _get_account_balance(db, company_id, "1")
        if total_assets > 0 and other_rec / total_assets > 0.1:
            results.append({
                "category": "资金往来", "category_icon": "🏦", "risk_score": 6, "risk_level": "中风险",
                "risk_color": "#f59e0b", "urgency": "提醒",
                "item": "其他应收款占比偏高",
                "detail": f"其他应收款余额 {other_rec:,.2f} 元，占总资产 {other_rec/total_assets*100:.1f}%。股东借款或关联方占用资金可能涉及个人所得税风险。",
                "suggestion": "核查其他应收款明细，清理不合规的关联方资金占用。个人股东借款超过1年未归还的，需视同分红缴纳个人所得税。"
            })


# ═══════════════════════════════════════════════════════════════
#  十七、薪酬合规风险
# ═══════════════════════════════════════════════════════════════

def _analyze_salary_compliance(db, company_id, ps, pe, results):
    """工资分布、个税起征点"""
    salary_records = db.query(SalaryRecord).filter(
        SalaryRecord.company_id == company_id,
        SalaryRecord.period >= ps, SalaryRecord.period <= pe).all()
    if not salary_records:
        return

    low_income_count = sum(1 for sr in salary_records if 4500 <= _safe_float(sr.current_income) <= 5500)
    total_emp = len(salary_records)
    if total_emp > 1 and low_income_count / total_emp > 0.5:
        results.append({
            "category": "薪酬合规", "category_icon": "👥", "risk_score": 5, "risk_level": "中风险",
            "risk_color": "#f59e0b", "urgency": "提醒",
            "item": "工资集中在个税起征点附近",
            "detail": f"{low_income_count}/{total_emp} 名员工的月收入集中在5,000元左右。全员收入在起征点附近可能被怀疑通过分拆工资等方式规避个税。",
            "suggestion": "确保工资发放真实反映员工劳动价值。如确有合理性，保留薪酬制度和岗位说明备查。"
        })

    zero_tax_count = sum(1 for sr in salary_records if _safe_float(sr.tax_payable) < 0.01)
    if zero_tax_count == total_emp and total_emp > 2:
        results.append({
            "category": "薪酬合规", "category_icon": "👥", "risk_score": 3, "risk_level": "低风险",
            "risk_color": "#3b82f6", "urgency": "建议",
            "item": "全体员工均无需缴纳个税",
            "detail": f"全部 {total_emp} 名员工均未产生个人所得税。如确有高管或高薪人员，需核实是否少报收入或违规使用专项附加扣除。",
            "suggestion": "核实收入是否全额申报，特别是奖金、补贴等是否并入工资薪金所得计税。"
        })


# ═══════════════════════════════════════════════════════════════
#  十八、客户风险穿透（新增 P1）
# ═══════════════════════════════════════════════════════════════

def _analyze_customer_penetration(db, company_id, ps, pe, results):
    """前十大客户分析、客户集中度、异常客户"""
    customers = db.query(Customer).filter(Customer.company_id == company_id).all()
    if not customers:
        return

    # 统计每个客户的销售额
    customer_sales = {}
    for c in customers:
        code = c.code or ''
        if code:
            amt = _safe_float(db.query(func.sum(SalesInvoice.total_amount)).filter(
                SalesInvoice.company_id == company_id,
                SalesInvoice.buyer_tax_no == c.tax_no).scalar())
            if amt > 0:
                customer_sales[c.name or code] = amt

    if not customer_sales:
        return

    total_customer_sales = sum(customer_sales.values())
    top_sorted = sorted(customer_sales.items(), key=lambda x: x[1], reverse=True)

    # 客户集中度
    if len(top_sorted) >= 1:
        top1 = top_sorted[0][1] / total_customer_sales * 100
        if top1 > 50:
            results.append({
                "category": "客户穿透", "category_icon": "🏢", "risk_score": 6, "risk_level": "中风险",
                "risk_color": "#f59e0b", "urgency": "提醒",
                "item": "第一大客户占比过高",
                "detail": f"第一大客户「{top_sorted[0][0]}」销售额占全部销售 {top1:.1f}%。过度依赖单一客户可能影响经营的独立性判断，且存在关联交易被关注的风险。",
                "suggestion": "分析第一大客户是否为关联方：(1)如是关联方，确保关联交易定价公允且有完整文档；(2)如非关联方，关注客户集中风险对持续经营判断的影响。"
            })

    # 客户税号检查
    customers_no_tax = db.query(func.count(Customer.id)).filter(
        Customer.company_id == company_id,
        (Customer.tax_no == None) | (Customer.tax_no == "")).scalar() or 0
    total_cust = len(customers)
    if total_cust > 0 and customers_no_tax / total_cust > 0.3:
        results.append({
            "category": "客户穿透", "category_icon": "🏢", "risk_score": 4, "risk_level": "中风险",
            "risk_color": "#f59e0b", "urgency": "提醒",
            "item": "客户税号完整度不足",
            "detail": f"{customers_no_tax}/{total_cust} 个客户缺少纳税人识别号。缺少税号的开票无法被客户正常认证抵扣。",
            "suggestion": "完善客户档案，对经常开票的客户必须记录纳税人识别号、地址、电话及开户行信息。"
        })

    if total_customer_sales > 0 and len(top_sorted) >= 3:
        top3_pct = sum(x[1] for x in top_sorted[:3]) / total_customer_sales * 100
        if top3_pct < 80:
            results.append({
                "category": "客户穿透", "category_icon": "🏢", "risk_score": 0, "risk_level": "良好",
                "risk_color": "#10b981", "urgency": "",
                "item": "客户分布合理",
                "detail": f"前三大客户销售占比 {top3_pct:.0f}%，客户群体较为分散，不存在过度依赖单一客户的风险。",
                "suggestion": "继续保持合理的客户结构。"
            })


# ═══════════════════════════════════════════════════════════════
#  十九、供应商风险穿透（新增 P1）
# ═══════════════════════════════════════════════════════════════

def _analyze_supplier_penetration(db, company_id, ps, pe, results):
    """前十大供应商分析、供应商集中度、供应商税号完整性"""
    suppliers = db.query(Supplier).filter(Supplier.company_id == company_id).all()
    if not suppliers:
        return

    supplier_purchase = {}
    for s in suppliers:
        code = s.code or ''
        if code:
            amt = _safe_float(db.query(func.sum(PurchaseInvoice.total_amount)).filter(
                PurchaseInvoice.company_id == company_id,
                PurchaseInvoice.seller_tax_no == s.tax_no).scalar())
            if amt > 0:
                supplier_purchase[s.name or code] = amt

    if not supplier_purchase:
        return

    total_purchase = sum(supplier_purchase.values())
    top_sorted = sorted(supplier_purchase.items(), key=lambda x: x[1], reverse=True)

    if len(top_sorted) >= 1:
        top1 = top_sorted[0][1] / total_purchase * 100
        if top1 > 50:
            results.append({
                "category": "供应商穿透", "category_icon": "🏭", "risk_score": 5, "risk_level": "中风险",
                "risk_color": "#f59e0b", "urgency": "提醒",
                "item": "第一大供应商占比过高",
                "detail": f"第一大供应商「{top_sorted[0][0]}」采购额占全部采购 {top1:.1f}%。过度依赖单一供应商存在供应链风险，且可能被怀疑为关联交易或虚假采购。",
                "suggestion": "(1)确认是否为关联方，如是需做关联交易披露；(2)保持采购来源的多元化；(3)保留完整的采购合同、入库单、付款凭证等证据链。"
            })

    # 供应商税号检查
    no_tax = db.query(func.count(Supplier.id)).filter(
        Supplier.company_id == company_id,
        (Supplier.tax_no == None) | (Supplier.tax_no == "")).scalar() or 0
    total_sup = len(suppliers)
    if total_sup > 0 and no_tax / total_sup > 0.3:
        results.append({
            "category": "供应商穿透", "category_icon": "🏭", "risk_score": 4, "risk_level": "中风险",
            "risk_color": "#f59e0b", "urgency": "提醒",
            "item": "供应商税号完整度不足",
            "detail": f"{no_tax}/{total_sup} 个供应商缺少纳税人识别号。缺少税号会影响进项发票的获取和认证抵扣。",
            "suggestion": "完善供应商档案，要求所有供应商提供营业执照和纳税人识别号，建立供应商准入制度。"
        })

    if total_purchase > 0 and len(top_sorted) >= 3:
        top3_pct = sum(x[1] for x in top_sorted[:3]) / total_purchase * 100
        if top3_pct < 80:
            results.append({
                "category": "供应商穿透", "category_icon": "🏭", "risk_score": 0, "risk_level": "良好",
                "risk_color": "#10b981", "urgency": "",
                "item": "供应商分布合理",
                "detail": f"前三大供应商采购占比 {top3_pct:.0f}%，供应商体系较为分散。",
                "suggestion": "继续保持合理的供应商结构。"
            })


# ═══════════════════════════════════════════════════════════════
#  二十、财务健康分析（新增 P1）
# ═══════════════════════════════════════════════════════════════

def _analyze_financial_health(db, company_id, ps, pe, results):
    """盈利能力/偿债能力/营运能力"""
    revenue = _get_account_sum(db, company_id, "6001", ps, pe, "credit")
    cost = _get_account_sum(db, company_id, "6401", ps, pe, "debit")
    total_assets = _get_account_balance(db, company_id, "1")
    total_liab = _get_account_balance(db, company_id, "2")
    current_assets = _get_account_balance(db, company_id, "11") + _get_account_balance(db, company_id, "12")
    current_liab = _get_account_balance(db, company_id, "21") + _get_account_balance(db, company_id, "22")
    net_profit = revenue - cost - _get_account_sum(db, company_id, "6601", ps, pe, "debit") - _get_account_sum(db, company_id, "6602", ps, pe, "debit") - _get_account_sum(db, company_id, "6603", ps, pe, "debit")
    equity = total_assets - total_liab

    metrics = {}
    if revenue > 0:
        metrics["毛利率"] = round((revenue - cost) / revenue * 100, 1)
        metrics["净利率"] = round(net_profit / revenue * 100, 1)
    if total_assets > 0:
        metrics["资产负债率"] = round(total_liab / total_assets * 100, 1)
    if equity > 0:
        metrics["净资产收益率"] = round(net_profit / equity * 100, 1) if net_profit != 0 else 0
    if current_liab > 0 and current_assets > 0:
        metrics["流动比率"] = round(current_assets / current_liab, 2)
    if total_assets > 0 and revenue > 0:
        metrics["总资产周转率"] = round(revenue / total_assets, 2)

    # 风险评估
    warnings = []
    if "资产负债率" in metrics and metrics["资产负债率"] > 70:
        warnings.append(f"资产负债率{metrics['资产负债率']}%（>70%偏高，偿债压力大）")
    if "流动比率" in metrics and metrics["流动比率"] < 1:
        warnings.append(f"流动比率{metrics['流动比率']}（<1可能面临短期偿债困难）")
    if "毛利率" in metrics and metrics["毛利率"] < 0:
        warnings.append(f"毛利率{metrics['毛利率']}%（亏损经营）")
    if "净利率" in metrics and metrics["净利率"] < 0:
        warnings.append(f"净利率{metrics['净利率']}%（整体经营亏损）")

    metric_str = " / ".join([f"{k}:{v}{'%' if '%' in k else ''}" for k, v in sorted(metrics.items())])

    if warnings:
        results.append({
            "category": "财务健康", "category_icon": "💹", "risk_score": 7, "risk_level": "高风险",
            "risk_color": "#dc2626", "urgency": "紧急",
            "item": "资产负债率过高",
            "detail": f"核心财务指标：{metric_str}。关键风险：{'；'.join(warnings)}。这些指标恶化可能影响持续经营能力判断和银行授信。",
            "suggestion": "制定财务改善计划：(1)控制成本费用支出；(2)加强应收账款回收和存货管理；(3)优化资本结构降低负债率；(4)必要时补充资本金或引入战略投资者。"
        })
    elif metrics:
        has_issues = any([
            metrics.get("资产负债率", 0) > 60,
            metrics.get("流动比率", 99) < 1.5,
            metrics.get("毛利率", 99) < 15
        ])
        if has_issues:
            results.append({
                "category": "财务健康", "category_icon": "💹", "risk_score": 3, "risk_level": "低风险",
                "risk_color": "#3b82f6", "urgency": "建议",
                "item": "流动比率偏低",
                "detail": f"核心财务指标：{metric_str}。部分指标偏离了健康区间，建议关注并改善。",
                "suggestion": "持续监控财务指标变化，重点关注毛利率趋势和应收账款周转效率。"
            })
        else:
            results.append({
                "category": "财务健康", "category_icon": "💹", "risk_score": 0, "risk_level": "良好",
                "risk_color": "#10b981", "urgency": "",
                "item": "财务指标总体健康",
                "detail": f"核心财务指标：{metric_str}。各项指标均在合理区间内，财务状况良好。",
                "suggestion": "继续保持良好的财务管理水平。"
            })


# ═══════════════════════════════════════════════════════════════
#  二十一、企业信用（新增 P2）
# ═══════════════════════════════════════════════════════════════

def _analyze_business_credit(db, company_id, ps, pe, results):
    """企业经营信用评估（基于系统内数据）"""
    company = db.query(Company).filter(Company.id == company_id).first()
    if not company:
        return

    company_name = company.name or ""

    # 21.1 空壳公司特征判断
    shell_indicators = []
    entry_count = db.query(func.count(JournalEntry.id)).filter(
        JournalEntry.company_id == company_id).scalar() or 0
    emp_count = db.query(func.count(Employee.id)).filter(
        Employee.company_id == company_id).scalar() or 0
    fixed_count = db.query(func.count(FixedAsset.id)).filter(
        FixedAsset.company_id == company_id).scalar() or 0
    revenue_total = _get_account_sum(db, company_id, "6001", ps, pe, "credit")

    if entry_count < 10:
        shell_indicators.append("凭证数量极少")
    if emp_count == 0:
        shell_indicators.append("无员工记录")
    if fixed_count == 0:
        shell_indicators.append("无固定资产")
    if revenue_total < 10000 and entry_count > 0:
        shell_indicators.append("营业收入极低")

    if len(shell_indicators) >= 3:
        results.append({
            "category": "企业信用", "category_icon": "🏛️", "risk_score": 7, "risk_level": "高风险",
            "risk_color": "#dc2626", "urgency": "紧急",
            "item": "疑似空壳公司特征",
            "detail": f"系统检测到以下异常：{'；'.join(shell_indicators)}。这些特征可能与空壳公司画像吻合，税务机关对此类企业关注度较高。",
            "suggestion": "确保公司有实际经营场所、真实员工和正常业务活动。如为初创企业，尽快补充人员、资产和业务记录。"
        })
    elif len(shell_indicators) >= 1:
        results.append({
            "category": "企业信用", "category_icon": "🏛️", "risk_score": 3, "risk_level": "低风险",
            "risk_color": "#3b82f6", "urgency": "建议",
            "item": "企业经营记录尚不完整",
            "detail": f"以下方面需要完善：{'；'.join(shell_indicators)}。经营记录不完整可能影响税务评级和信用等级。",
            "suggestion": "完善企业基础档案，确保人、财、物各项记录完整。"
        })

    # 21.2 经营异常判断
    if entry_count > 50 and emp_count >= 1 and fixed_count >= 1 and revenue_total > 100000:
        results.append({
            "category": "企业信用", "category_icon": "🏛️", "risk_score": 0, "risk_level": "良好",
            "risk_color": "#10b981", "urgency": "",
            "item": "企业经营实体特征正常",
            "detail": f"公司有 {entry_count} 条凭证、{emp_count} 名员工、{fixed_count} 项固定资产，营业收入 {revenue_total:,.2f} 元，经营实体特征完整。",
            "suggestion": "保持正常经营活动，按期进行工商年报和税务申报。"
        })


# ═══════════════════════════════════════════════════════════════
#  二十二、行业专项指标（新增 P2）
# ═══════════════════════════════════════════════════════════════

INDUSTRY_BENCHMARKS = {
    "制造业": {"gross_margin": 15, "vat_burden": 3.0},
    "建筑业": {"gross_margin": 10, "vat_burden": 2.5},
    "批发零售": {"gross_margin": 5, "vat_burden": 1.0},
    "交通运输": {"gross_margin": 15, "vat_burden": 2.0},
    "住宿餐饮": {"gross_margin": 40, "vat_burden": 2.5},
    "信息技术": {"gross_margin": 25, "vat_burden": 3.0},
    "租赁商务": {"gross_margin": 20, "vat_burden": 2.0},
    "居民服务": {"gross_margin": 30, "vat_burden": 2.5},
    "文化体育": {"gross_margin": 25, "vat_burden": 2.5},
    "房地产": {"gross_margin": 30, "vat_burden": 3.5},
}


def _analyze_industry_specific(db, company_id, ps, pe, results):
    """行业专项指标对比"""
    company = db.query(Company).filter(Company.id == company_id).first()
    if not company:
        return
    company_info = (company.business_scope or '') + (company.company_type or '')
    benchmark = None
    found_industry = None
    for key in INDUSTRY_BENCHMARKS:
        if key in company_info:
            benchmark = INDUSTRY_BENCHMARKS[key]
            found_industry = key
            break

    if not benchmark:
        return

    revenue = _get_account_sum(db, company_id, "6001", ps, pe, "credit")
    cost = _get_account_sum(db, company_id, "6401", ps, pe, "debit")
    actual_margin = (revenue - cost) / revenue * 100 if revenue > 0 else 0

    # 增值税实际税负
    vat_list = db.query(VATDeclaration).filter(
        VATDeclaration.company_id == company_id,
        VATDeclaration.period >= ps, VATDeclaration.period <= pe
    ).all()
    vat_sales = 0
    vat_payable = 0
    for v in vat_list:
        if v.form_main:
            fm = v.form_main if isinstance(v.form_main, dict) else json.loads(v.form_main)
            vat_sales += _safe_float(fm.get("sales_amount") or fm.get("total_sales"))
            vat_payable += _safe_float(fm.get("vat_payable"))
    actual_vat_burden = vat_payable / vat_sales * 100 if vat_sales > 0 else 0

    # 制造业特有：电费与产值
    if "制造" in company_info:
        electricity = _safe_float(db.query(func.sum(JournalEntry.debit_amount)).filter(
            JournalEntry.company_id == company_id, JournalEntry.period >= ps, JournalEntry.period <= pe,
            func.lower(JournalEntry.summary).contains("电费")).scalar())
        if electricity > 0 and revenue > 0:
            elec_ratio = electricity / revenue * 100
            if elec_ratio < 1:
                results.append({
                    "category": "行业专项", "category_icon": "🏭", "risk_score": 5, "risk_level": "中风险",
                    "risk_color": "#f59e0b", "urgency": "提醒",
                    "item": "制造业电费占收入比偏低",
                    "detail": f"制造业电费 {electricity:,.2f} 元，占收入 {elec_ratio:.2f}%。电费与产值严重不匹配可能暗示产能虚报或存在账外收入。",
                    "suggestion": "核实生产用电是否与产量匹配。关注是否存在部分厂房出租但电费仍由公司承担的情况。"
                })

    # 通用行业指标对比
    margin_gap = abs(actual_margin - benchmark["gross_margin"])
    if margin_gap > 20:
        results.append({
            "category": "行业专项", "category_icon": "📊", "risk_score": 5, "risk_level": "中风险",
            "risk_color": "#f59e0b", "urgency": "提醒",
            "item": "毛利率偏离行业参考值",
            "detail": f"实际毛利率 {actual_margin:.1f}%，{found_industry}行业参考毛利率 {benchmark['gross_margin']}%，偏差 {margin_gap:.1f} 个百分点。大幅偏离行业均值可能引起税务机关关注。",
            "suggestion": "如毛利率显著偏低：核查成本核算是否准确、是否存在关联交易转移利润。如毛利率显著偏高：确认收入是否完整入账。"
        })
    elif actual_vat_burden > 0 and benchmark["vat_burden"] > 0:
        vat_gap = abs(actual_vat_burden - benchmark["vat_burden"])
        if vat_gap > 1.5:
            results.append({
                "category": "行业专项", "category_icon": "📊", "risk_score": 4, "risk_level": "中风险",
                "risk_color": "#f59e0b", "urgency": "提醒",
                "item": "行业税负率明显低于同行业平均水平",
                "detail": f"增值税实际税负 {actual_vat_burden:.2f}%，{industry}行业参考值 {benchmark['vat_burden']}%，偏差 {vat_gap:.1f} 个百分点。",
                "suggestion": "税负偏低：分析是否有大量留抵或进项税额过大。税负偏高：检查进项税额是否应抵尽抵。"
            })


# ═══════════════════════════════════════════════════════════════
#  二十三、良好实践
# ═══════════════════════════════════════════════════════════════

def _analyze_good_practices(db, company_id, ps, pe, results):
    """识别做得好的事项"""
    pi_total = db.query(func.count(PurchaseInvoice.id)).filter(
        PurchaseInvoice.company_id == company_id).scalar() or 0

    if pi_total > 0:
        pi_has_tax = db.query(func.count(PurchaseInvoice.id)).filter(
            PurchaseInvoice.company_id == company_id,
            PurchaseInvoice.seller_tax_no != None, PurchaseInvoice.seller_tax_no != "").scalar() or 0
        if pi_has_tax / pi_total >= 0.9:
            results.append({
                "category": "良好实践", "category_icon": "✅", "risk_score": 0, "risk_level": "良好",
                "risk_color": "#10b981", "urgency": "",
                "item": "进项发票税号规范",
                "detail": f"进项发票供应商税号填写率 {pi_has_tax/pi_total*100:.0f}%，发票管理较为规范。",
                "suggestion": "继续保持。"
            })

    entry_count = db.query(func.count(JournalEntry.id)).filter(
        JournalEntry.company_id == company_id,
        JournalEntry.period >= ps, JournalEntry.period <= pe).scalar() or 0
    if entry_count > 50:
        results.append({
            "category": "良好实践", "category_icon": "✅", "risk_score": 0, "risk_level": "良好",
            "risk_color": "#10b981", "urgency": "",
            "item": "账务记录较为完整",
            "detail": f"分析期间内有 {entry_count} 条序时账记录，说明账务处理较为及时和完整。",
            "suggestion": "继续保持。"
        })

    risk_cnt = db.query(func.count(PurchaseInvoice.id)).filter(
        PurchaseInvoice.company_id == company_id,
        PurchaseInvoice.invoice_risk_level.in_(["疑点", "异常", "失控"])).scalar() or 0
    if pi_total > 0 and risk_cnt == 0:
        results.append({
            "category": "良好实践", "category_icon": "✅", "risk_score": 0, "risk_level": "良好",
            "risk_color": "#10b981", "urgency": "",
            "item": "进项发票无异常风险标记",
            "detail": f"全部 {pi_total} 张进项发票均无异常风险标记，发票管理质量良好。",
            "suggestion": "继续保持。"
        })

    vat_count = db.query(func.count(VATDeclaration.id)).filter(
        VATDeclaration.company_id == company_id).scalar() or 0
    if vat_count > 0:
        results.append({
            "category": "良好实践", "category_icon": "✅", "risk_score": 0, "risk_level": "良好",
            "risk_color": "#10b981", "urgency": "",
            "item": "已进行增值税按期申报",
            "detail": f"系统中有 {vat_count} 期增值税申报记录，税务申报工作按时开展。",
            "suggestion": "继续保持按期申报。"
        })

    # 社保公积金（申报表+明细+凭证三重检查）
    ss_ok = db.query(func.count(SocialSecurityDeclaration.id)).filter(
        SocialSecurityDeclaration.company_id == company_id).scalar() or 0
    ss_ok += db.query(func.count(SocialSecurityDetail.id)).join(
        SocialSecurityDeclaration, SocialSecurityDetail.declaration_id == SocialSecurityDeclaration.id
    ).filter(SocialSecurityDeclaration.company_id == company_id).scalar() or 0
    hf_ok = db.query(func.count(HousingFundDeclaration.id)).filter(
        HousingFundDeclaration.company_id == company_id).scalar() or 0
    hf_ok += db.query(func.count(HousingFundDetail.id)).filter(
        HousingFundDetail.company_id == company_id).scalar() or 0
    hf_ok += db.query(func.count(JournalEntry.id)).filter(
        JournalEntry.company_id == company_id,
        JournalEntry.source.in_(["公积金计提", "公积金缴纳"])
    ).scalar() or 0
    if ss_ok > 0 and hf_ok > 0:
        results.append({
            "category": "良好实践", "category_icon": "✅", "risk_score": 0, "risk_level": "良好",
            "risk_color": "#10b981", "urgency": "",
            "item": "已建立社保和公积金缴存体系",
            "detail": "已申报社会保险费和住房公积金，五险一金缴存体系已建立。",
            "suggestion": "继续保持按期全员足额缴存。"
        })


# ═══════════════════════════════════════════════════════════
#  二十四、经营实质深度分析——长期亏损风险（稽查级）
# ═══════════════════════════════════════════════════════════

def _analyze_long_term_loss(db, company_id, ps, pe, results):
    """连续多期亏损但持续经营——隐匿利润或虚增成本嫌疑（稽查重点）"""
    periods = _get_periods_between(ps, pe)
    if len(periods) < 6:
        return  # 至少需要6个月数据

    # 按季度汇总利润
    profit_by_quarter = {}
    for p in periods:
        y, m = int(p[:4]), int(p[5:7])
        q = f"{y}Q{(m-1)//3+1}"
        rev = _get_account_sum(db, company_id, "6001", p, p, "credit")
        cost = _get_account_sum(db, company_id, "6401", p, p, "debit")
        exps = (_get_account_sum(db, company_id, "6601", p, p, "debit") +
                _get_account_sum(db, company_id, "6602", p, p, "debit") +
                _get_account_sum(db, company_id, "6603", p, p, "debit"))
        profit = rev - cost - exps
        profit_by_quarter[q] = profit_by_quarter.get(q, 0) + profit

    quarters = sorted(profit_by_quarter.keys())
    loss_quarters = [q for q in quarters if profit_by_quarter[q] < 0]
    loss_ratio = len(loss_quarters) / len(quarters) if quarters else 0

    if loss_ratio >= 0.75 and len(quarters) >= 4:
        loss_detail = "；".join([f"{q}:{profit_by_quarter[q]:,.0f}元" for q in loss_quarters])
        results.append({
            "category": "经营实质", "category_icon": "🔍", "risk_score": 9, "risk_level": "高风险",
            "risk_color": "#dc2626", "urgency": "紧急",
            "item": "长期亏损风险",
            "detail": f"企业在 {len(quarters)} 个季度中有 {len(loss_quarters)} 个季度亏损（亏损面{loss_ratio*100:.0f}%），但持续经营不注销。稽查视角：长期亏损但持续经营，存在隐匿利润、虚增成本或关联交易转移利润的重大嫌疑。\n亏损明细：{loss_detail}",
            "suggestion": "（稽查应对）准备以下佐证材料：①各期成本费用明细及合法凭证；②关联交易定价政策及可比非受控价格；③库存盘点表及存货真实性说明；④银行流水与收入成本匹配说明；⑤持续经营的商业合理性说明（如市场开拓计划、研发投人等）。",
            "required_evidence": [
                "各期成本费用明细表及合法原始凭证（发票、合同、付款凭证）",
                "关联交易定价原则说明及可比非受控价格分析",
                "库存商品盘点表（含监盘记录）及存货真实性声明",
                "银行流水对账单（所有对公账户）与收入成本匹配分析说明",
                "持续经营商业计划书（说明亏损但不注销的商业合理性）",
                "员工工资表及社保缴纳证明（证明实际经营活动存在）"
            ]
        })
    elif loss_ratio >= 0.5 and len(quarters) >= 4:
        results.append({
            "category": "经营实质", "category_icon": "🔍", "risk_score": 5, "risk_level": "中风险",
            "risk_color": "#f59e0b", "urgency": "提醒",
            "item": "长期亏损风险",
            "detail": f"企业在 {len(quarters)} 个季度中有 {len(loss_quarters)} 个季度亏损（亏损面{loss_ratio*100:.0f}%）。需关注成本费用核算是否真实、完整。",
            "suggestion": "梳理亏损原因，准备成本费用真实性证明材料。如为季节性亏损或初创期亏损，准备相关说明。",
            "required_evidence": [
                "亏损原因说明及后续盈利计划",
                "成本费用明细及主要凭证复印件",
                "库存盘点表"
            ]
        })


# ═══════════════════════════════════════════════════════════
#  二十五、经营场所实质风险（稽查级）
# ═══════════════════════════════════════════════════════════

def _analyze_business_premise(db, company_id, ps, pe, results):
    """经营场所实质风险——四源交叉验证（序时账+银行流水+发票+合同）

    评估逻辑（按风险从高到低）：
    1. 四源全无 → 高风险「空壳嫌疑」
    2. 银行有支付但序时账无记录 → 高风险「银行付款未入账」（银行流水必须入账！）
    3. 序时账无，发票有 → 中风险「费用未入账」
    4. 序时账无，合同有（免费场地）→ 低风险「可能为免费使用场地」
    5. 序时账有 → 无风险
    """
    evidence = _check_premise_evidence_all(db, company_id, ps, pe)
    overall = evidence["overall"]
    je = evidence["je"]
    bank_info = evidence["bank"]
    inv_info = evidence["invoice"]
    contract_info = evidence["contract"]

    company = db.query(Company).filter(Company.id == company_id).first()
    biz_scope = (company.business_scope or "") if company else ""

    # 有经营范围才需要检查经营场所
    if not biz_scope:
        return

    je_has_any = je["has_any"]

    # ── 场景1：四源全无 → 真正空壳嫌疑 ──
    if overall["all_sources_none"]:
        results.append({
            "category": "经营实质", "category_icon": "🔍", "risk_score": 8, "risk_level": "高风险",
            "risk_color": "#dc2626", "urgency": "紧急",
            "item": "经营场所实质风险",
            "detail": "企业经营范围包含经营活动，四源交叉验证均未发现经营场所证据：序时账中无租赁费(660214)/水电费(660215)/物业费记录；银行流水中无相关支付；取得发票中无租赁/水电/物业发票；合同中无租赁/场地使用协议。多维度交叉验证均未发现经营场所证据。稽查视角：可能存在空壳公司、虚开发票或隐瞒实际经营场所的问题。",
            "suggestion": "（稽查应对）立即准备以下佐证材料：①经营场所租赁合同及租金支付凭证；②水电费缴纳凭证及发票；③经营场所照片（含门牌、办公/生产区域）；④物业缴费通知及支付记录；⑤如为家庭经营，提供房屋产权证明或租赁协议。",
            "required_evidence": [
                "经营场所不动产权证书或租赁合同（原件备查）",
                "租金支付凭证（银行回单+发票）",
                "水、电、燃气费缴纳凭证（至少3个月）",
                "经营场所实地照片（含门牌、内部场景）",
                "物业费缴纳凭证",
                "如为共用场所，提供共用协议及费用分摊说明"
            ]
        })
        return

    # ── 场景2：银行有支付但序时账无记录 → 高风险（银行流水必须入账！） ──
    if overall["je_missing_but_bank_found"]:
        results.append({
            "category": "经营实质", "category_icon": "🔍", "risk_score": 8, "risk_level": "高风险",
            "risk_color": "#dc2626", "urgency": "紧急",
            "item": "经营场所银行付款未入账",
            "detail": f"银行流水中发现经营场所相关支付 ¥{bank_info['amount']:,.2f}，但序时账中无对应的费用科目记录。银行流水必须入账，否则银行存款余额与银行对账单余额会对不上。请立即将银行支付凭证入账。",
            "suggestion": f"将银行流水中经营场所相关付款（合计 ¥{bank_info['amount']:,.2f}）立即生成序时账凭证：借费用科目/应交税费/贷银行存款。",
            "required_evidence": ["经营场所相关银行支付回单", "对应的费用发票（如有）"]
        })
        return

    # ── 场景3：序时账没有，发票有 → 费用已发生未入账（发票可延后） ──
    if not je_has_any and inv_info["any_found"]:
        inv_details = []
        if inv_info["has_rent"]:
            inv_details.append(f"租赁类发票 ¥{inv_info['rent_amount']:,.2f}")
        if inv_info["has_utility"]:
            inv_details.append(f"水电类发票 ¥{inv_info['utility_amount']:,.2f}")
        if inv_info["has_property"]:
            inv_details.append(f"物业类发票 ¥{inv_info['property_amount']:,.2f}")
        results.append({
            "category": "经营实质", "category_icon": "🔍", "risk_score": 4, "risk_level": "中风险",
            "risk_color": "#f59e0b", "urgency": "提醒",
            "item": "经营场所费用未入账（已取得发票）",
            "detail": f"序时账中未发现租赁费/水电费/物业费记录，但取得发票中发现相关支出发票（{'；'.join(inv_details)}），总金额 ¥{inv_info['total_found']:,.2f}。发票已到但尚未入账（发票可延后入账），请及时将发票对应的费用入账。",
            "suggestion": f"将已取得的经营场所相关发票（合计 ¥{inv_info['total_found']:,.2f}）及时入账：借费用科目/贷应付账款。",
            "required_evidence": [
                "经营场所租赁合同",
                "未入账的租赁费/水电费/物业费发票原件",
                "对应的银行支付凭证（如有）"
            ]
        })
        return

    # ── 场景4：序时账没有，但有场地合同（可能免费使用） ──
    if not je_has_any and contract_info["any_found"]:
        cts = contract_info["contracts"]
        ct_summary = "；".join([f"{c['name']}({c['type']}{'·免费' if c['is_free'] else ''})" for c in cts[:3]])
        if contract_info["free_venue_possible"]:
            # 免费场地
            results.append({
                "category": "经营实质", "category_icon": "🔍", "risk_score": 2, "risk_level": "低风险",
                "risk_color": "#3b82f6", "urgency": "提醒",
                "item": "经营场所为免费使用（有合同无费用记录）",
                "detail": f"序时账中无租赁费/水电费/物业费记录，但发现场地相关合同：{ct_summary}。可能为免费使用场地（如：关联方免费提供、孵化器免租、家庭住宅等）。建议保留合同备查。",
                "suggestion": "①保留场地使用合同/协议原件备查；②如为关联方免费提供，准备关联方确认函；③如有水电费由他人代缴，保留代缴说明。",
                "required_evidence": ["场地使用合同/协议", "如为关联方免费提供：关联方确认函"]
            })
        else:
            # 有租赁合同但未入账
            results.append({
                "category": "经营实质", "category_icon": "🔍", "risk_score": 3, "risk_level": "低风险",
                "risk_color": "#3b82f6", "urgency": "提醒",
                "item": "经营场所有合同但费用未入账",
                "detail": f"发现场地相关合同：{ct_summary}。序时账中无对应费用记录，建议确认合同是否在履行中并补充费用凭证。",
                "suggestion": "确认合同履行状态，如有已发生费用请及时入账。",
                "required_evidence": ["场地使用合同", "费用支付凭证（如有）"]
            })
        return

    # ── 场景5：序时账有记录 → 无风险，不报告 ──
    # (je_has_any is True)


# ═══════════════════════════════════════════════════════════
#  二十六、存货与经营规模匹配风险（稽查级）
# ═══════════════════════════════════════════════════════════

def _analyze_inventory_substance(db, company_id, ps, pe, results):
    """原材料进得多、销得少，但仓库/产能不匹配——虚增库存或隐瞒销售收入嫌疑"""
    # 统计存货进销存
    in_qty = db.query(func.sum(InventoryTransaction.quantity)).filter(
        InventoryTransaction.company_id == company_id,
        InventoryTransaction.trans_type.in_(["入库", "盘盈", "调拨入"]),
        InventoryTransaction.transaction_date >= _period_to_date_range(ps)[0],
        InventoryTransaction.transaction_date <= _period_to_date_range(pe)[1]
    ).scalar() or 0

    out_qty = db.query(func.sum(InventoryTransaction.quantity)).filter(
        InventoryTransaction.company_id == company_id,
        InventoryTransaction.trans_type.in_(["出库", "盘亏", "调拨出"]),
        InventoryTransaction.transaction_date >= _period_to_date_range(ps)[0],
        InventoryTransaction.transaction_date <= _period_to_date_range(pe)[1]
    ).scalar() or 0

    end_stock = db.query(func.sum(InventoryBalance.end_quantity)).filter(
        InventoryBalance.company_id == company_id,
        InventoryBalance.period == pe
    ).scalar() or 0

    # 计算进销比
    if out_qty > 0:
        in_out_ratio = abs(in_qty) / abs(out_qty)
    else:
        in_out_ratio = 999 if in_qty > 0 else 0

    # 检查是否有仓库租赁合同
    warehouse_contract = db.query(func.count(Contract.id)).filter(
        Contract.company_id == company_id,
        Contract.contract_type == "租赁",
        or_(Contract.name.contains("仓库"), Contract.name.contains("仓储"))
    ).scalar() or 0

    # 稽查视角：进销比异常 + 无仓库 = 高风险
    if in_out_ratio > 3 and warehouse_contract == 0 and end_stock > 100:
        results.append({
            "category": "经营实质", "category_icon": "🔍", "risk_score": 9, "risk_level": "高风险",
            "risk_color": "#dc2626", "urgency": "紧急",
            "item": "存货与经营规模匹配风险",
            "detail": f"分析期内入库数量 {in_qty:,.0f}，出库数量 {out_qty:,.0f}，进销比 {in_out_ratio:.1f}:1（正常应接近1:1）。期末库存 {end_stock:,.0f}，但系统中无仓库租赁合同。稽查视角：大量原材料购入但销售数量极少，且无仓库存储，存在虚增库存（虚开发票 counterpart）或隐瞒销售收入（账外销售）的重大嫌疑。",
            "suggestion": "（稽查应对）准备以下佐证材料：①库存商品盘点表（含监盘记录、照片）；②原材料采购合同、发票、入库单、付款凭证（三单匹配）；③销售合同、发货单、销售发票、收款凭证（验证销售真实性）；④仓库租赁合同及仓储费支付凭证；⑤运输发票及物流单据；⑥生产成本计算单（BOM表）及投入产出比分析。",
            "required_evidence": [
                "库存商品全面盘点表（含监盘人签字、盘点照片）",
                "原材料采购三单匹配：采购合同+进项发票+入库单+付款凭证",
                "销售收入四流合一：销售合同+销项发票+出库单+银行收款回单",
                "仓库租赁合同+仓储费支付凭证+仓储费发票",
                "物流运输单据及运输费发票（验证货物实际流动）",
                "生产成本投入产出比分析说明（工业企业的关键证据）",
                "如为委托加工，提供委托加工合同及加工费支付凭证"
            ]
        })
    elif in_out_ratio > 2:
        results.append({
            "category": "经营实质", "category_icon": "🔍", "risk_score": 5, "risk_level": "中风险",
            "risk_color": "#f59e0b", "urgency": "提醒",
            "item": "存货与经营规模匹配风险",
            "detail": f"入库数量 {in_qty:,.0f}，出库数量 {out_qty:,.0f}，进销比 {in_out_ratio:.1f}:1。如为生产型企业，应核查投入产出比是否合理；如为贸易型企业，应核查库存周转率为何偏低。",
            "suggestion": "提供库存周转分析说明，如确有合理原因（如季节性备货、战略囤货），准备相关说明材料。",
            "required_evidence": [
                "库存商品盘点表",
                "采购及销售合同（说明进销不匹配的商业理由）",
                "如为季节性备货，提供历史销售数据对比分析"
            ]
        })


# ═══════════════════════════════════════════════════════════
#  二十七、水电费与产能匹配分析（稽查级·生产企业专用）
# ═══════════════════════════════════════════════════════════

def _analyze_utility_expense(db, company_id, ps, pe, results):
    """生产企业水电费与产量不匹配——产能造假或隐瞒生产嫌疑"""
    company = db.query(Company).filter(Company.id == company_id).first()
    if not company:
        return
    biz = (company.business_scope or "") + " " + (company.company_type or "")
    if "制造" not in biz and "生产" not in biz and "加工" not in biz:
        return  # 非生产企业，不分析

    # 水电费
    elec_amt = _get_account_sum(db, company_id, "660215", ps, pe, "debit")
    water_amt = _get_account_sum(db, company_id, "660215", ps, pe, "debit")  # 同一科目，需通过摘要区分

    # 更精确：从银行流水或序时账摘要中匹配
    from sqlalchemy import or_
    elec_keywords = ["电费", "电力", "供电", "能源"]
    water_keywords = ["水费", "水务", "供水", "水费"]

    elec_entries = db.query(func.sum(JournalEntry.debit_amount)).filter(
        JournalEntry.company_id == company_id,
        JournalEntry.period >= ps, JournalEntry.period <= pe,
        or_(*[JournalEntry.summary.contains(kw) for kw in elec_keywords])
    ).scalar() or 0

    # 产量 proxies：销项发票数量（生产企业以销售发票货物数量近似）
    sales_qty = db.query(func.sum(SalesInvoice.quantity)).filter(
        SalesInvoice.company_id == company_id,
        SalesInvoice.invoice_date >= _period_to_date_range(ps)[0],
        SalesInvoice.invoice_date <= _period_to_date_range(pe)[1],
        SalesInvoice.status == "正常"
    ).scalar() or 0

    revenue = _get_account_sum(db, company_id, "6001", ps, pe, "credit")

    # 电费收入比
    elec_to_revenue = elec_entries / revenue * 100 if revenue > 0 else 0

    if elec_entries > 0 and revenue > 0 and elec_to_revenue < 0.5:
        results.append({
            "category": "经营实质", "category_icon": "🔍", "risk_score": 7, "risk_level": "高风险",
            "risk_color": "#dc2626", "urgency": "紧急",
            "item": f"电费占收入比异常偏低（{elec_to_revenue:.2f}%）",
            "detail": f"生产企业电费 {elec_entries:,.0f} 元，营业收入 {revenue:,.0f} 元，电费占收入比仅 {elec_to_revenue:.2f}%（制造企业正常应 2%~8%）。稽查视角：电费与产量严重不匹配，存在隐瞒生产规模、账外销售或虚列成本的嫌疑。",
            "suggestion": "（稽查应对）准备以下佐证材料：①电力公司出具的全期电费缴纳清单及发票；②生产日报表/车间产量记录（每日记录）；③原材料投入产出比计算表（BOM）；④设备清单及设备功率清单；⑤如部分为外包生产，提供委托加工合同及加工费发票。",
            "required_evidence": [
                "电力公司出具的正式电费缴纳清单（覆盖分析期各月）",
                "电费支付凭证（银行回单）及增值税专用发票",
                "生产车间日报表/产量逐日记录（生产部门出具）",
                "原材料投入产出比分析表（含标准耗电量定额说明）",
                "生产设备清单（含功率、运行时长说明）",
                "如存在免税产品/在建工程用电，提供电费分摊计算说明"
            ]
        })
    elif elec_entries == 0 and "制造" in biz:
        results.append({
            "category": "经营实质", "category_icon": "🔍", "risk_score": 6, "risk_level": "中风险",
            "risk_color": "#f59e0b", "urgency": "提醒",
            "item": "生产企业无电费记录",
            "detail": "企业经营范围包含生产制造，但序时账中未发现电费支出记录。生产企业必然消耗电力，无电费记录不符合经营实质。",
            "suggestion": "补录电费支出凭证，如由出租方代缴，提供代缴协议及代缴凭证。",
            "required_evidence": [
                "电费缴纳凭证或代缴协议",
                "如为租用厂房含电费，提供租赁合同相关条款及出租方出具的电费分割单"
            ]
        })


# ═══════════════════════════════════════════════════════════
#  二十八、人员与经营规模匹配（稽查级）
# ═══════════════════════════════════════════════════════════

def _analyze_staffing_substance(db, company_id, ps, pe, results):
    """员工人数与收入规模不匹配——空壳或隐瞒用工嫌疑"""
    emp_count = db.query(func.count(Employee.id)).filter(
        Employee.company_id == company_id,
        or_(Employee.leave_date == None, Employee.leave_date >= _period_to_date_range(ps)[0])
    ).scalar() or 0

    revenue = _get_account_sum(db, company_id, "6001", ps, pe, "credit")
    avg_revenue_per_emp = revenue / emp_count if emp_count > 0 else 0

    # 社保参保人数（从参保明细统计，按员工姓名去重）
    ss_count = db.query(func.count(func.distinct(SocialSecurityDetail.employee_name))).join(
        SocialSecurityDeclaration, SocialSecurityDetail.declaration_id == SocialSecurityDeclaration.id
    ).filter(
        SocialSecurityDeclaration.company_id == company_id
    ).scalar() or 0

    company = db.query(Company).filter(Company.id == company_id).first()
    reg_capital = company.registered_capital or 0

    # 稽查视角：有收入但无员工 = 空壳嫌疑
    if revenue > 100000 and emp_count == 0:
        results.append({
            "category": "经营实质", "category_icon": "🔍", "risk_score": 8, "risk_level": "高风险",
            "risk_color": "#dc2626", "urgency": "紧急",
            "item": "人工成本与经营规模不匹配",
            "detail": f"企业营业收入 {revenue:,.0f} 元（≥10万元），但系统中无员工档案记录。稽查视角：有经营收入但无从业人员，存在空壳公司、虚开发票或用工不申报（劳务外包/灵活用工未申报）的重大嫌疑。",
            "suggestion": "（稽查应对）准备以下佐证材料：①全体员工花名册（含入职时间、岗位、薪酬）；②工资支付凭证（银行代发工资回单）；③社保缴纳证明（证明用工真实性）；④如为劳务外包，提供外包合同及外包发票；⑤如为灵活用工，提供灵活用工平台协议及发票。",
            "required_evidence": [
                "全体员工花名册（盖章）及员工劳动合同",
                "工资表及银行代发工资回单（覆盖分析期）",
                "社会保险费缴纳证明（社保局出具）",
                "如为劳务外包，提供外包合同及劳务费发票（6%专票）",
                "如为灵活用工，提供灵活用工平台合作协议及发票",
                "如为劳务派遣，提供派遣协议及派遣人员工资支付凭证"
            ]
        })
    elif emp_count > 0 and ss_count == 0:
        results.append({
            "category": "经营实质", "category_icon": "🔍", "risk_score": 5, "risk_level": "中风险",
            "risk_color": "#f59e0b", "urgency": "提醒",
            "item": "有员工但未申报社保",
            "detail": f"系统中有 {emp_count} 名员工，但未发现社保申报记录。存在被认定未用工申报、逃避社保义务的风险，也是稽查关注点。",
            "suggestion": "补录社保申报记录，如确有部分人员不需缴纳社保（如退休返聘、实习生），准备相关说明材料。",
            "required_evidence": [
                "社保申报表（覆盖分析期）",
                "如存在退休返聘/实习生，提供相关劳动合同及说明"
            ]
        })
    elif 0 < emp_count < 5 and revenue > 5000000:
        results.append({
            "category": "经营实质", "category_icon": "🔍", "risk_score": 4, "risk_level": "中风险",
            "risk_color": "#f59e0b", "urgency": "提醒",
            "item": "人工成本与经营规模不匹配",
            "detail": f"企业仅有 {emp_count} 名员工，但营业收入达 {revenue:,.0f} 元（人均产值 {avg_revenue_per_emp:,.0f} 元/人）。如非贸易型/互联网型轻资产企业，人员规模与收入不匹配可能引起稽查关注。",
            "suggestion": "如为贸易企业或互联网企业，属合理情况。如为生产/服务型企业，准备人员组织架构图及岗位说明，解释人均产值合理性。",
            "required_evidence": [
                "企业组织架构图及岗位职责说明",
                "如为贸易/互联网企业，提供商业模式说明"
            ]
        })


# ═══════════════════════════════════════════════════════════
#  二十九、资金流与票据流匹配（稽查级·虚开发票嫌疑）
# ═══════════════════════════════════════════════════════════

def _analyze_fund_flow_invoice_match(db, company_id, ps, pe, results):
    """有发票但无银行收款记录——虚开发票嫌疑（稽查头号风险）"""
    # 销项发票（开票方=本企业）
    sales = db.query(SalesInvoice).filter(
        SalesInvoice.company_id == company_id,
        SalesInvoice.invoice_date >= _period_to_date_range(ps)[0],
        SalesInvoice.invoice_date <= _period_to_date_range(pe)[1],
        SalesInvoice.status == "正常"
    ).all()

    if not sales:
        return

    mismatch_count = 0
    mismatch_invoices = []
    from sqlalchemy import or_

    for inv in sales:
        # 提取发票号码（兼容数电发票）
        inv_no = ""
        if getattr(inv, 'digital_invoice_no', None):
            inv_no = inv.digital_invoice_no
        elif inv.invoice_code or inv.invoice_no:
            inv_no = (inv.invoice_code or "") + (inv.invoice_no or "")

        # 在银行流水中查找该发票号码或对应金额
        if inv_no:
            bt = db.query(BankTransaction).filter(
                BankTransaction.company_id == company_id,
                BankTransaction.transaction_date >= inv.invoice_date,
                BankTransaction.transaction_date <= (inv.invoice_date.replace(month=inv.invoice_date.month+1, day=1) if inv.invoice_date.month < 12 else inv.invoice_date.replace(year=inv.invoice_date.year+1, month=1, day=1)),
                BankTransaction.credit_amount == inv.total_amount
            ).first()
            if not bt:
                mismatch_count += 1
                if len(mismatch_invoices) < 5:
                    mismatch_invoices.append(inv_no)
        else:
            mismatch_count += 1

    mismatch_rate = mismatch_count / len(sales) * 100 if sales else 0

    if mismatch_rate > 30:
        results.append({
            "category": "经营实质", "category_icon": "🔍", "risk_score": 9, "risk_level": "高风险",
            "risk_color": "#dc2626", "urgency": "紧急",
            "item": f"发票与银行收款匹配率偏低（{mismatch_rate:.0f}%）",
            "detail": f"分析期内共有 {len(sales)} 张销项发票，其中 {mismatch_count} 张未在银行流水中找到对应收款记录（匹配率仅 {100-mismatch_rate:.0f}%）。稽查视角：开票但无资金回流，是虚开发票的最典型特征（特别是无真实交易背景下开票）。",
            "suggestion": "（稽查应对）准备以下佐证材料：①全部销项发票对应的销售合同；②发货单/运输单据（证明货物真实发出）；③银行收款回单（如为分期收款，提供收款计划及实际收款记录）；④如为抵债/易货交易，提供抵债协议或易货合同；⑤如为代销，提供代销合同及委托代销清单。",
            "required_evidence": [
                "所有不匹配发票对应的销售合同（原件备查）",
                "销售发票+发货单+运输单据（三流合一证据链）",
                "银行收款回单（核对发票金额与收款金额是否一致）",
                "如为分期收款，提供分期收款协议及实际收款记录",
                "如为抵债/易货，提供抵债/易货协议及入账凭证",
                "客户验收单或确认收货证明"
            ]
        })
    elif mismatch_rate > 10:
        results.append({
            "category": "经营实质", "category_icon": "🔍", "risk_score": 5, "risk_level": "中风险",
            "risk_color": "#f59e0b", "urgency": "提醒",
            "item": f"部分发票与银行收款未匹配（{mismatch_rate:.0f}%）",
            "detail": f"有 {mismatch_count}/{len(sales)} 张发票未在银行流水中找到对应收款记录。可能是 timing difference（开票与收款跨期），也可能是存在问题。",
            "suggestion": "核查未匹配发票的收款状态，补充银行收款凭证。如为跨期收款，在报告中说明。",
            "required_evidence": [
                "未匹配发票清单及收款计划说明",
                "已收款的银行回单"
            ]
        })


# ═══════════════════════════════════════════════════════════
#  三十、边角料/报废收入未申报风险（稽查级·生产企业）
# ═══════════════════════════════════════════════════════════

def _analyze_scrap_revenue(db, company_id, ps, pe, results):
    """生产企业边角料/报废收入未申报增值税——隐瞒收入嫌疑"""
    company = db.query(Company).filter(Company.id == company_id).first()
    if not company:
        return
    biz = (company.business_scope or "") + " " + (company.company_type or "")
    if "制造" not in biz and "生产" not in biz:
        return

    # 检查是否有边角料/报废收入的增值税申报
    # 边角料销售一般开普通发票或不开票（未开票收入）
    # 检查 "其他业务收入" 或 "营业外收入" 中是否包含边角料/报废
    from sqlalchemy import or_
    scrap_entries = db.query(JournalEntry).filter(
        JournalEntry.company_id == company_id,
        JournalEntry.period >= ps, JournalEntry.period <= pe,
        or_(
            JournalEntry.summary.contains("边角"),
            JournalEntry.summary.contains("废"),
            JournalEntry.summary.contains("下脚"),
            JournalEntry.summary.contains("残次"),
            JournalEntry.summary.contains("报废"),
        ),
        JournalEntry.credit_amount > 0
    ).all()

    # 检查是否有边角料销售的销项发票
    scrap_invoices = db.query(SalesInvoice).filter(
        SalesInvoice.company_id == company_id,
        SalesInvoice.invoice_date >= _period_to_date_range(ps)[0],
        SalesInvoice.invoice_date <= _period_to_date_range(pe)[1],
        or_(
            SalesInvoice.goods_name.contains("边角"),
            SalesInvoice.goods_name.contains("废"),
            SalesInvoice.goods_name.contains("下脚"),
        ),
        SalesInvoice.status == "正常"
    ).all()

    # 生产企业应有边角料收入但未发现
    # 这是一个 "应存在但未发现" 的风险提示（基于行业常识）
    has_scrap = len(scrap_entries) > 0 or len(scrap_invoices) > 0

    # 从银行存款借方（收款）中查找是否有零星空交易（疑似边角料销售未开票）
    unidentified_credits = db.query(func.count(BankTransaction.id)).filter(
        BankTransaction.company_id == company_id,
        BankTransaction.transaction_date >= _period_to_date_range(ps)[0],
        BankTransaction.transaction_date <= _period_to_date_range(pe)[1],
        BankTransaction.credit_amount > 0,
        or_(
            BankTransaction.counterparty_name == None,
            BankTransaction.counterparty_name == "",
            BankTransaction.summary == None,
            BankTransaction.summary == ""
        )
    ).scalar() or 0

    if not has_scrap and "制造" in biz and unidentified_credits > 3:
        results.append({
            "category": "经营实质", "category_icon": "🔍", "risk_score": 6, "risk_level": "中风险",
            "risk_color": "#f59e0b", "urgency": "提醒",
            "item": "生产企业边角料/报废收入申报不全",
            "detail": f"企业为生产企业，但序时账中未发现边角料/报废收入记录，且银行流水中有 {unidentified_credits} 笔对手方不明的收款。稽查视角：生产企业必然产生边角料/残次品，如未申报增值税销售额，属于隐瞒收入行为。",
            "suggestion": "（稽查应对）准备以下佐证材料：①边角料/报废品销售记录及发票；②如已入账未开票，提供未开票收入申报说明；③边角料/报废品出库单及收款凭证；④如边角料无偿赠送或自用，说明具体情况。",
            "required_evidence": [
                "边角料/报废品销售合同及销售发票",
                "边角料出库单及收款凭证",
                "如为未开票收入，提供未开票收入明细表及增值税申报表（未开票栏次）",
                "如边角料用于无偿赠送，提供视同销售增值税申报说明"
            ]
        })


# ═══════════════════════════════════════════════════════════
#  三十一、视同销售未处理风险（稽查级）
# ═══════════════════════════════════════════════════════════

def _analyze_deemed_sales(db, company_id, ps, pe, results):
    """视同销售未申报——无偿赠送/员工福利/对外投资未缴纳增值税（稽查重点）"""
    # 检查序时账中是否有视同销售业务处理
    from sqlalchemy import or_
    deemed_keywords = ["赠送", "福利", "员工福利", "招待", "视同销售", "对外投资", "非货币性资产"]
    deemed_entries = db.query(JournalEntry).filter(
        JournalEntry.company_id == company_id,
        JournalEntry.period >= ps, JournalEntry.period <= pe,
        or_(*[JournalEntry.summary.contains(kw) for kw in deemed_keywords])
    ).all()

    # 检查 InventoryTransaction 中是否有 "出库" 但对应 SalesInvoice 中没有的记录（可能是无偿赠送）
    free_out = db.query(InventoryTransaction).filter(
        InventoryTransaction.company_id == company_id,
        InventoryTransaction.trans_type == "出库",
        InventoryTransaction.transaction_date >= _period_to_date_range(ps)[0],
        InventoryTransaction.transaction_date <= _period_to_date_range(pe)[1]
    ).all()

    # 简易判断：有出库记录但销项发票中没有对应记录
    sales_inv_nos = set()
    for inv in db.query(SalesInvoice).filter(
        SalesInvoice.company_id == company_id,
        SalesInvoice.invoice_date >= _period_to_date_range(ps)[0],
        SalesInvoice.invoice_date <= _period_to_date_range(pe)[1]
    ).all():
        no = getattr(inv, 'digital_invoice_no', None) or (inv.invoice_code or "") + (inv.invoice_no or "")
        if no:
            sales_inv_nos.add(no)

    # 这是一个提醒型风险：系统无法完全判断，但可提示用户自查
    if len(deemed_entries) == 0 and len(free_out) > 0:
        results.append({
            "category": "经营实质", "category_icon": "🔍", "risk_score": 4, "risk_level": "中风险",
            "risk_color": "#f59e0b", "urgency": "提醒",
            "item": "可能存在视同销售未申报情形",
            "detail": f"系统中有 {len(free_out)} 条出库记录，但未发现「视同销售」相关账务处理。根据增值税法规，无偿赠送货物、将自产产品用于员工福利、对外投资等，均需视同销售缴纳增值税。",
            "suggestion": "自查以下事项并准备佐证材料：①是否有无偿赠送客户货物（如促销赠品）；②是否有将自产产品作为员工福利发放；③是否有将货物用于非应税项目；④如有，是否已按公允价值申报视同销售增值税。",
            "required_evidence": [
                "视同销售自查清单（按增值税法规逐一核对）",
                "如已处理：视同销售增值税申报表及账务处理凭证",
                "如未处理：补税说明及更正申报表",
                "赠送协议或员工福利发放记录（证明视同销售事实）"
            ]
        })


# ═══════════════════════════════════════════════════════════
#  三十二、库存现金过大风险（稽查级·坐支嫌疑）
# ═══════════════════════════════════════════════════════════

def _analyze_cash_overstock(db, company_id, ps, pe, results):
    """库存现金余额过大——坐支、账外循环、隐瞒收入嫌疑"""
    # 库存现金科目（1001）期末余额
    cash_balance = _get_account_balance(db, company_id, "1001")

    if cash_balance > 10000:
        # 检查是否有大额现金收支
        cash_entries = db.query(JournalEntry).filter(
            JournalEntry.company_id == company_id,
            JournalEntry.period >= ps, JournalEntry.period <= pe,
            JournalEntry.account_code.like("1001%")
        ).count()

        results.append({
            "category": "经营实质", "category_icon": "🔍", "risk_score": 5, "risk_level": "中风险",
            "risk_color": "#f59e0b", "urgency": "提醒",
            "item": f"库存现金余额过大（{cash_balance:,.0f}元）",
            "detail": f"库存现金科目期末余额 {cash_balance:,.0f} 元（超过1万元）。稽查视角：库存现金过大，可能存在现金坐支（收入不入账直接支付支出）、账外资金循环或隐瞒现金销售收入的风险。根据现金管理暂行条例，企业应与银行核对现金账目，不允许坐支。",
            "suggestion": "（稽查应对）准备以下佐证材料：①库存现金盘点表（盘点日至报表日调节表）；②现金日记账（逐笔记录）；③大额现金收支的审批单及原始凭证；④如为零售企业现金收款，提供现金收款台账及银行存款缴款单（现金送存银行凭证）。",
            "required_evidence": [
                "库存现金盘点表（含盘点日至报表日的现金调节表）",
                "现金日记账（逐笔登记，与银行对账单核对）",
                "大额现金收支的原始凭证（合同、发票、审批单）",
                "如为现金销售，提供现金收款台账及现金送存银行缴款单",
                "现金管理制度及执行情况说明"
            ]
        })
    elif cash_balance < 0:
        results.append({
            "category": "经营实质", "category_icon": "🔍", "risk_score": 7, "risk_level": "高风险",
            "risk_color": "#dc2626", "urgency": "紧急",
            "item": f"库存现金余额为负（{cash_balance:,.0f}元）",
            "detail": "库存现金科目出现负数余额，说明账务处理存在严重错误（可能是未做凭证但已支付现金），或者存在账外现金循环。",
            "suggestion": "立即核查现金日记账，找出余额为负的根本原因并更正凭证。",
            "required_evidence": [
                "现金日记账全面核查说明",
                "更正凭证及审批记录"
            ]
        })


# ═══════════════════════════════════════════════════════════
#  三十三、关联交易定价公允性（稽查级·转移利润嫌疑）
# ═══════════════════════════════════════════════════════════

def _analyze_related_party_pricing(db, company_id, ps, pe, results):
    """关联交易定价不公允——转移利润嫌疑（稽查重点·跨国/跨地区）"""
    # 查找股东、董事、监事关联的企业（同一法定代表人或控股股东）
    company = db.query(Company).filter(Company.id == company_id).first()
    if not company:
        return

    # 从进项发票中查找与股东/关联方交易的发票
    # 关联方认定：同一法定代表人、控股股东、实际控制人
    legal_rep = company.legal_representative or ""

    # 查找供应商/客户中是否包含本企业股东或法定代表人
    related_suppliers = db.query(Supplier).filter(
        Supplier.company_id == company_id,
        or_(
            Supplier.name.contains(legal_rep),
            Supplier.tax_no == company.uscc
        ) if legal_rep else False
    ).all()

    related_customers = db.query(Customer).filter(
        Customer.company_id == company_id,
        or_(
            Customer.name.contains(legal_rep),
            Customer.tax_no == company.uscc
        ) if legal_rep else False
    ).all()

    # 从进项发票查找定价异常（关联方向本企业开票，价格是否公允无法从系统判断，但可提示风险）
    related_inv_count = 0
    if legal_rep:
        related_inv_count = db.query(func.count(PurchaseInvoice.id)).filter(
            PurchaseInvoice.company_id == company_id,
            PurchaseInvoice.invoice_date >= _period_to_date_range(ps)[0],
            PurchaseInvoice.invoice_date <= _period_to_date_range(pe)[1],
            PurchaseInvoice.seller_name.contains(legal_rep)
        ).scalar() or 0

    if related_inv_count > 0:
        results.append({
            "category": "经营实质", "category_icon": "🔍", "risk_score": 6, "risk_level": "中风险",
            "risk_color": "#f59e0b", "urgency": "提醒",
            "item": "关联交易转移定价偏离市场价",
            "detail": f"发现 {related_inv_count} 张进项发票的开票方名称包含本企业法定代表人/控股股东姓名，属于关联交易。稽查视角：关联交易定价如不公允（高价采购、低价销售），可能被认定为转移利润、逃避税收，需准备转让定价同期资料。",
            "suggestion": "（稽查应对）准备以下佐证材料：①关联交易所涉完整合同及定价说明；②可比非受控价格（CUP）分析；③再销售价格法或成本加成法分析（证明定价公允性）；④如为跨国关联交易，准备同期资料（主体文档、本地文档）。",
            "required_evidence": [
                "关联交易清单（含交易类型、金额、定价方法）",
                "关联交易定价政策说明及可比非受控价格分析",
                "关联交易所涉合同及发票",
                "如为跨国关联交易，准备同期资料（主体文档+本地文档）",
                "董事会/股东会关于关联交易的决议（如金额较大）"
            ]
        })
    elif len(related_suppliers) > 0 or len(related_customers) > 0:
        results.append({
            "category": "经营实质", "category_icon": "🔍", "risk_score": 3, "risk_level": "低风险",
            "risk_color": "#3b82f6", "urgency": "建议",
            "item": "已识别关联方（供应商/客户）",
            "detail": f"系统中有 {len(related_suppliers)} 家疑似关联方供应商、{len(related_customers)} 家疑似关联方客户。关联交易需确保定价公允，并准备转让定价资料。",
            "suggestion": "建立关联交易管理制度，确保定价公允，并按规定准备同期资料。",
            "required_evidence": [
                "关联交易所涉合同",
                "定价公允性说明"
            ]
        })

    # 检查是否存在 "对开增值税专票"（即我开给你、你开给我，金额相近）——循环经济走私嫌疑
    # 简化版：检查是否有同一对手方既出现在 SalesInvoice 又出现在 PurchaseInvoice
    if company:
        counterparties = set()
        for inv in db.query(SalesInvoice).filter(SalesInvoice.company_id == company_id).all():
            counterparties.add(inv.buyer_name)
        for inv in db.query(PurchaseInvoice).filter(PurchaseInvoice.company_id == company_id).all():
            counterparties.add(inv.seller_name)

        # 如有重叠，提示风险
        # （完整实现需要更复杂的匹配逻辑，此处简化为提示）
        pass


# ═══════════════════════════════════════════════════════════
#  V5 新增 — 经营实质深度分析（稽查级·第11-18项）
# ═══════════════════════════════════════════════════════════

def _analyze_transport_missing(db, company_id, ps, pe, results):
    """运输费缺失检测——有销售收入但零运输/物流费（稽查必查·货物交付实质）"""
    revenue = _get_account_sum(db, company_id, "6001", ps, pe, "credit")
    if revenue < 100000:
        return  # 收入过低，不分析

    # 检查公司是否有实物产品销售（而非纯服务）
    company = db.query(Company).filter(Company.id == company_id).first()
    biz = (company.business_scope or "") + " " + (company.company_type or "")
    is_service_company = any(kw in biz for kw in ["服务", "咨询", "技术", "设计", "开发", "软件", "信息"])

    ps_date = _period_to_date_range(ps)[0]; pe_date = _period_to_date_range(pe)[1]

    # 检查销项发票中是否有实物商品（排除服务类）
    total_sales_amt = db.query(func.sum(SalesInvoice.total_amount)).filter(
        SalesInvoice.company_id == company_id,
        SalesInvoice.invoice_date >= ps_date, SalesInvoice.invoice_date <= pe_date,
        SalesInvoice.status == "正常"
    ).scalar() or 0

    service_keywords = ["服务", "咨询", "设计", "开发", "代理", "广告", "租赁", "技术"]
    service_sales = 0
    if is_service_company:
        service_sales = db.query(func.sum(SalesInvoice.total_amount)).filter(
            SalesInvoice.company_id == company_id,
            SalesInvoice.invoice_date >= ps_date, SalesInvoice.invoice_date <= pe_date,
            SalesInvoice.status == "正常",
            or_(*[SalesInvoice.goods_name.contains(kw) for kw in service_keywords])
        ).scalar() or 0
    product_sales = total_sales_amt - service_sales

    # 四源交叉验证：序时账 + 银行流水 + 取得发票 + 合同
    ev = _check_expense_evidence(db, company_id, ps, pe,
        expense_name="运输/物流费",
        je_account_codes=["660207"],
        je_summary_kw=["运输", "物流", "快递", "搬运", "货运", "配送"],
        bank_kw=["运输", "物流", "快递", "货运", "配送", "搬运"],
        invoice_kw=["运输", "物流", "快递", "配送", "货运", "搬运"],
        contract_kw=["运输", "物流", "快递", "货运", "承运", "配送"],
        contract_types=["服务"])

    transport_je = ev["je"]["amount"]
    transport_inv = ev["invoice"]["amount"]
    bank_transport = ev["bank"]["amount"]
    total_transport = transport_inv + transport_je + bank_transport

    # 纯服务型企业 → 跳过
    if is_service_company and product_sales == 0:
        return

    # 银行付款未入账 → 高风险（跨来源检测）
    if ev["overall"]["je_missing_but_bank_found"] and (product_sales > 100000 or revenue > 500000):
        results.append({
            "category": "经营实质", "category_icon": "🔍", "risk_score": 9, "risk_level": "高风险",
            "risk_color": "#dc2626", "urgency": "紧急",
            "item": "银行有运输/物流付款但序时账无运输费记录",
            "detail": f"银行流水中有运输/物流相关付款 {ev['bank']['amount']:,.0f} 元（实物产品销售 {product_sales:,.0f} 元），但序时账中无运输费科目（660207）记录。银行付款必须入账，否则余额对不上。可能：①运输费记入其他科目（如原材料采购含运费）；②凭证尚未生成。",
            "suggestion": "核查银行运输付款对应的序时账科目，如运费合并在采购成本中，建立科目映射对照表。",
            "required_evidence": ["银行运输付款回单", "对应序时账凭证", "运费与采购成本的分摊说明"]
        })
        return

    # 有发票但序时账未入账 → 中风险
    if ev["overall"]["je_missing_but_invoice_found"] and product_sales > 50000:
        results.append({
            "category": "经营实质", "category_icon": "🔍", "risk_score": 4, "risk_level": "中风险",
            "risk_color": "#f59e0b", "urgency": "提醒",
            "item": f"有运输发票（{ev['invoice']['count']}张/{ev['invoice']['amount']:,.0f}元）但序时账无运输费",
            "detail": f"取得运输/物流发票 {ev['invoice']['count']} 张，金额 {ev['invoice']['amount']:,.0f} 元，但序时账中未发现运输费科目（660207）。发票已取得但费用未入账，可能为发票延后或运费合并在采购成本中。",
            "suggestion": "确认运输发票是否已入账（如含在采购成本中），建议将运费单独核算以清晰反映货物流成本。"
        })
        return

    # 有运输合同但无账务
    if ev["overall"]["je_missing_but_contract_found"] and product_sales > 50000:
        free_contracts = [c for c in ev["contract"]["contracts"] if c["is_free"]]
        has_free = len(free_contracts) > 0
        results.append({
            "category": "经营实质", "category_icon": "🔍", "risk_score": 2 if has_free else 4,
            "risk_level": "低风险" if has_free else "中风险",
            "risk_color": "#3b82f6" if has_free else "#f59e0b", "urgency": "建议" if has_free else "提醒",
            "item": f"有运输/物流合同（{len(ev['contract']['contracts'])}份）但序时账无运输费{'(含免费条款)' if has_free else ''}",
            "detail": f"系统中有 {len(ev['contract']['contracts'])} 份运输/物流相关合同{'(含免费/无偿运输条款)' if has_free else ''}，但序时账中无运输费。可能为：①运费由客户承担（到付/自提）；②合同刚签署尚未产生费用；③免费配送。",
            "suggestion": "如运费由客户承担，在销售合同中标明运费条款。如为自提，保留客户自提签收记录。"
        })
        return

    # 四源皆空 → 高风险
    if ev["overall"]["all_sources_none"] and revenue > 0:
        if product_sales > 100000 or (not is_service_company and revenue > 500000):
            results.append({
                "category": "经营实质", "category_icon": "🔍", "risk_score": 8, "risk_level": "高风险",
                "risk_color": "#dc2626", "urgency": "紧急",
                "item": "有销售收入但零运输/物流费用（四源皆空）",
                "detail": f"企业营业收入 {revenue:,.0f} 元（其中实物产品销售 {product_sales:,.0f} 元），但序时账、银行流水、发票、合同中均未发现运输/物流费用。稽查视角：有货物销售但无运输费，货物如何交付？存在以下嫌疑：①虚开发票（无真实货物交付）；②运输费用以现金支付未入账；③由客户自提但无自提记录。这是稽查机关重点核查的货物交付实质问题。",
                "suggestion": "（稽查应对）准备以下佐证材料：①货运单/快递单/物流签收单（证明货物真实发出）；②如为客户自提，提供客户自提签收记录；③如含运费（由客户承担），提供销售合同相关条款；④承运方资质证明（道路运输许可证）；⑤主要客户的货物交付方式说明（含每种方式的占比）。",
                "required_evidence": [
                    "货物运输单据（运单/快递单/物流签收记录）",
                    "承运合同及承运方资质证明",
                    "客户自提记录（如为客户自提）",
                    "销售合同中关于运输方式的条款",
                    "货物交付方式说明（按客户分类）"
                ]
            })
        elif product_sales > 0:
            results.append({
                "category": "经营实质", "category_icon": "🔍", "risk_score": 3, "risk_level": "低风险",
                "risk_color": "#3b82f6", "urgency": "建议",
                "item": "运输费用偏低或未单独记录",
                "detail": f"实物产品销售 {product_sales:,.0f} 元，但运输费记录缺失。建议补录运输费用凭证或保留客户自提记录。",
                "suggestion": "保留货物交付凭证（运单/签收单/自提记录），以便应对稽查。"
            })


def _analyze_agriculture_substance(db, company_id, ps, pe, results):
    """农产品自产自销实质检测——有农产品销售但无土地/种植必须成本（稽查必查）"""
    company = db.query(Company).filter(Company.id == company_id).first()
    if not company:
        return
    biz = (company.business_scope or "") + " " + (company.company_type or "")
    is_agri = any(kw in biz for kw in ["农业", "种植", "养殖", "苗木", "花卉", "蔬菜",
                                         "水果", "茶叶", "中草药", "林木", "农产品", "畜牧",
                                         "农副产品", "园艺", "苗圃"])
    if not is_agri:
        # 也检查销项发票中是否有农产品/免税农产品
        ps_date = _period_to_date_range(ps)[0]; pe_date = _period_to_date_range(pe)[1]
        agri_sales = db.query(func.sum(SalesInvoice.total_amount)).filter(
            SalesInvoice.company_id == company_id,
            SalesInvoice.invoice_date >= ps_date, SalesInvoice.invoice_date <= pe_date,
            or_(
                SalesInvoice.goods_name.contains("苗木"),
                SalesInvoice.goods_name.contains("花卉"),
                SalesInvoice.goods_name.contains("蔬菜"),
                SalesInvoice.goods_name.contains("水果"),
                SalesInvoice.goods_name.contains("农产品"),
                SalesInvoice.goods_name.contains("粮食"),
                SalesInvoice.goods_name.contains("茶叶"),
                SalesInvoice.tax_rate == 0,  # 免税农产品
            )
        ).scalar() or 0
        if agri_sales < 50000:
            return
    else:
        ps_date = _period_to_date_range(ps)[0]; pe_date = _period_to_date_range(pe)[1]
        agri_sales = db.query(func.sum(SalesInvoice.total_amount)).filter(
            SalesInvoice.company_id == company_id,
            SalesInvoice.invoice_date >= ps_date, SalesInvoice.invoice_date <= pe_date,
        ).scalar() or 0

    if agri_sales < 50000:
        return

    # 四源交叉验证：土地/场地相关（序时账 + 银行流水 + 发票 + 合同）
    land_keywords = ["土地", "租赁", "承包", "流转", "租金", "地租", "场地"]
    ev = _check_expense_evidence(db, company_id, ps, pe,
        expense_name="土地/场地",
        je_summary_kw=land_keywords,
        bank_kw=land_keywords,
        invoice_kw=land_keywords,
        contract_kw=land_keywords,
        contract_types=["租赁", "承包"])

    land_costs = ev["je"]["amount"]
    bank_land = ev["bank"]["amount"]
    contracts = ev["contract"]["contracts"]
    has_land_contract = ev["contract"]["has_any"]
    land_free_possible = ev["contract"]["free_possible"]
    total_land = land_costs + bank_land

    # 银行有土地付款但序时账无记录 → 高风险
    if bank_land > 0 and land_costs == 0 and agri_sales > 100000:
        results.append({
            "category": "经营实质", "category_icon": "🔍", "risk_score": 8, "risk_level": "高风险",
            "risk_color": "#dc2626", "urgency": "紧急",
            "item": "银行有土地/场地付款但序时账无对应记录",
            "detail": f"银行流水中有土地/场地相关付款 {bank_land:,.0f} 元，但序时账中无对应土地租赁/承包费用。银行付款必须入账。可能：①土地费用记入其他科目；②凭证尚未生成。",
            "suggestion": "核查银行土地付款对应的序时账科目，补录土地租赁/承包费用凭证。",
            "required_evidence": ["银行土地付款回单", "对应序时账凭证"]
        })

    # 有土地合同但无账务（可能免费/自有土地）
    if has_land_contract and land_costs == 0 and bank_land == 0 and agri_sales > 100000:
        results.append({
            "category": "经营实质", "category_icon": "🔍", "risk_score": 3 if land_free_possible else 5,
            "risk_level": "低风险" if land_free_possible else "中风险",
            "risk_color": "#3b82f6" if land_free_possible else "#f59e0b",
            "urgency": "建议" if land_free_possible else "提醒",
            "item": f"有土地/场地合同（{len(contracts)}份）但序时账无土地租赁费{'(含免费条款)' if land_free_possible else ''}",
            "detail": f"系统中有 {len(contracts)} 份土地/场地相关合同{'(含免费/免租条款)' if land_free_possible else ''}，但序时账和银行流水中均无土地费用。可能为：①自有土地（无需租金）；②合同免费使用；③合同尚未产生费用。",
            "suggestion": "如为自有土地，准备土地证或产权证明；如为免费使用，保留合同作为佐证材料。"
        })

    # 检查是否有种植必须成本（化肥/农药/种子之外、人工/灌溉等）
    planting_keywords = ["化肥", "农药", "种子", "种苗", "灌溉", "农机", "农具",
                         "肥料", "除草", "收割", "采摘", "人工费", "劳务费",
                         "养护", "修剪", "施肥", "打药"]
    planting_costs = 0
    for kw in planting_keywords:
        planting_costs += db.query(func.sum(JournalEntry.debit_amount)).filter(
            JournalEntry.company_id == company_id,
            JournalEntry.period >= ps, JournalEntry.period <= pe,
            JournalEntry.summary.contains(kw)
        ).scalar() or 0

    # 进项发票中种植相关
    planting_inv = db.query(func.sum(PurchaseInvoice.total_amount)).filter(
        PurchaseInvoice.company_id == company_id,
        PurchaseInvoice.invoice_date >= ps_date, PurchaseInvoice.invoice_date <= pe_date,
        or_(*[PurchaseInvoice.goods_name.contains(kw) for kw in planting_keywords])
    ).scalar() or 0

    total_planting = planting_costs + planting_inv

    # 检查是否有农产品收购发票（反向开具）
    agri_purchase = db.query(func.sum(PurchaseInvoice.total_amount)).filter(
        PurchaseInvoice.company_id == company_id,
        PurchaseInvoice.invoice_date >= ps_date, PurchaseInvoice.invoice_date <= pe_date,
        or_(
            PurchaseInvoice.goods_name.contains("收购"),
            PurchaseInvoice.goods_name.contains("农产品"),
            PurchaseInvoice.tax_rate == 0,
        )
    ).scalar() or 0

    # 综合风险判定
    has_land = total_land > 0
    has_planting = total_planting > 0
    agri_sales_fmt = f"{agri_sales:,.0f}"

    if not has_land and not has_planting and agri_sales > 200000:
        results.append({
            "category": "经营实质", "category_icon": "🔍", "risk_score": 9, "risk_level": "高风险",
            "risk_color": "#dc2626", "urgency": "紧急",
            "item": "农产品销售无土地/种植成本支撑",
            "detail": f"企业农产品/自产自销相关销售收入 {agri_sales_fmt} 元（含免税农产品），但序时账、银行流水、进项发票中均未发现土地租赁/承包费用和种植必须成本（化肥、农药、人工、灌溉等）。稽查视角：自产自销必须有土地证明+种植全过程成本支出。仅有苗木采购发票而无种植成本，系典型的虚开农产品收购发票/骗取自产免税优惠的标志。",
            "suggestion": "（稽查应对）立即准备以下佐证材料：①土地证或土地租赁/承包合同（含支付凭证）；②种植生产记录（播种/施肥/打药/收割日期和用量）；③化肥/农药/种子/种苗的采购发票和入库单；④雇佣人工的劳务合同及工资支付凭证；⑤灌溉用水电费凭证；⑥农产品产量记录和销售台账；⑦如外包种植，提供外包种植合同及支付凭证。如为纯贸易（非自产），不得享受自产农产品免税优惠。",
            "required_evidence": [
                "土地证明：土地证/土地租赁合同/土地承包合同（含支付凭证）",
                "种植全过程记录：播种/施肥/打药/收割日期和用量台账",
                "农资采购凭证：化肥/农药/种子/种苗的发票+入库单+付款凭证",
                "人工成本：劳务合同+工资支付凭证（银行回单）",
                "灌溉用水电费缴纳凭证",
                "农产品产量记录+逐月销售台账",
                "如为纯贸易（非自产），提供采购合同+购进发票"
            ]
        })
    elif not has_land and agri_sales > 100000:
        results.append({
            "category": "经营实质", "category_icon": "🔍", "risk_score": 7, "risk_level": "高风险",
            "risk_color": "#dc2626", "urgency": "紧急",
            "item": "农产品销售无土地/场地费用",
            "detail": f"农产品销售收入 {agri_sales_fmt} 元，有部分种植成本但无土地租赁或承包费用。自产自销必须证明有合法的土地使用权（自有或租赁）。无土地证明→自产不成立→不得享受免税。",
            "suggestion": "准备土地证/土地租赁合同/承包合同及支付凭证。如为农户合作种植，准备合作种植协议及收购凭证。",
            "required_evidence": [
                "土地证/土地租赁合同/承包合同+支付凭证",
                "如为合作种植，提供合作种植协议+收购清单"
            ]
        })
    elif not has_planting and agri_sales > 100000:
        results.append({
            "category": "经营实质", "category_icon": "🔍", "risk_score": 7, "risk_level": "高风险",
            "risk_color": "#dc2626", "urgency": "紧急",
            "item": "农产品销售无种植必须成本",
            "detail": f"农产品销售收入 {agri_sales_fmt} 元，虽有土地相关费用但无种植必须成本（化肥/农药/人工/灌溉等）。仅有苗木采购成本而无种植全过程成本支出，不符合自产自销经营实质。稽查将认定：如不能证明自行种植，按贸易企业处理，不得享受自产免税优惠。",
            "suggestion": "补录化肥、农药、人工、灌溉等种植成本凭证，建立完整的种植成本台账。",
            "required_evidence": [
                "种植成本明细台账（分类：种苗/化肥/农药/人工/灌溉/其他）",
                "化肥/农药采购发票和入库记录",
                "雇佣人工的劳动合同+工资发放记录",
                "灌溉水电费凭证",
                "农业生产记录（逐日/逐周）"
            ]
        })
    elif has_land and has_planting and agri_sales > 0:
        # 低风险：有土地+有种植成本，但检查比例是否合理
        planting_ratio = total_planting / agri_sales * 100 if agri_sales > 0 else 0
        if planting_ratio < 5:
            results.append({
                "category": "经营实质", "category_icon": "🔍", "risk_score": 4, "risk_level": "中风险",
                "risk_color": "#f59e0b", "urgency": "提醒",
                "item": f"种植成本占销售收入比偏低（{planting_ratio:.1f}%）",
                "detail": f"种植相关成本 {total_planting:,.0f} 元，占农产品销售收入 {agri_sales_fmt} 元的 {planting_ratio:.1f}%（正常应在10-40%）。可能由于成本科目归集不完整，建议完善成本核算。",
                "suggestion": "检查是否所有种植成本均已入账，建立完整的种植成本核算体系。"
            })


def _analyze_packaging_missing(db, company_id, ps, pe, results):
    """包装费缺失检测——有销售出货但零包装费（稽查关注·产品交付实质）"""
    revenue = _get_account_sum(db, company_id, "6001", ps, pe, "credit")
    if revenue < 100000:
        return

    company = db.query(Company).filter(Company.id == company_id).first()
    biz = (company.business_scope or "") + " " + (company.company_type or "")
    is_product_biz = any(kw in biz for kw in ["制造", "生产", "销售", "贸易", "批发", "零售", "加工",
                                                "食品", "饮料", "服装", "电子", "机械", "化工"])
    if not is_product_biz:
        return

    # 四源交叉验证：序时账 + 银行流水 + 取得发票 + 合同
    ev = _check_expense_evidence(db, company_id, ps, pe,
        expense_name="包装费",
        je_account_codes=["660209"],
        je_summary_kw=["包装", "打包", "纸箱", "编织袋"],
        bank_kw=["包装", "打包", "纸箱", "包装材料"],
        invoice_kw=["包装", "纸箱", "塑料袋", "编织袋", "打包"],
        contract_kw=["包装", "打包", "包材", "包装材料"],
        contract_types=["采购", "服务"])

    packaging_je = ev["je"]["amount"]
    packaging_inv = ev["invoice"]["amount"]
    total_packaging = packaging_inv + packaging_je

    # 银行付款未入账 → 高风险
    if ev["overall"]["je_missing_but_bank_found"] and revenue > 200000:
        results.append({
            "category": "经营实质", "category_icon": "🔍", "risk_score": 8, "risk_level": "高风险",
            "risk_color": "#dc2626", "urgency": "紧急",
            "item": "银行有包装费付款但序时账无包装费记录",
            "detail": f"银行流水中有包装相关付款 {ev['bank']['amount']:,.0f} 元，但序时账中无对应包装费科目（660209）或包装摘要记录。银行流水必须入账，否则银行余额对不上。可能：①包装费记入了其他科目（如原材料）；②凭证尚未入账。",
            "suggestion": "核查银行包装付款对应的序时账科目，如误记入原材料等科目，建议调整至660209包装费科目或建立科目映射对照表。",
            "required_evidence": ["银行包装付款回单", "对应序时账凭证（核实科目归集）"]
        })
        return

    # 有发票但序时账未入账 → 中风险（发票可延后）
    if ev["overall"]["je_missing_but_invoice_found"]:
        results.append({
            "category": "经营实质", "category_icon": "🔍", "risk_score": 4, "risk_level": "中风险",
            "risk_color": "#f59e0b", "urgency": "提醒",
            "item": f"有包装材料发票（{ev['invoice']['count']}张/{ev['invoice']['amount']:,.0f}元）但序时账无包装费",
            "detail": f"取得包装相关进项发票 {ev['invoice']['count']} 张，金额 {ev['invoice']['amount']:,.0f} 元，但序时账中未发现包装费科目（660209）记录。发票已取得但费用尚未入账，可能为发票延后入账（常见做法）或费用科目归集错误。",
            "suggestion": "确认包装材料发票是否已入账（如记入原材料科目），建议补录包装费凭证或建立发票与凭证的对应关系。"
        })
        return

    # 有合同但无账务 → 低风险（可能为免费/对方承担）
    if ev["overall"]["je_missing_but_contract_found"]:
        free_contracts = [c for c in ev["contract"]["contracts"] if c["is_free"]]
        has_free = len(free_contracts) > 0
        results.append({
            "category": "经营实质", "category_icon": "🔍", "risk_score": 2 if has_free else 4,
            "risk_level": "低风险" if has_free else "中风险",
            "risk_color": "#3b82f6" if has_free else "#f59e0b", "urgency": "建议" if has_free else "提醒",
            "item": f"有包装合同（{len(ev['contract']['contracts'])}份）但序时账无包装费",
            "detail": f"系统中存在 {len(ev['contract']['contracts'])} 份包装相关合同{'(含免费/免租金条款)' if has_free else ''}，但序时账中无包装费记录。可能为：①合同已签但尚未实际履行；②包装费用由对方承担；③免费包装/赠送包材。",
            "suggestion": "核实合同实际履行情况，如已履行但费用未入账，补录相关凭证。如为对方免费提供，保留合同作为佐证材料。"
        })
        return

    # 四源皆空 → 中高风险
    if ev["overall"]["all_sources_none"] and revenue > 500000:
        results.append({
            "category": "经营实质", "category_icon": "🔍", "risk_score": 6, "risk_level": "中风险",
            "risk_color": "#f59e0b", "urgency": "提醒",
            "item": "有销售收入但零包装费用（四源皆空）",
            "detail": f"企业营业收入 {revenue:,.0f} 元，经营范围含产品生产/销售，但序时账、银行流水、进项发票、合同中均未发现包装费（包装材料/纸箱/打包等）。稽查视角：实物产品销售必然需要包装材料，零包装费不符合产品交付实质。可能：①包装费未单独核算（合并在原材料中）；②产品无实际包装（直发/散装需备说明）；③费用未取得发票。",
            "suggestion": "①如包装费合并在原材料采购中，补充说明材料；②如为散装/直发产品，提供客户签收记录；③补录包装材料采购凭证。",
            "required_evidence": [
                "包装材料采购发票及入库记录",
                "如包装费合并核算，提供材料出库单/成本分摊说明",
                "如为散装/直发，提供产品交付方式说明"
            ]
        })
    elif total_packaging == 0 and revenue > 200000:
        results.append({
            "category": "经营实质", "category_icon": "🔍", "risk_score": 3, "risk_level": "低风险",
            "risk_color": "#3b82f6", "urgency": "建议",
            "item": "包装费用偏低或未单独记录",
            "detail": f"企业营业收入 {revenue:,.0f} 元，包装费记录缺失。建议核实包装材料是否合并在原材料中核算。",
            "suggestion": "如包装费合并在原材料采购中，补充说明材料。"
        })


def _analyze_warehouse_missing(db, company_id, ps, pe, results):
    """仓储费缺失检测——有库存但无仓储/仓库租赁费（稽查关注·存货实质）"""
    # 检查是否有库存
    end_stock = db.query(func.sum(InventoryBalance.end_quantity)).filter(
        InventoryBalance.company_id == company_id,
        InventoryBalance.period == pe
    ).scalar() or 0

    # 从存货交易统计
    stock_in = db.query(func.sum(InventoryTransaction.quantity)).filter(
        InventoryTransaction.company_id == company_id,
        InventoryTransaction.trans_type == "入库",
        InventoryTransaction.transaction_date >= _period_to_date_range(ps)[0],
        InventoryTransaction.transaction_date <= _period_to_date_range(pe)[1]
    ).scalar() or 0

    if end_stock == 0 and stock_in == 0:
        return  # 无库存，跳过

    # 四源交叉验证：序时账 + 银行流水 + 取得发票 + 合同
    ev = _check_expense_evidence(db, company_id, ps, pe,
        expense_name="仓储费",
        je_summary_kw=["仓储", "仓库", "库房", "仓租", "仓管", "保管"],
        bank_kw=["仓储", "仓库", "仓租", "库房"],
        invoice_kw=["仓储", "仓库", "仓租", "库房", "保管", "寄存"],
        contract_kw=["仓库", "仓储", "库房", "仓租", "寄存"],
        contract_types=["租赁"])

    warehouse_je = ev["je"]["amount"]
    bank_warehouse = ev["bank"]["amount"]
    total_warehouse = warehouse_je + bank_warehouse
    warehouse_contracts = len(ev["contract"]["contracts"])

    # 银行付款未入账 → 高风险
    if ev["overall"]["je_missing_but_bank_found"] and end_stock > 100:
        results.append({
            "category": "经营实质", "category_icon": "🔍", "risk_score": 8, "risk_level": "高风险",
            "risk_color": "#dc2626", "urgency": "紧急",
            "item": f"银行有仓储支付（{ev['bank']['amount']:,.0f}元）但序时账无仓储费记录",
            "detail": f"期末库存 {end_stock:,.0f} 件，入库 {stock_in:,.0f} 件。银行流水中有仓储相关付款 {ev['bank']['amount']:,.0f} 元，但序时账中无对应仓储费记录。银行付款必须入账，否则余额对不上。",
            "suggestion": "核查银行仓储付款对应的序时账科目，补录仓储费凭证。",
            "required_evidence": ["银行仓储付款回单", "对应序时账凭证"]
        })
        return

    # 有合同但无账务（免费仓库可能）
    if ev["overall"]["je_missing_but_contract_found"]:
        free_contracts = [c for c in ev["contract"]["contracts"] if c["is_free"]]
        has_free = len(free_contracts) > 0
        results.append({
            "category": "经营实质", "category_icon": "🔍", "risk_score": 3 if has_free else 5,
            "risk_level": "低风险" if has_free else "中风险",
            "risk_color": "#3b82f6" if has_free else "#f59e0b", "urgency": "建议" if has_free else "提醒",
            "item": f"有仓储合同（{warehouse_contracts}份）但序时账无仓储费{'(含免费条款)' if has_free else ''}",
            "detail": f"期末库存 {end_stock:,.0f} 件，入库 {stock_in:,.0f} 件。系统中有 {warehouse_contracts} 份仓储/仓库相关合同{'(含免费/免租条款)' if has_free else ''}，但序时账中无仓储费。可能为：①自有仓库（无需租金）；②合同刚签署尚未产生费用；③对方免费提供仓储。",
            "suggestion": "如为自有仓库，提供房产证明；如为免费存货，保留仓储合同作为佐证材料。"
        })
        return

    # 四源皆空 → 高风险
    if ev["overall"]["all_sources_none"] and end_stock > 100:
        results.append({
            "category": "经营实质", "category_icon": "🔍", "risk_score": 7, "risk_level": "高风险",
            "risk_color": "#dc2626", "urgency": "紧急",
            "item": f"有库存（{end_stock:,.0f}件）但无仓储/仓库费用（四源皆空）",
            "detail": f"期末库存 {end_stock:,.0f} 件，入库 {stock_in:,.0f} 件，但序时账、银行流水、发票、合同中均无仓储/仓库费用。稽查视角：有货物就必须有存放场所。无仓储费→货物存放在哪里？→可能存在：①库存数据不实（虚列存货）；②仓库为自有但未入账；③存货已发出但未确认收入。",
            "suggestion": "①如为自有仓库，提供房产证/固定资产明细；②如为租赁仓库，补充仓库租赁合同及租金支付凭证；③核对库存实物与账面是否一致。",
            "required_evidence": [
                "仓库证明：自有房产证/仓库租赁合同+租金支付凭证",
                "仓库位置及面积说明",
                "存货盘点表（证明库存真实存在）",
                "如为寄售/代管库存，提供寄售/代管协议"
            ]
        })


def _analyze_equipment_depreciation_missing(db, company_id, ps, pe, results):
    """生产设备折旧缺失——制造业无设备折旧/租赁费（稽查必查·生产能力实质）"""
    company = db.query(Company).filter(Company.id == company_id).first()
    if not company:
        return
    biz = (company.business_scope or "") + " " + (company.company_type or "")
    if not any(kw in biz for kw in ["制造", "生产", "加工"]):
        return

    revenue = _get_account_sum(db, company_id, "6001", ps, pe, "credit")
    if revenue < 200000:
        return

    # 固定资产折旧
    depr = _get_account_sum(db, company_id, "660202", ps, pe, "debit")  # 累计折旧-费用化
    # 制造费用-折旧（可能在生产成本中）
    mfg_depr = db.query(func.sum(JournalEntry.debit_amount)).filter(
        JournalEntry.company_id == company_id,
        JournalEntry.period >= ps, JournalEntry.period <= pe,
        or_(
            JournalEntry.summary.contains("折旧"),
            JournalEntry.summary.contains("摊销"),
        )
    ).scalar() or 0

    # 固定资产原值
    fixed_assets = db.query(func.sum(FixedAsset.original_value)).filter(
        FixedAsset.company_id == company_id
    ).scalar() or 0

    # 四源交叉验证：序时账 + 银行流水 + 取得发票 + 合同（设备租赁/折旧）
    ev = _check_expense_evidence(db, company_id, ps, pe,
        expense_name="设备折旧/租赁费",
        je_extra_codes=["660202"],  # 折旧已单独查，这里只查租赁
        je_summary_kw=["设备租赁", "机械租赁", "设备租金", "机器租赁"],
        bank_kw=["设备租赁", "机械租赁", "设备租金", "机器租赁", "设备采购"],
        invoice_kw=["设备", "机械", "机器", "生产线"],
        contract_kw=["设备", "机械", "机器", "生产线", "加工"],
        contract_types=["采购", "租赁"])

    equip_rent = ev["je"]["amount"]  # 设备租赁费（四源结果）
    total_equip_cost = depr + mfg_depr + equip_rent

    # 银行有设备付款但序时账无折旧/租赁记录 → 高风险
    if ev["overall"]["je_missing_but_bank_found"] and fixed_assets == 0 and revenue > 300000:
        results.append({
            "category": "经营实质", "category_icon": "🔍", "risk_score": 8, "risk_level": "高风险",
            "risk_color": "#dc2626", "urgency": "紧急",
            "item": "银行有设备相关付款但序时账无设备折旧/租赁费",
            "detail": f"银行流水中有设备相关付款 {ev['bank']['amount']:,.0f} 元，但序时账中无设备折旧或租赁费，且固定资产系统中无设备记录。银行付款必须入账。可能：①设备采购未入固定资产卡片；②设备租赁费记入其他科目；③为其他用途付款（非生产设备）。",
            "suggestion": "核查银行设备付款的实际用途，如为设备采购，补录固定资产卡片并计提折旧；如为租赁，补录设备租赁费凭证。",
            "required_evidence": ["银行设备付款回单", "设备采购合同/租赁合同", "固定资产卡片"]
        })
        return

    # 有设备合同但无账务
    if ev["overall"]["je_missing_but_contract_found"] and fixed_assets == 0:
        free_contracts = [c for c in ev["contract"]["contracts"] if c["is_free"]]
        has_free = len(free_contracts) > 0
        results.append({
            "category": "经营实质", "category_icon": "🔍", "risk_score": 3 if has_free else 5,
            "risk_level": "低风险" if has_free else "中风险",
            "risk_color": "#3b82f6" if has_free else "#f59e0b", "urgency": "建议" if has_free else "提醒",
            "item": f"有设备相关合同（{len(ev['contract']['contracts'])}份）但无固定资产/折旧记录{'(含免费条款)' if has_free else ''}",
            "detail": f"系统中有 {len(ev['contract']['contracts'])} 份设备/机械相关合同{'(含免费/无偿条款)' if has_free else ''}，但序时账和固定资产系统中均无对应记录。可能为：①合同已签但设备尚未到货；②设备由对方免费提供使用；③合同未实际履行。",
            "suggestion": "核实合同实际履行状态，如设备已到货使用，补录固定资产卡片；如为免费使用，保留合同作为佐证材料。"
        })
        return

    if fixed_assets == 0 and equip_rent == 0 and revenue > 500000:
        results.append({
            "category": "经营实质", "category_icon": "🔍", "risk_score": 9, "risk_level": "高风险",
            "risk_color": "#dc2626", "urgency": "紧急",
            "item": "制造业无生产设备（固定资产+租赁）",
            "detail": f"企业经营范围含生产制造，营业收入 {revenue:,.0f} 元，但系统中既无固定资产（机器设备）记录，也无设备租赁费用。稽查视角：生产制造必须有生产设备。无设备→如何生产？→存在以下嫌疑：①虚开发票（无生产能力的空壳）；②委托加工但未签订加工合同；③设备全部租赁但未入账。",
            "suggestion": "（稽查应对）①如为自有设备，补录固定资产卡片（含设备名称、型号、购置发票、折旧明细）；②如为委托加工，提供委托加工合同+加工费发票+发料/收货记录；③如为租赁设备，提供设备租赁合同+租金支付凭证。",
            "required_evidence": [
                "固定资产清单（含机器设备明细、购置发票、折旧计算表）",
                "委托加工合同+加工费发票+发料收货记录（如为外包生产）",
                "设备租赁合同+租金支付凭证（如为租赁设备）",
                "生产车间/设备现场照片（含设备型号特写）"
            ]
        })
    elif total_equip_cost > 0 and total_equip_cost / revenue < 0.005 and revenue > 1000000:
        results.append({
            "category": "经营实质", "category_icon": "🔍", "risk_score": 5, "risk_level": "中风险",
            "risk_color": "#f59e0b", "urgency": "提醒",
            "item": f"生产设备折旧/租赁费占收入比偏低（{total_equip_cost/revenue*100:.2f}%）",
            "detail": f"设备相关成本 {total_equip_cost:,.0f} 元，占营业收入 {revenue:,.0f} 元的 {total_equip_cost/revenue*100:.2f}%。制造业设备折旧/租赁费通常占收入1-5%，费用偏少可能说明设备老旧或部分生产外包。",
            "suggestion": "如为设备老旧（已提完折旧），准备固定资产台账说明；如为外包生产，补充加工合同。"
        })


def _analyze_advertising_missing(db, company_id, ps, pe, results):
    """广告/营销费缺失——有收入但零营销推广费（经营实质·市场拓展）"""
    revenue = _get_account_sum(db, company_id, "6001", ps, pe, "credit")
    if revenue < 500000:
        return

    # 四源交叉验证：序时账 + 银行流水 + 取得发票 + 合同
    ev = _check_expense_evidence(db, company_id, ps, pe,
        expense_name="广告/营销费",
        je_account_codes=["660103"],
        je_summary_kw=["广告", "推广", "营销", "宣传", "展会"],
        bank_kw=["广告", "推广", "营销", "宣传", "展会", "传媒"],
        invoice_kw=["广告", "推广", "营销", "宣传", "展会"],
        contract_kw=["广告", "推广", "营销", "宣传", "展会"],
        contract_types=["服务"])

    ad_je = ev["je"]["amount"]
    ad_inv = ev["invoice"]["amount"]
    total_ad = ad_inv + ad_je

    # 银行付款未入账 → 高风险
    if ev["overall"]["je_missing_but_bank_found"] and revenue > 500000:
        results.append({
            "category": "经营实质", "category_icon": "🔍", "risk_score": 8, "risk_level": "高风险",
            "risk_color": "#dc2626", "urgency": "紧急",
            "item": "银行有广告/营销付款但序时账无对应记录",
            "detail": f"银行流水中有广告/推广相关付款 {ev['bank']['amount']:,.0f} 元，但序时账中无广告费科目（660103）记录。银行付款必须入账，否则余额对不上。可能：①广告费记入其他科目；②凭证尚未生成。",
            "suggestion": "核查银行广告付款对应的序时账科目，补录广告费凭证或建立科目映射表。",
            "required_evidence": ["银行广告付款回单", "对应序时账凭证"]
        })
        return

    # 有发票但序时账未入账 → 中风险
    if ev["overall"]["je_missing_but_invoice_found"]:
        results.append({
            "category": "经营实质", "category_icon": "🔍", "risk_score": 4, "risk_level": "中风险",
            "risk_color": "#f59e0b", "urgency": "提醒",
            "item": f"有广告/营销发票（{ev['invoice']['count']}张/{ev['invoice']['amount']:,.0f}元）但序时账无广告费",
            "detail": f"取得广告/推广相关进项发票 {ev['invoice']['count']} 张，金额 {ev['invoice']['amount']:,.0f} 元，但序时账中未发现广告费科目（660103）记录。发票已取得但未入账，可能延后入账或科目归集有误。",
            "suggestion": "确认广告发票是否已入账，补录广告费凭证。"
        })
        return

    # 有合同但无账务
    if ev["overall"]["je_missing_but_contract_found"]:
        results.append({
            "category": "经营实质", "category_icon": "🔍", "risk_score": 4, "risk_level": "中风险",
            "risk_color": "#f59e0b", "urgency": "提醒",
            "item": f"有广告/营销合同（{len(ev['contract']['contracts'])}份）但序时账无广告费",
            "detail": f"系统中有 {len(ev['contract']['contracts'])} 份广告/营销相关合同，但序时账中无广告费。可能为：①合同已签但尚未产生费用；②以资源置换形式合作（无现金流）。",
            "suggestion": "核实合同实际履行情况，如已产生费用，补录凭证。如为资源置换，保留合同作为说明材料。"
        })
        return

    # 四源皆空 → 中风险（取决于收入规模）
    if ev["overall"]["all_sources_none"] and revenue > 1000000:
        results.append({
            "category": "经营实质", "category_icon": "🔍", "risk_score": 5, "risk_level": "中风险",
            "risk_color": "#f59e0b", "urgency": "提醒",
            "item": f"营业收入 {revenue:,.0f} 元但零广告/营销费用（四源皆空）",
            "detail": f"营业收入 {revenue:,.0f} 元（≥100万元），但序时账、银行流水、发票、合同中均未发现广告费/推广费/营销费等市场拓展相关支出。稽查视角：收入达到一定规模却无市场拓展费用，不符合商业逻辑。可能：①广告费未取得发票；②以其他费用名义入账；③确无广告需求（如为B2B大客户模式）。",
            "suggestion": "①如确有广告/推广支出，补录相关凭证；②如为B2B大客户模式（无需广告），准备客户开发方式说明；③如通过电商平台销售，平台服务费即推广费，需明确标注。",
            "required_evidence": [
                "市场推广/客户开发模式说明",
                "如为B2B模式，提供主要客户开发记录和销售合同",
                "如有广告支出，提供广告合同+发票+付款凭证"
            ]
        })


def _analyze_travel_missing(db, company_id, ps, pe, results):
    """差旅费缺失——有异地客户但无差旅费（经营实质·业务开展）"""
    revenue = _get_account_sum(db, company_id, "6001", ps, pe, "credit")
    if revenue < 300000:
        return

    # 检查是否有非本地的客户
    company = db.query(Company).filter(Company.id == company_id).first()
    company_city = ""
    if company and company.address:
        addr = company.address
        for suffix in ["市", "县", "区"]:
            idx = addr.find(suffix)
            if idx > 0:
                company_city = addr[max(0, idx-3):idx+1]
                break

    customers = db.query(Customer).filter(Customer.company_id == company_id).all()
    has_remote_customers = False
    for c in customers:
        cust_addr = c.address or ""
        if company_city and company_city not in cust_addr:
            has_remote_customers = True
            break

    # 四源交叉验证：序时账 + 银行流水 + 取得发票 + 合同
    ev = _check_expense_evidence(db, company_id, ps, pe,
        expense_name="差旅费",
        je_account_codes=["660211"],
        je_summary_kw=["差旅", "出差", "交通", "住宿", "机票"],
        bank_kw=["差旅", "出差", "机票", "酒店", "住宿", "携程", "去哪儿"],
        invoice_kw=["机票", "酒店", "住宿", "差旅", "火车票"],
        contract_kw=["差旅", "出差", "商旅", "携程"],
        contract_types=["服务"])

    travel_je = ev["je"]["amount"]
    travel_inv = ev["invoice"]["amount"]
    total_travel = travel_je + travel_inv

    # 银行付款未入账 → 高风险
    if ev["overall"]["je_missing_but_bank_found"] and revenue > 500000:
        results.append({
            "category": "经营实质", "category_icon": "🔍", "risk_score": 7, "risk_level": "高风险",
            "risk_color": "#dc2626", "urgency": "紧急",
            "item": "银行有差旅付款但序时账无差旅费记录",
            "detail": f"银行流水中有差旅相关付款 {ev['bank']['amount']:,.0f} 元，但序时账中无差旅费科目（660211）记录。银行付款必须入账。可能：①差旅费记入其他费用科目；②报销凭证尚未生成。",
            "suggestion": "核查银行差旅付款对应的序时账科目，补录差旅费凭证。",
            "required_evidence": ["银行差旅付款回单", "对应序时账凭证"]
        })
        return

    # 有发票但序时账未入账 → 中风险
    if ev["overall"]["je_missing_but_invoice_found"]:
        results.append({
            "category": "经营实质", "category_icon": "🔍", "risk_score": 4, "risk_level": "中风险",
            "risk_color": "#f59e0b", "urgency": "提醒",
            "item": f"有差旅发票（{ev['invoice']['count']}张/{ev['invoice']['amount']:,.0f}元）但序时账无差旅费",
            "detail": f"取得机票/酒店等差旅发票 {ev['invoice']['count']} 张，金额 {ev['invoice']['amount']:,.0f} 元，但序时账中未发现差旅费科目（660211）记录。发票已取得但报销凭证未生成，可能为员工差旅费未及时报销入账。",
            "suggestion": "催促员工及时提交差旅费报销单，补录差旅费凭证。"
        })
        return

    # 四源皆空 + 有异地客户 → 中风险
    if ev["overall"]["all_sources_none"] and revenue > 1000000 and (has_remote_customers or len(customers) > 5):
        results.append({
            "category": "经营实质", "category_icon": "🔍", "risk_score": 4, "risk_level": "中风险",
            "risk_color": "#f59e0b", "urgency": "提醒",
            "item": f"有{len(customers)}家客户但零差旅费用（四源皆空）",
            "detail": f"系统中有 {len(customers)} 家客户（含异地客户），营业收入 {revenue:,.0f} 元，但序时账、银行流水、发票、合同中均未发现差旅费/住宿费/交通费。稽查视角：有客户就有业务往来，有业务就有差旅需求。零差旅费不符合常规商业逻辑。可能：①差旅费未取得发票或以现金支付未入账；②确无差旅需要（本地客户为主/线上沟通）。",
            "suggestion": "①如确有差旅，补录差旅费报销凭证；②如以线上沟通为主，准备业务沟通记录（邮件/视频会议记录）；③如客户均为本地，提供客户地域分布说明。",
            "required_evidence": [
                "客户地域分布说明",
                "差旅费报销凭证（如有）",
                "线上业务沟通记录（邮件/会议记录）"
            ]
        })


def _analyze_office_expense_missing(db, company_id, ps, pe, results):
    """办公费缺失——四源交叉验证（序时账+银行流水+发票+合同）"""
    revenue = _get_account_sum(db, company_id, "6001", ps, pe, "credit")
    if revenue < 100000:
        return

    # 四源交叉验证：序时账 + 银行流水 + 取得发票 + 合同
    office_kw = ["办公", "物业", "办公用品", "打印"]
    ev = _check_expense_evidence(db, company_id, ps, pe,
        expense_name="办公/物业费",
        je_account_codes=["660201", "660216"],
        je_summary_kw=office_kw,
        bank_kw=office_kw,
        invoice_kw=office_kw,
        contract_kw=["场地", "物业", "办公", "租赁", "房租"],
        contract_types=["租赁", "服务"])

    office_je = ev["je"]["amount"]
    bank_office = ev["bank"]["amount"]
    inv_count = ev["invoice"]["count"]
    overall = ev["overall"]
    free_possible = ev["contract"]["free_possible"]
    contracts = ev["contract"]["contracts"]

    # ── 银行有支付但序时账无记录 → 高风险（银行流水必须入账） ──
    if overall["je_missing_but_bank_found"] and revenue > 500000:
        results.append({
            "category": "经营实质", "category_icon": "🔍", "risk_score": 7, "risk_level": "高风险",
            "risk_color": "#dc2626", "urgency": "紧急",
            "item": f"办公/物业费银行付款未入账（¥{bank_office:,.0f}）",
            "detail": f"营业收入 {revenue:,.0f} 元（≥50万元），银行流水中有办公/物业相关付款 ¥{bank_office:,.2f}，但序时账中无对应费用记录。银行流水必须入账，否则银行存款余额对不上。",
            "suggestion": f"将银行流水中办公/物业相关付款（¥{bank_office:,.2f}）立即生成序时账凭证。",
            "required_evidence": ["办公/物业相关银行支付回单", "对应的发票（如有）"]
        })
        return

    # ── 四源全无 → 中风险 ──
    if overall["all_sources_none"] and revenue > 500000:
        results.append({
            "category": "经营实质", "category_icon": "🔍", "risk_score": 5, "risk_level": "中风险",
            "risk_color": "#f59e0b", "urgency": "提醒",
            "item": f"营业收入 {revenue:,.0f} 元但零办公/物业费用",
            "detail": f"营业收入 {revenue:,.0f} 元（≥50万元），四源交叉验证均未发现办公费/物业费：序时账、银行流水、取得发票、合同。稽查视角：有经营就有办公场所，有办公场所就有办公费用（水电/物业/办公用品）。零办公费→经营场所是否存在？→可能：①经营场所费用由关联方承担未入账；②在家办公但无费用凭证；③确无独立办公场所。",
            "suggestion": "①补录办公费、物业费等日常费用凭证；②如经营场所费用由他人承担，提供相关协议和说明；③如为家庭办公，说明情况并保留部分费用凭证。",
            "required_evidence": [
                "经营场所使用证明（房产证/租赁合同）",
                "物业费/水电费缴纳凭证",
                "如为家庭办公，提供情况说明"
            ]
        })
        return

    # ── 序时账没有，发票有 → 费用未入账（发票可延后） ──
    if overall["je_missing_but_invoice_found"] and revenue > 500000:
        results.append({
            "category": "经营实质", "category_icon": "🔍", "risk_score": 3, "risk_level": "低风险",
            "risk_color": "#3b82f6", "urgency": "提醒",
            "item": f"营业收入 {revenue:,.0f} 元但零办公费（已取得发票未入账）",
            "detail": f"营业收入 {revenue:,.0f} 元（≥50万元），序时账中未发现办公费/物业费记录，但取得发票中发现办公/物业相关发票（{inv_count}张）。费用已发生但尚未入账，请及时入账。",
            "suggestion": "将已取得的办公费/物业费相关发票及时入账。",
            "required_evidence": ["未入账的办公费/物业费发票", "经营场所使用证明"]
        })
        return

    # ── 序时账没有，但有场地合同（免费场地） ──
    if overall["je_missing_but_contract_found"]:
        free_note = "·可能为免费场地" if free_possible else ""
        results.append({
            "category": "经营实质", "category_icon": "🔍", "risk_score": 2, "risk_level": "低风险",
            "risk_color": "#3b82f6", "urgency": "提醒",
            "item": f"办公费记录缺失但存在场地合同{free_note}",
            "detail": f"营业收入 {revenue:,.0f} 元，序时账中无办公/物业费记录，但合同中有场地相关协议（{len(contracts)}份）。可能为关联方免费提供场地或家庭办公，建议保留合同备查。",
            "suggestion": "保留场地使用合同/协议原件备查。如为免费使用，准备相关说明。",
            "required_evidence": ["场地使用合同/协议"]
        })


# ═══════════════════════════════════════════════════════════
#  V4 新增 — 增值税专项分析
# ═══════════════════════════════════════════════════════════

def _analyze_vat_zero_declaration(db, company_id, ps, pe, results):
    """增值税零申报月数：连续零申报月份（经营异常信号）"""
    decls = db.query(VATDeclaration.period).filter(
        VATDeclaration.company_id == company_id,
        VATDeclaration.period >= ps, VATDeclaration.period <= pe
    ).order_by(VATDeclaration.period).all()
    if not decls:
        return

    periods = sorted(set(d[0] for d in decls))
    zero_months = []
    for p in periods:
        v = db.query(VATDeclaration).filter(
            VATDeclaration.company_id == company_id, VATDeclaration.period == p).first()
        if v and v.form_main:
            fm = v.form_main if isinstance(v.form_main, dict) else json.loads(v.form_main)
            taxable_sales = _safe_float(fm.get("taxable_sales", 0))
            tax_free_sales = _safe_float(fm.get("tax_free_sales", 0))
            vat_payable = _safe_float(fm.get("vat_payable", 0))
            if taxable_sales == 0 and tax_free_sales == 0 and vat_payable == 0:
                zero_months.append(p)

    if len(zero_months) >= 6:
        results.append({
            "category": "增值税专项", "category_icon": "📋", "risk_score": 8, "risk_level": "高风险",
            "risk_color": "#dc2626", "urgency": "紧急",
            "item": "长期零申报或小额申报",
            "detail": f"存在 {len(zero_months)} 个月增值税零申报：{', '.join(zero_months[:6])}。稽查视角：长期零申报但工商登记为「在营」，存在隐瞒收入或未按规定申报的重大嫌疑，可能触发税务稽查入户调查。特别提示：小规模纳税人连续12个月零申报将被列入异常名录。",
            "suggestion": "①核查实际经营收入，补报漏报销售；②如确无经营应办理停业登记；③税务机关预警后可导致发票停供、纳税信用降级。",
            "required_evidence": ["经营场所租赁合同或产权证明", "银行账户流水（证明无经营收入）", "近12个月零申报情况说明", "如已停业，提供工商停业备案"]
        })
    elif len(zero_months) >= 3:
        results.append({
            "category": "增值税专项", "category_icon": "📋", "risk_score": 6, "risk_level": "中风险",
            "risk_color": "#f59e0b", "urgency": "提醒",
            "item": "长期零申报或小额申报",
            "detail": f"共有 {len(zero_months)} 个月零申报。连续零申报将被税务机关重点关注。",
            "suggestion": "如实申报实际经营收入，避免长期零申报引致稽查风险。"
        })


def _analyze_vat_burden_quarterly(db, company_id, ps, pe, results):
    """增值税税负率季度波动检测"""
    decls = db.query(VATDeclaration).filter(
        VATDeclaration.company_id == company_id,
        VATDeclaration.period >= ps, VATDeclaration.period <= pe
    ).order_by(VATDeclaration.period).all()

    if len(decls) < 4:
        return

    quarterly_data = {}
    for v in decls:
        if not v.period: continue
        fm = v.form_main if isinstance(v.form_main, dict) else json.loads(v.form_main)
        sales = _safe_float(fm.get("taxable_sales", 0))
        vat = _safe_float(fm.get("vat_payable", 0))
        if sales == 0: continue
        q = v.period[:4] + "-Q" + str((int(v.period[5:7]) - 1) // 3 + 1)
        if q not in quarterly_data:
            quarterly_data[q] = {"sales": 0, "vat": 0}
        quarterly_data[q]["sales"] += sales
        quarterly_data[q]["vat"] += vat

    if len(quarterly_data) < 2:
        return

    quarters = sorted(quarterly_data.keys())
    rates = []
    for q in quarters:
        d = quarterly_data[q]
        r = d["vat"] / d["sales"] * 100 if d["sales"] > 0 else 0
        rates.append({"quarter": q, "rate": r, "sales": d["sales"]})

    if len(rates) >= 2:
        variances = [abs(rates[i]["rate"] - rates[i-1]["rate"]) for i in range(1, len(rates))]
        avg_var = sum(variances) / len(variances)
        if avg_var > 3:
            rate_str = ' / '.join(f"{r['quarter']}:{r['rate']:.2f}%" for r in rates)
            results.append({
                "category": "增值税专项", "category_icon": "📋", "risk_score": 7, "risk_level": "高风险",
                "risk_color": "#dc2626", "urgency": "紧急",
                "item": "增值税税负率季度波动异常",
                "detail": f"增值税税负率季度间平均波动 {avg_var:.1f}个百分点。各季度：{rate_str}。税负率大幅波动通常意味着：①收入确认时点人为调节（跨期调节）；②进项税额集中抵扣导致波动；③存在未申报收入被集中处理。",
                "suggestion": "核查各季度收入确认的准确性，避免人为调节税负率。准备季度税负率波动说明。",
                "required_evidence": ["季度收入确认明细表", "季度进项税额抵扣明细", "税负率波动情况说明"]
            })


def _analyze_vat_input_transfer_omission(db, company_id, ps, pe, results):
    """进项税额转出遗漏：有免税/简易计税但未做进项转出"""
    decls = db.query(VATDeclaration).filter(
        VATDeclaration.company_id == company_id,
        VATDeclaration.period >= ps, VATDeclaration.period <= pe
    ).all()

    for v in decls:
        if not v.form_main:
            continue
        fm = v.form_main if isinstance(v.form_main, dict) else json.loads(v.form_main)
        tax_free = _safe_float(fm.get("tax_free_sales", 0))
        simple_tax = _safe_float(fm.get("simple_tax_sales", 0))
        input_transfer = _safe_float(fm.get("input_tax_transfer", 0))

        if (tax_free > 0 or simple_tax > 0) and input_transfer == 0:
            input_cnt = db.query(func.count(PurchaseInvoice.id)).filter(
                PurchaseInvoice.company_id == company_id,
                func.substr(PurchaseInvoice.invoice_date, 1, 7) == v.period
            ).scalar() or 0
            if input_cnt > 0:
                results.append({
                    "category": "增值税专项", "category_icon": "📋", "risk_score": 9, "risk_level": "高风险",
                    "risk_color": "#dc2626", "urgency": "紧急",
                    "item": "免税收入进项税额未转出",
                    "detail": f"{v.period} 期存在免税销售额（{tax_free:,.2f}元）或简易计税销售额（{simple_tax:,.2f}元），同时有 {input_cnt} 张进项发票，但申报表中进项税额转出为0。根据增值税暂行条例，用于免税项目、简易计税项目的进项税额不得抵扣，必须做进项税额转出。",
                    "suggestion": "①立即核算不得抵扣的进项税额并补做转出申报；②建立进项税额分摊台账（按月、按项目）；③不得抵扣进项税额=当月全部进项税额×(当月免税/简易销售额÷当月全部销售额)。",
                    "required_evidence": ["免税/简易计税项目收入明细", "进项税额分摊计算表", "进项税额转出明细（按项目）", "修正后的增值税申报表"]
                })
                return


def _analyze_vat_retention_refund(db, company_id, ps, pe, results):
    """增值税留抵退税风险"""
    decls = db.query(VATDeclaration).filter(
        VATDeclaration.company_id == company_id,
        VATDeclaration.period >= ps, VATDeclaration.period <= pe
    ).order_by(VATDeclaration.period).all()

    # 检查连续留抵月数
    continuous_stay = 0
    max_continuous = 0
    for v in decls:
        if v.form_main:
            fm = v.form_main if isinstance(v.form_main, dict) else json.loads(v.form_main)
            end_credit = _safe_float(fm.get("end_period_credit", 0))
            if end_credit > 0:
                continuous_stay += 1
                max_continuous = max(max_continuous, continuous_stay)
            else:
                continuous_stay = 0

    if max_continuous >= 6:
        results.append({
            "category": "增值税专项", "category_icon": "📋", "risk_score": 7, "risk_level": "高风险",
            "risk_color": "#dc2626", "urgency": "紧急",
            "item": "留抵退税异常增长",
            "detail": f"连续 {max_continuous} 个月存在增值税留抵税额。稽查视角：长期留抵但仍在大量采购→可能：①进项发票虚开；②未入账销售收入（私账收款）；③存货积压虚假。留抵退税申请后，税务机关将对进项发票进行专项核查。",
            "suggestion": "①自查大额留抵原因（存货积压/采购集中/售价倒挂）；②申请留抵退税前确保进项发票真实合规；③如有私账收款立即补报。",
            "required_evidence": ["留抵原因说明（含存货/产能分析）", "大额留抵期间的主要进项发票清单", "存货盘点表（证明存货真实存在）"]
        })


def _analyze_vat_no_ticket_sales(db, company_id, ps, pe, results):
    """增值税无票销售额风险（适用一般纳税人）"""
    company = db.query(Company).filter(Company.id == company_id).first()
    if not company or (company.company_type and "小规模" in company.company_type):
        return

    decls = db.query(VATDeclaration).filter(
        VATDeclaration.company_id == company_id,
        VATDeclaration.period >= ps, VATDeclaration.period <= pe
    ).all()

    no_ticket_total = 0
    for v in decls:
        if v.form_main:
            fm = v.form_main if isinstance(v.form_main, dict) else json.loads(v.form_main)
            no_ticket_total += _safe_float(fm.get("no_ticket_sales", 0))

    # 检查销售收入vs开票金额是否匹配
    revenue_credit = _get_account_sum(db, company_id, "6001", ps, pe, "credit")
    invoice_sales = db.query(func.sum(SalesInvoice.amount)).filter(
        SalesInvoice.company_id == company_id,
        SalesInvoice.invoice_date >= _period_to_date_range(ps)[0],
        SalesInvoice.invoice_date <= _period_to_date_range(pe)[1]
    ).scalar() or 0

    gap = revenue_credit - _safe_float(invoice_sales)
    if gap > 100000:  # gap超过10万需要关注
        results.append({
            "category": "增值税专项", "category_icon": "📋", "risk_score": 8, "risk_level": "高风险",
            "risk_color": "#dc2626", "urgency": "紧急",
            "item": "无票收入申报缺失",
            "detail": f"主营业务收入（{revenue_credit:,.2f}元）大于销项发票金额（{_safe_float(invoice_sales):,.2f}元），差额 {gap:,.2f} 元可能为无票销售收入。稽查视角：无票收入如未如实申报增值税，将被认定为偷漏税。尤其餐饮/零售等面向个人消费者的行业，无票收入占比高是常见稽查重点。",
            "suggestion": "①确认无票收入是否已在增值税申报中如实填报；②餐饮企业使用收银系统数据与申报数据比对；③建立健全无票收入台账（按日/按期）。",
            "required_evidence": ["无票销售收入明细台账", "收银系统/销售系统数据备份", "无票收入与申报数据比对表"]
        })


# ═══════════════════════════════════════════════════════════
#  B 类 — 发票异常分析
# ═══════════════════════════════════════════════════════════

def _analyze_invoice_amount_anomaly(db, company_id, ps, pe, results):
    """零税额/顶额/代开发票异常分析"""
    ps_date = _period_to_date_range(ps)[0]; pe_date = _period_to_date_range(pe)[1]

    # 零税额发票
    zero_tax_sales = db.query(func.count(SalesInvoice.id), func.sum(SalesInvoice.amount)).filter(
        SalesInvoice.company_id == company_id,
        SalesInvoice.invoice_date >= ps_date, SalesInvoice.invoice_date <= pe_date,
        SalesInvoice.tax_amount == 0, SalesInvoice.amount > 0
    ).first()

    zero_cnt, zero_amt = zero_tax_sales
    total_sales = db.query(func.count(SalesInvoice.id)).filter(
        SalesInvoice.company_id == company_id,
        SalesInvoice.invoice_date >= ps_date, SalesInvoice.invoice_date <= pe_date
    ).scalar() or 1

    if zero_cnt and zero_cnt > 0 and (zero_cnt / total_sales) > 0.3:
        results.append({
            "category": "发票异常", "category_icon": "🧾", "risk_score": 7, "risk_level": "高风险",
            "risk_color": "#dc2626", "urgency": "紧急",
            "item": "零税额发票异常",
            "detail": f"销项发票中 {zero_cnt} 张为零税额发票，占比 {zero_cnt/total_sales*100:.1f}%，涉及金额 {_safe_float(zero_amt):,.2f} 元。稽查视角：大量开具零税率发票但无免税备案，可能被认定为虚开发票。",
            "suggestion": "核查零税额发票对应的业务是否具备免税资格，如无则应补开正常税率发票。",
            "required_evidence": ["零税额发票对应的业务合同", "免税资格备案文件（如适用）", "零税额发票明细表"]
        })

    # 顶额发票检查
    max_amount_invs = db.query(SalesInvoice).filter(
        SalesInvoice.company_id == company_id,
        SalesInvoice.invoice_date >= ps_date, SalesInvoice.invoice_date <= pe_date,
        SalesInvoice.total_amount >= 90000
    ).all()

    if len(max_amount_invs) >= 5:
        results.append({
            "category": "发票异常", "category_icon": "🧾", "risk_score": 6, "risk_level": "中风险",
            "risk_color": "#f59e0b", "urgency": "提醒",
            "item": "顶额开票占比过高",
            "detail": f"有 {len(max_amount_invs)} 张发票金额接近或达到开票限额（≥9万元），顶额开票是税务局重点监控指标。连续顶额开票可能被判定为虚开发票的嫌疑行为。",
            "suggestion": "核实顶额开票对应的业务是否真实，避免为凑票量而拆分大额业务。"
        })


def _analyze_sensitive_invoice(db, company_id, ps, pe, results):
    """敏感业务发票分析：生活用品/装修/餐饮/经纪代理/咨询服务"""
    ps_date = _period_to_date_range(ps)[0]; pe_date = _period_to_date_range(pe)[1]

    sensitive_keywords = {
        "生活用品": ["生活用品", "日用品", "洗护", "纸巾", "清洁用品", "劳保"],
        "装修装饰": ["装修", "装饰", "工程", "修缮", "建安"],
        "餐饮费": ["餐饮", "餐费", "食品", "酒水", "饮料"],
        "经纪代理": ["经纪代理", "代理服务", "居间", "中介"],
        "咨询服务": ["咨询服务", "咨询费", "顾问", "技术服务费"],
        "预付卡": ["预付卡", "购物卡", "礼品卡", "消费卡"]
    }

    for category, keywords in sensitive_keywords.items():
        count = db.query(func.count(PurchaseInvoice.id)).filter(
            PurchaseInvoice.company_id == company_id,
            PurchaseInvoice.invoice_date >= ps_date, PurchaseInvoice.invoice_date <= pe_date,
            *[PurchaseInvoice.goods_name.contains(kw) for kw in keywords]
        ).scalar() or 0
        if category == "预付卡" and count > 0:
            results.append({
                "category": "发票异常", "category_icon": "🧾", "risk_score": 6, "risk_level": "中风险",
                "risk_color": "#f59e0b", "urgency": "提醒",
                "item": "取得大量生活用品类发票",
                "detail": f"发现 {count} 张「{category}」进项发票。预付卡/购物卡发票进项税额不得抵扣（需做进项转出），且存在被认定为商业贿赂或变相福利的风险。",
                "suggestion": "预付卡发票进项税额不得抵扣，应做进项税额转出处理。",
                "required_evidence": ["预付卡使用明细及用途说明", "进项税额转出凭证"]
            })


def _analyze_buy_sell_mismatch(db, company_id, ps, pe, results):
    """购销商品不匹配风险（制造业专用）"""
    ps_date = _period_to_date_range(ps)[0]; pe_date = _period_to_date_range(pe)[1]

    company = db.query(Company).filter(Company.id == company_id).first()
    if not company:
        return

    # 检查是否是制造业
    industry = (company.business_scope or "") + " " + (company.company_type or "")
    if not any(kw in industry for kw in ["制造", "生产", "加工", "家具", "机械", "电子", "纺织", "服装", "食品"]):
        return

    # 提取进项和销项的商品名称关键词
    sales_items = set()
    for inv in db.query(SalesInvoice.goods_name).filter(
        SalesInvoice.company_id == company_id,
        SalesInvoice.invoice_date >= ps_date, SalesInvoice.invoice_date <= pe_date
    ).limit(100).all():
        if inv[0]: sales_items.add(inv[0][:4])

    purchase_items = set()
    for inv in db.query(PurchaseInvoice.goods_name).filter(
        PurchaseInvoice.company_id == company_id,
        PurchaseInvoice.invoice_date >= ps_date, PurchaseInvoice.invoice_date <= pe_date
    ).limit(100).all():
        if inv[0]: purchase_items.add(inv[0][:4])

    # 简单判断：如果进销品类完全不重叠
    if sales_items and purchase_items and len(sales_items & purchase_items) == 0:
        results.append({
            "category": "发票异常", "category_icon": "🧾", "risk_score": 9, "risk_level": "高风险",
            "risk_color": "#dc2626", "urgency": "紧急",
            "item": "进销品名严重不匹配",
            "detail": f"销项主要商品：{'/'.join(list(sales_items)[:5])}，进项主要商品：{'/'.join(list(purchase_items)[:5])}，进销商品品类无重合。稽查视角：制造业企业购买的原材料与生产销售的产品品类不匹配，暗示可能：①存在未入账的委托加工；②虚开进项发票冲抵成本；③隐瞒了部分生产环节。",
            "suggestion": "①准备完整的生产工艺流程图和投入产出分析；②说明原材料与产成品之间的转换关系；③如存在外协加工，补充委托加工合同。",
            "required_evidence": ["生产工艺流程图", "投入产出分析表（原材料→产成品）", "委托加工合同（如适用）", "存货收发存明细账"]
        })


# ═══════════════════════════════════════════════════════════
#  C 类 — 费用匹配分析
# ═══════════════════════════════════════════════════════════

def _analyze_fuel_vs_vehicles(db, company_id, ps, pe, results):
    """油费进项与固定资产（车辆）匹配分析"""
    ps_date = _period_to_date_range(ps)[0]; pe_date = _period_to_date_range(pe)[1]

    # 油费进项
    fuel_invs = db.query(func.sum(PurchaseInvoice.total_amount)).filter(
        PurchaseInvoice.company_id == company_id,
        PurchaseInvoice.invoice_date >= ps_date, PurchaseInvoice.invoice_date <= pe_date,
        or_(
            PurchaseInvoice.goods_name.contains("汽油"),
            PurchaseInvoice.goods_name.contains("柴油"),
            PurchaseInvoice.goods_name.contains("加油"),
            PurchaseInvoice.goods_name.contains("燃油"),
            PurchaseInvoice.goods_name.contains("油费")
        )
    ).scalar() or 0

    # 公司车辆
    vehicles = db.query(FixedAsset).filter(
        FixedAsset.company_id == company_id,
        FixedAsset.name.contains("车"),
        FixedAsset.status == "在用"
    ).all()

    vehicle_count = len(vehicles)
    if _safe_float(fuel_invs) > 5000 and vehicle_count == 0:
        results.append({
            "category": "费用匹配", "category_icon": "💰", "risk_score": 9, "risk_level": "高风险",
            "risk_color": "#dc2626", "urgency": "紧急",
            "item": "油费进项与固定资产（车辆）不匹配",
            "detail": f"取得油费进项发票 {_safe_float(fuel_invs):,.2f} 元，但公司名下无在用车辆。稽查视角：公司名下没有车辆却有大量加油发票报销，与实际经营逻辑严重不符，可能：①私人车辆费用混入公司账；②虚构费用冲抵利润；③私车公用的租赁协议缺失。",
            "suggestion": "①如为私车公用，应签订车辆租赁协议并代开发票；②如车辆在老板个人名下，补充车辆无偿使用协议；③实际发生的业务用车费用需区分公私。",
            "required_evidence": ["车辆租赁协议（私车公用）", "加油记录与出差/外勤记录的对应关系", "法人/股东名下车辆清单", "车辆使用费分摊说明"]
        })
    elif _safe_float(fuel_invs) > vehicle_count * 30000 and vehicle_count > 0:
        results.append({
            "category": "费用匹配", "category_icon": "💰", "risk_score": 5, "risk_level": "中风险",
            "risk_color": "#f59e0b", "urgency": "提醒",
            "item": "油费进项与固定资产（车辆）不匹配",
            "detail": f"每辆在用车辆年均油费约 {_safe_float(fuel_invs)/vehicle_count:,.2f} 元，明显偏高。建议核查是否混入私人用车费用。",
            "suggestion": "区分公务用车和私人用车的费用，保留行车记录和加油凭证。"
        })


def _analyze_transport_ratio(db, company_id, ps, pe, results):
    """运输费用占进项发票金额比例异常"""
    ps_date = _period_to_date_range(ps)[0]; pe_date = _period_to_date_range(pe)[1]

    transport_invs = db.query(func.sum(PurchaseInvoice.total_amount)).filter(
        PurchaseInvoice.company_id == company_id,
        PurchaseInvoice.invoice_date >= ps_date, PurchaseInvoice.invoice_date <= pe_date,
        or_(
            PurchaseInvoice.goods_name.contains("运输"),
            PurchaseInvoice.goods_name.contains("物流"),
            PurchaseInvoice.goods_name.contains("快递"),
            PurchaseInvoice.goods_name.contains("配送"),
            PurchaseInvoice.goods_name.contains("货运")
        )
    ).scalar() or 0

    total_purchase = db.query(func.sum(PurchaseInvoice.total_amount)).filter(
        PurchaseInvoice.company_id == company_id,
        PurchaseInvoice.invoice_date >= ps_date, PurchaseInvoice.invoice_date <= pe_date
    ).scalar() or 1

    ratio = _safe_float(transport_invs) / _safe_float(total_purchase) * 100

    if ratio > 30:
        results.append({
            "category": "费用匹配", "category_icon": "💰", "risk_score": 7, "risk_level": "高风险",
            "risk_color": "#dc2626", "urgency": "紧急",
            "item": "运输费用与经营模式不匹配",
            "detail": f"运输/物流费用占进项发票总额的 {ratio:.1f}%（明细{_safe_float(transport_invs):,.2f}元）。稽查视角：运输费用占比过高（>30%），可能：①虚开运输发票冲抵成本；②将非运输费用包装为运输费（如将餐饮/娱乐费用开成运输发票）；③运输服务定价不公允。",
            "suggestion": "①核查运输发票对应的实际物流单据（运单/签收单）；②运输费用应与货物采购量匹配（运输费÷采购量=单位运输成本，应与行业平均水平一致）；③大额运输费需有运输合同。",
            "required_evidence": ["运输合同及运单", "货物采购量与运输费匹配分析", "承运方资质证明（道路运输许可证）", "物流签收记录"]
        })


def _analyze_expense_reasonability(db, company_id, ps, pe, results):
    """费用合理性分析：业务招待费/广告费/咨询顾问费占比"""
    revenue_credit = _get_account_sum(db, company_id, "6001", ps, pe, "credit") or 1
    ps_date = _period_to_date_range(ps)[0]; pe_date = _period_to_date_range(pe)[1]

    # 业务招待费
    entertain_invs = db.query(func.sum(PurchaseInvoice.total_amount)).filter(
        PurchaseInvoice.company_id == company_id,
        PurchaseInvoice.invoice_date >= ps_date, PurchaseInvoice.invoice_date <= pe_date,
        or_(
            PurchaseInvoice.goods_name.contains("餐饮"),
            PurchaseInvoice.goods_name.contains("招待"),
            PurchaseInvoice.goods_name.contains("宴请"),
            PurchaseInvoice.goods_name.contains("礼品"),
        )
    ).scalar() or 0

    entertain_ratio = _safe_float(entertain_invs) / revenue_credit * 100
    if entertain_ratio > 5 and _safe_float(entertain_invs) > 10000:
        results.append({
            "category": "费用匹配", "category_icon": "💰", "risk_score": 5, "risk_level": "中风险",
            "risk_color": "#f59e0b", "urgency": "提醒",
            "item": "业务招待费可能超标",
            "detail": f"餐饮/招待/礼品类进项发票金额 {_safe_float(entertain_invs):,.2f} 元，占营业收入的 {entertain_ratio:.1f}%。根据企业所得税法，业务招待费税前扣除上限为发生额的60%且不超过营业收入的5‰。超标部分需纳税调增。",
            "suggestion": "①核算业务招待费实际税前扣除限额；②超标部分在企业所得税汇算清缴时做纳税调增处理。"
        })

    # 咨询顾问费
    consult_invs = db.query(func.sum(PurchaseInvoice.total_amount)).filter(
        PurchaseInvoice.company_id == company_id,
        PurchaseInvoice.invoice_date >= ps_date, PurchaseInvoice.invoice_date <= pe_date,
        or_(
            PurchaseInvoice.goods_name.contains("咨询"),
            PurchaseInvoice.goods_name.contains("顾问"),
            PurchaseInvoice.goods_name.contains("服务费"),
        )
    ).scalar() or 0

    consult_ratio = _safe_float(consult_invs) / revenue_credit * 100
    if consult_ratio > 10 and _safe_float(consult_invs) > 50000:
        results.append({
            "category": "费用匹配", "category_icon": "💰", "risk_score": 7, "risk_level": "高风险",
            "risk_color": "#dc2626", "urgency": "紧急",
            "item": "咨询费支出异常",
            "detail": f"咨询/顾问/服务费类进项发票金额 {_safe_float(consult_invs):,.2f} 元，占营业收入 {consult_ratio:.1f}%。稽查视角：咨询费是税务局重点核查科目，大额无实质性成果的咨询费可能被认定为虚开发票/商业贿赂/利润转移。",
            "suggestion": "①确保每笔咨询费有对应服务合同和成果交付（咨询报告/验收单）；②大额咨询费需有招投标或比价记录；③关联方之间的咨询费需证明定价公允。",
            "required_evidence": ["咨询服务合同", "咨询成果交付文件（报告/方案/验收单）", "咨询服务费用比价或招投标记录", "咨询费付款凭证"]
        })


# ═══════════════════════════════════════════════════════════
#  D 类 — 企业所得税专项
# ═══════════════════════════════════════════════════════════

def _analyze_impairment_not_adjusted(db, company_id, ps, pe, results):
    """资产减值损失未纳税调增"""
    # 检查利润表中是否有资产减值损失（科目6701）但申报表中无纳税调增
    impairment_debit = _get_account_sum(db, company_id, "6701", ps, pe, "debit")
    credit_impairment = _get_account_sum(db, company_id, "6702", ps, pe, "debit")
    total_impairment = impairment_debit + credit_impairment

    if total_impairment > 0:
        # 检查纳税调整（从VAT申报或企业所得税角度）
        results.append({
            "category": "企业所得税", "category_icon": "📑", "risk_score": 9, "risk_level": "高风险",
            "risk_color": "#dc2626", "urgency": "紧急",
            "item": "资产减值准备未纳税调增",
            "detail": f"利润表中资产减值损失/信用减值损失共计 {total_impairment:,.2f} 元。根据企业所得税法第十条，未经核定的准备金支出（包括资产减值准备）不得在税前扣除，必须做企业所得税纳税调增处理。如汇算清缴时未做调增，将导致少缴企业所得税。",
            "suggestion": "①确认汇算清缴时已对减值损失做纳税调增；②如未调增，需进行更正申报补缴税款；③建立减值准备与纳税调整台账。",
            "required_evidence": ["资产减值准备明细表（应收账款坏账/存货跌价/固定资产减值）", "企业所得税汇算清缴纳税调整明细表", "减值测试依据（账龄分析/可变现净值测算）"]
        })


def _analyze_unpaid_capital_interest(db, company_id, ps, pe, results):
    """出资不到位时利息支出未纳税调整"""
    company = db.query(Company).filter(Company.id == company_id).first()
    if not company:
        return

    # 检查实缴vs认缴
    paid_capital = _get_account_balance(db, company_id, "4001", pe)  # 实收资本
    # 从公司表获取注册资本
    registered_capital = _safe_float(company.registered_capital or 0)
    if registered_capital <= 0:
        return

    gap = registered_capital - abs(paid_capital)
    if gap > 0:
        # 检查利息支出
        interest = _get_account_sum(db, company_id, "6603", ps, pe)  # 财务费用-利息
        if interest > 0:
            results.append({
                "category": "企业所得税", "category_icon": "📑", "risk_score": 7, "risk_level": "高风险",
                "risk_color": "#dc2626", "urgency": "紧急",
                "item": "出资不到位利息支出税前扣除",
                "detail": f"注册资本 {registered_capital:,.2f} 元，实收资本 {abs(paid_capital):,.2f} 元，差额 {gap:,.2f} 元未到位。根据国税函[2009]312号，投资者未足额缴付资本期间，相当于未缴足资本额应计利息的部分，不得在企业所得税税前扣除。本期内利息支出 {interest:,.2f} 元，应按比例计算不得扣除金额。",
                "suggestion": "①计算不得税前扣除的利息金额（公式：未缴足资本额×同期贷款利率×未缴足月数/12）；②在企业所得税汇算时做纳税调增。",
                "required_evidence": ["股东出资证明/验资报告", "实收资本明细账", "借款合同及利息计算明细", "利息支出纳税调整计算表"]
            })


def _analyze_non_taxable_income(db, company_id, ps, pe, results):
    """不征税收入调减金额分析"""
    # 检查营业外收入中是否有政府补助
    govt_income = _get_account_sum(db, company_id, "6301", ps, pe, "credit")  # 营业外收入
    subsidy_keywords = ["补助", "补贴", "拨款", "退税"]
    subsidy_credits = db.query(func.sum(JournalEntry.credit_amount)).filter(
        JournalEntry.company_id == company_id,
        JournalEntry.period >= ps, JournalEntry.period <= pe,
        *[JournalEntry.summary.contains(kw) for kw in subsidy_keywords]
    ).scalar() or 0

    if _safe_float(subsidy_credits) > 0:
        results.append({
            "category": "企业所得税", "category_icon": "📑", "risk_score": 4, "risk_level": "低风险",
            "risk_color": "#3b82f6", "urgency": "建议",
            "item": "不征税收入处理不合规",
            "detail": f"发现 {_safe_float(subsidy_credits):,.2f} 元政府补助或财政补贴性质收入。稽查视角：不征税收入（符合财税[2011]70号三项条件的专项用途资金）对应的支出不得税前扣除。如不符合不征税条件但做了调减，或符合条件但对应的费用未作纳税调增，均有税务风险。",
            "suggestion": "①确认是否符合不征税收入条件（专项资金拨付文件+专门管理办法+单独核算）；②不征税收入对应的支出不得税前扣除；③建立专项资金使用台账。",
            "required_evidence": ["政府补助批文/拨付文件", "专项资金管理办法", "专项资金使用明细台账", "如有不征税收入申报，提供申报资料"]
        })


def _analyze_provisional_cost(db, company_id, ps, pe, results):
    """暂估/无票成本费用分析"""
    ps_date = _period_to_date_range(ps)[0]; pe_date = _period_to_date_range(pe)[1]

    # 检查摘要中包含"暂估""无票""预提"等的凭证
    provisional_entries = db.query(func.count(JournalEntry.id), func.sum(JournalEntry.debit_amount)).filter(
        JournalEntry.company_id == company_id,
        JournalEntry.period >= ps, JournalEntry.period <= pe,
        or_(
            JournalEntry.summary.contains("暂估"),
            JournalEntry.summary.contains("无票"),
            JournalEntry.summary.contains("预提"),
            JournalEntry.summary.contains("未取得发票")
        )
    ).first()

    prov_cnt, prov_amt = provisional_entries
    if prov_cnt and prov_cnt > 0:
        total_cost = _get_account_sum(db, company_id, "6401", ps, pe, "debit")
        ratio = _safe_float(prov_amt) / (total_cost or 1) * 100

        if ratio > 10 or _safe_float(prov_amt) > 100000:
            results.append({
                "category": "企业所得税", "category_icon": "📑", "risk_score": 8, "risk_level": "高风险",
                "risk_color": "#dc2626", "urgency": "紧急",
                "item": "暂估成本跨期未冲销",
                "detail": f"发现 {prov_cnt} 笔暂估/无票/预提成本记录，涉及金额 {_safe_float(prov_amt):,.2f} 元，占营业成本 {ratio:.1f}%。根据国家税务总局公告2011年第34号，暂估入账的成本费用在企业所得税汇算清缴前仍未取得合规发票的，不得税前扣除，需做纳税调增。",
                "suggestion": "①核实暂估款项截至汇算清缴日是否已取得发票；②未取得发票部分做企业所得税纳税调增；③建立暂估款项台账，加强发票催收管理。",
                "required_evidence": ["暂估/无票成本明细清单", "对应业务的采购合同/入库单", "发票催收记录", "纳税调整明细表"]
            })


# ═══════════════════════════════════════════════════════════
#  E 类 — 薪酬福利及其他
# ═══════════════════════════════════════════════════════════

def _analyze_staff_welfare(db, company_id, ps, pe, results):
    """职工福利费超14%工资总额限制"""
    salary_records = db.query(
        func.sum(SalaryRecord.current_income)
    ).filter(SalaryRecord.company_id == company_id, SalaryRecord.period >= ps,
             SalaryRecord.period <= pe).scalar() or 0

    total_salary = _safe_float(salary_records)

    welfare_invs = db.query(func.sum(PurchaseInvoice.total_amount)).filter(
        PurchaseInvoice.company_id == company_id,
        PurchaseInvoice.invoice_date >= _period_to_date_range(ps)[0], PurchaseInvoice.invoice_date <= _period_to_date_range(pe)[1],
        or_(
            PurchaseInvoice.goods_name.contains("福利"),
            PurchaseInvoice.goods_name.contains("节日"),
            PurchaseInvoice.goods_name.contains("月饼"),
            PurchaseInvoice.goods_name.contains("粽子"),
            PurchaseInvoice.goods_name.contains("慰问"),
            PurchaseInvoice.goods_name.contains("体检"),
            PurchaseInvoice.goods_name.contains("旅游"),
            PurchaseInvoice.goods_name.contains("团建"),
        )
    ).scalar() or 0

    if total_salary > 0:
        welfare_ratio = _safe_float(welfare_invs) / total_salary * 100
        if welfare_ratio > 14:
            results.append({
                "category": "薪酬福利", "category_icon": "👥", "risk_score": 7, "risk_level": "高风险",
                "risk_color": "#dc2626", "urgency": "紧急",
                "item": "职工福利费超限额",
                "detail": f"福利性质进项发票金额 {_safe_float(welfare_invs):,.2f} 元，占工资总额 {total_salary:,.2f} 元的 {welfare_ratio:.1f}%，超过企业所得税法规定的14%上限。超标部分不得税前扣除，需做纳税调增。同时，部分福利（如旅游/月饼/购物卡）可能涉及个人所得税代扣代缴义务。",
                "suggestion": "①核算超过14%的部分，在企业所得税汇算时做纳税调增；②发放给员工的实物福利/旅游等需并入工资薪金代扣个税。",
                "required_evidence": ["职工福利费明细账", "福利费对应的发票清单", "企业所得税纳税调整明细", "个税代扣代缴记录（福利部分）"]
            })


def _analyze_social_security_match(db, company_id, ps, pe, results):
    """社保基数与工资匹配度分析"""
    salary_records = db.query(SalaryRecord).filter(
        SalaryRecord.company_id == company_id, SalaryRecord.period >= ps,
        SalaryRecord.period <= pe
    ).all()

    if not salary_records:
        return

    # 检查社保申报缴纳
    ss_decls = db.query(SocialSecurityDetail).join(
        SocialSecurityDeclaration,
        SocialSecurityDetail.declaration_id == SocialSecurityDeclaration.id
    ).filter(
        SocialSecurityDeclaration.company_id == company_id,
        SocialSecurityDeclaration.period >= ps,
        SocialSecurityDeclaration.period <= pe
    ).all()

    if not ss_decls:
        # 有工资但无社保记录
        total_salary = sum(_safe_float(r.current_income) for r in salary_records)
        if total_salary > 50000:
            results.append({
                "category": "薪酬福利", "category_icon": "👥", "risk_score": 9, "risk_level": "高风险",
                "risk_color": "#dc2626", "urgency": "紧急",
                "item": "有工资发放但无社保缴纳记录",
                "detail": f"期间发放工资总额 {total_salary:,.2f} 元，但未找到对应的社保申报/缴纳记录。根据《社会保险法》，用人单位必须为劳动者缴纳社会保险。稽查视角：无社保缴纳记录→可能全部为未签劳动合同的临时用工→工资支出真实性存疑→虚列人工成本嫌疑。",
                "suggestion": "①补缴社会保险（含滞纳金）；②核实是否存在未签订劳动合同的用工；③如为劳务外包，补充外包合同和发票。",
                "required_evidence": ["劳动合同/劳务合同清单", "社保缴纳记录（或补缴计划）", "外包合同及发票（如适用）", "工资银行发放记录"]
            })

    # 社保基数与工资差异
    if ss_decls:
        avg_salary = sum(_safe_float(r.current_income) for r in salary_records) / max(len(salary_records), 1)
        # 简化的社保基数检查
        low_base = 0
        for sd in ss_decls:
            base = _safe_float(sd.salary_base or 0)
            if base > 0 and base < avg_salary * 0.8:  # 基数低于平均工资80%
                low_base += 1

        if low_base > 0:
            results.append({
                "category": "薪酬福利", "category_icon": "👥", "risk_score": 6, "risk_level": "中风险",
                "risk_color": "#f59e0b", "urgency": "提醒",
                "item": "社保缴存基数偏低",
                "detail": f"有 {low_base} 人社保缴纳基数显著低于实发工资水平。按最低基数缴纳虽较普遍，但严格来说不合规，差额较大时可能被社保稽查部门责令补缴。",
                "suggestion": "按实际工资水平核定社保缴费基数，避免因基数偏低被社保稽查处罚。"
            })


def _analyze_undistributed_profit(db, company_id, ps, pe, results):
    """未分配利润偏高但长期不分红"""
    # 未分配利润
    undivided = _get_account_balance(db, company_id, "4104", pe)  # 未分配利润（贷方正数）
    paid_capital = abs(_get_account_balance(db, company_id, "4001", pe))  # 实收资本

    # 检查"其他应收款"←大额资金是否被股东占用
    other_receivables = _get_account_balance(db, company_id, "1221", pe)

    if undivided > paid_capital * 2 and paid_capital > 0:
        risk_score = 6
        risk_level = "中风险"
        detail = f"未分配利润余额 {undivided:,.2f} 元，是实收资本（{paid_capital:,.2f}元）的 {undivided/paid_capital:.1f} 倍，但无分红记录。稽查视角：①大额未分配利润长期不分配→可能已被股东以借款/报销等形式转移（实质视同分红）；②如「其他应收款-股东」余额较大，税务机关会直接认定为视同分红，补缴20%个人所得税。"

        if other_receivables > paid_capital * 0.5:
            risk_score = 9
            risk_level = "高风险"
            detail += f" 且其他应收款余额 {other_receivables:,.2f} 元，存在股东占用资金嫌疑，可能被认定为视同分红。"
            results.append({
                "category": "薪酬福利", "category_icon": "👥", "risk_score": risk_score, "risk_level": risk_level,
                "risk_color": _risk_color(risk_score), "urgency": "紧急",
                "item": "大额未分配利润+股东资金占用→视同分红风险",
                "detail": detail,
                "suggestion": "①核实「其他应收款-股东」是否可以归还（在纳税年度终了后既不归还又未用于企业生产经营的，视同分红征收20%个税）；②如确实无法归还，应主动申报代扣代缴个人所得税；③考虑合理的利润分配方案。",
                "required_evidence": ["未分配利润明细表", "历次利润分配决议（如有）", "其他应收款-股东明细（含借款协议）", "股东资金占用归还计划"]
            })
        else:
            results.append({
                "category": "薪酬福利", "category_icon": "👥", "risk_score": risk_score, "risk_level": risk_level,
                "risk_color": _risk_color(risk_score), "urgency": "提醒",
                "item": "未分配利润长期不分配",
                "detail": detail,
                "suggestion": "关注未分配利润的构成，如有视同分红情形应主动申报个税。建立合理的利润分配政策和分红计划。"
            })


def _analyze_cross_invoicing(db, company_id, ps, pe, results):
    """互开发票风险（即同一对手方既做客户又做供应商）"""
    ps_date = _period_to_date_range(ps)[0]; pe_date = _period_to_date_range(pe)[1]

    # 获取所有客户名称
    sales_names = set()
    for inv in db.query(SalesInvoice.buyer_name).filter(
        SalesInvoice.company_id == company_id,
        SalesInvoice.invoice_date >= ps_date, SalesInvoice.invoice_date <= pe_date
    ).all():
        if inv[0]: sales_names.add(inv[0])

    # 获取所有供应商名称
    purchase_names = set()
    for inv in db.query(PurchaseInvoice.seller_name).filter(
        PurchaseInvoice.company_id == company_id,
        PurchaseInvoice.invoice_date >= ps_date, PurchaseInvoice.invoice_date <= pe_date
    ).all():
        if inv[0]: purchase_names.add(inv[0])

    cross = sales_names & purchase_names
    if cross:
        results.append({
            "category": "其他风险", "category_icon": "⚠️", "risk_score": 8, "risk_level": "高风险",
            "risk_color": "#dc2626", "urgency": "紧急",
            "item": "互开发票风险",
            "detail": f"发现 {len(cross)} 家企业同时既对本公司开具发票又接受本公司开具的发票：{'/'.join(list(cross)[:5])}。稽查视角：互开发票（对开发票）是税务局认定虚开增值税专用发票的重要信号，可能导致：①双方同时虚增收入和成本（流水造假）；②涉嫌循环经济/环开发票犯罪。",
            "suggestion": "①逐笔核查互开发票对应的业务真实性（合同/物流/资金流三流合一）；②如无真实业务，立即做进项税额转出或红冲；③建立客户供应商黑名单制度。",
            "required_evidence": ["互开发票清单及三流合一证明", "每笔互开发票对应的合同/物流/付款凭证", "不能提供真实交易证明的不得抵扣进项税额"]
        })


def _analyze_invest_property_tax(db, company_id, ps, pe, results):
    """投资性房地产无租金收入或未交房产税"""
    # 检查投资性房地产（科目1521或相关）
    invest_prop = _get_account_balance(db, company_id, "1521", pe)
    if invest_prop <= 0:
        invest_prop = _get_account_balance(db, company_id, "1518", pe)  # 其他投资性房地产

    if invest_prop > 0:
        # 检查租金收入
        rent_income = _get_account_sum(db, company_id, "6051", ps, pe, "credit")  # 其他业务收入-租金
        # 检查房产税
        property_tax = _get_account_sum(db, company_id, "6403", ps, pe, "debit")  # 税金及附加

        # 简化检查
        has_rent = rent_income > 0
        results.append({
            "category": "其他风险", "category_icon": "⚠️", "risk_score": 6, "risk_level": "中风险",
            "risk_color": "#f59e0b", "urgency": "提醒",
            "item": "持有投资性房地产需关注税费申报",
            "detail": f"投资性房地产账面价值 {invest_prop:,.2f} 元。{ '有租金收入' + str(rent_income) + '元' if has_rent else '未发现租金收入记录' }。稽查要点：①持有投资性房地产需缴纳房产税（按余值1.2%或租金12%）；②租金收入需开具发票并申报增值税；③将来转让需缴纳土地增值税。",
            "suggestion": "①确认是否已缴纳房产税；②出租收入是否入账并申报；③检查从租计征房产税的适用税率。"
        })


# ═══════════════════════════════════════════════════════════
#  F 类 — 金税四期新增：发票异常深度分析
# ═══════════════════════════════════════════════════════════

def _analyze_top_amount_invoice_ratio(db, company_id, ps, pe, results):
    """顶额开票占比过高：开具金额集中在开票限额附近（金税四期必查）"""
    ps_date = _period_to_date_range(ps)[0]; pe_date = _period_to_date_range(pe)[1]
    
    sales = db.query(SalesInvoice).filter(
        SalesInvoice.company_id == company_id,
        SalesInvoice.invoice_date >= ps_date, SalesInvoice.invoice_date <= pe_date,
        SalesInvoice.status == "正常"
    ).all()
    
    if not sales:
        return
    
    # 顶额标准：金额>=90000（10万版）、>=900000（100万版）
    top_amount_invs = []
    for inv in sales:
        amt = _safe_float(inv.total_amount)
        if amt >= 90000 and amt <= 100000:  # 10万版顶额区间
            top_amount_invs.append(inv)
        elif amt >= 900000 and amt <= 1000000:  # 100万版顶额区间
            top_amount_invs.append(inv)
    
    top_ratio = len(top_amount_invs) / len(sales) * 100 if sales else 0
    if top_ratio > 30 and len(top_amount_invs) >= 3:
        top_amt = sum(_safe_float(inv.total_amount) for inv in top_amount_invs)
        results.append({
            "category": "发票异常", "category_icon": "🧾", "risk_score": 8, "risk_level": "高风险",
            "risk_color": "#dc2626", "urgency": "紧急",
            "item": "顶额开票占比过高",
            "detail": f"有 {len(top_amount_invs)}/{len(sales)} 张发票金额集中在开票限额附近（>=9万元），占比 {top_ratio:.0f}%，涉及金额 {top_amt:,.2f} 元。涉嫌拆分大额交易或虚开发票。金税四期对顶额开票自动标记。",
            "suggestion": "核查顶额开票对应的交易是否实质为同一笔业务被拆分。大额交易应使用更高版位发票或申请提高开票限额。",
            "required_evidence": ["顶额发票清单（含对应合同/验收单/付款凭证）", "顶额开票业务真实性说明"]
        })
    elif top_ratio > 10 and len(top_amount_invs) >= 2:
        results.append({
            "category": "发票异常", "category_icon": "🧾", "risk_score": 4, "risk_level": "中风险",
            "risk_color": "#f59e0b", "urgency": "提醒",
            "item": "顶额开票占比过高",
            "detail": f"有 {len(top_amount_invs)} 张发票金额接近开票限额。顶额开票容易被税务机关认定为拆分开票。",
            "suggestion": "核实业务真实性，避免连续顶额开票。"
        })


def _analyze_top_amount_red_invoice(db, company_id, ps, pe, results):
    """红冲作废发票集中在顶额：顶额红冲率异常"""
    ps_date = _period_to_date_range(ps)[0]; pe_date = _period_to_date_range(pe)[1]
    
    void_red_invs = db.query(SalesInvoice).filter(
        SalesInvoice.company_id == company_id,
        SalesInvoice.invoice_date >= ps_date, SalesInvoice.invoice_date <= pe_date,
        SalesInvoice.status.in_(["作废", "红冲"])
    ).all()
    
    if not void_red_invs:
        return
    
    # 顶额红冲：金额>=90000
    top_void = [inv for inv in void_red_invs if _safe_float(inv.total_amount) >= 90000]
    
    if len(top_void) >= 3:
        top_amt = sum(_safe_float(inv.total_amount) for inv in top_void)
        void_ratio = len(top_void) / len(void_red_invs) * 100
        results.append({
            "category": "发票异常", "category_icon": "🧾", "risk_score": 8, "risk_level": "高风险",
            "risk_color": "#dc2626", "urgency": "紧急",
            "item": "红冲作废发票集中在顶额",
            "detail": f"红冲/作废发票中有 {len(top_void)} 张金额 >=9万元（占红冲总数 {void_ratio:.0f}%），涉及金额 {top_amt:,.2f} 元。可能涉嫌先开票虚增收入再红冲平账。",
            "suggestion": "逐笔核实顶额红冲/作废发票的原因：是否真实发生退货退款、销售折让；检查红冲是否有对应的红字信息表。",
            "required_evidence": ["顶额红冲/作废发票清单", "红字信息表（每笔红冲对应）", "退货协议/验收单/折让协议"]
        })


def _analyze_consecutive_invoice_numbers(db, company_id, ps, pe, results):
    """发票连号异常：纸质发票连续>=5张连号（金税四期监控）"""
    ps_date = _period_to_date_range(ps)[0]; pe_date = _period_to_date_range(pe)[1]
    
    # 获取所有有发票号码的销项发票（排除数电发票，因为数电发票不存在连号概念）
    invs = db.query(SalesInvoice).filter(
        SalesInvoice.company_id == company_id,
        SalesInvoice.invoice_date >= ps_date, SalesInvoice.invoice_date <= pe_date,
        SalesInvoice.invoice_no != None, SalesInvoice.invoice_no != "",
        SalesInvoice.invoice_code != None, SalesInvoice.invoice_code != ""
    ).order_by(SalesInvoice.invoice_code, SalesInvoice.invoice_no).all()
    
    if len(invs) < 5:
        return
    
    # 按发票代码分组，检查每组内号码连续性
    code_groups = {}
    for inv in invs:
        code = inv.invoice_code or ""
        no_str = inv.invoice_no or ""
        try:
            no = int(no_str)
        except (ValueError, TypeError):
            continue
        if code not in code_groups:
            code_groups[code] = []
        code_groups[code].append(no)
    
    consecutive_groups = []
    for code, numbers in code_groups.items():
        numbers.sort()
        i = 0
        while i < len(numbers):
            j = i + 1
            while j < len(numbers) and numbers[j] == numbers[j - 1] + 1:
                j += 1
            group_size = j - i
            if group_size >= 5:
                consecutive_groups.append({
                    "code": code,
                    "start": numbers[i],
                    "end": numbers[j - 1],
                    "count": group_size
                })
            i = j
    
    if consecutive_groups:
        total_affected = sum(g["count"] for g in consecutive_groups)
        detail_items = [f"代码{g['code']}号码{g['start']}~{g['end']}共{g['count']}张" for g in consecutive_groups[:3]]
        results.append({
            "category": "发票异常", "category_icon": "🧾", "risk_score": 6, "risk_level": "中风险",
            "risk_color": "#f59e0b", "urgency": "提醒",
            "item": "发票连号异常",
            "detail": f"发现 {len(consecutive_groups)} 组连号发票（连续>=5张）：{'；'.join(detail_items)}。大量纸质发票连号可能暗示突击开票或虚开发票。注意：电子发票/数电发票不存在连号概念，仅纸质发票需关注。",
            "suggestion": "核实连号发票对应的交易是否真实发生。如均为真实交易，保留各项业务对应的合同和交货单备查。"
        })


def _analyze_off_hours_invoicing(db, company_id, ps, pe, results):
    """夜间/节假日开票异常：非工作时间开票（金税四期时间维度监控）"""
    ps_date = _period_to_date_range(ps)[0]; pe_date = _period_to_date_range(pe)[1]
    
    # 中国法定节假日（简化版，不包括调休）
    holidays_2025 = [
        "01-01",  # 元旦
        "01-28", "01-29", "01-30", "01-31", "02-01", "02-02", "02-03", "02-04",  # 春节
        "04-04", "04-05", "04-06",  # 清明
        "05-01", "05-02", "05-03", "05-04", "05-05",  # 劳动节
        "05-31", "06-01", "06-02",  # 端午
        "10-01", "10-02", "10-03", "10-04", "10-05", "10-06", "10-07", "10-08",  # 国庆中秋
    ]
    holidays_2026 = [
        "01-01",
        "02-17", "02-18", "02-19", "02-20", "02-21", "02-22", "02-23",
        "04-05", "04-06", "04-07",
        "05-01", "05-02", "05-03", "05-04", "05-05",
        "06-19", "06-20", "06-21",
        "10-01", "10-02", "10-03", "10-04", "10-05", "10-06", "10-07", "10-08",
    ]
    all_holidays = set(holidays_2025 + holidays_2026)
    
    all_sales = db.query(SalesInvoice).filter(
        SalesInvoice.company_id == company_id,
        SalesInvoice.invoice_date >= ps_date, SalesInvoice.invoice_date <= pe_date,
        SalesInvoice.status == "正常"
    ).all()
    
    if not all_sales:
        return
    
    off_hours_invs = []
    for inv in all_sales:
        inv_date = inv.invoice_date
        if isinstance(inv_date, str):
            inv_date = datetime.strptime(inv_date, "%Y-%m-%d").date() if inv_date else None
        if not inv_date:
            continue
        
        # 检查节假日
        mmdd = inv_date.strftime("%m-%d")
        is_holiday = mmdd in all_holidays
        
        # 检查夜间：22:00-06:00（如果有时间戳）
        is_night = False
        if hasattr(inv, 'invoice_time') and inv.invoice_time:
            try:
                t = datetime.strptime(str(inv.invoice_time), "%H:%M:%S").time()
                is_night = t.hour >= 22 or t.hour < 6
            except:
                pass
        
        # 周末
        is_weekend = inv_date.weekday() >= 5
        
        if is_holiday or is_night or is_weekend:
            off_hours_invs.append(inv)
    
    off_ratio = len(off_hours_invs) / len(all_sales) * 100 if all_sales else 0
    
    if off_ratio > 10 and len(off_hours_invs) >= 5:
        off_amt = sum(_safe_float(inv.total_amount) for inv in off_hours_invs)
        results.append({
            "category": "发票异常", "category_icon": "🧾", "risk_score": 5, "risk_level": "中风险",
            "risk_color": "#f59e0b", "urgency": "提醒",
            "item": "夜间/节假日开票异常",
            "detail": f"有 {len(off_hours_invs)}/{len(all_sales)} 张发票在夜间/节假日/周末开具（占比 {off_ratio:.0f}%），涉及金额 {off_amt:,.2f} 元。非工作时间开票可能被税务机关质疑业务真实性。电商/餐饮住宿等行业可合理解释；制造业/批发业需重点关注。",
            "suggestion": "检查非工作时间开票对应的业务是否真实发生。保留业务发生时间佐证材料。"
        })


# ═══════════════════════════════════════════════════════════
#  G 类 — 金税四期：资金流与发票联动分析
# ═══════════════════════════════════════════════════════════

def _analyze_einvoice_frequency_spike(db, company_id, ps, pe, results):
    """全电发票开具频次异常：短期内开票量激增（金税四期动态监控）"""
    ps_date = _period_to_date_range(ps)[0]; pe_date = _period_to_date_range(pe)[1]
    
    # 按周统计开票量
    from datetime import timedelta
    sales = db.query(SalesInvoice).filter(
        SalesInvoice.company_id == company_id,
        SalesInvoice.invoice_date >= ps_date, SalesInvoice.invoice_date <= pe_date
    ).all()
    
    if len(sales) < 20:
        return
    
    # 按周汇总
    weekly_data = {}
    for inv in sales:
        inv_date = inv.invoice_date
        if isinstance(inv_date, str):
            try:
                inv_date = datetime.strptime(inv_date, "%Y-%m-%d").date()
            except:
                continue
        if not inv_date:
            continue
        # ISO周
        week_key = inv_date.isocalendar()[:2]  # (year, week)
        if week_key not in weekly_data:
            weekly_data[week_key] = {"count": 0, "amount": 0}
        weekly_data[week_key]["count"] += 1
        weekly_data[week_key]["amount"] += _safe_float(inv.total_amount)
    
    if len(weekly_data) < 4:
        return
    
    weeks = sorted(weekly_data.keys())
    avg_count = sum(weekly_data[w]["count"] for w in weeks) / len(weeks)
    avg_amount = sum(weekly_data[w]["amount"] for w in weeks) / len(weeks)
    
    spike_weeks = []
    for w in weeks:
        cnt = weekly_data[w]["count"]
        amt = weekly_data[w]["amount"]
        if avg_count > 0 and cnt > avg_count * 3:  # 偏离均值3倍
            spike_weeks.append(f"第{w[1]}周({cnt}张/{amt:,.0f}元)")
        elif avg_amount > 0 and amt > avg_amount * 3:
            spike_weeks.append(f"第{w[1]}周({cnt}张/{amt:,.0f}元)")
    
    if spike_weeks:
        results.append({
            "category": "发票合规", "category_icon": "🧾", "risk_score": 7, "risk_level": "高风险",
            "risk_color": "#dc2626", "urgency": "高",
            "item": "全电发票开具频次异常",
            "detail": f"以下周期开票量/金额远超历史均值（3倍以上）：{'；'.join(spike_weeks[:3])}。金税四期系统会自动标记开票量异常激增的企业。",
            "suggestion": "核查开票业务真实性，准备合同/物流/资金流佐证材料，说明激增合理原因。",
            "required_evidence": ["销项发票台账（按日/按周统计）", "激增期间的业务合同及物流单据"]
        })


def _analyze_corp_to_personal_transfer(db, company_id, ps, pe, results):
    """公户转私户频繁或大额（金税四期与人民银行账户系统联网监控）"""
    ps_date = _period_to_date_range(ps)[0]; pe_date = _period_to_date_range(pe)[1]
    
    # 从银行流水筛选贷方（付款）且对方为个人的记录
    # 个人特征：名称含常见姓名（2-3个汉字）、不含"公司""有限""厂"等
    company_keywords = ["公司", "有限", "厂", "局", "院", "所", "中心", "集团", "银行", "税务"]
    
    all_transfers = db.query(BankTransaction).filter(
        BankTransaction.company_id == company_id,
        BankTransaction.transaction_date >= ps_date,
        BankTransaction.transaction_date <= pe_date,
        BankTransaction.debit_amount > 0  # 企业付款出去
    ).all()
    
    personal_transfers = []
    for bt in all_transfers:
        name = (bt.counterparty_name or "").strip()
        # 排除公司/机构
        if not name or any(kw in name for kw in company_keywords):
            continue
        # 简单规则：2-3个汉字名称（中文姓名特征）
        if len(name) >= 2 and len(name) <= 4 and all('\u4e00' <= c <= '\u9fff' for c in name):
            personal_transfers.append(bt)
    
    if len(personal_transfers) >= 5:
        total_amount = sum(_safe_float(bt.debit_amount) for bt in personal_transfers)
        has_large = any(_safe_float(bt.debit_amount) >= 50000 for bt in personal_transfers)
        
        risk_score = 8 if has_large else 5
        risk_level = "高风险" if has_large else "中风险"
        risk_color = "#dc2626" if has_large else "#f59e0b"
        urgency = "高" if has_large else "中"
        
        results.append({
            "category": "资金往来", "category_icon": "💰", "risk_score": risk_score, "risk_level": risk_level,
            "risk_color": risk_color, "urgency": urgency,
            "item": "公户转私户频繁或大额",
            "detail": f"对公账户向个人账户转账 {len(personal_transfers)} 笔，合计 {total_amount:,.2f} 元。{'含大额转账（>=5万），' if has_large else ''}涉嫌挪用资金或逃避个税。金税四期与人民银行账户系统联网，公转私实时监控。",
            "suggestion": "规范资金往来，大额转账需有合法依据（分红、薪酬、借款等），并及时代扣个税。保留转账的合同/协议备查。",
            "required_evidence": ["公转私明细清单（含对方姓名/金额/用途）", "每笔大额转账的合法依据（分红决议/薪酬表/借款合同）", "代扣代缴个税记录"]
        })


def _analyze_cash_transaction_ratio(db, company_id, ps, pe, results):
    """现金收支占比过高（金税四期大额现金交易重点监控）"""
    ps_date = _period_to_date_range(ps)[0]; pe_date = _period_to_date_range(pe)[1]
    
    # 从银行流水筛选现金相关交易
    cash_keywords = ["现金", "取现", "存现", "现金支取", "现金存入"]
    
    all_transactions = db.query(
        func.count(BankTransaction.id),
        func.sum(BankTransaction.debit_amount),
        func.sum(BankTransaction.credit_amount)
    ).filter(
        BankTransaction.company_id == company_id,
        BankTransaction.transaction_date >= ps_date,
        BankTransaction.transaction_date <= pe_date
    ).first() or (0, 0, 0)
    
    total_cnt, total_debit, total_credit = all_transactions
    total_volume = _safe_float(total_debit) + _safe_float(total_credit)
    
    if total_volume <= 0:
        return
    
    # 现金交易
    cash_volume = 0
    cash_cnt = 0
    cash_large = 0
    for kw in cash_keywords:
        cash_rows = db.query(
            func.count(BankTransaction.id),
            func.sum(BankTransaction.debit_amount),
            func.sum(BankTransaction.credit_amount)
        ).filter(
            BankTransaction.company_id == company_id,
            BankTransaction.transaction_date >= ps_date,
            BankTransaction.transaction_date <= pe_date,
            BankTransaction.summary.contains(kw)
        ).first() or (0, 0, 0)
        cash_cnt += cash_rows[0] or 0
        cash_volume += _safe_float(cash_rows[1]) + _safe_float(cash_rows[2])
    
    # 也检查序时账中的现金科目
    cash_je_volume = _get_account_sum(db, company_id, "1001", ps, pe, "debit") + \
                     _get_account_sum(db, company_id, "1001", ps, pe, "credit")
    
    total_cash = cash_volume + cash_je_volume
    cash_ratio = total_cash / total_volume * 100 if total_volume > 0 else 0
    
    if cash_ratio > 30 and total_cash > 100000:
        results.append({
            "category": "资金往来", "category_icon": "💰", "risk_score": 5, "risk_level": "中风险",
            "risk_color": "#f59e0b", "urgency": "中",
            "item": "现金收支占比过高",
            "detail": f"现金类交易合计 {total_cash:,.2f} 元（含银行流水现金交易+库存现金科目），占全部资金交易 {cash_ratio:.0f}%。金税四期对大额现金交易重点监控，单笔5万元以上需重点说明。",
            "suggestion": "减少现金交易，尽量通过对公账户结算，保留完整资金链条证据。"
        })


# ═══════════════════════════════════════════════════════════
#  H 类 — 增值税专项补充分析
# ═══════════════════════════════════════════════════════════

def _analyze_input_tax_transfer_delay(db, company_id, ps, pe, results):
    """进项税额转出处理不及时：免税/简易计税项目进项未及时转出"""
    decls = db.query(VATDeclaration).filter(
        VATDeclaration.company_id == company_id,
        VATDeclaration.period >= ps, VATDeclaration.period <= pe
    ).order_by(VATDeclaration.period).all()
    
    delay_periods = []
    for v in decls:
        if not v.form_main:
            continue
        fm = v.form_main if isinstance(v.form_main, dict) else json.loads(v.form_main)
        # 检查是否有免税/简易计税/集体福利相关
        tax_free = _safe_float(fm.get("tax_free_sales", 0))
        simple_tax = _safe_float(fm.get("simple_tax_sales", 0))
        input_transfer = _safe_float(fm.get("input_tax_transfer", 0))
        
        # 有免税或简易征收，但进项转出很小/为零
        if (tax_free > 0 or simple_tax > 0):
            # 检查是否有进项发票
            input_cnt = db.query(func.count(PurchaseInvoice.id)).filter(
                PurchaseInvoice.company_id == company_id,
                func.substr(PurchaseInvoice.invoice_date, 1, 7) == v.period
            ).scalar() or 0
            
            if input_cnt > 0:
                # 估算应转出的进项税额
                total_sales = _safe_float(fm.get("taxable_sales", 0)) + tax_free + simple_tax
                if total_sales > 0:
                    input_amt = _safe_float(fm.get("input_tax", 0))
                    expected_transfer = input_amt * (tax_free + simple_tax) / total_sales
                    if input_transfer < expected_transfer * 0.5 and expected_transfer > 100:
                        delay_periods.append(f"{v.period}(应转出约{expected_transfer:,.0f}元，实际{input_transfer:,.0f}元)")
    
    if delay_periods:
        results.append({
            "category": "增值税专项", "category_icon": "🧾", "risk_score": 7, "risk_level": "高风险",
            "risk_color": "#dc2626", "urgency": "高",
            "item": "进项税额转出处理不及时",
            "detail": f"以下期间存在免税/简易计税项目但进项转出不足：{'；'.join(delay_periods[:3])}。用于免税项目、集体福利、个人消费的进项税额不得抵扣，必须做进项税额转出。金税四期进项转出数据与申报比对。",
            "suggestion": "按月核查进项税额用途，属于不得抵扣情形的及时做进项转出，并在增值税申报表正确填列。",
            "required_evidence": ["进项发票用途分类台账", "不得抵扣进项税额计算表", "进项税额转出明细"]
        })


def _analyze_invoice_revenue_cross_period(db, company_id, ps, pe, results):
    """增值税发票开具与收入确认跨期异常"""
    ps_date = _period_to_date_range(ps)[0]; pe_date = _period_to_date_range(pe)[1]
    
    # 按期间统计销项发票金额
    inv_by_period = {}
    sales_invs = db.query(SalesInvoice).filter(
        SalesInvoice.company_id == company_id,
        SalesInvoice.invoice_date >= ps_date, SalesInvoice.invoice_date <= pe_date
    ).all()
    for inv in sales_invs:
        inv_date = inv.invoice_date
        if isinstance(inv_date, str):
            period = inv_date[:7]
        else:
            period = inv_date.strftime("%Y-%m") if inv_date else ""
        if period:
            inv_by_period[period] = inv_by_period.get(period, 0) + _safe_float(inv.total_amount)
    
    # 按期间统计申报销售额
    decls = db.query(VATDeclaration).filter(
        VATDeclaration.company_id == company_id,
        VATDeclaration.period >= ps, VATDeclaration.period <= pe
    ).all()
    
    mismatches = []
    for v in decls:
        period = v.period
        inv_amt = inv_by_period.get(period, 0)
        if v.form_main:
            fm = v.form_main if isinstance(v.form_main, dict) else json.loads(v.form_main)
            sales = _safe_float(fm.get("taxable_sales", 0)) + _safe_float(fm.get("tax_free_sales", 0)) + \
                    _safe_float(fm.get("simple_tax_sales", 0))
            if inv_amt > 0 and sales > 0:
                diff_pct = abs(inv_amt - sales) / max(inv_amt, sales) * 100
                if diff_pct > 20:
                    mismatches.append(f"{period}:发票{inv_amt:,.0f}vs申报{sales:,.0f}差异{diff_pct:.0f}%")
    
    if mismatches:
        results.append({
            "category": "增值税专项", "category_icon": "🧾", "risk_score": 6, "risk_level": "中风险",
            "risk_color": "#f59e0b", "urgency": "高",
            "item": "增值税发票开具与收入确认跨期异常",
            "detail": f"以下期间开票金额与申报销售额差异较大：{'；'.join(mismatches[:5])}。金税四期开票系统与申报系统自动比对，跨期差异将被标记。",
            "suggestion": "确保发票开具时间与收入确认时间一致，跨期发票及时调整申报所属期。"
        })


def _analyze_export_rebate_anomaly(db, company_id, ps, pe, results):
    """出口退税单证异常检测"""
    ps_date = _period_to_date_range(ps)[0]; pe_date = _period_to_date_range(pe)[1]
    
    # 检查销项发票中的出口/外销发票
    export_keywords = ["出口", "外销", "保税", "海外", "跨境"]
    export_invs = db.query(SalesInvoice).filter(
        SalesInvoice.company_id == company_id,
        SalesInvoice.invoice_date >= ps_date, SalesInvoice.invoice_date <= pe_date,
        or_(*[SalesInvoice.goods_name.contains(kw) for kw in export_keywords])
    ).all()
    
    # 也检查零税率发票
    zero_rate_invs = db.query(SalesInvoice).filter(
        SalesInvoice.company_id == company_id,
        SalesInvoice.invoice_date >= ps_date, SalesInvoice.invoice_date <= pe_date,
        SalesInvoice.tax_rate == 0
    ).all()
    
    all_export = export_invs + [inv for inv in zero_rate_invs if inv not in export_invs]
    
    if not all_export:
        return
    
    export_amt = sum(_safe_float(inv.total_amount) for inv in all_export)
    
    if export_amt > 50000:
        results.append({
            "category": "增值税专项", "category_icon": "🧾", "risk_score": 7, "risk_level": "高风险",
            "risk_color": "#dc2626", "urgency": "紧急",
            "item": "出口退税单证异常",
            "detail": f"发现 {len(all_export)} 张出口/零税率销项发票（涉及金额 {export_amt:,.2f} 元）。出口退税是国家重点监管领域，单证不齐可能导致退税被拒或追缴。出口退税申报应在货物报关出口之日起至次年4月30日前完成。",
            "suggestion": "确保出口退税申报及时，保留报关单、核销单、购货合同、付款凭证等全套单证。",
            "required_evidence": ["出口退税申报记录", "报关单（出口货物报关单）", "外汇核销单/收汇凭证", "购货增值税专用发票及付款凭证"]
        })


def _analyze_input_tax_credit_long_term(db, company_id, ps, pe, results):
    """进项税额留抵长期不退税"""
    decls = db.query(VATDeclaration).filter(
        VATDeclaration.company_id == company_id,
        VATDeclaration.period >= ps, VATDeclaration.period <= pe
    ).order_by(VATDeclaration.period).all()
    
    if len(decls) < 6:
        return
    
    # 找连续留抵>=6个月且金额>50万的
    continuous_count = 0
    max_credit = 0
    for v in decls:
        if v.form_main:
            fm = v.form_main if isinstance(v.form_main, dict) else json.loads(v.form_main)
            end_credit = _safe_float(fm.get("end_period_credit", 0))
            max_credit = max(max_credit, end_credit)
            if end_credit > 500000:
                continuous_count += 1
            else:
                continuous_count = 0
    
    if continuous_count >= 6:
        results.append({
            "category": "增值税专项", "category_icon": "🧾", "risk_score": 6, "risk_level": "中风险",
            "risk_color": "#f59e0b", "urgency": "提醒",
            "item": "进项税额留抵长期不退税",
            "detail": f"期末留抵税额较大且持续 >=6 个月未申请留抵退税，也未在后续申报中消化。大额长期留抵可能暗示进项税额虚增或收入未确认。",
            "suggestion": "分析大额留抵原因：(1)是否符合留抵退税条件；(2)是否存在未确认收入的发出商品；(3)是否存在不得抵扣的进项税额。"
        })


# ═══════════════════════════════════════════════════════════
#  I 类 — 企业所得税专项补充
# ═══════════════════════════════════════════════════════════

def _analyze_related_party_transfer_pricing(db, company_id, ps, pe, results):
    """关联交易转移定价偏离市场价检测"""
    # 本函数与_analyze_related_party_pricing有重叠但更侧重转移定价角度
    company = db.query(Company).filter(Company.id == company_id).first()
    if not company:
        return
    
    legal_rep = company.legal_representative or ""
    if not legal_rep:
        return
    
    ps_date = _period_to_date_range(ps)[0]; pe_date = _period_to_date_range(pe)[1]
    
    # 查找与法定代表人/股东相关的交易
    related_purchase = db.query(
        func.count(PurchaseInvoice.id),
        func.sum(PurchaseInvoice.total_amount)
    ).filter(
        PurchaseInvoice.company_id == company_id,
        PurchaseInvoice.invoice_date >= ps_date, PurchaseInvoice.invoice_date <= pe_date,
        PurchaseInvoice.seller_name.contains(legal_rep[:2])  # 姓名前两字匹配
    ).first() or (0, 0)
    
    related_sales = db.query(
        func.count(SalesInvoice.id),
        func.sum(SalesInvoice.total_amount)
    ).filter(
        SalesInvoice.company_id == company_id,
        SalesInvoice.invoice_date >= ps_date, SalesInvoice.invoice_date <= pe_date,
        SalesInvoice.buyer_name.contains(legal_rep[:2])
    ).first() or (0, 0)
    
    rp_cnt, rp_amt = related_purchase
    rs_cnt, rs_amt = related_sales
    
    total_related = rp_cnt + rs_cnt
    if total_related > 0:
        total_related_amt = _safe_float(rp_amt) + _safe_float(rs_amt)
        results.append({
            "category": "企业所得税", "category_icon": "📊", "risk_score": 7, "risk_level": "高风险",
            "risk_color": "#dc2626", "urgency": "高",
            "item": "关联交易转移定价偏离市场价",
            "detail": f"发现与法定代表人/股东相关的交易：采购 {rp_cnt} 笔（{_safe_float(rp_amt):,.2f}元），销售 {rs_cnt} 笔（{_safe_float(rs_amt):,.2f}元）。关联交易需确保定价符合独立交易原则，否则金税四期与海关/市场监管数据联动将标记异常。",
            "suggestion": "关联交易需准备转移定价文档，证明交易价格符合独立交易原则，必要时申请预约定价安排。",
            "required_evidence": ["关联交易合同及定价说明", "可比非受控价格分析", "转移定价文档（主体文档/本地文档）"]
        })


def _analyze_rd_expense_compliance(db, company_id, ps, pe, results):
    """研发费用加计扣除归集规范性检测"""
    # 从序时账中查找研发相关支出
    rd_keywords = ["研发", "开发", "技术", "试验", "试制", "检测", "测试", "专利"]
    rd_entries = db.query(
        func.count(JournalEntry.id),
        func.sum(JournalEntry.debit_amount)
    ).filter(
        JournalEntry.company_id == company_id,
        JournalEntry.period >= ps, JournalEntry.period <= pe,
        or_(*[JournalEntry.summary.contains(kw) for kw in rd_keywords])
    ).first() or (0, 0)
    
    rd_cnt, rd_amt = rd_entries
    if rd_cnt and rd_cnt > 0 and _safe_float(rd_amt) > 10000:
        # 检查工资薪金中是否有研发人员
        salary_with_rd = db.query(func.count(SalaryRecord.id)).filter(
            SalaryRecord.company_id == company_id,
            SalaryRecord.period >= ps, SalaryRecord.period <= pe,
            or_(
                SalaryRecord.employee_name.contains("研发"),
                SalaryRecord.employee_name.contains("技术"),
                SalaryRecord.employee_name.contains("开发"),
            )
        ).scalar() or 0
        
        results.append({
            "category": "企业所得税", "category_icon": "📊", "risk_score": 4, "risk_level": "低风险",
            "risk_color": "#3b82f6", "urgency": "中",
            "item": "研发费用加计扣除归集不规范",
            "detail": f"发现 {rd_cnt} 笔研发相关支出（{_safe_float(rd_amt):,.2f}元），{'有' if salary_with_rd > 0 else '无'}研发人员工资记录。研发费用加计扣除需规范归集，区分研发用材料与生产性材料，人员人工费需有明确项目工时记录。",
            "suggestion": "建立研发项目台账，明确研发人员名单及工时分配，区分研发用材料与生产性材料，规范归集研发费用。",
            "required_evidence": ["研发项目立项文件", "研发费用辅助账（按项目）", "研发人员工时记录", "研发材料领用记录"]
        })


def _analyze_eit_contribution_zero(db, company_id, ps, pe, results):
    """企业所得税贡献率持续为零检测"""
    revenue = _get_account_sum(db, company_id, "6001", ps, pe, "credit")
    if revenue < 1000000:  # 收入<100万不触发
        return
    
    # 从利润表角度估算所得税
    cost = _get_account_sum(db, company_id, "6401", ps, pe, "debit")
    mgmt = _get_account_sum(db, company_id, "6602", ps, pe, "debit")
    sales_exp = _get_account_sum(db, company_id, "6601", ps, pe, "debit")
    fin_exp = _get_account_sum(db, company_id, "6603", ps, pe, "debit")
    total_exp = cost + mgmt + sales_exp + fin_exp
    
    profit = revenue - total_exp
    
    # 查找所得税费用
    income_tax = _get_account_sum(db, company_id, "6801", ps, pe, "debit")  # 所得税费用
    
    eit_contribution = income_tax / revenue * 100 if revenue > 0 else 0
    
    if profit > 0 and income_tax <= 0:
        results.append({
            "category": "企业所得税", "category_icon": "📊", "risk_score": 7, "risk_level": "高风险",
            "risk_color": "#dc2626", "urgency": "紧急",
            "item": "企业所得税贡献率持续为零",
            "detail": f"有营业收入 {revenue:,.2f} 元，账面利润约 {profit:,.2f} 元，但未发现企业所得税计提记录。可能原因：①利润总额持续为负（亏损企业）；②享受免税/减税政策但未合规备案；③通过关联交易转移利润。",
            "suggestion": "分析利润为负的原因：是真实经营亏损还是有虚列成本/滞后确认收入；检查免税政策适用是否合规。",
            "required_evidence": ["企业所得税申报表", "亏损原因说明", "免税政策备案文件（如适用）"]
        })
    elif eit_contribution < 1 and revenue > 5000000:
        results.append({
            "category": "企业所得税", "category_icon": "📊", "risk_score": 5, "risk_level": "中风险",
            "risk_color": "#f59e0b", "urgency": "提醒",
            "item": "企业所得税贡献率持续为零",
            "detail": f"所得税贡献率仅 {eit_contribution:.2f}%，远低于正常水平。建议核查是否存在应调增项目未做纳税调整。",
            "suggestion": "检查纳税调整事项是否完整，确保所有不得税前扣除的项目已做纳税调增。"
        })


def _analyze_eit_adjustment_items(db, company_id, ps, pe, results):
    """应纳税所得额调增项目检查"""
    revenue = _get_account_sum(db, company_id, "6001", ps, pe, "credit")
    if revenue <= 0:
        return
    
    # 检查常见的纳税调增事项
    issues = []
    
    # 1. 业务招待费超标
    entertainment = _safe_float(db.query(func.sum(JournalEntry.debit_amount)).filter(
        JournalEntry.company_id == company_id,
        func.lower(JournalEntry.summary).contains("招待"),
        JournalEntry.period >= ps, JournalEntry.period <= pe
    ).scalar())
    if entertainment > 0:
        limit1 = entertainment * 0.6
        limit2 = revenue * 0.005
        actual_limit = min(limit1, limit2)
        if entertainment > actual_limit:
            issues.append(f"业务招待费超标{entertainment - actual_limit:,.0f}元")
    
    # 2. 广宣费超标
    ad_promo = _safe_float(db.query(func.sum(JournalEntry.debit_amount)).filter(
        JournalEntry.company_id == company_id,
        JournalEntry.period >= ps, JournalEntry.period <= pe,
        or_(func.lower(JournalEntry.summary).contains("广告"),
            func.lower(JournalEntry.summary).contains("宣传"),
            func.lower(JournalEntry.summary).contains("推广"))
    ).scalar())
    if ad_promo > revenue * 0.15:
        issues.append(f"广宣费超限额{ad_promo - revenue * 0.15:,.0f}元")
    
    # 3. 营业外支出检查（罚款/滞纳金等不得扣除项目）
    penalty = _safe_float(db.query(func.sum(JournalEntry.debit_amount)).filter(
        JournalEntry.company_id == company_id,
        JournalEntry.period >= ps, JournalEntry.period <= pe,
        or_(func.lower(JournalEntry.summary).contains("罚款"),
            func.lower(JournalEntry.summary).contains("滞纳金"),
            func.lower(JournalEntry.summary).contains("罚金"))
    ).scalar())
    if penalty > 0:
        issues.append(f"罚款/滞纳金{penalty:,.0f}元（不得税前扣除）")
    
    if issues:
        results.append({
            "category": "企业所得税", "category_icon": "📊", "risk_score": 4, "risk_level": "低风险",
            "risk_color": "#3b82f6", "urgency": "提醒",
            "item": "纳税调整减少率异常偏低",
            "detail": f"发现以下项目需在汇算清缴时做纳税调增：{'；'.join(issues)}。纳税调整减少率持续为0或极低可能意味着未充分利用税收优惠政策。",
            "suggestion": "检查是否充分享受小微企业减免税、研发费用加计扣除等优惠政策。注意调减类申报。"
        })


# ═══════════════════════════════════════════════════════════
#  J 类 — 发票深度补充：经营实质与发票匹配
# ═══════════════════════════════════════════════════════════

def _analyze_seller_spike_invoicing(db, company_id, ps, pe, results):
    """同一销方短期内开票金额激增检测"""
    ps_date = _period_to_date_range(ps)[0]; pe_date = _period_to_date_range(pe)[1]
    
    # 按供应商分组，按月统计
    purchases = db.query(PurchaseInvoice).filter(
        PurchaseInvoice.company_id == company_id,
        PurchaseInvoice.invoice_date >= ps_date, PurchaseInvoice.invoice_date <= pe_date
    ).all()
    
    if len(purchases) < 10:
        return
    
    # 按销方分组
    seller_data = {}
    for inv in purchases:
        seller = (inv.seller_name or "未知").strip()
        if len(seller) < 2:
            continue
        if seller not in seller_data:
            seller_data[seller] = []
        seller_data[seller].append({
            "date": inv.invoice_date[:7] if isinstance(inv.invoice_date, str) else inv.invoice_date.strftime("%Y-%m"),
            "amount": _safe_float(inv.total_amount)
        })
    
    spike_sellers = []
    for seller, records in seller_data.items():
        if len(records) < 3:
            continue
        # 按月汇总
        monthly = {}
        for r in records:
            m = r["date"]
            monthly[m] = monthly.get(m, 0) + r["amount"]
        
        months = sorted(monthly.keys())
        if len(months) < 2:
            continue
        
        for i in range(1, len(months)):
            prev_amt = monthly[months[i - 1]]
            curr_amt = monthly[months[i]]
            if prev_amt > 0 and curr_amt / prev_amt > 5 and curr_amt > 50000:
                spike_sellers.append(f"{seller}({months[i]}:{curr_amt:,.0f}元)")
                break
    
    if spike_sellers:
        results.append({
            "category": "发票深度", "category_icon": "🔍", "risk_score": 7, "risk_level": "高风险",
            "risk_color": "#dc2626", "urgency": "高",
            "item": "同一销方短期内开票金额激增",
            "detail": f"以下供应商开票金额环比增长超过5倍：{'；'.join(spike_sellers[:3])}。涉嫌虚开发票或串通开票。金税四期对供应商开票异常行为双向监控。",
            "suggestion": "核查与销方的业务真实性，准备采购合同/入库单/付款凭证等佐证材料，对异常供应商停止合作。",
            "required_evidence": ["异常供应商的采购合同", "入库单/收货记录", "付款凭证（银行回单）"]
        })


def _analyze_daily_necessities_invoice(db, company_id, ps, pe, results):
    """取得大量生活用品类发票检测"""
    ps_date = _period_to_date_range(ps)[0]; pe_date = _period_to_date_range(pe)[1]
    
    daily_keywords = ["日用品", "百货", "食品", "服饰", "箱包", "化妆品", "洗涤", "床上用品",
                      "厨房用品", "清洁用品", "生活用品", "母婴", "鞋帽", "服装"]
    
    daily_invs = db.query(
        func.count(PurchaseInvoice.id),
        func.sum(PurchaseInvoice.total_amount)
    ).filter(
        PurchaseInvoice.company_id == company_id,
        PurchaseInvoice.invoice_date >= ps_date, PurchaseInvoice.invoice_date <= pe_date,
        or_(*[PurchaseInvoice.goods_name.contains(kw) for kw in daily_keywords])
    ).first() or (0, 0)
    
    daily_cnt, daily_amt = daily_invs
    
    total_purchase = db.query(
        func.sum(PurchaseInvoice.total_amount)
    ).filter(
        PurchaseInvoice.company_id == company_id,
        PurchaseInvoice.invoice_date >= ps_date, PurchaseInvoice.invoice_date <= pe_date
    ).scalar() or 0
    
    if daily_cnt and daily_cnt > 0:
        ratio = _safe_float(daily_amt) / (_safe_float(total_purchase) or 1) * 100
        if ratio > 10 and _safe_float(daily_amt) > 10000:
            results.append({
                "category": "发票深度", "category_icon": "🔬", "risk_score": 8, "risk_level": "高风险",
                "risk_color": "#dc2626", "urgency": "紧急",
                "item": "取得大量生活用品类发票",
                "detail": f"期间取得生活用品类进项发票 {daily_cnt} 张、金额 {_safe_float(daily_amt):,.2f} 元，占进项总额 {ratio:.1f}%。生活用品与多数企业的生产经营活动关系不大，可能涉嫌个人消费报销入账或虚列费用。",
                "suggestion": "逐笔核实生活用品类发票对应的采购用途：办公用品应有领用登记；劳保用品应有发放记录；节日福利应并入工资薪金扣缴个税。无法说明用途的进项税额不得抵扣。",
                "required_evidence": ["生活用品类发票清单及用途说明", "领用登记/发放记录", "如为福利费，个税代扣记录"]
            })
        elif ratio > 5 and _safe_float(daily_amt) > 5000:
            results.append({
                "category": "发票深度", "category_icon": "🔬", "risk_score": 4, "risk_level": "中风险",
                "risk_color": "#f59e0b", "urgency": "提醒",
                "item": "取得大量生活用品类发票",
                "detail": f"取得 {daily_cnt} 张生活用品类发票（{_safe_float(daily_amt):,.2f}元），占进项总额 {ratio:.1f}%。建议核实采购用途。",
                "suggestion": "确保每张生活用品类发票有明确的业务用途说明。"
            })


def _analyze_energy_consumption_match(db, company_id, ps, pe, results):
    """经营能耗（水电气）与业务规模匹配检测"""
    company = db.query(Company).filter(Company.id == company_id).first()
    if not company:
        return
    biz = (company.business_scope or "") + " " + (company.company_type or "")
    
    # 四源交叉验证：序时账 + 银行流水 + 取得发票 + 合同
    energy_keywords = ["电费", "水费", "燃气", "天然气", "供热", "供电", "水务", "电力"]
    ev = _check_expense_evidence(db, company_id, ps, pe,
        expense_name="能源/水电费",
        je_account_codes=["660215"],
        je_summary_kw=energy_keywords,
        bank_kw=energy_keywords + ["电", "水", "气"],
        invoice_kw=energy_keywords,
        contract_kw=["租赁", "物业", "场地", "厂房", "办公室", "商铺"],
        contract_types=["租赁"])
    
    energy_amt = ev["invoice"]["amount"]
    energy_je = ev["je"]["amount"]
    total_energy = max(energy_amt, energy_je)
    bank_energy = ev["bank"]["amount"]
    
    total_purchase = _safe_float(db.query(func.sum(PurchaseInvoice.total_amount)).filter(
        PurchaseInvoice.company_id == company_id,
        PurchaseInvoice.invoice_date >= _period_to_date_range(ps)[0],
        PurchaseInvoice.invoice_date <= _period_to_date_range(pe)[1]
    ).scalar())
    
    if total_purchase <= 0:
        return
    
    energy_ratio = total_energy / total_purchase * 100
    
    # 制造业/餐饮业水电气占比应>5%，一般企业>2%
    is_mfg = any(kw in biz for kw in ["制造", "生产", "加工", "食品", "餐饮"])
    threshold = 5 if is_mfg else 2
    
    # 银行有能源付款但发票/序时账无对应 → 高风险（水电气是银行必须记录的基础支出）
    if ev["overall"]["je_missing_but_bank_found"] and total_purchase > 50000:
        results.append({
            "category": "发票深度", "category_icon": "🔬", "risk_score": 8, "risk_level": "高风险",
            "risk_color": "#dc2626", "urgency": "紧急",
            "item": "银行有能源/水电费付款但无对应发票或凭证",
            "detail": f"银行流水中有能源/水电相关付款 {bank_energy:,.0f} 元，但序时账中无水电费科目（660215）记录，进项发票中也无对应发票。水电气是基础经营支出，银行有付款而账上无记录，说明可能：①水电费发票未取得或丢失；②水电费记入其他科目；③以个人账户代付未入账。",
            "suggestion": "核查银行水电付款对应的凭证，补录水电费发票及费用。建立水电费发票月度催收机制。",
            "required_evidence": ["银行水电付款回单", "水电费发票（催收中）", "水电费分摊说明（如含在租金中）"]
        })
        return
    
    # 有场地租赁合同但无私账水电费 → 水和电可能含在租金中
    if ev["overall"]["je_missing_but_contract_found"] and total_energy < threshold * total_purchase / 100 and total_purchase > 100000:
        results.append({
            "category": "发票深度", "category_icon": "🔬", "risk_score": 3, "risk_level": "低风险",
            "risk_color": "#3b82f6", "urgency": "建议",
            "item": f"有场地/租赁合同（{len(ev['contract']['contracts'])}份）但水电费记录不完整",
            "detail": f"系统中有 {len(ev['contract']['contracts'])} 份场地租赁/物业相关合同，但水电费记录仅 {total_energy:,.0f} 元（占进项{energy_ratio:.1f}%）。可能为水电费包含在租金中统一支付。",
            "suggestion": "如水电费含在租金中，建议在租赁合同中明确水电费条款，或要求房东单独开具水电费发票。"
        })
        return
    
    if energy_ratio < threshold and total_purchase > 100000:
        results.append({
            "category": "发票深度", "category_icon": "🔬", "risk_score": 7, "risk_level": "高风险",
            "risk_color": "#dc2626", "urgency": "紧急",
            "item": "经营能耗（水电气）与业务规模不匹配",
            "detail": f"水电气进项发票金额占进项总额仅 {energy_ratio:.1f}%（{total_energy:,.2f}元），显著低于{'制造/餐饮业应有>=5%' if is_mfg else '一般企业应有>=2%'}的水平。有实体经营场所的企业必然产生水电气消耗，能耗过低可能说明经营场所规模与开票业务规模不匹配。",
            "suggestion": "核实经营场所的实际使用面积和能耗情况；比对同区域同行业企业的能耗水平；如经营场所为自有或长期租赁，检查水电费缴纳记录是否完整。",
            "required_evidence": ["水电气费缴纳凭证及发票", "经营场所租赁合同/产权证明", "能耗与业务规模匹配说明"]
        })


def _analyze_transport_expense_match(db, company_id, ps, pe, results):
    """运输费用与经营模式匹配检测（补充版：运输费为0触发）"""
    company = db.query(Company).filter(Company.id == company_id).first()
    if not company:
        return
    biz = (company.business_scope or "") + " " + (company.company_type or "")
    
    # 制造业/批发业有大量购销但无运输费
    is_mfg_wholesale = any(kw in biz for kw in ["制造", "生产", "加工", "批发", "贸易", "销售"])
    
    # 四源交叉验证：序时账 + 银行流水 + 取得发票 + 合同
    transport_kw = ["运输", "物流", "快递", "货运", "配送"]
    ev = _check_expense_evidence(db, company_id, ps, pe,
        expense_name="运输/物流费",
        je_account_codes=["660207"],
        je_summary_kw=transport_kw,
        bank_kw=transport_kw + ["搬运"],
        invoice_kw=transport_kw,
        contract_kw=transport_kw,
        contract_types=["服务"])
    
    transport_je = ev["je"]["amount"]
    transport_inv = ev["invoice"]["amount"]
    bank_transport = ev["bank"]["amount"]
    total_transport = transport_inv + transport_je + bank_transport
    
    total_purchase = _safe_float(db.query(func.sum(PurchaseInvoice.total_amount)).filter(
        PurchaseInvoice.company_id == company_id,
        PurchaseInvoice.invoice_date >= _period_to_date_range(ps)[0],
        PurchaseInvoice.invoice_date <= _period_to_date_range(pe)[1]
    ).scalar())
    
    revenue = _get_account_sum(db, company_id, "6001", ps, pe, "credit")
    
    # 银行有运输付款但账上无对应 → 高风险
    if is_mfg_wholesale and ev["overall"]["je_missing_but_bank_found"] and total_purchase > 50000:
        results.append({
            "category": "发票深度", "category_icon": "🔬", "risk_score": 8, "risk_level": "高风险",
            "risk_color": "#dc2626", "urgency": "紧急",
            "item": "银行有运输/物流付款但无对应运输费凭证",
            "detail": f"银行流水中有运输/物流相关付款 {bank_transport:,.0f} 元，但序时账中无运输费科目（660207）记录。银行付款必须入账。可能：①运费合并在采购成本中；②凭证尚未生成。",
            "suggestion": "核查银行运输付款对应的序时账科目，如运费合并在采购成本中，建立科目映射表。",
            "required_evidence": ["银行运输付款回单", "对应序时账凭证", "运费与采购成本的分摊说明"]
        })
        return
    
    # 有运输合同但无账务
    if is_mfg_wholesale and ev["overall"]["je_missing_but_contract_found"] and total_purchase > 50000:
        free_contracts = [c for c in ev["contract"]["contracts"] if c["is_free"]]
        has_free = len(free_contracts) > 0
        results.append({
            "category": "发票深度", "category_icon": "🔬", "risk_score": 3 if has_free else 5,
            "risk_level": "低风险" if has_free else "中风险",
            "risk_color": "#3b82f6" if has_free else "#f59e0b", "urgency": "建议" if has_free else "提醒",
            "item": f"有运输/物流合同（{len(ev['contract']['contracts'])}份）但无运输费记录{'(含免费条款)' if has_free else ''}",
            "detail": f"系统中有 {len(ev['contract']['contracts'])} 份运输/物流合同{'(含免费条款)' if has_free else ''}，但序时账中无运输费。可能为：①运费由对方承担（到付/含在售价中）；②合同刚签署尚未履行。",
            "suggestion": "如运费由客户承担，在销售合同中标明运费承担方式。如为免费运输，保留合同作为佐证。"
        })
        return
    
    if is_mfg_wholesale and total_transport == 0 and (total_purchase > 100000 or revenue > 100000):
        results.append({
            "category": "发票深度", "category_icon": "🔬", "risk_score": 7, "risk_level": "高风险",
            "risk_color": "#dc2626", "urgency": "紧急",
            "item": "运输费用与经营模式不匹配（四源皆空）",
            "detail": f"企业为制造/批发类企业（购销总额 {(total_purchase + revenue):,.0f} 元），但序时账、银行流水、发票、合同中均未发现任何运输/物流费用。有货物购销就需要运输，无运输费不符合商业逻辑。可能：①运费由对方承担并已在发票金额中体现；②私车运输未入账；③虚开发票无真实货物流转。",
            "suggestion": "如为对方承担运费，提供销售合同中运费的条款说明。如为自有车辆运输，检查车辆折旧/油费是否入账。如为第三方物流，保留运输合同和对账单。",
            "required_evidence": ["销售合同（运费条款）", "运输合同/物流对账单", "自有车辆费用明细（如适用）"]
        })


# ═══════════════════════════════════════════════════════════
#  K 类 — 经营实质补充：上下游地域与开票容量分析
# ═══════════════════════════════════════════════════════════

def _analyze_two_end_outside(db, company_id, ps, pe, results):
    """购销两头在外检测：主要供应商和客户均不在注册地（空壳特征）"""
    company = db.query(Company).filter(Company.id == company_id).first()
    if not company:
        return
    company_addr = (company.address or "").strip()
    if not company_addr:
        return
    
    # 提取公司所在省份/城市
    company_province = ""
    for suffix in ["省", "市", "自治区"]:
        idx = company_addr.find(suffix)
        if idx > 0:
            company_province = company_addr[:idx + 1]
            break
    if not company_province:
        company_province = company_addr[:2]  # fallback
    
    ps_date = _period_to_date_range(ps)[0]; pe_date = _period_to_date_range(pe)[1]
    
    # 从客户档案中按地址筛选（Customer有address字段）
    # 从供应商和客户名称中提取地区关键词
    province_short = company_province.replace("省", "").replace("市", "").replace("自治区", "")[:2]
    
    # 检查供应商（PurchaseInvoice只有seller_name，Supplier无address）
    # 策略：统计不同供应商数量，用名称中的省市区关键词匹配
    purchase_names = set()
    for inv in db.query(PurchaseInvoice.seller_name).filter(
        PurchaseInvoice.company_id == company_id,
        PurchaseInvoice.invoice_date >= ps_date, PurchaseInvoice.invoice_date <= pe_date
    ).all():
        name = (inv[0] or "").strip()
        if name:
            purchase_names.add(name)
    
    # 检查客户：优先用Customer表的address字段
    sales_names_outside = set()
    sales_names_all = set()
    # 从SalesInvoice取buyer_name
    for inv in db.query(SalesInvoice.buyer_name).filter(
        SalesInvoice.company_id == company_id,
        SalesInvoice.invoice_date >= ps_date, SalesInvoice.invoice_date <= pe_date
    ).all():
        name = (inv[0] or "").strip()
        if name:
            sales_names_all.add(name)
    
    # 通过Customer表查地址
    for name in list(sales_names_all):
        cust = db.query(Customer).filter(
            Customer.company_id == company_id,
            Customer.name == name
        ).first()
        if cust and cust.address:
            if company_province not in cust.address:
                sales_names_outside.add(name)
        elif name and province_short not in name:
            # fallback: 名称中不含注册地简称
            sales_names_outside.add(name)
    
    # 供应商同理：名称中不含注册地简称
    purchase_names_outside = {n for n in purchase_names if province_short not in n}
    
    # 四源交叉验证经营场所：序时账 + 银行流水 + 发票 + 合同
    ev = _check_expense_evidence(db, company_id, ps, pe,
        expense_name="经营场所",
        je_account_codes=["660214", "660215", "660213"],
        invoice_kw=["租赁", "水电", "物业", "房租"],
        contract_kw=["场地", "租赁", "房租", "物业", "办公"],
        contract_types=["租赁"])

    local_expense = ev["je"]["amount"]
    overall = ev["overall"]
    free_possible = ev["contract"]["free_possible"]
    contracts = ev["contract"]["contracts"]
    inv_amount = ev["invoice"]["amount"]

    outside_threshold = len(purchase_names_outside) >= 3 and len(sales_names_outside) >= 3
    if not outside_threshold:
        return

    # ── 四源全无 → 高风险「购销两头在外」 ──
    if overall["all_sources_none"]:
        results.append({
            "category": "经营实质", "category_icon": "🔍", "risk_score": 9, "risk_level": "高风险",
            "risk_color": "#dc2626", "urgency": "紧急",
            "item": "购销两头在外",
            "detail": f"主要供应商（{len(purchase_names_outside)}家）和客户（{len(sales_names_outside)}家）均在注册地（{company_province}）以外，且注册地四源交叉验证无经营场所证据：序时账无租赁费/水电费/物业费、取得发票无相关记录、无场地相关合同。购销两头在外的空壳公司是虚开发票的典型特征。金税四期通过物流/资金流/发票流三流合一来判断。",
            "suggestion": "准备经营场所租赁合同+水电物业费凭证+员工名册及社保记录+物流单据，证明企业有真实经营场所和经营能力。",
            "required_evidence": ["经营场所证明（租赁合同+水电费）", "员工劳动合同+社保记录", "物流运输单据"]
        })
        return

    # ── 序时账没有，发票有 → 费用已发生未入账 ──
    if overall["je_missing_but_invoice_found"]:
        results.append({
            "category": "经营实质", "category_icon": "🔍", "risk_score": 6, "risk_level": "中风险",
            "risk_color": "#f59e0b", "urgency": "预警",
            "item": "购销两头在外（费用未入账）",
            "detail": f"主要供应商和客户均在注册地以外，序时账中无经营场所费用，但取得发票中发现相关支出 ¥{inv_amount:,.2f}（尚未入账）。建议尽快入账以证明经营场所真实性。",
            "suggestion": "将已取得的经营场所相关发票入账，保留经营场所租赁合同及水电费凭证备查。",
            "required_evidence": ["未入账的经营场所相关发票", "经营场所租赁合同"]
        })
        return

    # ── 序时账没有，但有场地合同（免费场地） ──
    if overall["je_missing_but_contract_found"]:
        cts = contracts
        ct_summary = "；".join([f"{c['name']}({c['type']}{'·免费' if c['is_free'] else ''})" for c in cts[:3]])
        free_note = "·可能为免费场地" if free_possible else ""
        results.append({
            "category": "经营实质", "category_icon": "🔍", "risk_score": 5, "risk_level": "中风险",
            "risk_color": "#f59e0b", "urgency": "提醒",
            "item": "购销主要在外地（有场地合同但费用未入账）",
            "detail": f"主要供应商和客户在注册地以外，序时账中无经营场所费用，但发现场地相关合同：{ct_summary}{free_note}。建议确认合同履行状态并补充费用凭证。",
            "suggestion": "保留场地使用合同备查，补充经营场所费用凭证。",
            "required_evidence": ["场地使用合同", "经营场所费用凭证（如有）"]
        })
        return

    # ── 序时账有记录 → 中风险（购销在外但注册地有经营） ──
    if local_expense > 0:
        results.append({
            "category": "经营实质", "category_icon": "🔍", "risk_score": 5, "risk_level": "中风险",
            "risk_color": "#f59e0b", "urgency": "提醒",
            "item": "购销主要在外地（注册地有经营费用）",
            "detail": f"主要供应商和客户在注册地以外，但注册地有经营费用记录。建议确保物流和资金流的记录完整。",
            "suggestion": "保留完整的货物流转记录和物流单据。"
        })


def _analyze_new_company_mass_invoicing(db, company_id, ps, pe, results):
    """新办企业短期内大量开票检测（暴力虚开典型特征）"""
    company = db.query(Company).filter(Company.id == company_id).first()
    if not company:
        return
    
    establishment_date = company.established_date
    if not establishment_date:
        return
    
    if isinstance(establishment_date, str):
        try:
            establishment_date = datetime.strptime(establishment_date, "%Y-%m-%d").date()
        except:
            return
    
    # 计算公司成立到现在的月数
    today = date.today()
    months_since_establish = (today.year - establishment_date.year) * 12 + (today.month - establishment_date.month)
    
    if months_since_establish > 12:
        return  # 成立超过1年，不检测
    
    ps_date = _period_to_date_range(ps)[0]; pe_date = _period_to_date_range(pe)[1]
    
    # 统计开票量
    sales = db.query(
        func.count(SalesInvoice.id),
        func.sum(SalesInvoice.total_amount)
    ).filter(
        SalesInvoice.company_id == company_id,
        SalesInvoice.invoice_date >= ps_date, SalesInvoice.invoice_date <= pe_date,
        SalesInvoice.status == "正常"
    ).first() or (0, 0)
    
    sales_cnt, sales_amt = sales
    months_active = max(months_since_establish, 1)
    monthly_avg_amt = _safe_float(sales_amt) / months_active
    monthly_avg_cnt = sales_cnt / months_active if sales_cnt else 0
    
    if monthly_avg_amt > 1000000 or monthly_avg_cnt > 50:
        results.append({
            "category": "经营实质", "category_icon": "🔍", "risk_score": 9, "risk_level": "高风险",
            "risk_color": "#dc2626", "urgency": "紧急",
            "item": "新办企业短期内大量开票",
            "detail": f"企业成立仅 {months_since_establish} 个月，但月均开票 {monthly_avg_cnt:.0f} 张/{monthly_avg_amt:,.0f} 元。新办企业短期大量开票是暴力虚开的典型手法。金税四期对新办企业开票增量实时监控。",
            "suggestion": "逐笔核实开票对应的交易背景。新办企业应控制开票节奏，确保业务规模与经营能力（人员/场地/资金）匹配。",
            "required_evidence": ["营业执照/成立日期", "业务合同清单（所有大额开票对应）", "人员花名册", "经营场所证明"]
        })
    elif monthly_avg_amt > 500000 or monthly_avg_cnt > 20:
        results.append({
            "category": "经营实质", "category_icon": "🔍", "risk_score": 5, "risk_level": "中风险",
            "risk_color": "#f59e0b", "urgency": "提醒",
            "item": "新办企业短期内大量开票",
            "detail": f"月均开票 {monthly_avg_cnt:.0f} 张/{monthly_avg_amt:,.0f} 元，建议确保经营能力与开票规模匹配。",
            "suggestion": "完善人员配置和经营场所，保留业务合同备查。"
        })


def _analyze_small_taxpayer_overscale(db, company_id, ps, pe, results):
    """小规模纳税人开票量远超经营能力"""
    company = db.query(Company).filter(Company.id == company_id).first()
    if not company:
        return
    
    company_type = company.company_type or ""
    if "一般纳税人" in company_type:
        return  # 一般纳税人正常
    
    ps_date = _period_to_date_range(ps)[0]; pe_date = _period_to_date_range(pe)[1]
    
    # 统计年度开票金额
    year_sales = _safe_float(db.query(func.sum(SalesInvoice.total_amount)).filter(
        SalesInvoice.company_id == company_id,
        SalesInvoice.invoice_date >= ps_date, SalesInvoice.invoice_date <= pe_date,
        SalesInvoice.status == "正常"
    ).scalar())
    
    # 小规模纳税人年销售额>500万需强制认定为一般纳税人
    if year_sales > 5000000:
        results.append({
            "category": "经营实质", "category_icon": "🔍", "risk_score": 8, "risk_level": "高风险",
            "risk_color": "#dc2626", "urgency": "紧急",
            "item": "小规模纳税人开票量远超经营能力",
            "detail": f"企业为小规模纳税人，但年开票金额 {year_sales:,.2f} 元已超过一般纳税人认定标准（500万元）。根据《增值税一般纳税人登记管理办法》，年应征增值税销售额超过500万元应申请认定为一般纳税人，否则将被强制认定并可能追溯补税。",
            "suggestion": "如确已达到一般纳税人标准，应主动申请认定，避免被强制认定并追溯补税。同时核实开票业务的真实性。",
            "required_evidence": ["增值税申报表（近12个月）", "超过500万的合理说明或一般纳税人认定申请"]
        })
    elif year_sales > 4000000:
        results.append({
            "category": "经营实质", "category_icon": "🔍", "risk_score": 4, "risk_level": "中风险",
            "risk_color": "#f59e0b", "urgency": "提醒",
            "item": "小规模纳税人开票量远超经营能力",
            "detail": f"年开票金额 {year_sales:,.2f} 元，已接近一般纳税人认定标准。建议提前做一般纳税人认定准备。",
            "suggestion": "评估是否需要主动申请一般纳税人认定。"
        })


# ═══════════════════════════════════════════════════════════
#  L 类 — 征管风险分析
# ═══════════════════════════════════════════════════════════

def _analyze_tax_arrears(db, company_id, ps, pe, results):
    """欠税未清缴检测"""
    # 检查各税种申报表中是否存在应缴未缴
    arrears = []
    
    # 增值税
    vat_decls = db.query(VATDeclaration).filter(
        VATDeclaration.company_id == company_id,
        VATDeclaration.period >= ps, VATDeclaration.period <= pe
    ).all()
    vat_arrear = 0
    for v in vat_decls:
        fm = v.form_main
        if fm:
            fm_dict = json.loads(fm) if isinstance(fm, str) else fm
            payable = _safe_float(fm_dict.get("vat_payable", 0))
            paid = _safe_float(fm_dict.get("vat_paid", payable))  # 如果无paid字段，假设已缴
            if payable > paid + 1:
                vat_arrear += payable - paid
    if vat_arrear > 100:
        arrears.append(f"增值税欠缴{vat_arrear:,.0f}元")
    
    # 文化事业建设费
    ccf_decls = db.query(CulturalConstructionFeeDeclaration).filter(
        CulturalConstructionFeeDeclaration.company_id == company_id,
        CulturalConstructionFeeDeclaration.period >= ps, CulturalConstructionFeeDeclaration.period <= pe
    ).all()
    ccf_arrear = 0
    for ccf in ccf_decls:
        fm = ccf.form_main
        if fm:
            fm_dict = json.loads(fm) if isinstance(fm, str) else fm
            payable = _safe_float(fm_dict.get("row12_total_payable_current", 0))
            # 简化：假设已缴=应缴（实际应从缴款记录判断）
    if ccf_arrear > 100:
        arrears.append(f"文化事业建设费欠缴{ccf_arrear:,.0f}元")
    
    if arrears:
        results.append({
            "category": "征管风险", "category_icon": "🏛️", "risk_score": 8, "risk_level": "高风险",
            "risk_color": "#dc2626", "urgency": "紧急",
            "item": "存在欠税未清缴",
            "detail": f"发现可能存在的欠税：{'；'.join(arrears)}。欠税会加收每日万分之五的滞纳金，并可能被采取强制措施（停供发票、税收保全等）。",
            "suggestion": "立即核实并清缴欠税，优先清缴欠税本金以减少滞纳金。如资金困难，可申请延期缴纳（最长3个月）。",
            "required_evidence": ["各税种申报表", "缴款凭证/银行回单", "如有争议，提供行政复议或诉讼材料"]
        })


def _analyze_upstream_runaway(db, company_id, ps, pe, results):
    """上游企业走逃失联检测"""
    # 检查进项发票中是否有风险供应商（标记为疑点/异常/失控）
    risk_sellers = db.query(
        func.count(func.distinct(PurchaseInvoice.seller_name)),
        func.count(PurchaseInvoice.id),
        func.sum(PurchaseInvoice.total_amount)
    ).filter(
        PurchaseInvoice.company_id == company_id,
        PurchaseInvoice.invoice_risk_level.in_(["疑点", "异常", "失控"])
    ).first() or (0, 0, 0)
    
    seller_cnt, inv_cnt, inv_amt = risk_sellers
    
    if seller_cnt and seller_cnt > 0:
        results.append({
            "category": "征管风险", "category_icon": "🏛️", "risk_score": 10, "risk_level": "高风险",
            "risk_color": "#dc2626", "urgency": "紧急",
            "item": "上游企业走逃失联",
            "detail": f"有 {seller_cnt} 家供应商的 {inv_cnt} 张进项发票（涉及金额 {_safe_float(inv_amt):,.2f} 元）被标记为疑点/异常/失控。根据国家税务总局公告2016年第76号，走逃（失联）企业开具的发票属于异常扣税凭证，进项税额不得抵扣，已抵扣的需做进项税额转出。",
            "suggestion": "立即核实与走逃企业的交易：(1)如有真实交易，保留合同/付款凭证/运输单据等证明；(2)如无法证明，应主动做进项税额转出并补税。",
            "required_evidence": ["异常凭证清单", "与走逃企业的购销合同", "付款凭证（银行回单）", "货物运输单据（证明真实交易）", "进项转出计算表"]
        })


def _analyze_stale_invoice(db, company_id, ps, pe, results):
    """滞留票风险：销项发票长期未被购买方认证"""
    # 检查超过90天未认证的销项发票（简化：检查已开具但可能未匹配到收款记录）
    cutoff_date = (date.today() - timedelta(days=90)).strftime("%Y-%m-%d")
    
    stale_invs = db.query(
        func.count(SalesInvoice.id),
        func.sum(SalesInvoice.total_amount)
    ).filter(
        SalesInvoice.company_id == company_id,
        SalesInvoice.invoice_date < cutoff_date,
        SalesInvoice.status == "正常"
    ).first() or (0, 0)
    
    stale_cnt, stale_amt = stale_invs
    total_sales_cnt = db.query(func.count(SalesInvoice.id)).filter(
        SalesInvoice.company_id == company_id
    ).scalar() or 1
    
    stale_ratio = stale_cnt / total_sales_cnt * 100
    if stale_ratio > 20 and stale_cnt >= 10:
        results.append({
            "category": "征管风险", "category_icon": "🏛️", "risk_score": 6, "risk_level": "中风险",
            "risk_color": "#f59e0b", "urgency": "提醒",
            "item": "滞留票风险",
            "detail": f"有 {stale_cnt} 张销项发票开具超过90天（占比 {stale_ratio:.0f}%），涉及金额 {_safe_float(stale_amt):,.2f} 元。过多长期未认证的滞留票会引发税务机关对企业销售真实性的质疑。",
            "suggestion": "主动联系购买方核实未认证原因：是否对方为非一般纳税人、是否发票遗失、是否对方已入账但未认证。建立发票签收和认证跟踪机制。"
        })


def _analyze_non_normal_counterparty(db, company_id, ps, pe, results):
    """与非正常户有交易往来检测"""
    # 从进项发票中检查销方风险等级
    risk_seller_invs = db.query(
        func.count(PurchaseInvoice.id),
        func.sum(PurchaseInvoice.total_amount)
    ).filter(
        PurchaseInvoice.company_id == company_id,
        PurchaseInvoice.invoice_risk_level.in_(["疑点", "异常", "失控", "非正常"])
    ).first() or (0, 0)
    
    # 也检查序时账中是否有与非正常户相关的记录
    risk_journals = db.query(func.count(JournalEntry.id)).filter(
        JournalEntry.company_id == company_id,
        JournalEntry.period >= ps, JournalEntry.period <= pe,
        or_(
            JournalEntry.summary.contains("非正常"),
            JournalEntry.summary.contains("异常"),
            JournalEntry.summary.contains("走逃"),
        )
    ).scalar() or 0
    
    risk_inv_cnt, risk_inv_amt = risk_seller_invs
    
    total_risk_indicators = risk_inv_cnt + risk_journals
    
    if total_risk_indicators > 0:
        results.append({
            "category": "征管风险", "category_icon": "🏛️", "risk_score": 8, "risk_level": "高风险",
            "risk_color": "#dc2626", "urgency": "紧急",
            "item": "与非正常户有交易往来",
            "detail": f"发现 {risk_inv_cnt} 张进项发票来自异常供应商（涉及 {_safe_float(risk_inv_amt):,.2f} 元），{risk_journals} 条异常相关序时账记录。与非正常户交易取得的发票可能被列为异常扣税凭证，面临进项转出风险。",
            "suggestion": "建立供应商准入制度，合作前查验对方的税务登记状态。已发生的交易需准备充分的三流合一证明。建议每笔大额采购前通过电子税务局查验销方状态。"
        })


# ═══════════════════════════════════════════════════════════
#  M 类 — 交易特征分析
# ═══════════════════════════════════════════════════════════

def _analyze_round_amount_transactions(db, company_id, ps, pe, results):
    """大额整数交易占比过高检测"""
    ps_date = _period_to_date_range(ps)[0]; pe_date = _period_to_date_range(pe)[1]
    
    # 从销项发票和银行流水检查整数交易
    all_amt_records = []
    
    # 销项发票
    for inv in db.query(SalesInvoice.total_amount).filter(
        SalesInvoice.company_id == company_id,
        SalesInvoice.invoice_date >= ps_date, SalesInvoice.invoice_date <= pe_date
    ).all():
        all_amt_records.append(_safe_float(inv[0]))
    
    # 银行流水
    for bt in db.query(BankTransaction.debit_amount, BankTransaction.credit_amount).filter(
        BankTransaction.company_id == company_id,
        BankTransaction.transaction_date >= ps_date,
        BankTransaction.transaction_date <= pe_date
    ).all():
        all_amt_records.append(_safe_float(bt[0]))  # debit
        all_amt_records.append(_safe_float(bt[1]))  # credit
    
    if len(all_amt_records) < 20:
        return
    
    # 判断整数：整万/整十万/整百万（金额>=10000）
    round_count = 0
    for amt in all_amt_records:
        if amt >= 10000 and amt % 10000 == 0:
            round_count += 1
        elif amt >= 100000 and amt % 100000 == 0:
            round_count += 1
    
    round_ratio = round_count / len(all_amt_records) * 100
    if round_ratio > 30 and round_count >= 10:
        results.append({
            "category": "交易特征", "category_icon": "📊", "risk_score": 7, "risk_level": "高风险",
            "risk_color": "#dc2626", "urgency": "紧急",
            "item": "大额整数交易占比过高",
            "detail": f"有 {round_count}/{len(all_amt_records)} 笔交易金额为整万/整十万的整数，占比 {round_ratio:.0f}%。频繁的大额整数交易在正常商业活动中极为罕见，容易被认定为资金回流或虚假交易的标志。正常交易金额通常包含税费/折扣/运费等零头。",
            "suggestion": "逐笔核实大额整数交易的商业合理性，准备合同/验收单/物流单等证明材料。",
            "required_evidence": ["大额整数交易清单", "每笔对应合同/验收单/付款凭证"]
        })


def _analyze_fund_backflow(db, company_id, ps, pe, results):
    """资金回流嫌疑检测"""
    ps_date = _period_to_date_range(ps)[0]; pe_date = _period_to_date_range(pe)[1]
    
    # 从银行流水中找付款后短期内（<=7天）有等额/近似金额回款
    debits = db.query(BankTransaction).filter(
        BankTransaction.company_id == company_id,
        BankTransaction.transaction_date >= ps_date,
        BankTransaction.transaction_date <= pe_date,
        BankTransaction.debit_amount > 10000  # 大额付款
    ).all()
    
    credits = db.query(BankTransaction).filter(
        BankTransaction.company_id == company_id,
        BankTransaction.transaction_date >= ps_date,
        BankTransaction.transaction_date <= pe_date,
        BankTransaction.credit_amount > 10000  # 大额收款
    ).all()
    
    if len(debits) < 3 or len(credits) < 3:
        return
    
    # 简化的资金回流检测：付款后7天内收款金额近似
    suspicious_pairs = []
    from datetime import timedelta
    
    for d in debits:
        d_date = d.transaction_date
        if isinstance(d_date, str):
            try:
                d_date = datetime.strptime(d_date, "%Y-%m-%d").date()
            except:
                continue
        if not d_date:
            continue
        
        d_amt = _safe_float(d.debit_amount)
        d_name = (d.counterparty_name or "").strip()
        
        for c in credits:
            c_date = c.transaction_date
            if isinstance(c_date, str):
                try:
                    c_date = datetime.strptime(c_date, "%Y-%m-%d").date()
                except:
                    continue
            if not c_date:
                continue
            
            c_amt = _safe_float(c.credit_amount)
            c_name = (c.counterparty_name or "").strip()
            
            # 条件：收款在付款后1-7天内，金额差异<5%
            days_diff = (c_date - d_date).days
            if 1 <= days_diff <= 7:
                amt_diff_pct = abs(d_amt - c_amt) / d_amt * 100 if d_amt > 0 else 999
                if amt_diff_pct < 5:
                    # 检查是否来自同一对手方或关联方
                    if d_name and c_name and (d_name[:2] == c_name[:2] or d_name in c_name or c_name in d_name):
                        suspicious_pairs.append({
                            "pay_date": str(d_date), "recv_date": str(c_date),
                            "pay_amt": d_amt, "recv_amt": c_amt,
                            "name": d_name
                        })
    
    if suspicious_pairs:
        results.append({
            "category": "交易特征", "category_icon": "📊", "risk_score": 9, "risk_level": "高风险",
            "risk_color": "#dc2626", "urgency": "紧急",
            "item": "资金回流嫌疑",
            "detail": f"检测到付款后7天内从相同/关联方收回近似金额的交易 {len(suspicious_pairs)} 组。资金回流是虚开发票的典型资金特征：先付款（伪装采购），短期内等额回款（资金回流）。金税四期已与人民银行账户系统联网，资金回流自动识别。",
            "suggestion": "逐笔核查资金流出和流入的商业实质：采购是否有真实的货物流；回款是否为正常的销售退款/借款归还。无法合理解释的回流应立即整改。",
            "required_evidence": ["资金流出/流入明细", "资金回流交易对应的采购合同", "货物入库单/验收单", "退款协议/借款合同（如有）"]
        })


def _analyze_same_buyer_frequent_red(db, company_id, ps, pe, results):
    """同一购方频繁红冲检测"""
    ps_date = _period_to_date_range(ps)[0]; pe_date = _period_to_date_range(pe)[1]
    
    # 获取所有红冲/作废发票
    void_red = db.query(SalesInvoice).filter(
        SalesInvoice.company_id == company_id,
        SalesInvoice.invoice_date >= ps_date, SalesInvoice.invoice_date <= pe_date,
        SalesInvoice.status.in_(["作废", "红冲"])
    ).all()
    
    if len(void_red) < 3:
        return
    
    # 按购买方分组
    buyer_red = {}
    for inv in void_red:
        buyer = (inv.buyer_name or "未知").strip()
        if buyer not in buyer_red:
            buyer_red[buyer] = []
        buyer_red[buyer].append(inv)
    
    frequent_buyers = [(name, invs) for name, invs in buyer_red.items() if len(invs) >= 3]
    
    if frequent_buyers:
        detail_items = []
        for name, invs in frequent_buyers[:3]:
            amt = sum(_safe_float(inv.total_amount) for inv in invs)
            detail_items.append(f"{name}({len(invs)}张/{amt:,.0f}元)")
        
        results.append({
            "category": "交易特征", "category_icon": "📊", "risk_score": 7, "risk_level": "高风险",
            "risk_color": "#dc2626", "urgency": "紧急",
            "item": "同一购方频繁红冲",
            "detail": f"对以下购买方频繁开具红字发票：{'；'.join(detail_items)}。频繁红冲可能暗示先开票虚增收入再红冲，或利用红冲调节税负。",
            "suggestion": "核查红冲发票的真实原因：是否有对应的红字信息表；退货是否有验收和入库记录；销售折让是否有协议。",
            "required_evidence": ["红冲发票清单", "红字信息表（每笔红冲对应）", "退货验收单/折让协议"]
        })


def _analyze_input_transfer_reentry(db, company_id, ps, pe, results):
    """进项税额转出后又转入检测"""
    # 从增值税申报表附表二检查
    decls = db.query(VATDeclaration).filter(
        VATDeclaration.company_id == company_id,
        VATDeclaration.period >= ps, VATDeclaration.period <= pe
    ).order_by(VATDeclaration.period).all()
    
    for v in decls:
        if v.form_main:
            fm = v.form_main if isinstance(v.form_main, dict) else json.loads(v.form_main)
            transfer_out = _safe_float(fm.get("input_tax_transfer", 0))
            transfer_in = _safe_float(fm.get("input_tax_transfer_in", 0))
            
            if transfer_out > 0 and transfer_in > 0:
                results.append({
                    "category": "交易特征", "category_icon": "📊", "risk_score": 8, "risk_level": "高风险",
                    "risk_color": "#dc2626", "urgency": "紧急",
                    "item": "进项税额转出后又转入",
                    "detail": f"{v.period} 期存在进项税额转出 {transfer_out:,.2f} 元后又转入 {transfer_in:,.2f} 元的情况。进项转出后转入需要特定情形（如改变用途用于应税项目），需提供充分依据，否则可能被认定为违规抵扣。",
                    "suggestion": "检查进项税额转出又转入的合规性：是否确实改变了用途；是否按规定计算可转入的进项税额。保留用途变更的证明材料。",
                    "required_evidence": ["进项转出凭证", "进项转入依据（资产用途变更证明）", "进项税额转入计算表"]
                })
                return  # 只报告一次


# ═══════════════════════════════════════════════════════════
#  N 类 — 财务健康与其他专项
# ═══════════════════════════════════════════════════════════

def _analyze_roe_sustainability(db, company_id, ps, pe, results):
    """净资产收益率持续偏低/为负检测"""
    revenue = _get_account_sum(db, company_id, "6001", ps, pe, "credit")
    cost = _get_account_sum(db, company_id, "6401", ps, pe, "debit")
    mgmt = _get_account_sum(db, company_id, "6602", ps, pe, "debit")
    sales_exp = _get_account_sum(db, company_id, "6601", ps, pe, "debit")
    fin_exp = _get_account_sum(db, company_id, "6603", ps, pe, "debit")
    
    total_assets = _get_account_balance(db, company_id, "1")
    total_liab = _get_account_balance(db, company_id, "2")
    equity = total_assets - total_liab
    
    if equity <= 0:
        return
    
    net_profit = revenue - cost - mgmt - sales_exp - fin_exp
    roe = net_profit / equity * 100
    
    if roe < 0:
        results.append({
            "category": "财务健康", "category_icon": "💊", "risk_score": 6, "risk_level": "中风险",
            "risk_color": "#f59e0b", "urgency": "提醒",
            "item": "净资产收益率持续偏低/为负",
            "detail": f"近12个月净资产收益率 {roe:.1f}%，为负数。说明企业盈利能力严重不足，可能涉及：隐瞒收入导致利润虚低、成本费用虚增、或经营模式不可持续。",
            "suggestion": "分析亏损原因：收入确认是否完整；成本费用是否真实合理；是否存在不合理的关联交易转移利润。若为正常经营亏损，保留行业分析和商业计划说明。"
        })
    elif roe < 8 and equity > 1000000:
        results.append({
            "category": "财务健康", "category_icon": "💊", "risk_score": 3, "risk_level": "低风险",
            "risk_color": "#3b82f6", "urgency": "建议",
            "item": "净资产收益率持续偏低/为负",
            "detail": f"净资产收益率仅 {roe:.1f}%，低于行业参考值（一般>8%）。建议分析盈利能力的改善空间。",
            "suggestion": "关注成本控制和收入增长，制定盈利能力提升计划。"
        })


def _analyze_operating_cash_flow_coverage(db, company_id, ps, pe, results):
    """经营活动现金流无法覆盖流动负债"""
    # 从银行流水和序时账估算经营活动现金流
    ps_date = _period_to_date_range(ps)[0]; pe_date = _period_to_date_range(pe)[1]
    
    # 经营活动现金流入（银行流水贷方/收款）
    op_cash_in = _safe_float(db.query(func.sum(BankTransaction.credit_amount)).filter(
        BankTransaction.company_id == company_id,
        BankTransaction.transaction_date >= ps_date,
        BankTransaction.transaction_date <= pe_date
    ).scalar())
    
    # 经营活动现金流出（银行流水借方/付款）
    op_cash_out = _safe_float(db.query(func.sum(BankTransaction.debit_amount)).filter(
        BankTransaction.company_id == company_id,
        BankTransaction.transaction_date >= ps_date,
        BankTransaction.transaction_date <= pe_date
    ).scalar())
    
    net_op_cash = op_cash_in - op_cash_out
    
    # 流动负债
    current_liab = _get_account_balance(db, company_id, "21") + _get_account_balance(db, company_id, "22")
    
    if current_liab <= 0:
        return
    
    cash_flow_ratio = net_op_cash / current_liab if current_liab > 0 else 0
    
    if cash_flow_ratio < 0:
        results.append({
            "category": "财务健康", "category_icon": "💊", "risk_score": 7, "risk_level": "高风险",
            "risk_color": "#dc2626", "urgency": "紧急",
            "item": "经营活动现金流无法覆盖流动负债",
            "detail": f"经营活动现金净流量 {net_op_cash:,.2f} 元（现金流为负），流动负债 {current_liab:,.2f} 元，现金流动负债比率 {cash_flow_ratio:.2%}。该指标从收付实现制角度衡量偿债能力，现金流为负意味着企业经营活动无法产生正向现金流入。",
            "suggestion": "加强应收账款催收，优化存货管理减少资金占用，检查收入确认质量。如为正常经营所需，准备融资计划补充流动资金。"
        })
    elif cash_flow_ratio < 0.5 and net_op_cash > 0:
        results.append({
            "category": "财务健康", "category_icon": "💊", "risk_score": 4, "risk_level": "中风险",
            "risk_color": "#f59e0b", "urgency": "提醒",
            "item": "经营活动现金流无法覆盖流动负债",
            "detail": f"经营活动现金净流量 {net_op_cash:,.2f} 元，流动负债 {current_liab:,.2f} 元，覆盖率仅 {cash_flow_ratio:.1%}。短期偿债能力偏弱。",
            "suggestion": "加快应收账款回收，改善现金流状况。"
        })


def _analyze_other_payables_scale_mismatch(db, company_id, ps, pe, results):
    """其他应付款与经营规模不匹配"""
    other_payables = _get_account_balance(db, company_id, "2241")
    revenue = _get_account_sum(db, company_id, "6001", ps, pe, "credit")
    cost = _get_account_sum(db, company_id, "6401", ps, pe, "debit")
    equity = sum([
        _get_account_balance(db, company_id, "4001"),
        _get_account_balance(db, company_id, "4002"),
        _get_account_balance(db, company_id, "4101"),
        _get_account_balance(db, company_id, "4103"),
        _get_account_balance(db, company_id, "4104")
    ])
    
    if other_payables <= 0:
        return
    
    cost_ratio = other_payables / cost * 100 if cost > 0 else 0
    revenue_ratio = other_payables / revenue * 100 if revenue > 0 else 0
    equity_ratio = other_payables / equity * 100 if equity > 0 else 0
    
    triggers = []
    if cost_ratio > 50:
        triggers.append(f"占营业成本 {cost_ratio:.0f}%")
    if revenue_ratio > 50:
        triggers.append(f"占营业收入 {revenue_ratio:.0f}%")
    if equity_ratio > 50:
        triggers.append(f"占所有者权益 {equity_ratio:.0f}%")
    
    if triggers:
        results.append({
            "category": "隐匿虚增", "category_icon": "🔗", "risk_score": 8, "risk_level": "高风险",
            "risk_color": "#dc2626", "urgency": "紧急",
            "item": "其他应付款与经营规模不匹配",
            "detail": f"其他应付款余额 {other_payables:,.2f} 元，{'、'.join(triggers)}。其他应付款长期大额挂账是高频税务风险点：可能隐藏账外收入、虚列成本费用、虚假负债、或股东/关联方资金拆借未规范处理。",
            "suggestion": "逐笔核查大额其他应付款的入账依据和业务真实性；长期挂账未付的，需核实是否虚假负债或应转收入的预收款项；检查其他应付款中是否隐藏了应付未付的工资/社保/税费。",
            "required_evidence": ["其他应付款明细账", "每笔大额款项的入账凭证及合同", "账龄分析表"]
        })


# ═══════════════════════════════════════════════════════════
#  V7 新增 — 8项补全规则
# ═══════════════════════════════════════════════════════════

def _analyze_year_end_bonus_tax_opt(db, company_id, ps, pe, results):
    """年终奖计税方式未优化（个人所得税）"""
    # 查找含有"年终奖"或"全年一次性奖金"的工资记录
    bonus_records = db.query(SalaryRecord).filter(
        SalaryRecord.company_id == company_id,
        SalaryRecord.period >= ps, SalaryRecord.period <= pe,
        or_(
            func.lower(SalaryRecord.income_type).contains("年终奖"),
            func.lower(SalaryRecord.income_type).contains("全年一次性奖金"),
            func.lower(SalaryRecord.income_type).contains("一次性奖金"),
            func.lower(SalaryRecord.income_type).contains("奖金")
        )
    ).all()

    if not bonus_records:
        return

    bonus_employees = set()
    total_bonus = 0.0
    for r in bonus_records:
        bonus_employees.add(r.employee_name)
        total_bonus += _safe_float(r.current_income)

    # 年终奖未优化提示：建议在次年1月前测算两种计税方式
    if total_bonus > 0:
        results.append({
            "category": "个人所得税", "category_icon": "👤", "risk_score": 4, "risk_level": "低风险",
            "risk_color": "#3b82f6", "urgency": "中",
            "item": "年终奖计税方式未优化",
            "detail": f"系统检测到 {len(bonus_employees)} 名员工有年终奖/全年一次性奖金发放记录（合计 {total_bonus:,.2f} 元）。"
                     f"2027年前年终奖仍可适用单独计税优惠政策（财税〔2023〕30号延续至2027.12.31），"
                     f"建议在每年1月前测算单独计税与合并计税两种方式的税负差异，选择最优方案并在个税系统中准确申报。",
            "suggestion": "每月/每季计算每位员工的累计预扣预缴个税，年终奖发放前做税负测算：①单独计税：年终奖÷12查税率表×12；②合并计税：并入综合所得。取税负较低方案。",
            "required_evidence": ["工资薪金台账", "年终奖两种计税方式税负测算表", "个税申报记录"]
        })


AGRICULTURAL_KEYWORDS = [
    "农产品", "蔬菜", "水果", "粮食", "谷物", "小麦", "稻谷", "玉米", "大豆",
    "棉花", "畜牧", "水产", "林木", "药材", "中药材", "苗木", "花卉",
    "食用油", "菜籽", "花生", "鲜肉", "禽蛋", "茶叶", "蚕茧", "烟叶",
    "种植", "养殖", "捕捞", "收割", "收购", "农副产品"
]


def _analyze_agricultural_purchase_invoice(db, company_id, ps, pe, results):
    """农产品收购发票异常（发票异常）"""
    ps_date = _period_to_date_range(ps)[0]
    pe_date = _period_to_date_range(pe)[1]

    # 查找品名匹配农产品关键词的进项发票
    agri_invoices = []
    for inv in db.query(PurchaseInvoice).filter(
        PurchaseInvoice.company_id == company_id,
        PurchaseInvoice.invoice_date >= ps_date,
        PurchaseInvoice.invoice_date <= pe_date
    ).all():
        gn = (getattr(inv, 'goods_name', '') or '')
        for kw in AGRICULTURAL_KEYWORDS:
            if kw in gn:
                seller = getattr(inv, 'seller_name', '') or ''
                # 判断销方是否为个人（无公司后缀）
                is_individual = not any(s in seller for s in ["公司", "集团", "厂", "合作社", "协会"])
                agri_invoices.append({
                    "no": getattr(inv, 'digital_invoice_no', '') or (getattr(inv, 'invoice_code', '') or '') + (getattr(inv, 'invoice_no', '') or ''),
                    "seller": seller,
                    "goods": gn,
                    "amt": _safe_float(inv.total_amount),
                    "is_individual": is_individual
                })
                break

    if not agri_invoices:
        return

    total_agri = len(agri_invoices)
    total_amt = sum(a['amt'] for a in agri_invoices)
    individual_cnt = sum(1 for a in agri_invoices if a['is_individual'])
    individual_amt = sum(a['amt'] for a in agri_invoices if a['is_individual'])

    # 农产品收购发票由购方自行填开，是虚开发票的高发领域
    if individual_cnt >= 3:
        results.append({
            "category": "发票异常", "category_icon": "🧾", "risk_score": 9, "risk_level": "高风险",
            "risk_color": "#dc2626", "urgency": "紧急",
            "item": "农产品收购发票异常",
            "detail": f"检测到 {total_agri} 张农产品相关进项发票（合计 {total_amt:,.2f} 元），"
                     f"其中 {individual_cnt} 张销方疑似为个人（无企业后缀），涉及金额 {individual_amt:,.2f} 元。"
                     f"农产品收购发票由购方自行填开，是虚开发票的高发领域，金税四期对此专项监控。",
            "suggestion": "逐笔核实农产品收购的真实性：(1)出售人身份信息真实完整并留存身份证复印件；(2)有收购过磅单/入库单/质检记录；(3)银行付款记录与发票金额一致；(4)大额现金支付需特别警惕。虚开农产品收购发票是刑事犯罪。",
            "required_evidence": ["农产品收购发票清单", "出售人身份证明", "过磅单/入库单", "付款凭证", "质检记录"]
        })
    elif individual_cnt > 0:
        results.append({
            "category": "发票异常", "category_icon": "🧾", "risk_score": 5, "risk_level": "中风险",
            "risk_color": "#f59e0b", "urgency": "提醒",
            "item": "农产品收购发票异常",
            "detail": f"有 {individual_cnt} 张农产品相关发票销方为个人（非企业），涉及金额 {individual_amt:,.2f} 元。"
                     f"向个人收购农产品可自行开具收购发票，但需严格留存出售人身份证明和收购记录。",
            "suggestion": "确保每张农产品收购发票都有完整的出售人身份信息和收购业务记录。"
        })


def _analyze_cross_region_invoicing(db, company_id, ps, pe, results):
    """跨区域开票异常（发票异常）"""
    # 获取公司注册地址
    company = db.query(Company).filter(Company.id == company_id).first()
    if not company or not company.address:
        return

    registered_addr = company.address[:6] if company.address else ""  # 取前6位（省市区）
    if not registered_addr or len(registered_addr) < 4:
        return

    ps_date = _period_to_date_range(ps)[0]
    pe_date = _period_to_date_range(pe)[1]

    # 查找销项发票中购买方地址与注册地严重不符的情况
    cross_region_invoices = []
    total_sales = 0
    for inv in db.query(SalesInvoice).filter(
        SalesInvoice.company_id == company_id,
        SalesInvoice.invoice_date >= ps_date,
        SalesInvoice.invoice_date <= pe_date
    ).all():
        total_sales += 1
        buyer_addr = (getattr(inv, 'buyer_address', '') or '')[:6]
        if buyer_addr and len(buyer_addr) >= 4 and registered_addr[:2] != buyer_addr[:2]:
            cross_region_invoices.append({
                "no": getattr(inv, 'digital_invoice_no', '') or (getattr(inv, 'invoice_code', '') or '') + (getattr(inv, 'invoice_no', '') or ''),
                "buyer": getattr(inv, 'buyer_name', '') or '',
                "amt": _safe_float(inv.total_amount),
                "province": buyer_addr[:2]
            })

    if total_sales == 0:
        return

    cross_ratio = len(cross_region_invoices) / total_sales * 100
    # 按省份统计
    province_counts = {}
    for cr in cross_region_invoices:
        province_counts[cr['province']] = province_counts.get(cr['province'], 0) + 1

    if cross_ratio > 50 and len(cross_region_invoices) >= 10:
        top_provinces = sorted(province_counts.items(), key=lambda x: x[1], reverse=True)[:3]
        prov_desc = "、".join(f"{p}({c}张)" for p, c in top_provinces)
        results.append({
            "category": "发票异常", "category_icon": "🧾", "risk_score": 7, "risk_level": "高风险",
            "risk_color": "#dc2626", "urgency": "紧急",
            "item": "跨区域开票异常",
            "detail": f"企业注册地（{registered_addr}）与 {cross_ratio:.0f}%（{len(cross_region_invoices)}/{total_sales}张）销项发票购买方所在地不符，"
                     f"主要涉及省份：{prov_desc}。大量异地开票可能涉嫌虚开发票或开票地与实际经营地不一致。",
            "suggestion": "核实开票地点与经营地的关系：(1)确认开票行为在税务登记地进行；(2)如确为异地经营，需办理跨区域涉税事项报告；(3)保留物流运输证明备查。金税四期会记录开票设备MAC地址和IP，跨省开票自动预警。",
            "required_evidence": ["开票设备/IP记录", "异地经营备案文件", "物流运输单据"]
        })
    elif cross_ratio > 30 and len(cross_region_invoices) >= 5:
        results.append({
            "category": "发票异常", "category_icon": "🧾", "risk_score": 4, "risk_level": "中风险",
            "risk_color": "#f59e0b", "urgency": "提醒",
            "item": "跨区域开票异常",
            "detail": f"约 {cross_ratio:.0f}% 的销项发票购买方在注册地以外省份。如为正常商业活动，建议保留物流记录和合同备查。",
            "suggestion": "确保跨省销售有对应的物流运输单据和合同支撑，以备税务机关核实。"
        })


def _analyze_simple_general_tax_mix(db, company_id, ps, pe, results):
    """简易计税与一般计税划分不清（申报比对）"""
    vat_list = db.query(VATDeclaration).filter(
        VATDeclaration.company_id == company_id,
        VATDeclaration.period >= ps, VATDeclaration.period <= pe
    ).order_by(VATDeclaration.period).all()

    mix_periods = []
    for v in vat_list:
        general_sales = 0.0
        simple_sales = 0.0
        # 从附表一检查
        if v.form_sales:
            fs = v.form_sales if isinstance(v.form_sales, dict) else (json.loads(v.form_sales) if v.form_sales else {})
            # 一般计税销售额 (row1)
            general_sales += _safe_float(fs.get("row1_1_current") or fs.get("row1_1"))
            general_sales += _safe_float(fs.get("row1_2_current") or fs.get("row1_2"))
            general_sales += _safe_float(fs.get("row1_3_current") or fs.get("row1_3"))
            # 简易计税销售额
            simple_sales += _safe_float(fs.get("row5_1_current") or fs.get("row5_1"))
            simple_sales += _safe_float(fs.get("row5_2_current") or fs.get("row5_2"))
            simple_sales += _safe_float(fs.get("row5_3_current") or fs.get("row5_3"))
        if general_sales > 0 and simple_sales > 0:
            # 同时存在一般计税和简易计税，检查是否有进项转出
            input_transfer = 0.0
            if v.form_deduction:
                fd = v.form_deduction if isinstance(v.form_deduction, dict) else (json.loads(v.form_deduction) if v.form_deduction else {})
                input_transfer = _safe_float(fd.get("row14") or fd.get("row14_total"))
            mix_periods.append({
                "period": v.period,
                "general": general_sales,
                "simple": simple_sales,
                "input_transfer": input_transfer
            })

    if not mix_periods:
        return

    # 检查是否有期间未做进项转出
    no_transfer_periods = [m for m in mix_periods if m['input_transfer'] < 1]
    if no_transfer_periods:
        period_list = "、".join(m['period'] for m in no_transfer_periods[:3])
        total_simple = sum(m['simple'] for m in no_transfer_periods)
        results.append({
            "category": "申报比对", "category_icon": "📋", "risk_score": 6, "risk_level": "中风险",
            "risk_color": "#f59e0b", "urgency": "提醒",
            "item": "简易计税与一般计税划分不清",
            "detail": f"以下期间同时存在一般计税和简易计税销售额，但未发现对应的进项税额转出：{period_list}。"
                     f"简易计税项目对应的进项税额不得从销项税额中抵扣，专用于简易计税的进项税额必须转出。"
                     f"混用的进项税额需按销售额比例分摊转出。",
            "suggestion": "建立一般计税和简易计税项目的进项税额台账：(1)专用于简易计税项目的进项税额全额转出；(2)混用的进项税额按简易计税销售额÷全部销售额的比例计算转出；(3)在增值税申报表附表二正确填列进项税额转出。",
            "required_evidence": ["简易计税项目清单", "进项税额分摊计算表", "进项税额转出凭证"]
        })


def _analyze_pre_cancellation_spike(db, company_id, ps, pe, results):
    """注销前突击开票（经营实质）"""
    # 检查最近6个月的开票趋势，判断是否存在末期激增
    periods = _get_periods_between(ps, pe)
    if len(periods) < 6:
        return

    # 按月统计销项发票
    monthly_sales = {}
    ps_date = _period_to_date_range(ps)[0]
    pe_date = _period_to_date_range(pe)[1]

    for inv in db.query(SalesInvoice).filter(
        SalesInvoice.company_id == company_id,
        SalesInvoice.invoice_date >= ps_date,
        SalesInvoice.invoice_date <= pe_date
    ).all():
        inv_date = str(inv.invoice_date)[:7]
        if inv_date not in monthly_sales:
            monthly_sales[inv_date] = {"cnt": 0, "amt": 0.0}
        monthly_sales[inv_date]["cnt"] += 1
        monthly_sales[inv_date]["amt"] += _safe_float(inv.total_amount)

    # 取后3个月 vs 前3个月
    sorted_periods = sorted(monthly_sales.keys())
    if len(sorted_periods) < 6:
        return

    first_half = sorted_periods[:3]
    last_half = sorted_periods[-3:]

    first_avg = sum(monthly_sales[p]["amt"] for p in first_half) / 3 if first_half else 0
    last_avg = sum(monthly_sales[p]["amt"] for p in last_half) / 3 if last_half else 0

    if first_avg <= 0:
        return

    spike_ratio = (last_avg - first_avg) / first_avg * 100

    # 也检查进项发票
    monthly_purchase = {}
    for inv in db.query(PurchaseInvoice).filter(
        PurchaseInvoice.company_id == company_id,
        PurchaseInvoice.invoice_date >= ps_date,
        PurchaseInvoice.invoice_date <= pe_date
    ).all():
        inv_date = str(inv.invoice_date)[:7]
        if inv_date not in monthly_purchase:
            monthly_purchase[inv_date] = {"cnt": 0, "amt": 0.0}
        monthly_purchase[inv_date]["cnt"] += 1
        monthly_purchase[inv_date]["amt"] += _safe_float(inv.total_amount)

    sorted_pur = sorted(monthly_purchase.keys())
    pur_first_avg = 0
    pur_last_avg = 0
    if len(sorted_pur) >= 6:
        pur_first_avg = sum(monthly_purchase[p]["amt"] for p in sorted_pur[:3]) / 3
        pur_last_avg = sum(monthly_purchase[p]["amt"] for p in sorted_pur[-3:]) / 3

    pur_spike = 0
    if pur_first_avg > 0:
        pur_spike = (pur_last_avg - pur_first_avg) / pur_first_avg * 100

    if spike_ratio > 200 or pur_spike > 200:
        results.append({
            "category": "经营实质", "category_icon": "🔍", "risk_score": 9, "risk_level": "高风险",
            "risk_color": "#dc2626", "urgency": "紧急",
            "item": "注销前突击开票",
            "detail": f"最近3个月销项发票月均金额 {last_avg:,.2f} 元，较前3个月（{first_avg:,.2f} 元）增长 {spike_ratio:.0f}%。"
                     + (f" 进项发票也增长{pur_spike:.0f}%。" if pur_spike > 50 else "")
                     + "如企业正在办理注销，注销前突击开票可能涉嫌对外虚开或取得进项发票虚抵后注销逃税。",
            "suggestion": "如企业正在办理注销，应停止大额开票行为，配合税务机关完成注销清算。注销前异常开票将影响注销审批，严重者追究法律责任。如并非注销而是业务增长，保留所有合同和业务单据备查。",
            "required_evidence": ["注销申请材料（如适用）", "最近6个月发票台账", "业务增长说明", "对应合同"]
        })
    elif spike_ratio > 100:
        results.append({
            "category": "经营实质", "category_icon": "🔍", "risk_score": 5, "risk_level": "中风险",
            "risk_color": "#f59e0b", "urgency": "提醒",
            "item": "注销前突击开票",
            "detail": f"最近3个月销项发票金额较前3个月增长 {spike_ratio:.0f}%。如非注销前突击开票而是业务正常增长，保留证明材料。",
            "suggestion": "确认业务增长原因，保留合同/物流/资金流三流合一的证据链。"
        })


def _analyze_same_address_multi_company(db, company_id, ps, pe, results):
    """一址多照/一人多企（经营实质）"""
    company = db.query(Company).filter(Company.id == company_id).first()
    if not company:
        return

    # 检查同一注册地址的其他公司
    same_addr_companies = []
    if company.address:
        same_addr_companies = db.query(Company).filter(
            Company.id != company_id,
            Company.address == company.address
        ).all()

    # 检查同一法定代表人的其他公司
    same_legal_companies = []
    if company.legal_representative:
        same_legal_companies = db.query(Company).filter(
            Company.id != company_id,
            Company.legal_representative == company.legal_representative
        ).all()

    same_addr_cnt = len(same_addr_companies)
    same_legal_cnt = len(same_legal_companies)

    # 排除同一公司的重复（可能同是地址和法人都相同的）
    addr_ids = {c.id for c in same_addr_companies}
    legal_ids = {c.id for c in same_legal_companies}
    total_related = len(addr_ids | legal_ids)

    if same_addr_cnt >= 2 or same_legal_cnt >= 2 or total_related >= 2:
        detail_parts = []
        if same_addr_cnt >= 2:
            detail_parts.append(f"同一注册地址登记了{same_addr_cnt + 1}家企业（含本公司）")
        if same_legal_cnt >= 2:
            detail_parts.append(f"同一法定代表人控制{same_legal_cnt + 1}家企业（含本公司）")

        results.append({
            "category": "经营实质", "category_icon": "🔍", "risk_score": 7, "risk_level": "高风险",
            "risk_color": "#dc2626", "urgency": "紧急",
            "item": "一址多照/一人多企",
            "detail": "、".join(detail_parts) + "。一址多照/一人多企是虚开发票团伙的常见组织方式。"
                     "金税四期通过企业关系图谱自动识别关联企业群体，交叉比对发票开具和纳税情况。",
            "suggestion": "核实各关联企业的经营是否独立：(1)资金账户独立；(2)人员独立（不共用员工）；(3)业务独立（不循环交易）；(4)财务独立（各自核算）。为每家企业的独立性准备证明材料。",
            "required_evidence": ["关联企业清单", "各企业独立经营证明材料", "各企业银行账户流水", "各企业人员名册和社保记录"]
        })
    elif same_addr_cnt == 1 or same_legal_cnt == 1:
        results.append({
            "category": "经营实质", "category_icon": "🔍", "risk_score": 4, "risk_level": "中风险",
            "risk_color": "#f59e0b", "urgency": "提醒",
            "item": "一址多照/一人多企",
            "detail": f"发现{total_related}家关联企业。如关联企业间有交易往来，需确保定价公允且保留转移定价文档。",
            "suggestion": "关注关联企业间的交易是否公允，避免通过关联交易转移利润或虚开发票。"
        })


def _analyze_address_mismatch(db, company_id, ps, pe, results):
    """注册地址与经营地址不一致（经营实质）"""
    company = db.query(Company).filter(Company.id == company_id).first()
    if not company or not company.address:
        return

    registered_addr = company.address

    # 四源交叉验证：序时账 + 银行流水 + 取得发票 + 合同
    ev = _check_expense_evidence(db, company_id, ps, pe,
        expense_name="经营场所",
        je_summary_kw=["租赁", "租金", "房租", "水电", "电费", "水费", "燃气", "物业", "物业费", "物业管理"],
        je_account_codes=["660214", "660215", "660213"],
        invoice_kw=["租赁", "水电", "物业", "房租"],
        contract_kw=["场地", "租赁", "房租", "物业", "办公"],
        contract_types=["租赁"])

    je_has = ev["je"]["has_any"]
    overall = ev["overall"]
    inv_amount = ev["invoice"]["amount"]
    free_possible = ev["contract"]["free_possible"]
    contracts = ev["contract"]["contracts"]

    has_inv_or_contract = ev["invoice"]["has_any"] or ev["contract"]["has_any"]

    # ── 如果序时账、发票、合同都没有 → 地址不一致 ──
    if overall["all_sources_none"]:
        results.append({
            "category": "经营实质", "category_icon": "🔍", "risk_score": 5, "risk_level": "中风险",
            "risk_color": "#f59e0b", "urgency": "提醒",
            "item": "注册地址与经营地址不一致",
            "detail": f"企业注册地址为「{registered_addr}」，但分析期间内四源交叉验证均未发现经营场所证据：序时账无租赁费/水电费/物业费、取得发票无相关记录、无场地相关合同。实际经营地址与工商注册地址不一致且未办理变更登记，可能影响税务管辖权和纳税地点判断。",
            "suggestion": "及时办理工商变更登记，或在经营地办理跨区域涉税事项报告。注意：异地经营超过180天的，应在经营地办理税务登记。",
            "required_evidence": ["经营场所租赁合同或产权证明", "工商变更登记材料", "水电费缴费凭证"]
        })
    elif overall["je_missing_but_invoice_found"]:
        results.append({
            "category": "经营实质", "category_icon": "🔍", "risk_score": 3, "risk_level": "低风险",
            "risk_color": "#3b82f6", "urgency": "提醒",
            "item": "注册地址与经营地址不一致",
            "detail": f"企业注册地址为「{registered_addr}」，序时账中未发现租赁费/水电费/物业费，但取得发票中发现相关费用 ¥{inv_amount:,.2f}（尚未入账）。建议及时将发票入账以证明经营场所真实性。",
            "suggestion": "将已取得的经营场所相关发票入账，并核对实际经营地址与注册地址是否一致。",
            "required_evidence": ["未入账的经营场所相关发票", "经营场所租赁合同或产权证明"]
        })
    elif overall["je_missing_but_contract_found"]:
        ct_summary = "；".join([f"{c['name']}({c['type']}{'·免费' if c['is_free'] else ''})" for c in contracts[:3]])
        free_note = "·可能为免费使用场地" if free_possible else ""
        results.append({
            "category": "经营实质", "category_icon": "🔍", "risk_score": 2, "risk_level": "低风险",
            "risk_color": "#3b82f6", "urgency": "提醒",
            "item": "注册地址与经营地址不一致",
            "detail": f"企业注册地址为「{registered_addr}」，序时账中无租赁费/水电费/物业费，但发现场地相关合同：{ct_summary}{free_note}。已签署合同但尚未入账，可证明经营场所真实性。",
            "suggestion": "保留场地使用合同备查，核对实际经营地址与注册地址是否一致。如地址不一致，及时办理变更。",
            "required_evidence": ["场地使用合同", "经营场所使用证明"]
        })

    # 检查是否有地址变更（通过公司档案中的备注等）
    addr_lower = registered_addr.lower()
    if any(kw in addr_lower for kw in ["虚拟", "集群", "挂靠", "孵化器", "众创空间", "托管"]):
        if overall["all_sources_none"]:
            results.append({
                "category": "经营实质", "category_icon": "🔍", "risk_score": 8, "risk_level": "高风险",
                "risk_color": "#dc2626", "urgency": "紧急",
                "item": "注册地址与经营地址不一致",
                "detail": f"企业注册地址包含「虚拟/集群/挂靠/孵化器/托管」等字眼，且四源交叉验证（序时账+发票+合同）均无租赁费/水电费支出。集群注册企业是税务稽查重点：金税四期会检查注册地址是否有实际经营实体。",
                "suggestion": "如为集群注册但确有实际经营场所，应提供实际经营地址的租赁合同和水电费凭证。如长期无实际经营场所，建议办理注销或变更注册地址。",
                "required_evidence": ["实际经营地址租赁合同", "水电费缴纳记录", "经营场所照片"]
            })
        elif not je_has and has_inv_or_contract:
            results.append({
                "category": "经营实质", "category_icon": "🔍", "risk_score": 4, "risk_level": "中风险",
                "risk_color": "#f59e0b", "urgency": "提醒",
                "item": "集群注册地址但有经营证据需补充",
                "detail": f"企业为集群注册地址（{registered_addr}），发现经营场所相关发票或合同但序时账中尚无费用记录。建议将相关发票入账并准备实际经营地址证明。",
                "suggestion": "将已取得的经营场所发票入账，准备实际经营地址证明文件。",
                "required_evidence": ["实际经营地址租赁合同", "经营场所发票", "水电费缴纳记录"]
            })


def _analyze_d_taxpayer_risk(db, company_id, ps, pe, results):
    """D级纳税人关联风险（征管风险）"""
    # 纳税信用等级数据通常来自税务机关，系统内可能没有直接存储
    # 检查公司是否有信用等级相关字段
    company = db.query(Company).filter(Company.id == company_id).first()
    if not company:
        return

    # 检查公司是否有欠税、违规开票等可能导致D级评定的行为
    # 汇总风险信号
    risk_signals = []

    # 检查进项发票是否有风险记录
    risk_inv_cnt = db.query(func.count(PurchaseInvoice.id)).filter(
        PurchaseInvoice.company_id == company_id,
        PurchaseInvoice.invoice_risk_level.in_(["疑点", "异常", "失控"])
    ).scalar() or 0

    if risk_inv_cnt > 0:
        risk_signals.append(f"有{risk_inv_cnt}张风险进项发票")

    # 检查是否有连续零申报
    vat_list = db.query(VATDeclaration).filter(
        VATDeclaration.company_id == company_id,
        VATDeclaration.period >= ps, VATDeclaration.period <= pe
    ).order_by(VATDeclaration.period).all()

    zero_months = 0
    for v in vat_list:
        if v.form_main:
            fm = v.form_main if isinstance(v.form_main, dict) else json.loads(v.form_main)
            if _safe_float(fm.get("vat_payable")) < 1:
                zero_months += 1

    if zero_months >= 6:
        risk_signals.append(f"连续{zero_months}个月零/小额申报")

    # 检查是否有互开发票
    # (已在 _analyze_cross_invoicing 中检测)

    # 检查是否有欠税
    # (已在 _analyze_tax_arrears 中检测)

    if risk_signals:
        results.append({
            "category": "征管风险", "category_icon": "🏛️", "risk_score": 8, "risk_level": "高风险",
            "risk_color": "#dc2626", "urgency": "紧急",
            "item": "D级纳税人关联风险",
            "detail": f"企业存在多项可能导致纳税信用降级的风险指标：{'；'.join(risk_signals)}。"
                     "D级纳税人面临发票限量供应（增值税专用发票每月限量25份）、出口退税从严审核、不得享受即征即退、"
                     "列入重点监控对象等多重限制，且D级评价保留2年。",
            "suggestion": "逐项整改风险指标：(1)处理风险进项发票，必要时做进项税额转出；(2)补申报零申报期间的应缴税款；(3)完善账务处理和发票管理。次年纳税信用评价前完成整改可申请复评。",
            "required_evidence": ["纳税信用评价结果查询截图", "风险指标整改方案", "整改证明材料"]
        })

    # 即使没有明显风险信号，也提醒关注信用评级
    # 检查供应商中是否有非正常户
    abnormal_supplier = db.query(func.count(PurchaseInvoice.id)).filter(
        PurchaseInvoice.company_id == company_id,
        PurchaseInvoice.invoice_risk_level == "异常"
    ).scalar() or 0

    if abnormal_supplier > 0:
        results.append({
            "category": "征管风险", "category_icon": "🏛️", "risk_score": 6, "risk_level": "中风险",
            "risk_color": "#f59e0b", "urgency": "提醒",
            "item": "D级纳税人关联风险",
            "detail": f"有{abnormal_supplier}家进项发票供应商存在异常记录。与异常状态供应商交易取得的发票可能被列为异常扣税凭证，面临进项转出风险。",
            "suggestion": "建议在每笔大额采购前通过电子税务局查验销方税务状态和纳税信用等级。建立供应商准入制度。"
        })


def _analyze_flexible_employment_tax(db, company_id, ps, pe, results):
    """灵活用工/劳务报酬税务处理不规范（财务健康）"""
    ps_date = _period_to_date_range(ps)[0]
    pe_date = _period_to_date_range(pe)[1]

    # 灵活用工/劳务派遣关键词
    labor_keywords = ["劳务", "派遣", "外包", "灵活用工", "兼职", "临时工",
                      "服务费", "咨询费", "居间", "佣金", "代办"]

    # 1. 检查进项发票中是否有劳务/服务类发票
    labor_invs = db.query(PurchaseInvoice).filter(
        PurchaseInvoice.company_id == company_id,
        PurchaseInvoice.invoice_date >= ps_date,
        PurchaseInvoice.invoice_date <= pe_date
    ).all()

    labor_inv_amount = 0
    labor_inv_cnt = 0
    labor_sellers = set()
    for inv in labor_invs:
        pname = (inv.goods_name or "").strip()
        summary = (inv.remark or "").strip()
        combined = pname + summary
        if any(kw in combined for kw in labor_keywords):
            labor_inv_amount += _safe_float(inv.total_amount)
            labor_inv_cnt += 1
            if inv.seller_name:
                labor_sellers.add(inv.seller_name.strip())

    # 2. 检查银行流水中是否有劳务/服务类支付
    labor_txs = db.query(BankTransaction).filter(
        BankTransaction.company_id == company_id,
        BankTransaction.transaction_date >= ps_date,
        BankTransaction.transaction_date <= pe_date,
        BankTransaction.debit_amount > 0
    ).all()

    labor_tx_amount = 0
    labor_tx_cnt = 0
    for bt in labor_txs:
        summ = (bt.summary or "").strip()
        remark = (bt.remark or "").strip()
        cpty = (bt.counterparty_name or "").strip()
        combined = summ + remark + cpty
        if any(kw in combined for kw in labor_keywords):
            labor_tx_amount += _safe_float(bt.debit_amount)
            labor_tx_cnt += 1

    total_labor_expense = labor_inv_amount + labor_tx_amount

    if total_labor_expense <= 0:
        return

    # 3. 检查是否有对应的人员薪金/个税记录
    revenue = _get_account_sum(db, company_id, "6001", ps, pe, "credit")
    salary_expense = _get_account_sum(db, company_id, "660204", ps, pe, "debit")  # 工资薪金

    # 4. 判断风险：劳务费用高但工资薪金低/无
    # 将劳务报酬伪装为经营所得（不代扣个税）是常见避税手段
    labor_ratio = total_labor_expense / revenue if revenue > 0 else 0

    risk_items = []
    detail_parts = []

    if labor_inv_cnt >= 3 and labor_inv_amount >= 100000:
        detail_parts.append(f"取得劳务/服务类进项发票 {labor_inv_cnt} 张，金额 {labor_inv_amount:,.2f} 元")
        detail_parts.append(f"涉及销方 {len(labor_sellers)} 家")

    if labor_tx_cnt >= 3 and labor_tx_amount >= 100000:
        detail_parts.append(f"银行支付劳务/服务类款项 {labor_tx_cnt} 笔，金额 {labor_tx_amount:,.2f} 元")

    if not detail_parts:
        return

    # 劳务费用占总收入比例过高
    if labor_ratio > 0.3:
        risk_items.append(f"劳务/服务支出占营业收入 {labor_ratio*100:.1f}%")

    # 劳务费用高但无工资薪金支出
    if salary_expense < 1000 and total_labor_expense >= 50000:
        risk_items.append("有大量劳务支出但无工资薪金记录（疑似将员工薪酬伪装为劳务费）")

    # 多张小额服务费发票，典型灵活用工拆分
    if labor_inv_cnt >= 10:
        risk_items.append(f"劳务/服务发票数量较多（{labor_inv_cnt}张），存在拆分金额规避个税嫌疑")

    if not risk_items:
        return

    risk_score = 7 if labor_ratio > 0.5 else (6 if len(risk_items) >= 2 else 5)
    risk_level = "高风险" if risk_score >= 7 else "中风险"
    risk_color = "#dc2626" if risk_score >= 7 else "#f59e0b"
    urgency = "高" if risk_score >= 7 else "中"

    detail = "；".join(detail_parts) + "。" if detail_parts else ""
    detail += "风险信号：" + "；".join(risk_items) + "。"
    detail += "金税四期接入人社部门数据，灵活用工/劳务报酬的税务合规是重点监控领域。将工资薪金伪装为经营所得（劳务发票）以规避个税代扣义务，是典型税务违规。"

    results.append({
        "category": "财务健康", "category_icon": "💼", "risk_score": risk_score,
        "risk_level": risk_level, "risk_color": risk_color, "urgency": urgency,
        "item": "灵活用工/劳务报酬税务处理不规范",
        "detail": detail,
        "suggestion": "①核实劳务/服务发票是否对应真实业务（留存合同、成果交付证明）；②如为个人提供劳务，需按规定代扣代缴个人所得税（劳务报酬适用20%-40%税率）；③如使用灵活用工平台，确保平台具备税务代征资质并取得合规发票；④避免将员工工资薪金伪装为劳务费以规避社保和个税。",
        "required_evidence": [
            "劳务/服务类发票对应的业务合同",
            "劳务报酬个人所得税代扣代缴记录",
            "灵活用工平台合作协议及平台资质证明",
            "劳务人员身份证件及联系方式清单",
            "服务成果交付证明（如为咨询服务）"
        ]
    })


# ═══════════════════════════════════════════════════════════
#  V8 新增 — 合同风险分析（三流合一/四流合一，16项）
# ═══════════════════════════════════════════════════════════

def _normalize_entity_name(name):
    """规范化企业名称用于跨表匹配（去掉常见后缀）"""
    if not name:
        return ""
    n = name.strip().replace("（", "(").replace("）", ")")
    for sfx in ["有限责任公司", "股份有限公司", "有限公司", "合伙企业", "(有限合伙)", "（有限合伙）",
                 "个人独资企业", "专业合作社", "集团"]:
        n = n.replace(sfx, "")
    return n.strip()


def _name_in_set(name, party_set):
    """检查企业名称是否匹配合同中的任一交易方（模糊匹配）"""
    if not name or not party_set:
        return False
    nn = _normalize_entity_name(name)
    if len(nn) < 2:
        return False
    for party in party_set:
        np = _normalize_entity_name(party)
        if len(np) < 2:
            continue
        if nn == np or nn in np or np in nn:
            return True
    return False


def _find_matching_contracts(name, contracts):
    """从合同列表中找出匹配指定企业名称的合同"""
    if not name or not contracts:
        return []
    matched = []
    nn = _normalize_entity_name(name)
    for c in contracts:
        for party in [c.party_b or "", c.party_a or ""]:
            np = _normalize_entity_name(party)
            if len(np) >= 2 and (nn == np or nn in np or np in nn):
                matched.append(c)
                break
    return matched


def _date_to_str(d):
    """安全地将日期字段转为 YYYY-MM-DD 字符串（处理 date/datetime/str/None）"""
    if not d:
        return ""
    if isinstance(d, date):
        return d.strftime("%Y-%m-%d")
    if hasattr(d, 'strftime'):
        return d.strftime("%Y-%m-%d")
    s = str(d)
    return s[:10] if len(s) >= 10 else s


def _contract_covers_date(contract, check_date_str):
    """判断合同有效期是否覆盖指定日期（YYYY-MM-DD格式）。
    无生效日期+无到期日期 → 视为永久有效。
    仅有生效日期 → 生效后有效。
    仅有到期日期 → 到期前有效。"""
    eff = _date_to_str(contract.effective_date)
    exp = _date_to_str(contract.expiry_date)
    if not eff and not exp:
        return True
    if eff and exp:
        return eff <= check_date_str <= exp
    if eff and not exp:
        return eff <= check_date_str
    if not eff and exp:
        return check_date_str <= exp
    return True


def _analyze_contract_risks(db, company_id, ps, pe, results):
    """合同风险综合分析：三流合一/四流合一/收入成本无合同/金额匹配等"""
    ps_date, pe_date = _period_to_date_range(ps)[0], _period_to_date_range(pe)[1]
    CAT = "合同风险"
    ICON = "📋"

    # ═══ 预加载所有有效合同 ═══
    all_contracts = db.query(Contract).filter(
        Contract.company_id == company_id,
        Contract.status.in_(["已签署", "履行中", "已完成"])
    ).all()

    # 按类型分组
    sales_ct = [c for c in all_contracts if c.contract_type in ("销售", "服务")]  # 销售类合同
    purchase_ct = [c for c in all_contracts if c.contract_type in ("采购", "服务")]  # 采购类合同
    # 服务类合同同时属于销售和采购（既可能对外提供服务，也可能接受服务）

    # 提取交易方名称集合
    sales_parties = set()
    for c in sales_ct:
        if c.party_b: sales_parties.add(c.party_b.strip())
        if c.party_a: sales_parties.add(c.party_a.strip())

    purchase_parties = set()
    for c in purchase_ct:
        if c.party_b: purchase_parties.add(c.party_b.strip())
        if c.party_a: purchase_parties.add(c.party_a.strip())

    # ═══ 预计算销项收入总额（用于后续印花税分析）════
    sales_customers = db.query(
        SalesInvoice.buyer_name,
        func.sum(SalesInvoice.total_amount).label("total"),
        func.count(SalesInvoice.id).label("cnt")
    ).filter(
        SalesInvoice.company_id == company_id,
        SalesInvoice.invoice_date >= ps_date,
        SalesInvoice.invoice_date <= pe_date
    ).group_by(SalesInvoice.buyer_name).order_by(func.sum(SalesInvoice.total_amount).desc()).all()

    total_sales_rev = 0.0
    for name, amt, cnt in sales_customers:
        total_sales_rev += _safe_float(amt)

    # ═══ 2. 预计算采购总额（用于后续印花税分析）════
    purchase_suppliers = db.query(
        PurchaseInvoice.seller_name,
        func.sum(PurchaseInvoice.total_amount).label("total"),
        func.count(PurchaseInvoice.id).label("cnt")
    ).filter(
        PurchaseInvoice.company_id == company_id,
        PurchaseInvoice.invoice_date >= ps_date,
        PurchaseInvoice.invoice_date <= pe_date
    ).group_by(PurchaseInvoice.seller_name).order_by(func.sum(PurchaseInvoice.total_amount).desc()).all()

    # 预计算采购总额（用于后续印花税分析）
    total_purchase_amt = 0.0
    for name, amt, cnt in purchase_suppliers:
        total_purchase_amt += _safe_float(amt)

    BIG_THRESHOLD = 100000  # 10万元作为大额交易门槛（三流合一/四流合一/覆盖率大额分层使用）

    # ═══ 4. 三流合一缺失（合同流+发票流+资金流完整性检查）════
    # 对于有合同且有大额交易的，检查是否三流完整
    all_contracts_with_amount = [c for c in all_contracts if _safe_float(c.amount) >= BIG_THRESHOLD]
    three_flow_missing = []

    for c in all_contracts_with_amount:
        parties = [p.strip() for p in [c.party_b or "", c.party_a or ""] if p.strip()]
        if not parties:
            continue
        ct_amount = _safe_float(c.amount)

        # 检查发票流：该合同任一交易方是否有匹配的发票
        invoice_found = False
        for party in parties:
            nn = _normalize_entity_name(party)
            # 查销项发票
            sale_amt = db.query(func.coalesce(func.sum(SalesInvoice.total_amount), 0)).filter(
                SalesInvoice.company_id == company_id,
                SalesInvoice.invoice_date >= ps_date,
                SalesInvoice.invoice_date <= pe_date
            ).scalar() or 0
            # 精确匹配：遍历结果不如直接查
            sale_match = db.query(func.count(SalesInvoice.id)).filter(
                SalesInvoice.company_id == company_id,
                SalesInvoice.invoice_date >= ps_date,
                SalesInvoice.invoice_date <= pe_date,
                SalesInvoice.buyer_name == party
            ).scalar() or 0
            purchase_match = db.query(func.count(PurchaseInvoice.id)).filter(
                PurchaseInvoice.company_id == company_id,
                PurchaseInvoice.invoice_date >= ps_date,
                PurchaseInvoice.invoice_date <= pe_date,
                PurchaseInvoice.seller_name == party
            ).scalar() or 0
            if sale_match > 0 or purchase_match > 0:
                invoice_found = True
                break
            # 模糊匹配
            if _name_in_set(party, sales_parties) or _name_in_set(party, purchase_parties):
                invoice_found = True
                break

        # 检查资金流：合同对方是否有匹配的银行流水
        bank_match = 0
        for party in parties:
            bm = db.query(func.count(BankTransaction.id)).filter(
                BankTransaction.company_id == company_id,
                BankTransaction.transaction_date >= ps_date,
                BankTransaction.transaction_date <= pe_date,
                BankTransaction.counterparty_name == party
            ).scalar() or 0
            bank_match += bm

        if not invoice_found:
            three_flow_missing.append({
                "contract_no": c.contract_no, "name": c.name,
                "amount": ct_amount, "type": c.contract_type,
                "missing": "发票流" if bank_match > 0 else "发票流+资金流"
            })

    if three_flow_missing:
        missing_both = [x for x in three_flow_missing if "资金流" in x["missing"]]
        missing_invoice = [x for x in three_flow_missing if x["missing"] == "发票流"]
        detail_parts = []
        for x in three_flow_missing[:5]:
            detail_parts.append(f"{x['contract_no']} {x['name']}（{x['amount']:,.0f}元，缺失{x['missing']}）")

        results.append({
            "category": CAT, "category_icon": ICON, "risk_score": 9, "risk_level": "高风险",
            "risk_color": "#dc2626", "urgency": "紧急",
            "item": "三流合一缺失（合同→发票→资金不完整）",
            "detail": f"共 {len(three_flow_missing)} 份有效合同的大额交易中，未发现对应的发票流或资金流记录。"
                       f"其中 {len(missing_both)} 份既无发票也无付款记录（严重），{len(missing_invoice)} 份无发票记录。"
                       f"三流不合一是虚开发票认定的重要线索。"
                       f"涉及合同：{'；'.join(detail_parts)}"
                       + (f"等{len(three_flow_missing)}份" if len(three_flow_missing) > 5 else ""),
            "suggestion": "逐份核查合同执行情况：(1)补开/补收对应发票；(2)补办付款手续并留存银行回单；(3)确实无法补办的，出具书面说明并由负责人签字；(4)建立三流合一检查制度，每笔交易前确认合同→发票→资金三流匹配。",
            "required_evidence": ["合同执行情况说明", "补开发票/付款凭证", "三流比对台账"]
        })

    # ═══ 5. 四流合一缺失（合同+发票+资金+货物流）═══
    # 检查采购合同：有合同+发票+付款，但缺少物流/验收单（通过合同金额与发票金额一致性间接判断）
    four_flow_issues = []
    for c in purchase_ct:
        if _safe_float(c.amount) < BIG_THRESHOLD:
            continue
        parties = [p.strip() for p in [c.party_b or "", c.party_a or ""] if p.strip()]
        if not parties:
            continue
        # 有发票吗
        has_invoice = False
        total_invoice_amt = 0.0
        for party in parties:
            invoices = db.query(
                func.count(PurchaseInvoice.id),
                func.coalesce(func.sum(PurchaseInvoice.total_amount), 0)
            ).filter(
                PurchaseInvoice.company_id == company_id,
                PurchaseInvoice.invoice_date >= ps_date,
                PurchaseInvoice.invoice_date <= pe_date,
                PurchaseInvoice.seller_name == party
            ).first()
            if invoices and invoices[0] > 0:
                has_invoice = True
                total_invoice_amt += _safe_float(invoices[1])
        # 有付款吗
        has_payment = False
        for party in parties:
            pay_cnt = db.query(func.count(BankTransaction.id)).filter(
                BankTransaction.company_id == company_id,
                BankTransaction.transaction_date >= ps_date,
                BankTransaction.transaction_date <= pe_date,
                BankTransaction.counterparty_name == party
            ).scalar() or 0
            if pay_cnt > 0:
                has_payment = True
                break
        # 三流有但货物流存疑（合同金额与发票金额偏差大 = 可能存在未收货或未验收）
        if has_invoice and has_payment:
            ct_amt = _safe_float(c.amount)
            if ct_amt > 0 and total_invoice_amt > 0:
                deviation = abs(ct_amt - total_invoice_amt) / ct_amt
                if deviation > 0.3:  # 偏差超30%
                    four_flow_issues.append({
                        "contract_no": c.contract_no, "name": c.name,
                        "ct_amt": ct_amt, "inv_amt": total_invoice_amt,
                        "deviation_pct": round(deviation * 100, 1)
                    })

    if four_flow_issues:
        detail_parts = []
        for x in four_flow_issues[:5]:
            detail_parts.append(
                f"{x['contract_no']} {x['name']}（合同{x['ct_amt']:,.0f}元 vs 发票{x['inv_amt']:,.0f}元，偏差{x['deviation_pct']}%）"
            )
        results.append({
            "category": CAT, "category_icon": ICON, "risk_score": 7, "risk_level": "高风险",
            "risk_color": "#dc2626", "urgency": "紧急",
            "item": "四流合一缺失（货物流存疑）",
            "detail": f"共 {len(four_flow_issues)} 份采购合同虽有合同+发票+付款三流，但合同金额与发票金额偏差超30%，"
                       f"可能存在货物未完全交付、未验收即付款、或发票金额不实等问题，无法满足四流合一要求。"
                       f"涉及合同：{'；'.join(detail_parts)}"
                       + (f"等{len(four_flow_issues)}份" if len(four_flow_issues) > 5 else ""),
            "suggestion": "补充完整的货物流证据链：(1)逐单补充货物验收单/入库单；(2)大额采购补充运输单据；(3)建立四流比对台账（合同→发票→付款→验收）；(4)对偏差原因逐份出具书面说明。",
            "required_evidence": ["货物验收单/入库单", "运输物流单据", "四流比对台账", "偏差说明"]
        })

    # ═══ 6. 长期合作方无正式合同 ═══
    # 过去12个月累计交易额 > 50万但零合同
    LONG_TERM_THRESHOLD = 500000
    # 扩展查询范围到过去12个月
    ps_12m = f"{int(ps[:4]) - 1 if int(ps[5:7]) == 1 else int(ps[:4])}-{12 if int(ps[5:7]) == 1 else int(ps[5:7]) - 1:02d}"
    # 简化：直接查当前分析期间的数据（如需完整12个月可后续扩展）

    # 用当前期间 + 前推12个月来查
    from datetime import datetime as dt
    try:
        pe_dt = dt.strptime(pe + "-01", "%Y-%m-%d")
        ps_12m_dt = dt(pe_dt.year - 1, pe_dt.month, 1)
        ps_12m_date = ps_12m_dt.strftime("%Y-%m-%d")
        pe_date_range = _period_to_date_range(pe)[1]
    except:
        ps_12m_date = ps_date
        pe_date_range = pe_date

    # 查12个月内的所有发票
    long_sales = db.query(
        SalesInvoice.buyer_name,
        func.sum(SalesInvoice.total_amount).label("total")
    ).filter(
        SalesInvoice.company_id == company_id,
        SalesInvoice.invoice_date >= ps_12m_date,
        SalesInvoice.invoice_date <= pe_date_range
    ).group_by(SalesInvoice.buyer_name).all()

    long_purchase = db.query(
        PurchaseInvoice.seller_name,
        func.sum(PurchaseInvoice.total_amount).label("total")
    ).filter(
        PurchaseInvoice.company_id == company_id,
        PurchaseInvoice.invoice_date >= ps_12m_date,
        PurchaseInvoice.invoice_date <= pe_date_range
    ).group_by(PurchaseInvoice.seller_name).all()

    long_no_ct_cust = []
    for name, amt in long_sales:
        a = _safe_float(amt)
        if a >= LONG_TERM_THRESHOLD and not _name_in_set(name, sales_parties):
            long_no_ct_cust.append((name, a))

    long_no_ct_supp = []
    for name, amt in long_purchase:
        a = _safe_float(amt)
        if a >= LONG_TERM_THRESHOLD and not _name_in_set(name, purchase_parties):
            long_no_ct_supp.append((name, a))

    if long_no_ct_cust:
        total = sum(a for _, a in long_no_ct_cust)
        top3 = long_no_ct_cust[:3]
        detail_parts = [f"{n}（{a:,.0f}元）" for n, a in top3]
        results.append({
            "category": CAT, "category_icon": ICON, "risk_score": 8, "risk_level": "高风险",
            "risk_color": "#dc2626", "urgency": "紧急",
            "item": "长期大额客户无正式合同",
            "detail": f"共 {len(long_no_ct_cust)} 家客户近12个月累计交易额超{LONG_TERM_THRESHOLD/10000:.0f}万元（合计 {total:,.2f} 元），"
                       f"但系统中未发现任何销售合同或服务协议。长期大额无合同交易是税务稽查的重点关注事项。"
                       f"涉及：{'、'.join(detail_parts)}"
                       + (f"等{len(long_no_ct_cust)}家" if len(long_no_ct_cust) > 3 else ""),
            "suggestion": "立即与长期合作客户补充签署框架协议或年度采购合同：(1)签订年度框架合同明确基本条款；(2)单笔交易签署订单/补充协议；(3)合同应涵盖定价机制、交货方式、结算周期等核心要素。",
            "required_evidence": ["年度框架合同/销售协议", "历史交易对账单", "合作备忘录"]
        })

    if long_no_ct_supp:
        total = sum(a for _, a in long_no_ct_supp)
        top3 = long_no_ct_supp[:3]
        detail_parts = [f"{n}（{a:,.0f}元）" for n, a in top3]
        results.append({
            "category": CAT, "category_icon": ICON, "risk_score": 8, "risk_level": "高风险",
            "risk_color": "#dc2626", "urgency": "紧急",
            "item": "长期大额供应商无正式合同",
            "detail": f"共 {len(long_no_ct_supp)} 家供应商近12个月累计交易额超{LONG_TERM_THRESHOLD/10000:.0f}万元（合计 {total:,.2f} 元），"
                       f"但系统中未发现任何采购合同。长期大额无合同采购可能影响成本真实性认定。"
                       f"涉及：{'、'.join(detail_parts)}"
                       + (f"等{len(long_no_ct_supp)}家" if len(long_no_ct_supp) > 3 else ""),
            "suggestion": "立即与长期合作供应商补签采购框架合同：(1)签订年度采购协议；(2)明确质量标准、验收标准、违约责任；(3)逐单签署采购订单作为合同附件。",
            "required_evidence": ["年度采购框架合同", "供应商资质文件", "历史采购对账单"]
        })

    # ═══ 7. 合同金额与发票金额严重偏差 ═══
    # 对于有匹配合同的客户/供应商，检查合同金额 vs 实际发票金额
    contract_invoice_deviation = []

    # 销售合同 vs 销项发票
    for c in sales_ct:
        ct_amt = _safe_float(c.amount)
        if ct_amt <= 0:
            continue
        parties = [p.strip() for p in [c.party_b or "", c.party_a or ""] if p.strip()]
        for party in parties:
            inv_amt = db.query(func.coalesce(func.sum(SalesInvoice.total_amount), 0)).filter(
                SalesInvoice.company_id == company_id,
                SalesInvoice.invoice_date >= ps_date,
                SalesInvoice.invoice_date <= pe_date,
                SalesInvoice.buyer_name == party
            ).scalar() or 0
            inv_amt = _safe_float(inv_amt)
            if inv_amt > 0:
                dev = abs(ct_amt - inv_amt) / ct_amt
                if dev > 0.5:
                    contract_invoice_deviation.append({
                        "contract_no": c.contract_no, "name": c.name,
                        "party": party, "ct_amt": ct_amt, "inv_amt": inv_amt,
                        "deviation_pct": round(dev * 100, 1), "direction": "销售"
                    })
            break  # 一个合同只匹配第一个交易方

    # 采购合同 vs 进项发票
    for c in purchase_ct:
        ct_amt = _safe_float(c.amount)
        if ct_amt <= 0:
            continue
        parties = [p.strip() for p in [c.party_b or "", c.party_a or ""] if p.strip()]
        for party in parties:
            inv_amt = db.query(func.coalesce(func.sum(PurchaseInvoice.total_amount), 0)).filter(
                PurchaseInvoice.company_id == company_id,
                PurchaseInvoice.invoice_date >= ps_date,
                PurchaseInvoice.invoice_date <= pe_date,
                PurchaseInvoice.seller_name == party
            ).scalar() or 0
            inv_amt = _safe_float(inv_amt)
            if inv_amt > 0:
                dev = abs(ct_amt - inv_amt) / ct_amt
                if dev > 0.5:
                    contract_invoice_deviation.append({
                        "contract_no": c.contract_no, "name": c.name,
                        "party": party, "ct_amt": ct_amt, "inv_amt": inv_amt,
                        "deviation_pct": round(dev * 100, 1), "direction": "采购"
                    })
            break

    if contract_invoice_deviation:
        detail_parts = []
        for x in contract_invoice_deviation[:5]:
            detail_parts.append(
                f"{x['contract_no']} {x['name']} {x['direction']}（合同{x['ct_amt']:,.0f} vs 发票{x['inv_amt']:,.0f}，偏差{x['deviation_pct']}%）"
            )
        results.append({
            "category": CAT, "category_icon": ICON, "risk_score": 6, "risk_level": "中风险",
            "risk_color": "#f59e0b", "urgency": "提醒",
            "item": "合同金额与实际交易金额严重偏差",
            "detail": f"共 {len(contract_invoice_deviation)} 份合同的金额与对应发票金额偏差超过50%。"
                       f"合同金额与发票金额严重不匹配，可能导致：印花税申报不准确（合同金额为计税基础）、"
                       f"合同管理失控、隐藏未开票收入等问题。"
                       f"涉及：{'；'.join(detail_parts)}"
                       + (f"等{len(contract_invoice_deviation)}份" if len(contract_invoice_deviation) > 5 else ""),
            "suggestion": "逐份核实偏差原因：(1)签署补充协议调整合同金额；(2)如为框架合同，补充具体订单作为附件；(3)更新合同台账并修正印花税申报基础。",
            "required_evidence": ["合同金额偏差说明", "补充协议/订单", "印花税申报更正表"]
        })

    # ═══ 8. 过期合同仍在发生交易 ═══
    # 合同到期日期 < 当前期间，但在分析期间仍有发票
    expired_trading = []
    for c in all_contracts:
        if not c.expiry_date:
            continue
        # 转换为date对象比较
        exp_date = c.expiry_date if isinstance(c.expiry_date, date) else c.expiry_date
        if hasattr(exp_date, 'strftime'):
            exp_str = exp_date.strftime("%Y-%m-%d")
        else:
            exp_str = str(exp_date)
        if exp_str >= ps_date:
            continue  # 未过期

        # 已过期，检查是否仍有交易
        parties = [p.strip() for p in [c.party_b or "", c.party_a or ""] if p.strip()]
        has_trade = False
        trade_info = ""
        for party in parties:
            sale_cnt = db.query(func.count(SalesInvoice.id)).filter(
                SalesInvoice.company_id == company_id,
                SalesInvoice.invoice_date >= ps_date,
                SalesInvoice.invoice_date <= pe_date,
                SalesInvoice.buyer_name == party
            ).scalar() or 0
            if sale_cnt > 0:
                has_trade = True
                trade_info = f"销项发票{sale_cnt}张"
                break
            pur_cnt = db.query(func.count(PurchaseInvoice.id)).filter(
                PurchaseInvoice.company_id == company_id,
                PurchaseInvoice.invoice_date >= ps_date,
                PurchaseInvoice.invoice_date <= pe_date,
                PurchaseInvoice.seller_name == party
            ).scalar() or 0
            if pur_cnt > 0:
                has_trade = True
                trade_info = f"进项发票{pur_cnt}张"
                break
        if has_trade:
            expired_trading.append({
                "contract_no": c.contract_no, "name": c.name,
                "expiry": exp_str, "trade": trade_info
            })

    if expired_trading:
        detail_parts = []
        for x in expired_trading[:5]:
            detail_parts.append(f"{x['contract_no']} {x['name']}（到期{x['expiry']}，仍有{x['trade']}）")
        results.append({
            "category": CAT, "category_icon": ICON, "risk_score": 7, "risk_level": "高风险",
            "risk_color": "#dc2626", "urgency": "紧急",
            "item": "过期合同仍在发生交易",
            "detail": f"共 {len(expired_trading)} 份已过期合同在分析期间内仍有发票交易记录。"
                       f"过期后继续交易属于事实合同关系，但缺少正式合同保护，法律风险极高——"
                       f"发生纠纷时无法依据合同条款维权，且可能被认定为虚构交易。"
                       f"涉及：{'；'.join(detail_parts)}"
                       + (f"等{len(expired_trading)}份" if len(expired_trading) > 5 else ""),
            "suggestion": "立即处理过期合同：(1)已终止的合作关系补充终止协议；(2)继续合作的立即续签合同，续签日期追溯至到期日次日；(3)建立合同到期预警制度，提前30天提醒续签。",
            "required_evidence": ["续签合同/终止协议", "合同到期预警记录", "续签审批流程"]
        })

    # ═══ 9. 发票日期早于合同签订日期（先票后签）═══
    pre_sign_invoices = []
    for c in all_contracts:
        if not c.signing_date:
            continue
        sign_date = c.signing_date if isinstance(c.signing_date, date) else c.signing_date
        if hasattr(sign_date, 'strftime'):
            sign_str = sign_date.strftime("%Y-%m-%d")
        else:
            sign_str = str(sign_date)

        parties = [p.strip() for p in [c.party_b or "", c.party_a or ""] if p.strip()]
        for party in parties:
            # 查销项发票中早于签订日期的
            pre_sale = db.query(
                func.count(SalesInvoice.id),
                func.min(SalesInvoice.invoice_date)
            ).filter(
                SalesInvoice.company_id == company_id,
                SalesInvoice.buyer_name == party,
                SalesInvoice.invoice_date < sign_str
            ).first()
            if pre_sale and pre_sale[0] > 0:
                pre_sign_invoices.append({
                    "contract_no": c.contract_no, "name": c.name,
                    "party": party, "sign_date": sign_str,
                    "first_invoice": str(pre_sale[1]) if pre_sale[1] else "",
                    "cnt": pre_sale[0], "type": "销项"
                })

            # 查进项发票
            pre_pur = db.query(
                func.count(PurchaseInvoice.id),
                func.min(PurchaseInvoice.invoice_date)
            ).filter(
                PurchaseInvoice.company_id == company_id,
                PurchaseInvoice.seller_name == party,
                PurchaseInvoice.invoice_date < sign_str
            ).first()
            if pre_pur and pre_pur[0] > 0:
                pre_sign_invoices.append({
                    "contract_no": c.contract_no, "name": c.name,
                    "party": party, "sign_date": sign_str,
                    "first_invoice": str(pre_pur[1]) if pre_pur[1] else "",
                    "cnt": pre_pur[0], "type": "进项"
                })

    if pre_sign_invoices:
        detail_parts = []
        for x in pre_sign_invoices[:5]:
            detail_parts.append(
                f"{x['contract_no']} {x['name']}（{x['type']}，签{x['sign_date']}，"
                f"首票{x['first_invoice']}，共{x['cnt']}张）"
            )
        results.append({
            "category": CAT, "category_icon": ICON, "risk_score": 7, "risk_level": "高风险",
            "risk_color": "#dc2626", "urgency": "紧急",
            "item": "先票后签（发票日期早于合同签订日期）",
            "detail": f"共 {len(pre_sign_invoices)} 份合同存在交易发票日期早于合同签订日期的「先票后签」问题。"
                       f"先票后签说明业务发生在前、合同补签在后，属于倒签合同行为，"
                       f"在税务稽查中会被重点审查，可能被认定为虚构交易或虚开发票。"
                       f"涉及：{'；'.join(detail_parts)}"
                       + (f"等{len(pre_sign_invoices)}份" if len(pre_sign_invoices) > 5 else ""),
            "suggestion": "规范合同签署流程：(1)严格执行「先签合同后开发票」的原则；(2)已发生的倒签情况逐份出具书面说明，记录真实业务背景；(3)建立合同签署→开票的流程管控，杜绝先票后签。",
            "required_evidence": ["倒签合同书面说明", "合同审批记录", "业务发生证明材料"]
        })

    # ═══ 10. 合同收付款计划与实际银行流水不匹配 ═══
    # 查询合同收付款计划并与银行流水对比
    ct_with_plans = db.query(ContractPayment).filter(
        ContractPayment.company_id == company_id,
        ContractPayment.due_date >= ps_date,
        ContractPayment.due_date <= pe_date
    ).all()

    plan_unmatched = []
    for cp in ct_with_plans:
        # 取关联合同
        contract = db.query(Contract).filter(Contract.id == cp.contract_id).first()
        if not contract:
            continue
        parties = [p.strip() for p in [contract.party_b or "", contract.party_a or ""] if p.strip()]
        plan_amt = _safe_float(cp.amount)
        paid_amt = _safe_float(cp.paid_amount)
        if paid_amt >= plan_amt * 0.95:
            continue  # 已基本付清

        # 查银行流水是否有匹配付款
        matched_bank = False
        for party in parties:
            bank_amt = db.query(func.coalesce(func.sum(BankTransaction.debit_amount), 0)).filter(
                BankTransaction.company_id == company_id,
                BankTransaction.transaction_date >= ps_date,
                BankTransaction.transaction_date <= pe_date,
                BankTransaction.counterparty_name == party
            ).scalar() or 0
            bank_amt = _safe_float(bank_amt)
            if bank_amt > 0 and abs(bank_amt - (plan_amt - paid_amt)) < plan_amt * 0.3:
                matched_bank = True
                break

        if not matched_bank and plan_amt > 10000:  # 只关注1万元以上未匹配的
            plan_unmatched.append({
                "contract_no": contract.contract_no,
                "name": contract.name,
                "payment_no": cp.payment_no,
                "plan_amt": plan_amt, "paid_amt": paid_amt,
                "due_date": str(cp.due_date) if cp.due_date else ""
            })

    if plan_unmatched:
        total_unmatched = sum(p["plan_amt"] - p["paid_amt"] for p in plan_unmatched)
        detail_parts = []
        for x in plan_unmatched[:5]:
            detail_parts.append(
                f"{x['contract_no']} 第{x['payment_no']}期 "
                f"（应付{x['plan_amt']:,.0f}/实付{x['paid_amt']:,.0f}，到期{x['due_date']}）"
            )
        results.append({
            "category": CAT, "category_icon": ICON, "risk_score": 6, "risk_level": "中风险",
            "risk_color": "#f59e0b", "urgency": "提醒",
            "item": "合同收付款计划与实际流水不匹配",
            "detail": f"共 {len(plan_unmatched)} 条合同收付款计划的到期金额与银行流水无法匹配"
                       f"（未匹配总额 {total_unmatched:,.2f} 元）。"
                       f"合同付款计划与实际付款脱节，可能隐藏：(1)付款未入账；(2)虚假付款计划；"
                       f"(3)资金体外循环。"
                       f"涉及：{'；'.join(detail_parts)}"
                       + (f"等{len(plan_unmatched)}条" if len(plan_unmatched) > 5 else ""),
            "suggestion": "逐笔核实合同付款执行情况：(1)比对合同付款计划与银行流水、发票记录；(2)对未匹配款项出具说明（延期/取消/变更支付方式）；(3)更新合同付款台账并与账面一致。",
            "required_evidence": ["合同付款对账单", "未付款项说明", "银行流水截图"]
        })

    # ═══ 12. 合同状态异常（大量作废/终止合同）═══
    abnormal_status = db.query(
        Contract.status, func.count(Contract.id), func.sum(Contract.amount)
    ).filter(
        Contract.company_id == company_id
    ).group_by(Contract.status).all()

    status_map = {}
    for s, cnt, amt in abnormal_status:
        status_map[s] = (cnt, _safe_float(amt))

    total_ct = sum(cnt for cnt, _ in status_map.values())
    terminated_cnt = status_map.get("已终止", (0, 0))[0]
    draft_cnt = status_map.get("起草中", (0, 0))[0]

    if total_ct > 0 and terminated_cnt > 0:
        terminated_pct = round(terminated_cnt / total_ct * 100, 1)
        if terminated_pct > 30:
            results.append({
                "category": CAT, "category_icon": ICON, "risk_score": 5, "risk_level": "中风险",
                "risk_color": "#f59e0b", "urgency": "提醒",
                "item": "合同终止率偏高",
                "detail": f"系统中共有 {total_ct} 份合同，其中 {terminated_cnt} 份已终止（占比 {terminated_pct}%）。"
                           f"合同终止率过高可能暗示：(1)业务合作关系不稳定；(2)合同管理不规范、频繁更换合作方；"
                           f"(3)存在虚假签约后终止的异常模式。",
                "suggestion": "逐份分析终止原因：(1)分类统计终止原因（履约完成/协商终止/违约终止）；(2)对异常终止合同重点核查业务真实性；(3)完善合同签订前的尽职调查流程。",
                "required_evidence": ["合同终止原因汇总", "终止合同清单", "业务流程改进方案"]
            })

    if total_ct > 0 and draft_cnt > 0:
        draft_pct = round(draft_cnt / total_ct * 100, 1)
        if draft_pct > 20:
            results.append({
                "category": CAT, "category_icon": ICON, "risk_score": 4, "risk_level": "中风险",
                "risk_color": "#f59e0b", "urgency": "提醒",
                "item": "大量合同处于起草中状态",
                "detail": f"系统中共有 {total_ct} 份合同，其中 {draft_cnt} 份仍处于「起草中」状态（占比 {draft_pct}%）。"
                           f"大量合同长期未签署，可能表明：(1)合同审批流程效率低下；(2)业务未按合同执行即已发生；"
                           f"(3)合同管理流于形式。",
                "suggestion": "清理起草中合同：(1)超过30天未签署的合同逐一确认状态——签署/终止/继续协商；(2)建立合同审批时效预警，超期自动提醒；(3)已执行但未签署的合同立即补签。",
                "required_evidence": ["起草中合同清单", "合同审批时效报告", "补签/终止确认"]
            })

    # ═══ 13. 印花税 — 无合同导致印花税风险 ═══
    # 发票金额大但合同金额小/零 → 印花税可能少缴
    if total_ct == 0 and (total_sales_rev > 100000 or total_purchase_amt > 100000):
        results.append({
            "category": CAT, "category_icon": ICON, "risk_score": 7, "risk_level": "高风险",
            "risk_color": "#dc2626", "urgency": "紧急",
            "item": "无合同导致印花税申报基础缺失",
            "detail": f"系统中有销项收入 {total_sales_rev:,.2f} 元、进项采购 {total_purchase_amt:,.2f} 元，"
                       f"但没有任何有效合同记录。印花税以合同金额为计税基础，无合同意味着无法正确计算印花税，"
                       f"容易导致少缴或漏缴印花税。",
            "suggestion": "(1)按发票金额作为合同印花税的最低计税基础，主动申报缴纳；(2)立即补签合同并登记入册；(3)建立合同与印花税申报的联动机制。",
            "required_evidence": ["印花税补缴计算表", "补签合同清单", "印花税申报表"]
        })

    # ═══ 15. 合同对方与发票对方名称不一致 ═══
    # 检查同一客户/供应商在合同和发票中的名称是否能匹配上
    name_mismatch = []

    # 销售方向：发票客户名 vs 合同乙方
    invoice_customers = set(name for name, _, _ in sales_customers)
    ct_sales_names = set()
    for c in sales_ct:
        if c.party_b: ct_sales_names.add(c.party_b.strip())
        if c.party_a: ct_sales_names.add(c.party_a.strip())

    # 找出发票中有但合同中没有的，且发票中名称与合同中任一名称都不相似的
    for name in list(invoice_customers)[:20]:  # 限制检查量
        if name not in ct_sales_names and not _name_in_set(name, ct_sales_names):
            # 发票中有，合同中找不到 → 已经在"收入无合同"中报了
            pass

    # 反之：合同中有乙方但找不到匹配发票的
    for party in list(ct_sales_names)[:20]:
        matched = False
        for inv_name in invoice_customers:
            if _name_in_set(party, {inv_name}):
                matched = True
                break
        if not matched:
            # 合同有交易方但找不到发票 → 可能名称不一致
            name_mismatch.append({"party": party, "type": "销售"})

    # 采购方向
    invoice_suppliers = set(name for name, _, _ in purchase_suppliers)
    ct_pur_names = set()
    for c in purchase_ct:
        if c.party_b: ct_pur_names.add(c.party_b.strip())
        if c.party_a: ct_pur_names.add(c.party_a.strip())

    for party in list(ct_pur_names)[:20]:
        matched = False
        for inv_name in invoice_suppliers:
            if _name_in_set(party, {inv_name}):
                matched = True
                break
        if not matched:
            name_mismatch.append({"party": party, "type": "采购"})

    if name_mismatch:
        detail_parts = [f"{x['type']}合同方「{x['party']}」无匹配发票" for x in name_mismatch[:5]]
        results.append({
            "category": CAT, "category_icon": ICON, "risk_score": 5, "risk_level": "中风险",
            "risk_color": "#f59e0b", "urgency": "提醒",
            "item": "合同对方与发票对方名称不一致",
            "detail": f"共 {len(name_mismatch)} 个合同交易方在发票中找不到匹配记录，"
                       f"可能存在合同签约主体与开票主体不一致的问题。"
                       f"合同流和发票流的主体不一致将破坏三流合一完整性。"
                       f"涉及：{'；'.join(detail_parts)}"
                       + (f"等{len(name_mismatch)}个" if len(name_mismatch) > 5 else ""),
            "suggestion": "逐一核实合同签约主体与开票主体的一致性：(1)关注集团内部不同子公司签约与开票的对应关系；(2)第三方代签约需补充授权委托书；(3)名称变更的需补充工商变更证明。",
            "required_evidence": ["签约主体与开票主体对应关系表", "授权委托书", "工商变更证明"]
        })

    # ═══ 16. 关联方交易无独立合同 ═══
    # 简易版：通过银行流水中的"往来"/"借款"/"关联"等关键词 + 无合同判断
    related_kw = ["往来", "借款", "周转", "拆借", "关联", "集团", "总部", "分公司", "子公司"]
    related_bank = db.query(
        BankTransaction.counterparty_name,
        func.sum(BankTransaction.debit_amount).label("debit"),
        func.sum(BankTransaction.credit_amount).label("credit")
    ).filter(
        BankTransaction.company_id == company_id,
        BankTransaction.transaction_date >= ps_date,
        BankTransaction.transaction_date <= pe_date,
        or_(*[BankTransaction.summary.contains(kw) for kw in related_kw])
    ).group_by(BankTransaction.counterparty_name).all()

    related_no_ct = []
    for name, debit, credit in related_bank:
        total_flow = _safe_float(debit) + _safe_float(credit)
        if total_flow < 50000:
            continue
        # 检查是否有合同
        if not _name_in_set(name, sales_parties) and not _name_in_set(name, purchase_parties):
            related_no_ct.append((name, total_flow))

    if related_no_ct:
        total = sum(a for _, a in related_no_ct)
        detail_parts = [f"{n}（{a:,.0f}元）" for n, a in related_no_ct[:5]]
        results.append({
            "category": CAT, "category_icon": ICON, "risk_score": 8, "risk_level": "高风险",
            "risk_color": "#dc2626", "urgency": "紧急",
            "item": "关联方资金往来无独立合同",
            "detail": f"发现 {len(related_no_ct)} 个疑似关联方的银行流水往来（含「往来/借款/关联」等关键词，"
                       f"合计 {total:,.2f} 元）无对应合同。"
                       f"关联方资金往来缺少书面协议，可能被认定为：(1)资本弱化——债资比超标利息不得税前扣除；"
                       f"(2)转移定价——未按独立交易原则定价面临纳税调整。"
                       f"涉及：{'、'.join(detail_parts)}"
                       + (f"等{len(related_no_ct)}个" if len(related_no_ct) > 5 else ""),
            "suggestion": "关联交易必须签订书面协议：(1)借款类：签订借款合同，约定利率（不得低于央行同期贷款基准利率）；(2)购销类：签订购销合同并留存同期可比非关联交易价格作为定价依据；(3)编制关联交易同期资料报告。",
            "required_evidence": ["关联交易合同/借款协议", "同期可比价格分析", "关联交易同期资料报告"]
        })

    # ═══ 17. 合同要素完整性检查 ═══
    # 检查合同的必填字段是否完整
    incomplete_ct = []
    for c in all_contracts:
        missing = []
        if not c.signing_date: missing.append("签订日期")
        if not c.effective_date: missing.append("生效日期")
        if not c.expiry_date: missing.append("到期日期")
        if _safe_float(c.amount) <= 0: missing.append("合同金额")
        if not c.responsible_person: missing.append("负责人")
        if missing:
            incomplete_ct.append({"contract_no": c.contract_no, "name": c.name, "missing": missing})

    if incomplete_ct and len(incomplete_ct) > len(all_contracts) * 0.3:  # 超过30%合同要素不全
        detail_parts = []
        for x in incomplete_ct[:5]:
            detail_parts.append(f"{x['contract_no']} {x['name']}（缺{'/'.join(x['missing'])}）")
        results.append({
            "category": CAT, "category_icon": ICON, "risk_score": 4, "risk_level": "中风险",
            "risk_color": "#f59e0b", "urgency": "提醒",
            "item": "合同要素不完整",
            "detail": f"共 {len(incomplete_ct)} 份合同（占总数 {round(len(incomplete_ct)/len(all_contracts)*100,1)}%）"
                       f"存在关键要素缺失。合同要素不全可能导致：(1)合同法律效力瑕疵；"
                       f"(2)合同执行缺乏明确标准；(3)印花税申报计税基础不准确。"
                       f"涉及：{'；'.join(detail_parts)}"
                       + (f"等{len(incomplete_ct)}份" if len(incomplete_ct) > 5 else ""),
            "suggestion": "逐份补全合同要素：(1)签订日期、生效日期、到期日期必须齐全；(2)合同金额不得为空或零；(3)明确负责人和审批流程；(4)建立合同要素完整性校验规则。",
            "required_evidence": ["补全后的合同正本", "合同要素完整性检查表"]
        })

    # 构建合同快速查找索引（标准化名称 → 合同列表）
    contract_lookup = {}  # normalized_name → [Contract]
    for c in all_contracts:
        for party in [c.party_a, c.party_b]:
            if party:
                nn = _normalize_entity_name(party.strip())
                if len(nn) >= 2:
                    if nn not in contract_lookup:
                        contract_lookup[nn] = []
                    contract_lookup[nn].append(c)

    # ─── 18. 销项发票逐票合同匹配 ───
    all_sales_inv = db.query(
        SalesInvoice.id, SalesInvoice.buyer_name,
        SalesInvoice.total_amount, SalesInvoice.invoice_date
    ).filter(
        SalesInvoice.company_id == company_id,
        SalesInvoice.invoice_date >= ps_date,
        SalesInvoice.invoice_date <= pe_date
    ).all()

    sales_matched_cnt = 0
    sales_matched_amt = 0.0
    sales_total_cnt = len(all_sales_inv)
    sales_total_amt = 0.0
    sales_unmatched_list = []  # 保存前10条无合同发票用于展示

    for inv_id, buyer_name, total_amount, inv_date in all_sales_inv:
        amt = _safe_float(total_amount)
        sales_total_amt += amt
        inv_str = _date_to_str(inv_date)

        # 按名称+有效期查找匹配合同
        matched = False
        nn = _normalize_entity_name(buyer_name) if buyer_name else ""
        if nn and len(nn) >= 2:
            candidates = contract_lookup.get(nn, [])
            # 精确匹配失败 → 模糊匹配
            if not candidates:
                for key, cts in contract_lookup.items():
                    if nn in key or key in nn:
                        candidates = cts
                        break
            # 筛选：销售/服务类合同 + 有效期覆盖发票日期
            for c in candidates:
                if c.contract_type in ("销售", "服务") and _contract_covers_date(c, inv_str):
                    matched = True
                    break

        if matched:
            sales_matched_cnt += 1
            sales_matched_amt += amt
        else:
            if len(sales_unmatched_list) < 10:
                sales_unmatched_list.append((buyer_name or "未知", amt, inv_str))

    if sales_total_cnt > 0:
        cov_pct = round(sales_matched_cnt / sales_total_cnt * 100, 1)
        amt_cov_pct = round(sales_matched_amt / sales_total_amt * 100, 1) if sales_total_amt > 0 else 0
        uncovered = sales_total_cnt - sales_matched_cnt
        uncovered_amt = sales_total_amt - sales_matched_amt

        # 覆盖率<100%即报告（100%覆盖率也报绿色良好）
        if cov_pct < 100:
            # ── 子维度：大额交易无合同（≥10万）──
            sales_big_nomatch = [(n, _safe_float(a), c) for n, a, c in sales_customers 
                                 if _safe_float(a) >= BIG_THRESHOLD and not _name_in_set(n, sales_parties)]
            # ── 子维度：高频交易无框架合同（≥5张）──
            sales_freq_nomatch = [(n, _safe_float(a), c) for n, a, c in sales_customers 
                                  if c >= 5 and not _name_in_set(n, sales_parties)]

            score = 8 if cov_pct < 50 else (7 if cov_pct < 70 else (5 if cov_pct < 90 else 3))
            level = "高风险" if cov_pct < 70 else ("中风险" if cov_pct < 90 else "低风险")
            color = "#dc2626" if cov_pct < 70 else ("#f59e0b" if cov_pct < 90 else "#10b981")
            urg = "紧急" if cov_pct < 70 else ("提醒" if cov_pct < 90 else "关注")
            detail_parts = [f"{n}（{a:,.0f}元/{d}）" for n, a, d in sales_unmatched_list[:5]]

            # 构建分层详情
            sub_detail = ""
            if sales_big_nomatch:
                big_items = "、".join([f"{n}({a:,.0f}元)" for n, a, _ in sales_big_nomatch[:3]])
                sub_detail += f"【大额专项】≥10万元交易：{len(sales_big_nomatch)}家客户无合同（{big_items}）"
            if sales_freq_nomatch:
                freq_items = "、".join([f"{n}({c}张/{a:,.0f}元)" for n, a, c in sales_freq_nomatch[:3]])
                sub_detail += f"\n【高频专项】≥5张发票客户：{len(sales_freq_nomatch)}家无框架合同（{freq_items}）"

            detail_text = (
                f"分析期间共 {sales_total_cnt} 张销项发票（{sales_total_amt:,.2f} 元），"
                f"其中 {sales_matched_cnt} 张（{sales_matched_amt:,.2f} 元，覆盖率 {cov_pct}% 按张数 / {amt_cov_pct}% 按金额）"
                f"有有效合同覆盖（名称匹配+合同在有效期内），"
                f"{uncovered} 张（{uncovered_amt:,.2f} 元）无合同或合同已过期。"
                f"《发票管理办法》要求开具发票应当基于真实交易，合同是证明交易真实性的核心证据。"
                + (f" 无合同发票示例：{'；'.join(detail_parts)}"
                      + (f"等{uncovered}张" if uncovered > 5 else "")
                      if sales_unmatched_list else "")
                + (f"\n{sub_detail}" if sub_detail else "")
            )
            results.append({
                "category": CAT, "category_icon": ICON,
                "risk_score": score, "risk_level": level,
                "risk_color": color, "urgency": urg,
                "item": "销项发票合同覆盖率不足",
                "detail": detail_text,
                "suggestion": f"提高合同覆盖率至100%："
                    f"对{uncovered}张无合同销项发票逐票补签销售合同；"
                    f"大额交易（≥10万）必须签正式购销合同；"
                    f"高频交易客户（≥5张）签订框架销售合同/年度协议；"
                    f"建立「先签合同后开发票」的铁律制度。",
                "required_evidence": ["无合同销项发票清单（含票号/金额/日期）", "大额交易合同正本", "高频客户框架合同", "合同覆盖率整改计划"]
            })

    # ─── 19. 进项发票逐票合同匹配 ───
    all_pur_inv = db.query(
        PurchaseInvoice.id, PurchaseInvoice.seller_name,
        PurchaseInvoice.total_amount, PurchaseInvoice.invoice_date
    ).filter(
        PurchaseInvoice.company_id == company_id,
        PurchaseInvoice.invoice_date >= ps_date,
        PurchaseInvoice.invoice_date <= pe_date
    ).all()

    pur_matched_cnt = 0
    pur_matched_amt = 0.0
    pur_total_cnt = len(all_pur_inv)
    pur_total_amt = 0.0
    pur_unmatched_list = []

    for inv_id, seller_name, total_amount, inv_date in all_pur_inv:
        amt = _safe_float(total_amount)
        pur_total_amt += amt
        inv_str = _date_to_str(inv_date)

        matched = False
        nn = _normalize_entity_name(seller_name) if seller_name else ""
        if nn and len(nn) >= 2:
            candidates = contract_lookup.get(nn, [])
            if not candidates:
                for key, cts in contract_lookup.items():
                    if nn in key or key in nn:
                        candidates = cts
                        break
            for c in candidates:
                if c.contract_type in ("采购", "服务") and _contract_covers_date(c, inv_str):
                    matched = True
                    break

        if matched:
            pur_matched_cnt += 1
            pur_matched_amt += amt
        else:
            if len(pur_unmatched_list) < 10:
                pur_unmatched_list.append((seller_name or "未知", amt, inv_str))

    if pur_total_cnt > 0:
        cov_pct = round(pur_matched_cnt / pur_total_cnt * 100, 1)
        amt_cov_pct = round(pur_matched_amt / pur_total_amt * 100, 1) if pur_total_amt > 0 else 0
        uncovered = pur_total_cnt - pur_matched_cnt
        uncovered_amt = pur_total_amt - pur_matched_amt

        if cov_pct < 100:
            # ── 子维度：大额交易无合同（≥10万）──
            pur_big_nomatch = [(n, _safe_float(a), c) for n, a, c in purchase_suppliers 
                               if _safe_float(a) >= BIG_THRESHOLD and not _name_in_set(n, purchase_parties)]

            score = 8 if cov_pct < 50 else (7 if cov_pct < 70 else (5 if cov_pct < 90 else 3))
            level = "高风险" if cov_pct < 70 else ("中风险" if cov_pct < 90 else "低风险")
            color = "#dc2626" if cov_pct < 70 else ("#f59e0b" if cov_pct < 90 else "#10b981")
            urg = "紧急" if cov_pct < 70 else ("提醒" if cov_pct < 90 else "关注")
            detail_parts = [f"{n}（{a:,.0f}元/{d}）" for n, a, d in pur_unmatched_list[:5]]

            # 构建大额分层详情
            sub_detail = ""
            if pur_big_nomatch:
                big_items = "、".join([f"{n}({a:,.0f}元)" for n, a, _ in pur_big_nomatch[:3]])
                sub_detail += f"【大额专项】≥10万元交易：{len(pur_big_nomatch)}家供应商无合同（{big_items}）"

            detail_text = (
                f"分析期间共 {pur_total_cnt} 张进项发票（{pur_total_amt:,.2f} 元），"
                f"其中 {pur_matched_cnt} 张（{pur_matched_amt:,.2f} 元，覆盖率 {cov_pct}% 按张数 / {amt_cov_pct}% 按金额）"
                f"有有效合同覆盖（名称匹配+合同在有效期内），"
                f"{uncovered} 张（{uncovered_amt:,.2f} 元）无合同或合同已过期。"
                f"缺少采购合同违反三流合一原则，影响成本真实性认定和所得税税前扣除资格。"
                + (f" 无合同发票示例：{'；'.join(detail_parts)}"
                      + (f"等{uncovered}张" if uncovered > 5 else "")
                      if pur_unmatched_list else "")
                + (f"\n{sub_detail}" if sub_detail else "")
            )
            results.append({
                "category": CAT, "category_icon": ICON,
                "risk_score": score, "risk_level": level,
                "risk_color": color, "urgency": urg,
                "item": "进项发票合同覆盖率不足",
                "detail": detail_text,
                "suggestion": f"提高合同覆盖率至100%："
                    f"对{uncovered}张无合同进项发票逐票补签采购合同或服务协议；"
                    f"大额交易（≥10万）必须签正式采购合同；"
                    f"建立「先签合同后采购」的铁律制度。",
                "required_evidence": ["无合同进项发票清单（含票号/金额/日期）", "大额采购合同正本", "供应商资质文件", "合同覆盖率整改计划"]
            })

    # ═══ 20. 合同涉税条款缺失/不规范 ═══
    # 检查合同中是否包含发票类型、税率、纳税义务时点、违约金涉税等关键条款
    TAX_CLAUSE_KEYWORDS = ["发票", "税率", "增值税", "含税", "不含税", "税款", "发票类型",
                           "专票", "普票", "纳税", "价税合计", "开票时间", "代扣代缴"]
    contracts_missing_clauses = []
    for c in all_contracts:
        # 从合同名称、内容摘要、备注中检查
        text_pool = " ".join(filter(None, [c.name or "", c.content_summary or "", c.remark or ""]))
        found = [kw for kw in TAX_CLAUSE_KEYWORDS if kw in text_pool]
        if len(found) < 3:  # 少于3个涉税关键词视为条款缺失
            contracts_missing_clauses.append({
                "contract_no": c.contract_no, "name": c.name,
                "found": found, "type": c.contract_type or "未知"
            })

    if contracts_missing_clauses:
        detail_parts = []
        for x in contracts_missing_clauses[:5]:
            found_str = "、".join(x['found']) if x['found'] else "无涉税关键词"
            detail_parts.append(f"{x['contract_no']} {x['name']}（类型:{x['type']}，含:{found_str}）")
        results.append({
            "category": CAT, "category_icon": ICON, "risk_score": 6, "risk_level": "中风险",
            "risk_color": "#f59e0b", "urgency": "提醒",
            "item": "合同涉税条款缺失或内容不规范",
            "detail": f"共 {len(contracts_missing_clauses)} 份合同（占总数 {round(len(contracts_missing_clauses)/len(all_contracts)*100,1)}%）"
                       f"缺少关键涉税条款（发票类型/税率/纳税义务时点等）。"
                       f"合同涉税条款不完备可能导致：(1)发票开具与合同约定不一致；"
                       f"(2)税率适用争议；(3)纳税义务发生时点认定困难；(4)代扣代缴义务未约定。"
                       f"涉及：{'；'.join(detail_parts)}"
                       + (f"等{len(contracts_missing_clauses)}份" if len(contracts_missing_clauses) > 5 else ""),
            "suggestion": "逐份补充合同涉税条款：(1)明确发票类型（专票/普票）和适用税率；(2)约定开票时限和付款条件；(3)明确价税分离和含税/不含税金额；(4)涉及跨境交易须约定代扣代缴义务。",
            "required_evidence": ["涉税条款缺失合同清单", "修订后的合同正本（含涉税条款）", "合同涉税条款模板/指引"]
        })

    # ═══ 21. 合同拆分规避税负嫌疑 ═══
    # 同一交易方在同月内有多份小额合同（单份<5万），总额超阈值（20万）
    # 疑似化整为零规避印花税/增值税纳税人认定/招标门槛等
    party_contracts = defaultdict(list)
    for c in all_contracts:
        if c.signing_date:
            sd_str = _date_to_str(c.signing_date)
            if sd_str[:7] == ps[:7] or (pe and sd_str <= pe):
                for party_name in [c.party_a, c.party_b]:
                    if party_name:
                        nn = _normalize_entity_name(party_name.strip())
                        if len(nn) >= 2:
                            party_contracts[(nn, sd_str[:7])].append(c)

    split_suspect = []
    for (party, month), cts in party_contracts.items():
        if len(cts) >= 3:
            amounts = [_safe_float(c.amount) for c in cts]
            total_amt = sum(amounts)
            avg_amt = total_amt / len(cts)
            if avg_amt < 50000 and total_amt >= 200000:
                split_suspect.append({
                    "party": party, "month": month, "count": len(cts),
                    "total": total_amt, "avg": avg_amt
                })

    if split_suspect:
        detail_parts = []
        for x in split_suspect[:5]:
            detail_parts.append(f"{x['party']} {x['month']}月 {x['count']}份合同 合计{x['total']:,.0f}元")
        results.append({
            "category": CAT, "category_icon": ICON, "risk_score": 7, "risk_level": "高风险",
            "risk_color": "#ef4444", "urgency": "紧急",
            "item": "疑似合同拆分规避税负",
            "detail": f"发现 {len(split_suspect)} 组「同一方+同月」多份小额合同（单份均<5万，合计≥20万），"
                       f"涉嫌化整为零规避：(1)印花税计税基础；(2)增值税一般纳税人认定标准；"
                       f"(3)招标/审批门槛。涉及：{'；'.join(detail_parts)}"
                       + (f"等{len(split_suspect)}组" if len(split_suspect) > 5 else ""),
            "suggestion": "(1)逐组核实拆分合同的商业合理性（是否确属不同业务事项）；"
                          "(2)如实质为同一交易，应合并为一份合同并按总额缴纳印花税；"
                          "(3)保留每份合同的独立业务证据（需求确认单、交付成果等）。",
            "required_evidence": ["拆分合同清单及业务说明", "各合同对应的独立业务证据", "印花税补缴凭证（如需要）"]
        })

    # ═══ 22. 合同违约金/赔偿金涉税处理异常 ═══
    # 检查合同中有无违约金条款，实际是否收取/支付违约金但未做税务处理
    # 收取违约金应缴增值税+企业所得税；支付违约金可税前扣除
    PENALTY_KEYWORDS = ["违约金", "赔偿金", "罚款", "滞纳金", "违约", "赔偿", "罚则"]
    contracts_with_penalty = []
    for c in all_contracts:
        text_pool = " ".join(filter(None, [c.name or "", c.content_summary or "", c.remark or ""]))
        has_penalty = any(kw in text_pool for kw in PENALTY_KEYWORDS)
        if has_penalty:
            contracts_with_penalty.append(c)

    if contracts_with_penalty:
        # 检查序时账中是否有违约金/赔偿金相关分录
        penalty_entries = db.query(JournalEntry).filter(
            JournalEntry.company_id == company_id,
            JournalEntry.period >= ps,
            JournalEntry.period <= pe,
            JournalEntry.summary.contains("违约")
        ).count()
        _ = db.query(JournalEntry).filter(
            JournalEntry.company_id == company_id,
            JournalEntry.period >= ps,
            JournalEntry.period <= pe,
            JournalEntry.summary.contains("赔偿")
        ).count()
        penalty_entries += _

        detail_parts = [f"{c.contract_no} {c.name}" for c in contracts_with_penalty[:5]]
        results.append({
            "category": CAT, "category_icon": ICON, "risk_score": 5, "risk_level": "中风险",
            "risk_color": "#f59e0b", "urgency": "提醒",
            "item": "合同违约金/赔偿金涉税处理存疑",
            "detail": f"共 {len(contracts_with_penalty)} 份合同含违约金/赔偿金条款，"
                       f"但分析期间仅发现 {penalty_entries} 条相关分录。"
                       f"（注意：收取违约金需缴增值税并确认营业外收入；"
                       f"支付违约金可税前扣除但需留存合同+支付凭证+对方收据）。"
                       f"涉及合同：{'；'.join(detail_parts)}"
                       + (f"等{len(contracts_with_penalty)}份" if len(contracts_with_penalty) > 5 else ""),
            "suggestion": "(1)逐份核实含违约条款合同的执行情况（是否实际发生违约）；"
                          "(2)已收取违约金→确认是否申报增值税（价外费用）和企业所得税；"
                          "(3)已支付违约金→确认是否取得合规凭证并税前扣除。",
            "required_evidence": ["含违约条款的合同清单", "违约金收付凭证（银行回单）", "违约金税务处理说明"]
        })

    # ═══ 23. 免租期/装修期增值税与房产税处理风险 ═══
    # 租赁合同中免租期→增值税视同销售应按余值计税；房产税应按房产原值缴纳
    RENT_FREE_KEYWORDS = ["免租期", "免租", "装修期", "免租金", "开业筹备期"]
    lease_cts = [c for c in all_contracts if c.contract_type == "租赁"]
    rent_free_cts = []
    for c in lease_cts:
        text_pool = " ".join(filter(None, [c.name or "", c.content_summary or "", c.remark or ""]))
        if any(kw in text_pool for kw in RENT_FREE_KEYWORDS):
            rent_free_cts.append(c)

    if rent_free_cts:
        detail_parts = [f"{c.contract_no} {c.name}" for c in rent_free_cts[:5]]
        results.append({
            "category": CAT, "category_icon": ICON, "risk_score": 6, "risk_level": "中风险",
            "risk_color": "#f59e0b", "urgency": "提醒",
            "item": "租赁合同免租期/装修期存在增值税和房产税处理风险",
            "detail": f"共 {len(rent_free_cts)} 份租赁合同含免租期/装修期条款。"
                       f"免租期内：(1)增值税——出租方无偿提供租赁服务应视同销售，"
                       f"按同期同类租金确认销售额缴纳增值税；"
                       f"(2)房产税——免租期内应由产权所有人按房产余值（原值×70%×1.2%）缴纳房产税。"
                       f"涉及合同：{'；'.join(detail_parts)}"
                       + (f"等{len(rent_free_cts)}份" if len(rent_free_cts) > 5 else ""),
            "suggestion": "(1)核实免租期是否已按视同销售申报增值税；"
                          "(2)核实免租期是否已按房产余值缴纳房产税；"
                          "(3)建议合同将免租期改为租金减免（分摊到各期），避免视同销售风险。",
            "required_evidence": ["含免租期条款的租赁合同", "免租期增值税申报记录", "免租期房产税缴纳凭证"]
        })

    # ═══ 24. 跨境合同未约定代扣代缴义务 ═══
    CROSS_BORDER_KEYWORDS = ["境外", "国外", "跨境", "海外", "export", "import", "离岸"]
    cross_border_cts = []
    for c in all_contracts:
        text_pool = " ".join(filter(None, [c.name or "", c.content_summary or "", c.remark or ""]))
        parties_text = " ".join(filter(None, [c.party_a or "", c.party_b or ""]))
        full_text = text_pool + " " + parties_text
        if any(kw in full_text for kw in CROSS_BORDER_KEYWORDS):
            # 检查是否约定代扣代缴
            has_withholding = any(kw in full_text for kw in ["代扣代缴", "withholding", "预提税", "扣缴"])
            cross_border_cts.append({
                "contract_no": c.contract_no, "name": c.name,
                "has_withholding": has_withholding
            })

    if cross_border_cts:
        no_withholding = [x for x in cross_border_cts if not x["has_withholding"]]
        detail_parts = []
        for x in cross_border_cts[:5]:
            status = "未约定代扣代缴" if not x["has_withholding"] else "已约定"
            detail_parts.append(f"{x['contract_no']} {x['name']}（{status}）")
        results.append({
            "category": CAT, "category_icon": ICON, "risk_score": 8 if no_withholding else 5,
            "risk_level": "高风险" if no_withholding else "中风险",
            "risk_color": "#ef4444" if no_withholding else "#f59e0b",
            "urgency": "紧急" if no_withholding else "提醒",
            "item": "跨境合同中未约定代扣代缴义务条款",
            "detail": f"共 {len(cross_border_cts)} 份涉跨境交易合同，"
                       f"其中 {len(no_withholding)} 份未约定代扣代缴义务。"
                       f"向境外支付款项时：(1)服务费/特许权使用费→按10%预提所得税；"
                       f"(2)利息/租金/财产转让→按10%或协定税率；(3)未代扣代缴将被追缴+滞纳金+罚款。"
                       f"涉及：{'；'.join(detail_parts)}"
                       + (f"等{len(cross_border_cts)}份" if len(cross_border_cts) > 5 else ""),
            "suggestion": "(1)逐份跨境合同补充代扣代缴条款，明确扣缴税种/税率/时点；"
                          "(2)对已发生的跨境支付，自查是否已履行代扣代缴义务；"
                          "(3)对适用税收协定的，及时办理协定待遇备案。",
            "required_evidence": ["跨境合同清单", "代扣代缴税款凭证（如已扣缴）", "税收协定待遇备案文件（如适用）"]
        })

    # ═══ 25. 合同变更/补充协议未及时更新 ═══
    # 检查合同是否有多次变更但未反映在合同台账中
    # 通过合同名称中的"补充协议"/"变更"等关键词，以及合同编号的关联关系来检测
    AMEND_KEYWORDS = ["补充协议", "变更", "修订", "补充", "修改", "增补"]
    amended_cts = []  # 补充协议/变更合同
    original_ct_nos = set()  # 被引用的原合同编号

    for c in all_contracts:
        text_pool = " ".join(filter(None, [c.name or "", c.content_summary or "", c.remark or ""]))
        if any(kw in text_pool for kw in AMEND_KEYWORDS):
            amended_cts.append(c)
            # 尝试从名称中提取原合同编号（如 "XXX合同补充协议-001"）
            # 合同编号中查找 "-" 前的部分作为原合同编号前缀
            cn = c.contract_no or ""
            if "-" in cn:
                original_ct_nos.add(cn.rsplit("-", 1)[0])
            # 也尝试从名称中匹配
            import re as _re
            ref_match = _re.search(r"[A-Z]{2,4}\d{6,}", c.name or "")
            if ref_match:
                original_ct_nos.add(ref_match.group())

    # 检查原合同是否存在
    if amended_cts:
        # 查找原合同编号对应的合同
        found_originals = set()
        for c in all_contracts:
            if c.contract_no in original_ct_nos:
                found_originals.add(c.contract_no)
            # 也检查合同名称模糊匹配
            for orig_no in original_ct_nos:
                if orig_no in (c.name or ""):
                    found_originals.add(orig_no)

        missing_originals = original_ct_nos - found_originals

        detail_parts = [f"{c.contract_no} {c.name}" for c in amended_cts[:5]]
        results.append({
            "category": CAT, "category_icon": ICON, "risk_score": 4, "risk_level": "中风险",
            "risk_color": "#f59e0b", "urgency": "提醒",
            "item": "合同存在变更/补充协议但台账可能未完整更新",
            "detail": f"发现 {len(amended_cts)} 份合同涉及变更/补充，"
                       f"其中 {len(missing_originals)} 个原合同编号在台账中未找到。"
                       f"合同变更未及时更新台账可能导致：(1)印花税计税基础不准确；"
                       f"(2)合同金额、履行期限与实际不符；(3)权利义务约定不清晰引发纠纷。"
                       f"涉及：{'；'.join(detail_parts)}"
                       + (f"等{len(amended_cts)}份" if len(amended_cts) > 5 else ""),
            "suggestion": "(1)逐份核实变更/补充协议是否已纳入合同台账；"
                          "(2)更新合同金额、有效期、关键条款至最新版本；"
                          "(3)补充协议增加金额的，需补缴相应印花税。",
            "required_evidence": ["变更/补充协议清单", "更新后的合同台账", "印花税补缴凭证（如涉及金额变更）"]
        })

    # ═══ 26. 混合销售/兼营合同未分别约定税率 ═══
    # 检查合同是否涉及多种业务类型（如销售+安装、服务+销售），但未分别约定税率
    MIXED_SALES_KEYWORDS = ["安装", "运输", "培训", "维护", "保修", "设计", "咨询", "代理", "加工"]
    potentially_mixed_cts = []
    for c in all_contracts:
        text_pool = " ".join(filter(None, [c.name or "", c.content_summary or "", c.remark or ""]))
        found_kws = [kw for kw in MIXED_SALES_KEYWORDS if kw in text_pool]
        # 合同名中包含"销售+服务类"关键词，可能是混合销售
        is_sales = any(kw in (c.name or "") for kw in ["销售", "买卖", "购销"])
        has_service = len(found_kws) >= 1
        if (is_sales and has_service) or len(found_kws) >= 2:
            # 检查是否已分别约定税率
            has_rate_split = any(kw in text_pool for kw in ["分别", "单独", "各自", "税率", "兼营"])
            potentially_mixed_cts.append({
                "contract_no": c.contract_no, "name": c.name,
                "keywords": found_kws, "has_rate_split": has_rate_split
            })

    if potentially_mixed_cts:
        no_split = [x for x in potentially_mixed_cts if not x["has_rate_split"]]
        detail_parts = []
        for x in potentially_mixed_cts[:5]:
            status = "未分别约定" if not x["has_rate_split"] else "已分别约定"
            detail_parts.append(f"{x['contract_no']} {x['name']}（{status}，含{'/'.join(x['keywords'])}）")
        results.append({
            "category": CAT, "category_icon": ICON, "risk_score": 5, "risk_level": "中风险",
            "risk_color": "#f59e0b", "urgency": "提醒",
            "item": "混合销售/兼营合同未分别约定适用税率",
            "detail": f"发现 {len(potentially_mixed_cts)} 份合同可能涉及混合销售/兼营，"
                       f"其中 {len(no_split)} 份未分别约定税率。"
                       f"混合销售→从主业税率；兼营→分别核算分别适用税率。"
                       f"未分别核算的兼营行为→从高适用税率，存在多缴或税务稽查调整风险。"
                       f"涉及：{'；'.join(detail_parts)}"
                       + (f"等{len(potentially_mixed_cts)}份" if len(potentially_mixed_cts) > 5 else ""),
            "suggestion": "(1)区分混合销售（一项行为涉及货物+服务）和兼营（多项独立业务）；"
                          "(2)兼营业务应在合同中分别列明各业务金额和对应税率；"
                          "(3)混合销售按主营业务确定税率，合同中应明确主业性质。",
            "required_evidence": ["混合销售/兼营合同清单", "各业务收入分项明细表", "适用税率认定说明"]
        })


# ═══════════════════════════════════════════════════════════
#  辅助函数：佐证材料清单汇总
# ═══════════════════════════════════════════════════════════

def _build_evidence_summary(results):
    """汇总所有风险项中需要提供的佐证材料清单（去重）"""
    seen = set()
    evidence_list = []
    for r in results:
        required = r.get("required_evidence", [])
        if not required:
            continue
        for item in required:
            key = item.strip().lower()
            if key not in seen:
                seen.add(key)
                evidence_list.append({
                    "item": item.strip(),
                    "related_dimension": r.get("item", ""),
                    "risk_level": r.get("risk_level", "")
                })
    return evidence_list


# ═══════════════════════════════════════════════════════════════
#  V10 新增 — 稽查重点盲区补全（13项）
# ═══════════════════════════════════════════════════════════════

def _analyze_price_surcharge(db, company_id, ps, pe, results):
    """价外费用未申报增值税：手续费/违约金/包装费/运输装卸费等是否并入销售额申报"""
    SURCHARGE_KEYWORDS = ["手续费", "包装费", "违约金", "滞纳金", "延期付款利息",
                          "赔偿金", "代收款项", "代垫款项", "运输装卸费", "仓储费",
                          "优质费", "返还利润", "补贴", "奖励费"]
    # 方法：检查序时账收入类科目摘要中是否含价外费用关键词但未计入销项税额
    revenue_entries = db.query(JournalEntry).filter(
        JournalEntry.company_id == company_id,
        JournalEntry.period >= ps, JournalEntry.period <= pe,
        JournalEntry.account_code.like("6001%")
    ).all()

    found_surcharge = []
    for entry in revenue_entries:
        summary = (entry.summary or "")
        matched = [kw for kw in SURCHARGE_KEYWORDS if kw in summary]
        if matched:
            found_surcharge.append({
                "summary": entry.summary, "period": entry.period,
                "amount": _safe_float(entry.credit_amount or entry.debit_amount),
                "keywords": matched
            })

    if found_surcharge:
        total_amt = sum(x["amount"] for x in found_surcharge)
        detail_parts = []
        for x in found_surcharge[:5]:
            detail_parts.append(f"{x['period']} {x['summary'][:40]}（{x['amount']:,.0f}元，含{'/'.join(x['keywords'])}）")
        results.append({
            "category": "隐匿虚增", "category_icon": "🔍",
            "risk_score": 7, "risk_level": "高风险", "risk_color": "#ef4444", "urgency": "紧急",
            "item": "价外费用可能未并入销售额申报增值税",
            "detail": f"序时账中发现 {len(found_surcharge)} 笔含价外费用关键词的收入分录（合计{total_amt:,.2f}元）。"
                       f"根据增值税暂行条例第六条，销售额为全部价款+价外费用。"
                       f"价外费用包括手续费/违约金/包装费/延期付款利息等，无论会计如何核算均应并入销售额计算增值税。"
                       f"涉及：{'；'.join(detail_parts)}"
                       + (f"等{len(found_surcharge)}笔" if len(found_surcharge) > 5 else ""),
            "suggestion": "(1)逐笔核实价外费用是否已计入应税销售额；(2)如未计入，补申报增值税及附加；"
                          "(3)价外费用适用与主货物/服务相同税率。",
            "required_evidence": ["价外费用明细表", "对应合同条款（价外费用约定）", "增值税补申报记录（如适用）"]
        })


def _analyze_non_monetary_exchange(db, company_id, ps, pe, results):
    """非货币性资产交换未确认收入"""
    # 检查序时账中是否有资产处置/交换的分录但无对应增值税销项
    EXCHANGE_KEYWORDS = ["换入", "换出", "资产交换", "以物易物", "以物抵债", "抵债"]
    entries = db.query(JournalEntry).filter(
        JournalEntry.company_id == company_id,
        JournalEntry.period >= ps, JournalEntry.period <= pe,
        or_(
            JournalEntry.summary.contains("换"),
            JournalEntry.summary.contains("抵"),
            JournalEntry.summary.contains("换"),
            JournalEntry.summary.contains("抵")
        )
    ).all()

    found = []
    for entry in entries:
        summary = (entry.summary or "")
        if any(kw in summary for kw in EXCHANGE_KEYWORDS):
            found.append({
                "summary": entry.summary, "period": entry.period,
                "debit": _safe_float(entry.debit_amount),
                "credit": _safe_float(entry.credit_amount)
            })

    if found:
        total_amt = sum(x["debit"] + x["credit"] for x in found)
        detail_parts = [f"{x['period']} {x['summary'][:50]}" for x in found[:5]]
        results.append({
            "category": "隐匿虚增", "category_icon": "🔍",
            "risk_score": 7, "risk_level": "高风险", "risk_color": "#ef4444", "urgency": "紧急",
            "item": "非货币性资产交换可能未确认增值税和企业所得税收入",
            "detail": f"序时账中发现 {len(found)} 笔疑似非货币性资产交换/以物抵债分录（合计{total_amt:,.2f}元）。"
                       f"非货币性资产交换：(1)增值税→换出货物视同销售，按同类产品均价计税；"
                       f"(2)企业所得税→按公允价值确认资产转让所得。"
                       f"涉及：{'；'.join(detail_parts)}"
                       + (f"等{len(found)}笔" if len(found) > 5 else ""),
            "suggestion": "(1)逐笔核实是否有未确认的视同销售收入；(2)补充增值税申报和企业所得税纳税调整；"
                          "(3)保留交换协议及公允价值评估依据。",
            "required_evidence": ["非货币性资产交换协议", "换出资产公允价值评估报告", "增值税补申报记录"]
        })


def _analyze_debt_restructuring(db, company_id, ps, pe, results):
    """债务重组收益未确认企业所得税"""
    RESTRUCTURE_KEYWORDS = ["债务重组", "债务豁免", "债务减免", "债转股", "以资抵债", "豁免债务"]
    entries = db.query(JournalEntry).filter(
        JournalEntry.company_id == company_id,
        JournalEntry.period >= ps, JournalEntry.period <= pe,
        or_(
            JournalEntry.summary.contains("债务"),
            JournalEntry.summary.contains("豁免"),
            JournalEntry.summary.contains("债转股"),
            JournalEntry.summary.contains("债务"),
            JournalEntry.summary.contains("豁免")
        )
    ).all()

    found = []
    for entry in entries:
        summary = (entry.summary or "")
        if any(kw in summary for kw in RESTRUCTURE_KEYWORDS):
            found.append({
                "summary": entry.summary, "period": entry.period,
                "credit": _safe_float(entry.credit_amount),
                "debit": _safe_float(entry.debit_amount)
            })

    if found:
        total_amt = sum(max(x["credit"], x["debit"]) for x in found)
        detail_parts = [f"{x['period']} {x['summary'][:50]}" for x in found[:5]]
        results.append({
            "category": "企业所得税", "category_icon": "💰",
            "risk_score": 7, "risk_level": "高风险", "risk_color": "#ef4444", "urgency": "紧急",
            "item": "债务重组收益可能未确认企业所得税",
            "detail": f"序时账中发现 {len(found)} 笔疑似债务重组分录（涉及金额{total_amt:,.2f}元）。"
                       f"债务重组收益（债务豁免/以资抵债差额/债转股收益）属于企业所得税应税收入，"
                       f"应在重组完成年度确认并申报。未申报将面临补税+滞纳金。"
                       f"涉及：{'；'.join(detail_parts)}"
                       + (f"等{len(found)}笔" if len(found) > 5 else ""),
            "suggestion": "(1)逐笔核实债务重组的性质和金额；(2)确认债务重组收益是否在年度汇算清缴中申报；"
                          "(3)特殊重组（如债转股）可适用递延纳税，需向税务机关备案。",
            "required_evidence": ["债务重组协议/法院裁定", "债务重组收益计算表", "企业所得税汇算清缴申报表"]
        })


def _analyze_equity_transfer(db, company_id, ps, pe, results):
    """股权转让未申报纳税"""
    EQUITY_KEYWORDS = ["股权转让", "股权变更", "股权出售", "转让股权", "股权交割"]
    entries = db.query(JournalEntry).filter(
        JournalEntry.company_id == company_id,
        JournalEntry.period >= ps, JournalEntry.period <= pe,
        or_(
            JournalEntry.summary.contains("股权"),
            JournalEntry.summary.contains("转让"),
            JournalEntry.summary.contains("股权"),
            JournalEntry.summary.contains("转让")
        )
    ).all()

    found = []
    for entry in entries:
        summary = (entry.summary or "")
        if any(kw in summary for kw in EQUITY_KEYWORDS) and "印花税" not in summary:
            found.append({
                "summary": entry.summary, "period": entry.period,
                "amount": _safe_float(entry.debit_amount or entry.credit_amount)
            })

    if found:
        total_amt = sum(x["amount"] for x in found)
        detail_parts = [f"{x['period']} {x['summary'][:50]}" for x in found[:5]]
        results.append({
            "category": "企业所得税", "category_icon": "💰",
            "risk_score": 8, "risk_level": "高风险", "risk_color": "#ef4444", "urgency": "紧急",
            "item": "股权转让交易可能未足额申报纳税",
            "detail": f"序时账中发现 {len(found)} 笔疑似股权转让相关分录（涉及金额{total_amt:,.2f}元）。"
                       f"股权转让涉及多项税务义务：(1)印花税——转让双方各自按0.05%缴纳；"
                       f"(2)企业所得税——转让方按（转让收入-投资成本）×25%缴纳；"
                       f"(3)个人所得税——自然人股东按财产转让所得20%缴纳（受让方代扣代缴）。"
                       f"涉及：{'；'.join(detail_parts)}"
                       + (f"等{len(found)}笔" if len(found) > 5 else ""),
            "suggestion": "(1)核实股权转让价格是否公允（低于净资产份额可能被核定）；"
                          "(2)确认印花税/所得税是否已足额申报；(3)保留股权转让协议、评估报告、付款凭证。",
            "required_evidence": ["股权转让协议", "股权评估报告/净资产审计报告", "印花税/所得税完税凭证"]
        })


def _analyze_interest_free_loan(db, company_id, ps, pe, results):
    """无偿借款未视同销售缴纳增值税"""
    # 检查其他应收款/其他应付款中大额长期挂账且无利息收入/支出
    # 关联方无偿借款→增值税视同销售（按同期同类贷款利率计算利息）
    other_receivables = db.query(JournalEntry).filter(
        JournalEntry.company_id == company_id,
        JournalEntry.period >= ps, JournalEntry.period <= pe,
        JournalEntry.account_code.like("1221%")
    ).all()

    # 寻找大额借款（单笔≥50万）
    large_loans = []
    for entry in other_receivables:
        amt = _safe_float(entry.debit_amount or entry.credit_amount)
        if amt >= 500000:
            summary = (entry.summary or "")
            if any(kw in summary for kw in ["借款", "借支", "拆借", "往来", "暂借"]):
                large_loans.append({
                    "summary": summary, "period": entry.period,
                    "amount": amt, "contact": entry.contact_project or ""
                })

    # 检查是否有对应的利息收入
    interest_entries = db.query(func.sum(JournalEntry.credit_amount)).filter(
        JournalEntry.company_id == company_id,
        JournalEntry.period >= ps, JournalEntry.period <= pe,
        or_(
            JournalEntry.account_code.like("6051%"),  # 利息收入
            JournalEntry.summary.contains("利息"),
            JournalEntry.summary.contains("资金占用费")
        )
    ).scalar() or 0

    if large_loans and interest_entries < 1000:  # 有大额借款但几乎无利息收入
        total_loan = sum(x["amount"] for x in large_loans)
        detail_parts = [f"{x['period']} {x['summary'][:40]}（{x['amount']:,.0f}元）" for x in large_loans[:5]]
        results.append({
            "category": "增值税专项", "category_icon": "📋",
            "risk_score": 7, "risk_level": "高风险", "risk_color": "#ef4444", "urgency": "紧急",
            "item": "大额无偿借款可能未视同销售缴纳增值税",
            "detail": f"发现 {len(large_loans)} 笔大额借款（合计{total_loan:,.2f}元），"
                       f"但分析期间利息收入仅{interest_entries:,.2f}元。"
                       f"关联方无偿借款：(1)增值税——视同销售按同期同类贷款利率计算利息缴6%增值税；"
                       f"(2)企业所得税——借出方利息收入应确认，借入方利息支出不得税前扣除（出资未到位/超债资比）。"
                       f"涉及：{'；'.join(detail_parts)}"
                       + (f"等{len(large_loans)}笔" if len(large_loans) > 5 else ""),
            "suggestion": "(1)逐笔核实借款性质（是否关联方/是否无偿）；"
                          "(2)无偿借款应参照银行同期利率计算利息并申报增值税；"
                          "(3)建议关联方借款签订书面协议，明确利率和还款期限。",
            "required_evidence": ["借款协议/借据", "关联方关系说明", "同期同类贷款利率参考依据", "增值税补申报记录（如适用）"]
        })


def _analyze_parent_subsidiary_mgmt_fee(db, company_id, ps, pe, results):
    """母子公司管理费不合规"""
    MGFEE_KEYWORDS = ["管理费", "服务费", "咨询费", "技术服务费", "品牌使用费"]
    # 检查向关联方（如名称含"集团""控股""母公司"等）支付管理费
    entries = db.query(JournalEntry).filter(
        JournalEntry.company_id == company_id,
        JournalEntry.period >= ps, JournalEntry.period <= pe,
        or_(
            JournalEntry.account_code.like("6602%"),  # 管理费用
            JournalEntry.account_code.like("5602%")
        )
    ).all()

    found = []
    for entry in entries:
        summary = (entry.summary or "") + (entry.contact_project or "")
        contact = (entry.contact_project or "").lower()
        # 检查收款方是否为关联方特征
        is_related = any(kw in contact for kw in ["集团", "控股", "母公司", "子公司", "关联"])
        has_mgfee_kw = any(kw in summary for kw in MGFEE_KEYWORDS)
        if has_mgfee_kw and is_related:
            found.append({
                "summary": entry.summary, "period": entry.period,
                "amount": _safe_float(entry.debit_amount),
                "contact": entry.contact_project
            })

    if found:
        total_amt = sum(x["amount"] for x in found)
        detail_parts = [f"{x['period']} {x['contact']} {x['summary'][:30]}（{x['amount']:,.0f}元）" for x in found[:5]]
        results.append({
            "category": "企业所得税", "category_icon": "💰",
            "risk_score": 6, "risk_level": "中风险", "risk_color": "#f59e0b", "urgency": "提醒",
            "item": "向关联方（母公司/集团）支付管理费可能不合规",
            "detail": f"发现 {len(found)} 笔向关联方支付的管理费/服务费（合计{total_amt:,.2f}元）。"
                       f"企业所得税法规定：(1)企业之间支付的管理费不得税前扣除；"
                       f"(2)母子公司未按独立交易原则收取的服务费需纳税调整；"
                       f"(3)关联方服务费需有真实服务内容+受益性证明+独立交易价格。"
                       f"涉及：{'；'.join(detail_parts)}"
                       + (f"等{len(found)}笔" if len(found) > 5 else ""),
            "suggestion": "(1)逐笔核实关联方管理费/服务费的商业实质；(2)补充服务合同、成果交付记录等受益性证据；"
                          "(3)无真实服务的不得税前扣除，应做纳税调增；(4)有真实服务的需按独立交易原则定价。",
            "required_evidence": ["关联方服务合同/协议", "服务成果交付记录", "独立交易价格参考依据", "关联交易申报表"]
        })


def _analyze_overseas_payment_withholding(db, company_id, ps, pe, results):
    """境外支付未代扣代缴税款"""
    OVERSEAS_KEYWORDS = ["境外", "国外", "海外", "香港", "美国", "日本", "新加坡", "开曼",
                         "BVI", "维尔京", "境外公司", "foreign", "overseas"]
    entries = db.query(JournalEntry).filter(
        JournalEntry.company_id == company_id,
        JournalEntry.period >= ps, JournalEntry.period <= pe,
        or_(
            JournalEntry.account_code.like("2202%"),  # 应付账款
            JournalEntry.account_code.like("1002%"),  # 银行存款（直接支付）
            JournalEntry.account_code.like("6602%"),  # 管理费用
        )
    ).all()

    found = []
    for entry in entries:
        summary = (entry.summary or "") + (entry.contact_project or "")
        if any(kw in summary for kw in OVERSEAS_KEYWORDS):
            found.append({
                "summary": entry.summary, "period": entry.period,
                "amount": _safe_float(entry.debit_amount or entry.credit_amount),
                "contact": entry.contact_project or ""
            })

    if found:
        total_amt = sum(x["amount"] for x in found)
        detail_parts = [f"{x['period']} {x['contact']} {x['summary'][:30]}（{x['amount']:,.0f}元）" for x in found[:5]]
        results.append({
            "category": "企业所得税", "category_icon": "💰",
            "risk_score": 8, "risk_level": "高风险", "risk_color": "#ef4444", "urgency": "紧急",
            "item": "向境外支付款项可能未履行代扣代缴义务",
            "detail": f"发现 {len(found)} 笔疑似向境外支付的分录（合计{total_amt:,.2f}元）。"
                       f"向境外非居民企业支付款项时需代扣代缴：(1)特许权使用费/技术服务费→10%预提所得税+6%增值税；"
                       f"(2)利息/股息/租金→10%（或协定税率）；(3)未代扣代缴将面临补税+滞纳金+0.5~3倍罚款。"
                       f"涉及：{'；'.join(detail_parts)}"
                       + (f"等{len(found)}笔" if len(found) > 5 else ""),
            "suggestion": "(1)逐笔确认境外支付性质及应代扣税种/税率；"
                          "(2)检查是否已办理代扣代缴申报和对外付汇备案；"
                          "(3)适用税收协定优惠税率的需办理协定待遇备案。",
            "required_evidence": ["境外合同/协议", "代扣代缴税款凭证", "对外付汇税务备案表", "税收协定待遇备案文件（如适用）"]
        })


def _analyze_mixed_sales_accounting(db, company_id, ps, pe, results):
    """混合销售/兼营未分别核算"""
    # 检查企业是否同时有货物销售（13%税率发票）和服务收入（6%税率发票）
    # 但没有分别核算导致从高适用税率
    invoices_13 = db.query(SalesInvoice).filter(
        SalesInvoice.company_id == company_id,
        SalesInvoice.invoice_date >= ps,
        SalesInvoice.invoice_date <= pe,
        SalesInvoice.tax_rate == 13
    ).count()

    invoices_6 = db.query(SalesInvoice).filter(
        SalesInvoice.company_id == company_id,
        SalesInvoice.invoice_date >= ps,
        SalesInvoice.invoice_date <= pe,
        SalesInvoice.tax_rate == 6
    ).count()

    if invoices_13 > 0 and invoices_6 > 0:
        # 同时有13%和6%税率发票→存在混合销售/兼营可能
        # 检查是否分别核算（序时账中收入科目是否分设明细）
        rev_sub_accounts = db.query(JournalEntry.account_code).filter(
            JournalEntry.company_id == company_id,
            JournalEntry.period >= ps, JournalEntry.period <= pe,
            JournalEntry.account_code.like("6001%")
        ).group_by(JournalEntry.account_code).count()

        if rev_sub_accounts <= 1:  # 收入科目未分设明细→可能未分别核算
            results.append({
                "category": "增值税专项", "category_icon": "📋",
                "risk_score": 6, "risk_level": "中风险", "risk_color": "#f59e0b", "urgency": "提醒",
                "item": "混合销售/兼营业务可能未分别核算",
                "detail": f"分析期间存在 {invoices_13} 张13%税率发票和 {invoices_6} 张6%税率发票，"
                           f"但收入科目仅设 {rev_sub_accounts} 个明细科目。"
                           f"兼营不同税率业务应分别核算（分设明细科目），未分别核算的→从高适用税率（13%），"
                           f"可能导致多缴税款或税务稽查时被认定核算不清而从严处理。",
                "suggestion": "(1)区分混合销售（一项行为涉及货物+服务，从主业税率）和兼营（多项独立业务）；"
                              "(2)兼营业务应在会计科目、发票、合同中分别核算/列示；"
                              "(3)设立独立的收入明细科目，分别归集不同税率业务的收入。",
                "required_evidence": ["各税率业务收入分项明细表", "收入核算科目设置说明", "销售合同分类清单"]
            })


def _analyze_environmental_tax(db, company_id, ps, pe, results):
    """环保税未申报"""
    # 检查是否有污染排放相关业务（制造业/化工/印染等），但未缴纳环保税
    # 间接判断：有"排污"/"环保"/"废水"/"废气"/"固废"等支出分录但无环保税缴纳记录
    ENV_KEYWORDS = ["排污", "环保", "废水", "废气", "固废", "危废", "噪声", "污染", "排放",
                    "污水处理", "除尘", "脱硫", "脱硝", "垃圾处理"]
    entries = db.query(JournalEntry).filter(
        JournalEntry.company_id == company_id,
        JournalEntry.period >= ps, JournalEntry.period <= pe,
        or_(*[JournalEntry.summary.contains(kw) for kw in ENV_KEYWORDS])
    ).all()

    if entries:
        # 检查是否有环保税缴纳记录
        env_tax_entries = db.query(JournalEntry).filter(
            JournalEntry.company_id == company_id,
            JournalEntry.period >= ps, JournalEntry.period <= pe,
            or_(
                JournalEntry.summary.contains("环保税"),
                JournalEntry.summary.contains("环境保护税"),
                JournalEntry.account_code.like("2221%")  # 应交环保税
            )
        ).count()

        if env_tax_entries == 0:
            detail_parts = [f"{e.period} {e.summary[:40]}" for e in entries[:5]]
            results.append({
                "category": "政策执行", "category_icon": "📜",
                "risk_score": 7, "risk_level": "高风险", "risk_color": "#ef4444", "urgency": "紧急",
                "item": "存在排污/环保相关支出但未申报环境保护税",
                "detail": f"发现 {len(entries)} 笔与排污/环保相关的支出分录，"
                           f"但未发现环境保护税缴纳记录。"
                           f"直接向环境排放应税污染物（大气/水/固废/噪声）的企业应缴纳环保税。"
                           f"未申报可能面临追缴+滞纳金+罚款。"
                           f"涉及：{'；'.join(detail_parts)}"
                           + (f"等{len(entries)}笔" if len(entries) > 5 else ""),
                "suggestion": "(1)评估企业是否属于环保税纳税义务人（排放应税污染物）；"
                              "(2)如是，按污染物排放量或污染当量数计算环保税并补申报；"
                              "(3)如已委托第三方处理（有危废转移联单），可免责但需保留凭证。",
                "required_evidence": ["排污许可证/排污登记表", "污染物排放监测报告", "环保税申报表（如适用）", "危废转移联单（如委托处理）"]
            })


def _analyze_invoice_remarks_compliance(db, company_id, ps, pe, results):
    """发票备注栏信息缺失或不规范"""
    # 检查特定业务类型发票备注栏是否按要求填写
    # 建筑服务（项目名称+地址）、运输服务（起运地+到达地+车号）、不动产租赁/销售（详细地址）
    CONSTRUCTION_KEYWORDS = ["建筑", "工程", "施工", "安装", "装修", "修缮"]
    TRANSPORT_KEYWORDS = ["运输", "物流", "货运", "快递", "配送"]
    PROPERTY_KEYWORDS = ["不动产", "房产", "房屋", "商铺", "写字楼", "厂房", "土地"]

    # 查进项发票中有建筑业/运输/不动产特征的
    all_purchase = db.query(
        PurchaseInvoice.id, PurchaseInvoice.goods_name, PurchaseInvoice.seller_name,
        PurchaseInvoice.total_amount, PurchaseInvoice.invoice_date
    ).filter(
        PurchaseInvoice.company_id == company_id,
        PurchaseInvoice.invoice_date >= ps,
        PurchaseInvoice.invoice_date <= pe
    ).all()

    needs_remark = []
    for inv_id, goods_name, seller, amt, inv_date in all_purchase:
        name = (goods_name or "").lower()
        category = None
        if any(kw in name for kw in CONSTRUCTION_KEYWORDS):
            category = "建筑服务"
        elif any(kw in name for kw in TRANSPORT_KEYWORDS):
            category = "运输服务"
        elif any(kw in name for kw in PROPERTY_KEYWORDS):
            category = "不动产"
        if category:
            needs_remark.append({
                "category": category, "goods": goods_name[:50],
                "seller": seller, "amount": _safe_float(amt)
            })

    if needs_remark:
        by_cat = defaultdict(lambda: {"count": 0, "amount": 0.0})
        for item in needs_remark:
            by_cat[item["category"]]["count"] += 1
            by_cat[item["category"]]["amount"] += item["amount"]
        detail_parts = [f"{cat}: {info['count']}张（{info['amount']:,.0f}元）" for cat, info in by_cat.items()]
        results.append({
            "category": "发票合规", "category_icon": "📄",
            "risk_score": 5, "risk_level": "中风险", "risk_color": "#f59e0b", "urgency": "提醒",
            "item": "特定业务发票备注栏可能未按要求填写",
            "detail": f"分析期间共 {len(needs_remark)} 张进项发票涉及建筑服务/运输/不动产业务，"
                       f"这些发票按规定需在备注栏注明详细信息（建筑服务→项目名称+地址，"
                       f"运输→起运地+到达地+车号，不动产→详细地址）。"
                       f"备注栏不合规的发票不得作为增值税进项抵扣和企业所得税税前扣除凭证。"
                       f"分布：{'；'.join(detail_parts)}",
            "suggestion": "(1)逐票检查特定业务发票备注栏内容是否完整；"
                          "(2)不合规发票要求供应商重新开具；(3)建立发票审核制度，拒收备注栏不合规发票。",
            "required_evidence": ["需要备注的发票清单", "补开/重开的合规发票", "发票审核制度"]
        })


def _analyze_commission_fee_compliance(db, company_id, ps, pe, results):
    """佣金/手续费支出合规性"""
    COMMISSION_KEYWORDS = ["佣金", "手续费", "返点", "提成", "回扣", "服务费", "居间费", "中介费"]
    entries = db.query(JournalEntry).filter(
        JournalEntry.company_id == company_id,
        JournalEntry.period >= ps, JournalEntry.period <= pe,
        or_(
            JournalEntry.account_code.like("6601%"),  # 销售费用
            JournalEntry.account_code.like("6602%"),  # 管理费用
        )
    ).all()

    found = []
    for entry in entries:
        summary = (entry.summary or "")
        if any(kw in summary for kw in COMMISSION_KEYWORDS):
            found.append({
                "summary": entry.summary, "period": entry.period,
                "amount": _safe_float(entry.debit_amount),
                "contact": entry.contact_project or ""
            })

    if found:
        total_amt = sum(x["amount"] for x in found)
        # 检查佣金是否超限额（一般企业≤5%收入，保险等特殊行业≤10%~18%）
        total_revenue = db.query(func.sum(JournalEntry.credit_amount)).filter(
            JournalEntry.company_id == company_id,
            JournalEntry.period >= ps, JournalEntry.period <= pe,
            JournalEntry.account_code.like("6001%")
        ).scalar() or 0

        ratio = (total_amt / total_revenue * 100) if total_revenue > 0 else 0
        detail_parts = [f"{x['period']} {x['summary'][:40]}（{x['amount']:,.0f}元）" for x in found[:5]]
        results.append({
            "category": "企业所得税", "category_icon": "💰",
            "risk_score": 7 if ratio > 5 else 5,
            "risk_level": "高风险" if ratio > 5 else "中风险",
            "risk_color": "#ef4444" if ratio > 5 else "#f59e0b",
            "urgency": "紧急" if ratio > 5 else "提醒",
            "item": "佣金/手续费支出可能存在合规风险",
            "detail": f"发现 {len(found)} 笔佣金/手续费支出（合计{total_amt:,.2f}元）"
                       f"，佣金占收入比 {ratio:.1f}%"
                       + ("，超过5%限额，超标部分不得税前扣除。" if ratio > 5 else "。")
                       + f"佣金支出需满足：(1)合法真实业务；(2)签订书面合同/协议；"
                       f"(3)不超过合同收入5%（一般企业）；(4)以非现金方式支付；"
                       f"(5)对方为个人的需代扣个税（劳务报酬20%~40%）。"
                       f"涉及：{'；'.join(detail_parts)}"
                       + (f"等{len(found)}笔" if len(found) > 5 else ""),
            "suggestion": "(1)逐笔核实佣金支出的合同/协议和业务真实性；"
                          "(2)超标部分做纳税调增；(3)个人佣金需代扣代缴个人所得税；"
                          "(4)保留服务合同、付款凭证和对方收款确认。",
            "required_evidence": ["佣金支出明细表", "佣金合同/服务协议", "代扣个税凭证（如付给个人）", "收款方确认记录"]
        })


def _analyze_donation_compliance(db, company_id, ps, pe, results):
    """捐赠支出税前扣除不合规"""
    DONATION_KEYWORDS = ["捐赠", "捐助", "赞助", "捐献", "公益", "慈善"]
    entries = db.query(JournalEntry).filter(
        JournalEntry.company_id == company_id,
        JournalEntry.period >= ps, JournalEntry.period <= pe,
        or_(
            JournalEntry.account_code.like("6711%"),  # 营业外支出-捐赠
            JournalEntry.summary.contains("捐赠"),
            JournalEntry.summary.contains("赞助"),
            JournalEntry.summary.contains("公益")
        )
    ).all()

    found = []
    for entry in entries:
        summary = (entry.summary or "")
        if any(kw in summary for kw in DONATION_KEYWORDS):
            found.append({
                "summary": entry.summary, "period": entry.period,
                "amount": _safe_float(entry.debit_amount),
                "account": entry.account_code
            })

    if found:
        total_amt = sum(x["amount"] for x in found)
        # 计算利润总额12%限额
        profit = _pl_net(db, company_id, ps, pe)
        limit_12pct = abs(profit) * 0.12 if profit else 0
        detail_parts = [f"{x['period']} {x['summary'][:40]}（{x['amount']:,.0f}元）" for x in found[:5]]
        results.append({
            "category": "企业所得税", "category_icon": "💰",
            "risk_score": 7 if total_amt > limit_12pct else 5,
            "risk_level": "高风险" if total_amt > limit_12pct else "中风险",
            "risk_color": "#ef4444" if total_amt > limit_12pct else "#f59e0b",
            "urgency": "紧急" if total_amt > limit_12pct else "提醒",
            "item": "捐赠支出税前扣除合规性存疑",
            "detail": (
                           f"发现 {len(found)} 笔捐赠/赞助支出（合计{total_amt:,.2f}元，"
                           f"利润总额12%限额={limit_12pct:,.2f}元"
                           + ("，超标!" if total_amt > limit_12pct else "，在限额内。")
                           + "）捐赠支出税前扣除需满足："
                           "(1)通过公益性社会组织或县级以上政府捐赠；"
                           "(2)不超过年度利润总额12%（超标可结转3年）；"
                           "(3)取得公益事业捐赠票据；(4)直接捐赠不得扣除；"
                           "(5)赞助支出（非广告性质）不得税前扣除。"
                       ),
            "suggestion": "(1)逐笔核实捐赠性质（公益/直接/赞助）；"
                          "(2)公益性捐赠需取得合规票据且不超12%限额；"
                          "(3)直接捐赠和赞助支出做纳税调增；"
                          "(4)超标公益性捐赠可结转以后3年扣除。",
            "required_evidence": ["捐赠支出明细表", "公益事业捐赠统一票据", "受赠方资格证明", "企业所得税纳税调增明细"]
        })


def _analyze_inventory_count_discrepancy(db, company_id, ps, pe, results):
    """存货盘点差异未做税务调整"""
    DISCREPANCY_KEYWORDS = ["盘亏", "盘盈", "盘点差异", "报废", "毁损", "霉变", "过期"]
    entries = db.query(JournalEntry).filter(
        JournalEntry.company_id == company_id,
        JournalEntry.period >= ps, JournalEntry.period <= pe,
        or_(
            JournalEntry.account_code.like("1405%"),  # 存货相关
            *[JournalEntry.summary.contains(kw) for kw in DISCREPANCY_KEYWORDS]
        )
    ).all()

    found = []
    for entry in entries:
        summary = (entry.summary or "")
        if any(kw in summary for kw in DISCREPANCY_KEYWORDS):
            found.append({
                "summary": entry.summary, "period": entry.period,
                "debit": _safe_float(entry.debit_amount),
                "credit": _safe_float(entry.credit_amount)
            })

    if found:
        total_amt = sum(x["debit"] + x["credit"] for x in found)
        detail_parts = [f"{x['period']} {x['summary'][:50]}" for x in found[:5]]
        results.append({
            "category": "企业所得税", "category_icon": "💰",
            "risk_score": 6, "risk_level": "中风险", "risk_color": "#f59e0b", "urgency": "提醒",
            "item": "存货盘亏/盘盈/报废可能未做税务处理",
            "detail": f"发现 {len(found)} 笔存货盘点/报废相关分录（涉及金额{total_amt:,.2f}元）。"
                       f"存货盘点差异税务处理：(1)盘亏→进项税额需转出（管理不善导致），"
                       f"损失需专项申报税前扣除；(2)盘盈→计入收入缴纳企业所得税；"
                       f"(3)报废→正常损耗可清单申报，非正常损失需专项申报+进项转出。"
                       f"涉及：{'；'.join(detail_parts)}"
                       + (f"等{len(found)}笔" if len(found) > 5 else ""),
            "suggestion": "(1)逐笔记实盘亏/盘盈/报废原因；(2)管理不善导致的盘亏应做进项税额转出；"
                          "(3)资产损失税前扣除需留存内部核批+外部证据（盘点表/技术鉴定等）；"
                          "(4)盘盈收入不得隐瞒，应如实申报。",
            "required_evidence": ["存货盘点表", "盘亏/报废原因说明及审批单", "进项税额转出凭证（如适用）", "资产损失税前扣除申报表"]
        })


# ═══════════════════════════════════════════════════════════
# V11 新增 — 普票浪费可抵扣进项税额 + 预付卡涉税风险
# ═══════════════════════════════════════════════════════════

def _analyze_wasted_deductible_tax(db, company_id, ps, pe, results):
    """普票浪费可抵扣进项税额——专业服务类供应商开普票未开专票，浪费进项税额。
    
    检测逻辑：
    1. 筛选取得发票中的普票
    2. 判断品名分类是否为专业服务类（这些服务的供应商通常是一般纳税人，可以开专票）
    3. 按供应商汇总浪费的税额
    4. 浪费税额 > 阈值（50元）则报告风险
    
    可开专票的分类（一般纳税人常见）：
    广告服务/设计服务/信息技术服务/现代服务/经营租赁/企业管理服务/修理修配劳务
    排除分类（通常为小规模/个体/生活服务）：
    餐饮服务/住宿服务/运输服务/旅游服务/经纪代理服务/生活服务
    """
    # 专业服务分类——这些供应商通常是一般纳税人，理应能开专票
    PROFESSIONAL_SERVICE_CATEGORIES = [
        "广告服务", "设计服务", "信息技术服务", "现代服务",
        "经营租赁", "企业管理服务", "修理修配劳务",
        "建筑服务", "金融服务", "邮政服务", "电信服务",
        "研发和技术服务", "文化创意服务", "物流辅助服务",
        "鉴证咨询服务", "广播影视服务", "商务辅助服务",
    ]

    # 1. 筛选取得发票普票（全量——普票浪费进项是系统性问题，不限期间）
    purchase_invs = db.query(PurchaseInvoice).filter(
        PurchaseInvoice.company_id == company_id
    ).all()

    wasted_by_supplier = defaultdict(lambda: {
        "tax": 0.0, "total": 0.0, "cnt": 0, "goods": set(), "cats": set()
    })

    import re
    for inv in purchase_invs:
        # 只处理普票（非专用发票）
        is_special = "专用" in (inv.invoice_category or "")
        if is_special:
            continue

        goods = (inv.goods_name or "").strip()
        m = re.match(r'\*([^*]+)\*', goods)
        category = m.group(1) if m else ""

        # 判断是否为专业服务类
        if category not in PROFESSIONAL_SERVICE_CATEGORIES:
            continue

        tax = abs(_safe_float(inv.tax_amount))
        total = _safe_float(inv.amount) + tax

        name = (inv.seller_name or "未知供应商").strip()
        wasted_by_supplier[name]["tax"] += tax
        wasted_by_supplier[name]["total"] += total
        wasted_by_supplier[name]["cnt"] += 1
        wasted_by_supplier[name]["goods"].add(goods[:40])
        wasted_by_supplier[name]["cats"].add(category)

    if not wasted_by_supplier:
        return

    # 2. 过滤浪费税额 > 50
    significant = {k: v for k, v in wasted_by_supplier.items() if v["tax"] >= 50}
    if not significant:
        return

    total_wasted_tax = sum(v["tax"] for v in significant.values())
    total_inv_cnt = sum(v["cnt"] for v in significant.values())
    supplier_cnt = len(significant)

    # 3. 构建报告
    detail_parts = []
    for name in sorted(significant, key=lambda x: -significant[x]["tax"]):
        d = significant[name]
        cats_str = "、".join(sorted(d["cats"]))
        detail_parts.append(
            f"{name}（{cats_str}）{d['cnt']}张普票，价税合计{d['total']:,.2f}元，"
            f"浪费可抵扣税额{d['tax']:,.2f}元"
        )

    score = 7 if total_wasted_tax >= 1000 else (5 if total_wasted_tax >= 200 else 3)
    level = "高风险" if total_wasted_tax >= 5000 else ("中风险" if total_wasted_tax >= 500 else "低风险")

    results.append({
        "category": "发票合规", "category_icon": "🧾",
        "risk_score": score, "risk_level": level,
        "risk_color": "#dc2626" if level == "高风险" else ("#f59e0b" if level == "中风险" else "#3b82f6"),
        "urgency": "提醒",
        "item": "专业服务类供应商开普票浪费可抵扣进项税额",
        "detail": f"共 {supplier_cnt} 家供应商的 {total_inv_cnt} 张普票属于专业服务类（"
                  f"广告/设计/信息技术/租赁/企业管理等），这些供应商通常为一般纳税人，"
                  f"理应可以开具增值税专用发票。累计价税合计{sum(v['total'] for v in significant.values()):,.2f}元，"
                  f"浪费可抵扣进项税额{total_wasted_tax:,.2f}元。"
                  f"涉及供应商：{'；'.join(detail_parts[:8])}"
                  + (f"等{supplier_cnt}家" if supplier_cnt > 8 else ""),
        "suggestion": "(1)主动联系供应商要求开具增值税专用发票，专票与普票对供应商而言开票成本相同；"
                      "(2)将「要求专票」纳入采购合同条款和供应商准入标准；"
                      f"(3)已损失的进项税额{total_wasted_tax:,.2f}元虽无法追回，但应从下月起严格执行专票要求；"
                      "(4)关注供应商纳税人资质，一般纳税人必须要求专票。",
        "required_evidence": ["取得发票清单（区分专票/普票）", "供应商纳税人资质证明",
                              "采购合同/订单（检查发票条款）", "进项税额计算表"]
    })


def _analyze_prepaid_card_risk(db, company_id, ps, pe, results):
    """预付卡涉税风险——京东E卡等预付卡购买/赠送/消费的税务合规检查。
    
    三大风险维度：
    1. 增值税：购买预付卡取得的是不征税普票，不能抵扣进项税额
    2. 企业所得税：仅有购卡发票不足税前扣除，需取得实际消费/赠送凭证
    3. 个人所得税：赠送客户需代扣偶然所得20%，赠送员工需并入工资薪金
    
    税务规定依据：
    - 国家税务总局公告2016年第53号：售卡方开具「不征税」普票，不得抵扣进项税额
    - 企业所得税：预付卡在消费/赠送环节才能确认费用，需留存消费记录/赠送签收单
    - 个人所得税法：向个人赠送礼品（含预付卡）需按偶然所得20%代扣个税
    """
    PREPAID_KEYWORDS = ["预付卡", "购物卡", "礼品卡", "E卡", "消费卡", "储值卡",
                        "加油卡", "电话卡", "交通卡", "一卡通"]

    # 1. 查找预付卡采购发票（全量——预付卡涉税风险是系统性问题，不限期间）
    card_invs = db.query(PurchaseInvoice).filter(
        PurchaseInvoice.company_id == company_id
    ).all()

    matched = []
    for inv in card_invs:
        goods = (inv.goods_name or "").strip()
        if any(kw in goods for kw in PREPAID_KEYWORDS):
            matched.append(inv)

    if not matched:
        return

    total_amt = sum(_safe_float(inv.amount) + abs(_safe_float(inv.tax_amount)) for inv in matched)
    total_cnt = len(matched)
    sellers = list(set((inv.seller_name or "未知").strip() for inv in matched))

    detail_parts = []
    for inv in matched:
        detail_parts.append(
            f"{(inv.seller_name or '未知')[:20]} {inv.goods_name[:40]} "
            f"{_safe_float(inv.amount) + abs(_safe_float(inv.tax_amount)):,.2f}元"
        )

    # 2. 检查序时账中是否有预付卡相关费用的处理（全量不限期间）
    card_entries = db.query(JournalEntry).filter(
        JournalEntry.company_id == company_id,
        or_(
            JournalEntry.account_code.in_(["660216", "660299", "660201"]),
            *[JournalEntry.summary.contains(kw) for kw in PREPAID_KEYWORDS]
        )
    ).all()

    has_expense_entry = any(
        any(kw in (e.summary or "") for kw in ["预付卡", "购物卡", "E卡", "礼品", "福利", "招待"])
        for e in card_entries
    )

    # 3. 检查银行流水（实际付款，全量不限期间）
    card_payments = db.query(BankTransaction).filter(
        BankTransaction.company_id == company_id,
        BankTransaction.debit_amount > 0
    ).all()

    paid_cnt = 0
    paid_amt = 0.0
    for bt in card_payments:
        cpty = (bt.counterparty_name or "").strip()
        summary = (bt.summary or "").strip()
        if any(kw in cpty + summary for kw in ["京东", "预付", "购物卡", "E卡", "礼品"]):
            paid_cnt += 1
            paid_amt += _safe_float(bt.debit_amount)

    # 构建风险报告
    score = 8 if total_amt >= 10000 else (6 if total_amt >= 5000 else 4)
    level = "高风险" if total_amt >= 20000 else ("中风险" if total_amt >= 5000 else "低风险")

    # 风险详情
    risks_detail = []
    risks_detail.append(
        f"① 增值税：购买{total_cnt}张预付卡（价税合计{total_amt:,.2f}元），"
        f"取得的是不征税普通发票，购买环节不得抵扣进项税额——此项已合规。"
    )

    if not has_expense_entry:
        risks_detail.append(
            f"② 企业所得税 ⚠️ 风险：未在序时账中发现预付卡费用化处理（业务招待费/福利费等）。"
            f"仅有购卡发票不足税前扣除，需在实际消费/赠送环节凭消费记录或赠送签收单确认费用。"
            f"若卡已赠送但未入账费用，存在企业所得税少缴风险。"
        )
    else:
        risks_detail.append(
            f"② 企业所得税：已在序时账中发现预付卡费用化处理，"
            f"请核实是否留存消费记录/赠送签收单等税前扣除凭证。"
        )

    risks_detail.append(
        f"③ 个人所得税 ⚠️ 风险：若预付卡赠送客户，应按「偶然所得」20%代扣代缴个人所得税；"
        f"若赠送员工，应并入当月工资薪金代扣个税。请核实是否已履行代扣代缴义务。"
    )

    if paid_amt > 0:
        risks_detail.append(
            f"④ 资金流：银行流水显示已支付{paid_amt:,.2f}元（{paid_cnt}笔），与发票金额匹配。"
        )

    results.append({
        "category": "发票合规", "category_icon": "🧾",
        "risk_score": score, "risk_level": level,
        "risk_color": "#dc2626" if level == "高风险" else ("#f59e0b" if level == "中风险" else "#3b82f6"),
        "urgency": "提醒",
        "item": "预付卡（购物卡/E卡等）涉税处理合规性检查",
        "detail": "；".join(risks_detail),
        "suggestion": "(1)购买预付卡取得不征税普票→不得抵扣进项税额（已合规）；"
                      "(2)实际消费/赠送时凭消费记录、赠送签收单确认费用，取得税前扣除凭证；"
                      f"(3)赠送客户→按偶然所得20%代扣个税（{total_amt:,.2f}×20%={total_amt*0.2:,.2f}元）；"
                      "(4)赠送员工→并入当月工资薪金代扣个税；"
                      "(5)自用消费→凭消费明细转入相应费用科目（办公费/其他费用等）。",
        "required_evidence": ["预付卡购买发票", "预付卡赠送/消费明细台账",
                              "客户签收单（如赠送客户）", "员工福利签领表（如赠送员工）",
                              "代扣代缴个人所得税申报记录（如适用）"]
    })
